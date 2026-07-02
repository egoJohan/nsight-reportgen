import pytest

from nsight import config
from nsight.coding import top_words
from nsight.store.survey_store import SurveyStore
from nsight.tabulate import awareness_by_brand


@pytest.fixture(scope="module", autouse=True)
def _isolated_survey_db(tmp_path_factory):
    """Isolate these goldens' SurveyStore in a FRESH per-module DuckDB. They otherwise
    share the fixed `work/survey.duckdb`, which another test can leave holding a
    different dataset (the goldens only ingest `if frame().empty`) — making them fail
    purely on test-collection order. A private DB makes them order-independent."""
    db = tmp_path_factory.mktemp("golden_db") / "survey.duckdb"
    orig = config.SURVEY_DB
    config.SURVEY_DB = db
    try:
        yield
    finally:
        config.SURVEY_DB = orig


@pytest.mark.integration
def test_aided_awareness_matches_deck():
    if not config.ATTENDO_SAV.exists():
        pytest.skip("Attendo files not present")
    store = SurveyStore(db_path=config.SURVEY_DB, codebook_path=config.CODEBOOK_JSON)
    if store.frame().empty:
        store.ingest(config.ATTENDO_SAV)
    from nsight.attendo_bindings import (
        AIDED_AWARENESS_POSITIVE,
        AIDED_AWARENESS_VARS,
        DECK_AIDED_AWARENESS,
        WEIGHT_VAR,
    )

    frame = store.frame()
    res = awareness_by_brand(
        frame,
        brand_vars=AIDED_AWARENESS_VARS,
        positive_values=AIDED_AWARENESS_POSITIVE,
        weight=WEIGHT_VAR,
    )
    for brand, expected in DECK_AIDED_AWARENESS.items():
        got = res[brand].pct
        assert abs(got - expected) <= 1, f"{brand}: got {got}, want {expected}"


@pytest.mark.integration
def test_brand_image_top10_overlaps_deck():
    if not config.ATTENDO_SAV.exists():
        pytest.skip("Attendo files not present")
    store = SurveyStore(db_path=config.SURVEY_DB, codebook_path=config.CODEBOOK_JSON)
    if store.frame().empty:
        store.ingest(config.ATTENDO_SAV)
    from nsight.attendo_bindings import IMAGE_WORD_VARS, IMAGE_SYNONYMS, DECK_IMAGE_TOP10
    res = top_words(store.frame(), text_vars=IMAGE_WORD_VARS, top_n=10, synonyms=IMAGE_SYNONYMS)
    got = [w for w, _ in res]
    overlap = len(set(got) & set(DECK_IMAGE_TOP10))
    assert overlap >= 7, f"only {overlap}/10 overlap; got {got}; want {DECK_IMAGE_TOP10}"


@pytest.mark.integration
def test_general_opinion_matches_deck():
    """M-2: general opinion (var20 private / var21 public) vs deck slide 17."""
    if not config.ATTENDO_SAV.exists():
        pytest.skip("Attendo files not present")
    import pandas as pd
    from nsight.attendo_bindings import (
        DECK_OPINION_DIST,
        DECK_OPINION_POSITIVE,
        OPINION_NEGATIVE,
        OPINION_NEUTRAL,
        OPINION_POSITIVE,
        OPINION_PRIVATE_VAR,
        OPINION_PUBLIC_VAR,
        OPINION_SERIES_CODES,
    )

    store = SurveyStore(db_path=config.SURVEY_DB, codebook_path=config.CODEBOOK_JSON)
    if store.frame().empty:
        store.ingest(config.ATTENDO_SAV)
    frame = store.frame()

    print("\nM-2 general opinion (computed vs deck):")
    for key, var in (("private", OPINION_PRIVATE_VAR), ("public", OPINION_PUBLIC_VAR)):
        s = pd.to_numeric(frame[var], errors="coerce")
        tot = s.notna().sum()
        # per-series chart proportions
        for name, code in OPINION_SERIES_CODES.items():
            got = round((s == code).sum() / tot, 2)
            want = DECK_OPINION_DIST[key][name]
            print(f"  {key} {name}: got {got} deck {want}")
            assert abs(got - want) <= 0.01, f"{key} {name}: {got} vs {want}"
        pos = round(((s.isin(OPINION_POSITIVE)).sum()) / tot * 100)
        want_pos = DECK_OPINION_POSITIVE[key]
        neg = round((s.isin(OPINION_NEGATIVE)).sum() / tot * 100)
        neu = round((s.isin(OPINION_NEUTRAL)).sum() / tot * 100)
        print(f"  {key} POSITIVE: got {pos} deck {want_pos} (neg {neg}, neutral {neu})")
        assert abs(pos - want_pos) <= 1, f"{key} positive: {pos} vs {want_pos}"


@pytest.mark.integration
def test_spontaneous_awareness_matches_deck():
    """M-3: spontaneous any-mention + top-of-mind (var17 open list) vs deck slide 13."""
    if not config.ATTENDO_SAV.exists():
        pytest.skip("Attendo files not present")
    from nsight.attendo_bindings import (
        DECK_SPONTANEOUS,
        SPONTANEOUS_FIRST_VAR,
        SPONTANEOUS_MENTION_VARS,
        SPONTANEOUS_PATTERNS,
    )
    from nsight.tabulate import spontaneous_any_mention, top_of_mind_patterns

    store = SurveyStore(db_path=config.SURVEY_DB, codebook_path=config.CODEBOOK_JSON)
    if store.frame().empty:
        store.ingest(config.ATTENDO_SAV)
    frame = store.frame()

    print("\nM-3 spontaneous awareness (computed vs deck):")
    for brand, (deck_any, deck_tom) in DECK_SPONTANEOUS.items():
        pats = SPONTANEOUS_PATTERNS[brand]
        am = spontaneous_any_mention(frame, mention_vars=SPONTANEOUS_MENTION_VARS, patterns=pats).pct
        tom = top_of_mind_patterns(frame, first_mention_var=SPONTANEOUS_FIRST_VAR, patterns=pats).pct
        print(f"  {brand}: any={am} (deck {deck_any}) | tom={tom} (deck {deck_tom})")
        assert abs(am - deck_any) <= 1, f"{brand} any: {am} vs {deck_any}"
        assert abs(tom - deck_tom) <= 1, f"{brand} tom: {tom} vs {deck_tom}"


@pytest.mark.integration
def test_image_words_rendered():
    """Slide-24 brand-image TOP-10 word list is data-driven (regenerated from SPSS).

    Counts differ from the deck (deck used manual thematic merging), so we assert on
    lemma OVERLAP: the top lemma "Kallis" is present and >=7 of the 10 deck lemmas
    appear in the rendered current-wave word box.
    """
    if not config.ATTENDO_SAV.exists():
        pytest.skip("Attendo files not present")
    from pptx import Presentation

    from nsight.attendo_bindings import DECK_IMAGE_TOP10
    from nsight.generate import generate_deck

    out = generate_deck(
        sav=config.ATTENDO_SAV,
        brief_path=config.BRIEFS_DIR / "attendo.md",
        template=config.ATTENDO_TEMPLATE,
        out=config.GENERATED_PPTX,
    )

    shape = [
        s
        for s in Presentation(str(out)).slides[24].shapes
        if s.name == "Rectangle: Rounded Corners 7"
    ][0]
    text = shape.text_frame.text
    print("\nSlide-24 current-wave (Marras 25) rendered TOP-10:\n" + text)

    # Header preserved.
    assert "TOP 10" in text
    assert "Marras 25" in text
    # Top lemma present (capitalized like the deck).
    assert "Kallis" in text
    # Overlap with deck lemmas (counts differ; compare lemma membership).
    lowered = text.casefold()
    overlap = sum(1 for lemma in DECK_IMAGE_TOP10 if lemma.casefold() in lowered)
    assert overlap >= 7, f"only {overlap}/10 deck lemmas in rendered text:\n{text}"


@pytest.mark.integration
def test_full_deck_fidelity():
    # Plumbing / no-corruption check: verifies that shapes survive round-trip without
    # structural damage. Real numeric correctness is guarded by the unit tests and
    # test_aided_awareness_matches_deck + test_generated_slide14_waves.
    if not config.ATTENDO_SAV.exists():
        pytest.skip("Attendo files not present")
    from nsight.generate import generate_deck
    from nsight.fidelity.extract import extract_deck
    from nsight.fidelity.compare import compare_decks
    out = generate_deck(sav=config.ATTENDO_SAV, brief_path=config.BRIEFS_DIR / "attendo.md",
                        template=config.ATTENDO_TEMPLATE, out=config.GENERATED_PPTX)
    rep = compare_decks(extract_deck(out), extract_deck(config.ATTENDO_TEMPLATE))
    print("chart fidelity:", rep.chart_score, "matched", rep.charts_matched, "/", rep.charts_total)
    print("first mismatches:", rep.mismatches[:10])
    assert rep.chart_score >= 99.0


@pytest.mark.integration
def test_generated_slide14_waves():
    """Verify slide-14 current-wave values come from the SPSS computation, not a template copy.

    The assertion ties generated proportions to independently computed tabulation
    results, so a no-op template copy would produce different numbers (the prior wave's
    values) and fail the anti-copy guard.
    """
    if not config.ATTENDO_SAV.exists():
        pytest.skip("Attendo files not present")
    from pptx import Presentation
    from nsight.generate import generate_deck
    from nsight.store.survey_store import SurveyStore
    from nsight.attendo_bindings import AIDED_AWARENESS_VARS, AIDED_AWARENESS_POSITIVE, WEIGHT_VAR
    from nsight.tabulate import awareness_by_brand

    generate_deck(sav=config.ATTENDO_SAV, brief_path=config.BRIEFS_DIR / "attendo.md",
                  template=config.ATTENDO_TEMPLATE, out=config.GENERATED_PPTX)

    def series_of(path):
        sh = [s for s in Presentation(str(path)).slides[14].shapes
              if s.name == "Content Placeholder 9"][0]
        chart_cats = [str(c) for c in sh.chart.plots[0].categories]
        series = {s.name: list(s.values) for s in sh.chart.plots[0].series}
        return chart_cats, series

    gen_cats, gen_series = series_of(config.GENERATED_PPTX)
    _, tmpl_series = series_of(config.ATTENDO_TEMPLATE)

    # --- Independently compute expected proportions from SPSS frame ---
    store = SurveyStore(db_path=config.SURVEY_DB, codebook_path=config.CODEBOOK_JSON)
    if store.frame().empty:
        store.ingest(config.ATTENDO_SAV)
    frame = store.frame()
    res = awareness_by_brand(
        frame,
        brand_vars=AIDED_AWARENESS_VARS,
        positive_values=AIDED_AWARENESS_POSITIVE,
        weight=WEIGHT_VAR,
    )

    # Build expected proportions aligned to the chart's category order.
    # gen_cats holds the brand labels as they appear in the chart.
    expected_proportions = [round(res[cat].pct) / 100 for cat in gen_cats]

    # Assert each generated "Marraskuu 2025" proportion matches tabulation.
    for i, (cat, exp) in enumerate(zip(gen_cats, expected_proportions)):
        got = gen_series["Marraskuu 2025"][i]
        assert abs(got - exp) < 1e-9, (
            f"Marraskuu 2025 [{cat}] (idx {i}): generated {got!r} != expected {exp!r} "
            f"(from tabulation pct={res[cat].pct})"
        )

    # --- Anti-copy guard ---
    # A no-op template copy would leave "Marraskuu 2025" equal to its template value,
    # which equals the prior-wave "Toukokuu 2025" template series — they should differ.
    # More directly: the generated current wave must NOT equal the template's "Toukokuu 2024"
    # series (an entirely different wave), confirming the values were actually written.
    assert gen_series["Marraskuu 2025"] != tmpl_series["Toukokuu 2024"], (
        "Anti-copy guard: generated 'Marraskuu 2025' equals template 'Toukokuu 2024' — "
        "the current wave was NOT written (or the data is identical, which is implausible)"
    )

    # --- Prior-wave preservation ---
    for wave in ("Toukokuu 2025", "Marraskuu 2024", "Toukokuu 2024"):
        assert gen_series[wave] == tmpl_series[wave], f"prior wave {wave} changed"
