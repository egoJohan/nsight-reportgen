"""Integration tests for the pure-python PPTX build layer (SOFFICE-FREE).

`reportbuilder.export.pptx_build.build_pptx` is pure python-pptx + matplotlib —
no LibreOffice — so these tests run in the standard suite. They assert real
behavior: a valid deck that python-pptx can reopen, one slide per chart spec,
image mode places PICTUREs while native mode places `c:chart` graphic frames,
and a spec that fails to compute still yields a deck (empty-series fallback).
"""
from __future__ import annotations

import dataclasses

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from reportbuilder.export.pptx_build import build_pptx
from reportbuilder.model.report import (
    ChartSpec,
    ElementToggles,
    NumberFormat,
    Report,
    SortSpec,
)
from reportbuilder.testing.fixtures import (
    one_chart_report,
    tiny_model_and_data,
    two_chart_report,
)


def _has_chart(slide) -> bool:
    """True if the slide carries a native `c:chart` graphic frame."""
    return any(getattr(sh, "has_chart", False) for sh in slide.shapes)


def _has_picture(slide) -> bool:
    return any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in slide.shapes)


# ---------------------------------------------------------------------------
# Valid decks reopen and are complete (one slide per chart spec)
# ---------------------------------------------------------------------------

@pytest.mark.parametrize("render_mode", ["native", "image"])
def test_build_pptx_one_chart_reopens(tmp_path, render_mode):
    """A one-chart Report builds a .pptx that python-pptx reopens; 1 slide == 1 spec."""
    model, data = tiny_model_and_data()
    report = dataclasses.replace(one_chart_report(), render_mode=render_mode)

    out = build_pptx(report, model, data, str(tmp_path / "one.pptx"))

    assert out.endswith("one.pptx")
    prs = Presentation(out)  # reopens without exception
    assert len(prs.slides._sldIdLst) == len(report.charts) == 1


@pytest.mark.parametrize("render_mode", ["native", "image"])
def test_build_pptx_two_charts_complete(tmp_path, render_mode):
    """A two-chart Report yields exactly two slides (completeness)."""
    model, data = tiny_model_and_data()
    report = dataclasses.replace(two_chart_report(), render_mode=render_mode)

    out = build_pptx(report, model, data, str(tmp_path / "two.pptx"))

    prs = Presentation(out)
    assert len(prs.slides._sldIdLst) == len(report.charts) == 2


# ---------------------------------------------------------------------------
# Render-mode shape kinds: native -> c:chart frames, image -> PICTUREs
# ---------------------------------------------------------------------------

def test_native_mode_places_chart_frames(tmp_path):
    """Native render mode places a `c:chart` graphic frame (has_chart) and no picture."""
    model, data = tiny_model_and_data()
    report = dataclasses.replace(one_chart_report(), render_mode="native")

    prs = Presentation(build_pptx(report, model, data, str(tmp_path / "n.pptx")))
    slide = next(iter(prs.slides))

    assert _has_chart(slide), "native mode must place a c:chart graphic frame"
    assert not _has_picture(slide), "native mode must not rasterize to a picture"


def test_image_mode_places_pictures(tmp_path):
    """Image render mode places a PICTURE and no native chart frame."""
    model, data = tiny_model_and_data()
    report = dataclasses.replace(one_chart_report(), render_mode="image")

    prs = Presentation(build_pptx(report, model, data, str(tmp_path / "i.pptx")))
    slide = next(iter(prs.slides))

    assert _has_picture(slide), "image mode must place a rasterized PICTURE"
    assert not _has_chart(slide), "image mode must not place a native chart frame"


# ---------------------------------------------------------------------------
# Failing compute falls back to an empty series (no raise, deck still built)
# ---------------------------------------------------------------------------

def test_failing_compute_still_yields_deck(tmp_path):
    """A spec whose compute() raises still produces a complete deck (empty-series fallback).

    A wordcloud over the categorical q1 has no text answers, so
    `compute` raises ValueError; `build_pptx` must swallow it, fall back to an
    empty series, and still emit one slide. Image mode is used because a
    wordcloud is not supported by the native renderer.
    """
    model, data = tiny_model_and_data()
    bad_spec = ChartSpec(
        question_ref="q1",
        chart_type="wordcloud",
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"),
        template_slot="slot1",
        elements=ElementToggles(),
    )
    report = Report(
        name="Fallback",
        render_mode="image",
        template_ref="t.pptx",
        charts=(bad_spec,),
    )

    out = build_pptx(report, model, data, str(tmp_path / "fallback.pptx"))

    prs = Presentation(out)  # deck is valid despite the compute failure
    assert len(prs.slides._sldIdLst) == len(report.charts) == 1
