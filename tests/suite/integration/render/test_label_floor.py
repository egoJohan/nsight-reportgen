"""The percentage under which a data label is not drawn is the author's to set.

A 100%-stacked scale puts the whole tail of the scale — the 1 %, the 2 %, the
3 % — into slivers narrower than the number that belongs in them. Below 1 % of
the axis the label was dropped silently; above it, it was drawn anyway and
collided with its neighbour and with the bar's own name, which is what a reader
of the finished deck sees as "the small percentages are unreadable". There is no
one right cut-off: it depends on how wide the chart is drawn and how much the
author is willing to lose, so it belongs to the author.

Default is what the renderer did before: 1 % of the value axis for a stack, 4 %
for a pie wedge, whose labels sit on a curve and collide sooner.
"""
from __future__ import annotations

import pytest

from reportbuilder.model.report import NumberFormat
from reportbuilder.render.image._mpl import label_floor


def test_the_default_is_what_the_renderer_already_did():
    assert label_floor(NumberFormat(), default_pct=1.0, axis_max=100.0) == 1.0
    assert label_floor(NumberFormat(), default_pct=4.0, axis_max=100.0) == 4.0


def test_the_author_s_own_cut_off_wins_over_the_chart_s_default():
    fmt = NumberFormat(hide_below_pct=5.0)
    assert label_floor(fmt, default_pct=1.0, axis_max=100.0) == 5.0
    assert label_floor(fmt, default_pct=4.0, axis_max=100.0) == 5.0


def test_the_cut_off_is_a_share_of_the_axis_not_a_number_of_units():
    """A true-width chart's axis does not stop at 100.

    On a 0-465 axis a "1 unit" sliver is invisible but would still be labelled,
    and the labels pile onto each other — the reason the floor was written as a
    proportion in the first place.
    """
    assert label_floor(NumberFormat(hide_below_pct=2.0),
                       default_pct=1.0, axis_max=465.0) == pytest.approx(9.3)


def test_zero_shows_every_value_the_chart_has():
    assert label_floor(NumberFormat(hide_below_pct=0.0),
                       default_pct=4.0, axis_max=100.0) == 0.0


# ── through the renderer ─────────────────────────────────────────────────────

def _stacked_png(hide_below_pct):
    """A 7-point scale whose tail is 1 %, 2 % and 3 %, drawn as a 100% stack."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches

    from reportbuilder.model.report import ChartSpec, ElementToggles, SortSpec
    from reportbuilder.render.base import RenderContext, Slot, StyleSpec
    from reportbuilder.render.image.bars import build_image_bar_stacked
    from reportbuilder.stats.series import Cell, SeriesResult

    cats = ("1", "2", "3", "4", "5", "6", "7")
    vals = [1.0, 2.0, 3.0, 17.0, 30.0, 30.0, 17.0]
    cells = {(c, "Suomi"): Cell(pct=v, count=v * 10, mean=None)
             for c, v in zip(cats, vals)}
    series = SeriesResult(categories=cats, segments=("Suomi",), cells=cells,
                          base_n={"Suomi": 1000}, statistic="pct")
    spec = ChartSpec(
        question_ref="q1", chart_type="stacked_horizontal_bar", statistic="pct",
        classifying_var="maa", number_format=NumberFormat(hide_below_pct=hide_below_pct),
        sort=SortSpec(basis="data_order"), template_slot="slot1",
        elements=ElementToggles(title=False, legend=True, data_labels=True))
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slot = Slot(slide_index=0, left=Inches(0.5), top=Inches(0.5),
                width=Inches(12.0), height=Inches(5.0), name="slot1")
    build_image_bar_stacked(RenderContext(
        slide=slide, slot=slot, style=StyleSpec(), spec=spec, series=series,
        fmt=spec.number_format))
    pic = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE][0]
    return pic.image.blob


def test_raising_the_cut_off_takes_the_small_numbers_off_the_chart():
    """The 3 % label is drawn by default and gone at 5 %."""
    assert _stacked_png(None) != _stacked_png(5.0)


def test_the_default_draws_exactly_what_1_percent_draws():
    """`None` must mean the renderer's own long-standing floor, to the pixel —
    otherwise adding the setting quietly redraws every existing deck."""
    assert _stacked_png(None) == _stacked_png(1.0)


def test_zero_puts_every_value_back_on_the_chart():
    assert _stacked_png(0.0) != _stacked_png(None)
