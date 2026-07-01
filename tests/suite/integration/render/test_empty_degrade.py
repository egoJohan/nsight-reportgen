"""Empty / degenerate series degrade to a placeholder picture instead of crashing.

Covers series_is_empty detection, render_empty_chart placement, the deck's
empty branch, and build_pptx's compute-failure fallback (all soffice-free).
"""
from __future__ import annotations

import dataclasses

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from suite._helpers import make_ctx, assert_single_picture, picture_shapes
from suite.integration.render._series import empty_series, all_zero_series

from reportbuilder.render.base import StyleSpec
from reportbuilder.render.deck import render_report
from reportbuilder.render.image._mpl import series_is_empty, render_empty_chart
from reportbuilder.export.pptx_build import build_pptx
from reportbuilder.model.report import Report
from reportbuilder.testing.fixtures import one_chart_report, _chart, tiny_model_and_data


def _pics(prs):
    return sum(1 for s in prs.slides for sh in s.shapes
               if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)


# ---------------------------------------------------------------------------
# series_is_empty detection
# ---------------------------------------------------------------------------

def test_series_is_empty_true_for_no_categories():
    assert series_is_empty(empty_series()) is True


def test_series_is_empty_true_for_all_zero():
    assert series_is_empty(all_zero_series()) is True


def test_series_is_empty_false_for_real_data():
    from suite.integration.render._series import partition_series
    assert series_is_empty(partition_series()) is False


# ---------------------------------------------------------------------------
# render_empty_chart places a placeholder picture
# ---------------------------------------------------------------------------

def test_render_empty_chart_places_one_picture():
    prs, slide, slot, ctx = make_ctx("vertical_bar", empty_series())
    ret = render_empty_chart(ctx)
    assert ret is None
    assert_single_picture(slide, slot)


# ---------------------------------------------------------------------------
# Deck degrades an empty series to a placeholder (image mode)
# ---------------------------------------------------------------------------

def test_deck_empty_series_yields_placeholder_picture():
    report = dataclasses.replace(one_chart_report(), render_mode="image")
    prs = render_report(report, {"q1": empty_series()}, StyleSpec())
    # Exactly one picture (the placeholder), no crash, no native chart.
    assert _pics(prs) == 1
    assert not any(getattr(sh, "has_chart", False)
                   for s in prs.slides for sh in s.shapes)


# ---------------------------------------------------------------------------
# build_pptx: a chart whose compute fails/yields empty still produces a deck
# ---------------------------------------------------------------------------

def test_build_pptx_wordcloud_no_words_falls_back(tmp_path):
    """A wordcloud over a non-text question yields no words -> empty-series fallback,
    so build_pptx still writes a deck (one placeholder picture) without raising."""
    model, df = tiny_model_and_data()
    report = Report(name="W", render_mode="image", template_ref="t.pptx",
                    charts=(_chart(chart_type="wordcloud"),))
    out = tmp_path / "wc.pptx"

    path = build_pptx(report, model, df, str(out))
    assert path == str(out)
    assert out.exists()

    prs = Presentation(str(out))
    assert _pics(prs) == 1


def test_build_pptx_normal_chart_still_works(tmp_path):
    """Sanity: a computable chart also produces exactly one picture in image mode."""
    model, df = tiny_model_and_data()
    report = Report(name="N", render_mode="image", template_ref="t.pptx",
                    charts=(_chart(chart_type="vertical_bar"),))
    out = tmp_path / "n.pptx"
    build_pptx(report, model, df, str(out))
    prs = Presentation(str(out))
    assert _pics(prs) == 1
