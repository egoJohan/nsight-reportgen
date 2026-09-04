"""A bullet that wraps: every line of it starts at the same x.

Reported: "when creating a special slide the multi row bullets are not indented
correctly. The second+ line(s) are not exactly on the same indentation level as
the first" — and it has to hold whatever template the slide is drawn on.

The slide itself is right: the paragraph carries a hanging indent (marL where
the text sits, a negative indent pulling the glyph back) and a left tab stop at
marL that snaps the first line's text to the same x as its wrapped lines, which
is what LibreOffice draws. The PREVIEW is what the author looks at, and the
compositor that draws it replaced the tab with two spaces — so the first line's
text began wherever the glyph happened to end, about ten pixels short of its own
continuation.

Measured in ink, on the rendered pixels, because that is the thing being
reported: not "is the XML right" (it was) but "do the lines line up".
"""
from __future__ import annotations

import numpy as np
import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from reportbuilder.render.base import StyleSpec
from reportbuilder.render.image import fast_preview
from reportbuilder.render.image.special_slide import _bullet_box

pytestmark = pytest.mark.integration

# Same first letter on every line, so where the ink starts IS where the pen
# started — different letters carry different left side bearings and would make
# a correct layout look a pixel or two out.
WORD = "Hnnn"
LONG = " ".join([WORD] * 40)


def _line_starts(img) -> list[list[int]]:
    """The x of each ink run, per row of text, top to bottom."""
    a = np.asarray(img.convert("L"))
    dark = a < 170
    rows = dark.any(axis=1)
    out, start = [], None
    for y, on in enumerate(rows):
        if on and start is None:
            start = y
        elif not on and start is not None:
            if y - start > 4:
                band = dark[start:y]
                xs = np.where(band.any(axis=0))[0]
                starts, prev = [], None
                for x in xs:
                    if prev is None or x - prev > 2:
                        starts.append(int(x))
                    prev = x
                out.append(starts)
            start = None
    return out


def _composited(bullets):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style = StyleSpec()
    _bullet_box(slide, int(prs.slide_width), int(prs.slide_height), bullets)
    img = fast_preview.compose_from_slide(style, slide)
    if img is None:
        pytest.skip("no ground image available for the compositor")
    return img


def test_a_wrapped_bullet_lines_up_with_its_own_first_line():
    lines = _line_starts(_composited([(0, LONG)]))
    assert len(lines) >= 2, "the bullet did not wrap; nothing to compare"
    first_line_text = lines[0][1]      # [0] is the glyph, [1] is where text starts
    for cont in lines[1:]:
        assert abs(cont[0] - first_line_text) <= 1, (
            f"wrapped line starts at {cont[0]}, first line's text at "
            f"{first_line_text}")


def test_a_nested_bullet_lines_up_too():
    lines = _line_starts(_composited([(1, LONG)]))
    assert len(lines) >= 2
    first_line_text = lines[0][1]
    for cont in lines[1:]:
        assert abs(cont[0] - first_line_text) <= 1, (
            f"wrapped line starts at {cont[0]}, first line's text at "
            f"{first_line_text}")
