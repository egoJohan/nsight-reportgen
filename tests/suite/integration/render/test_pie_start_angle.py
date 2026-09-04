"""Where a pie starts drawing, so the legend reads the way the circle does.

Reported: "kuvaajassa 'Kyllä' on oikealla puolella mutta legendissä se on
vasemmalla puolella" — the first category was drawn from 12 o'clock clockwise,
which puts it on the RIGHT of the circle, while the legend lists it first, on
the LEFT.

Starting at 6 o'clock instead, still clockwise, sends the first category up the
LEFT side: first in the legend, first on the left of the picture. The categories
then run left to right across the top, the way the legend does.
"""
from __future__ import annotations

import pytest

from reportbuilder.model.report import ChartSpec, ElementToggles, NumberFormat, SortSpec
from reportbuilder.render.base import RenderContext, Slot, StyleSpec
from reportbuilder.stats.series import Cell, SeriesResult

pytestmark = pytest.mark.integration


def _wedges(pcts, donut=False):
    from pptx import Presentation
    from pptx.util import Inches
    from matplotlib.patches import Wedge

    from reportbuilder.render.image.pie import _build_pie_figure

    cats = tuple(f"c{i}" for i in range(len(pcts)))
    cells = {(c, "Total"): Cell(pct=p, count=p * 10, mean=None)
             for c, p in zip(cats, pcts)}
    series = SeriesResult(categories=cats, segments=("Total",), cells=cells,
                          base_n={"Total": 1000}, statistic="pct")
    spec = ChartSpec(question_ref="q", chart_type="pie", statistic="pct",
                     classifying_var=None, number_format=NumberFormat(),
                     sort=SortSpec(basis="data_order"), template_slot="s",
                     elements=ElementToggles(title=False))
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slot = Slot(slide_index=0, left=Inches(0.7), top=Inches(1.4),
                width=Inches(12.0), height=Inches(4.6), name="s")
    fig = _build_pie_figure(
        RenderContext(slide=slide, slot=slot, style=StyleSpec(), spec=spec,
                      series=series, fmt=spec.number_format), donut=donut)
    return [w for ax in fig.axes for w in ax.patches if isinstance(w, Wedge)]


def test_the_first_category_starts_at_six_o_clock():
    """270° in matplotlib's reckoning — anticlockwise from 3 o'clock."""
    first = _wedges([25.0, 75.0])[0]
    assert first.theta2 == pytest.approx(270.0)


def test_it_runs_clockwise_from_there_so_the_first_slice_is_on_the_LEFT():
    first = _wedges([25.0, 75.0])[0]
    # Sweeping clockwise from the bottom means falling angles: a quarter of the
    # circle lands the wedge between 9 o'clock (180°) and 6 o'clock (270°) —
    # the lower LEFT, which is the side its legend entry sits on.
    assert 180.0 <= first.theta1 < 270.0
    assert first.theta1 == pytest.approx(180.0)


def test_the_categories_follow_each_other_clockwise():
    a, b, c = _wedges([25.0, 25.0, 50.0])
    assert a.theta1 == pytest.approx(b.theta2)
    assert b.theta1 == pytest.approx(c.theta2)


def test_a_doughnut_starts_in_the_same_place():
    assert _wedges([25.0, 75.0], donut=True)[0].theta2 == pytest.approx(270.0)
