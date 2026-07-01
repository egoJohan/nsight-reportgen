"""Deck assembly: render_report / render_to_file in BOTH image and native modes.

Asserts the completeness + native-purity guards, the correct kind/number of
chart objects, the titles= override, multi-chart reports, and that a saved deck
reopens with the same tallies.  Deterministic — no soffice.
"""
from __future__ import annotations

import dataclasses

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from reportbuilder.render.base import StyleSpec
from reportbuilder.render.deck import (
    render_report, render_to_file,
    assert_complete, assert_no_pictures_in_chart_slots,
    CompletenessError, NativePurityError,
)
from reportbuilder.testing.fixtures import (
    one_chart_report, two_chart_report, known_series,
)


def _series_by_ref(report):
    return {spec.question_ref: known_series() for spec in report.charts}


def _counts(prs):
    charts = sum(1 for s in prs.slides for sh in s.shapes if getattr(sh, "has_chart", False))
    pics = sum(1 for s in prs.slides for sh in s.shapes
               if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
    return charts, pics


def _chart_titles(prs):
    out = []
    for s in prs.slides:
        for sh in s.shapes:
            if getattr(sh, "has_chart", False) and sh.chart.has_title:
                out.append(sh.chart.chart_title.text_frame.text)
    return out


# ---------------------------------------------------------------------------
# Native mode
# ---------------------------------------------------------------------------

def test_native_render_one_chart():
    report = one_chart_report()           # render_mode="native"
    assert report.render_mode == "native"
    prs = render_report(report, _series_by_ref(report), StyleSpec(), titles={"q1": "Sat"})

    charts, pics = _counts(prs)
    assert charts == 1
    assert pics == 0
    # Guards embedded in render_report already ran; re-assert explicitly.
    assert_complete(prs, report)
    assert_no_pictures_in_chart_slots(prs, report)
    # titles= override flowed onto the native chart title.
    assert "Sat" in _chart_titles(prs)


def test_native_purity_would_fail_for_image_shapes():
    """A native report whose deck actually contains a picture trips NativePurityError."""
    report = one_chart_report()
    # Build an IMAGE deck (has a picture), then assert purity as-if native.
    img_report = dataclasses.replace(report, render_mode="image")
    prs = render_report(img_report, _series_by_ref(img_report), StyleSpec())
    charts, pics = _counts(prs)
    assert pics == 1 and charts == 0
    with pytest.raises(NativePurityError):
        assert_no_pictures_in_chart_slots(prs, report)  # report.render_mode == native


# ---------------------------------------------------------------------------
# Image mode
# ---------------------------------------------------------------------------

def test_image_render_one_chart():
    report = dataclasses.replace(one_chart_report(), render_mode="image")
    prs = render_report(report, _series_by_ref(report), StyleSpec())

    charts, pics = _counts(prs)
    assert pics == 1
    assert charts == 0
    assert_complete(prs, report)
    # Purity guard is a no-op for image mode (does not raise).
    assert_no_pictures_in_chart_slots(prs, report)


# ---------------------------------------------------------------------------
# Multi-chart
# ---------------------------------------------------------------------------

@pytest.mark.parametrize("mode,expect_charts,expect_pics", [
    ("native", 2, 0),
    ("image", 0, 2),
])
def test_multi_chart_report(mode, expect_charts, expect_pics):
    report = dataclasses.replace(two_chart_report(), render_mode=mode)
    prs = render_report(report, _series_by_ref(report), StyleSpec())
    charts, pics = _counts(prs)
    assert charts == expect_charts
    assert pics == expect_pics
    assert_complete(prs, report)


# ---------------------------------------------------------------------------
# render_to_file: saved deck reopens with the same tallies
# ---------------------------------------------------------------------------

@pytest.mark.parametrize("mode,expect_charts,expect_pics", [
    ("native", 1, 0),
    ("image", 0, 1),
])
def test_render_to_file_reopens(tmp_path, mode, expect_charts, expect_pics):
    report = dataclasses.replace(one_chart_report(), render_mode=mode)
    out = tmp_path / f"deck_{mode}.pptx"

    returned = render_to_file(report, _series_by_ref(report), StyleSpec(), str(out),
                              titles={"q1": "Reopened"})
    assert returned == str(out)
    assert out.exists()

    reopened = Presentation(str(out))
    charts, pics = _counts(reopened)
    assert charts == expect_charts
    assert pics == expect_pics
    if mode == "native":
        assert "Reopened" in _chart_titles(reopened)
