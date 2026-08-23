"""Rendering a slide into a customer's template.

Johan's rule: the TITLE follows the template — its font, size, colour and box —
and the furniture that repeats on every slide comes with it. nSight positions
everything else. There are two ways a template states its design, and a deck has
to come out right from both:

  * Attendo and Holiday Club keep it in a LAYOUT their own slides are built on,
    so a slide built from that layout inherits it.
  * Synsam and the agent deck draw it on the SLIDES (98 of 147 on "Tom"; all 20
    on `Blank`), so it is harvested and redrawn.

Real client files, like test_template_check and test_template_profile: the
difficulty is entirely in what real templates do.
"""
from __future__ import annotations

import dataclasses
import logging
import pathlib
from unittest import mock

import pytest
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from reportbuilder.render.deck import render_report
from reportbuilder.render.image.slide_chrome import content_floor
from reportbuilder.render.style_spec import load_style_spec
from reportbuilder.model.report import (
    ChartSpec, ElementToggles, NumberFormat, Report, SortSpec,
)
from reportbuilder.testing.fixtures import one_chart_report, known_series

_TEMPLATES = {
    "attendo": "input/Attendo Bränditutkimus Marraskuu 2025.pptx",
    "synsam": "input/Synsam_Segmentointitutkimus_30.4.2025_nSight.pptx",
    "holidayclub": "input/Holiday Club_Loyalty tutkimus_raportti_19.2.2026.pptx",
    "agent_deck": "work/attendo_agent_deck.pptx",   # not in git; skipped if absent
}

_LONG = ("Attendo liitetään vahvimmin attribuutteihin luotettava ja "
         "ammattitaitoinen")


#: The question, which becomes the subtitle when a distinct headline is set.
_QUESTION = "Mikä seuraavista vastaa työtilannettasi tällä hetkellä?"


def _render(name, title=_LONG, headline=""):
    """Render one chart into *name*'s template. With *headline* set, the title is
    the headline and the QUESTION becomes the subtitle line under it."""
    path = pathlib.Path(_TEMPLATES[name])
    if not path.exists():
        pytest.skip(f"{path} not available locally")
    style = load_style_spec(str(path))
    report = dataclasses.replace(one_chart_report(), render_mode="image")
    if headline:
        charts = tuple(dataclasses.replace(c, slide_title=headline)
                       for c in report.charts)
        report = dataclasses.replace(report, charts=charts)
        title = _QUESTION
    prs = render_report(report, {"q1": known_series()}, style, titles={"q1": title})
    return style, prs, prs.slides[0]


def _picture(slide):
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    return pics[0]


def _subtitle(slide):
    return next(s for s in slide.shapes
                if s.has_text_frame and s.text_frame.text.strip() == _QUESTION)


def _title_shape(slide, text):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() == text:
            return shape
    return None


class TestWhichWayTheSlideIsBuilt:
    @pytest.mark.parametrize("name", ["attendo", "holidayclub"])
    def test_a_layout_template_builds_from_its_own_layout(self, name):
        style, _prs, slide = _render(name)
        assert style.chart_layout_index is not None
        assert slide.slide_layout.name == style.profile.source.split(":", 1)[1]

    @pytest.mark.parametrize("name", ["synsam", "agent_deck"])
    def test_a_hand_drawn_template_is_not_built_from_a_layout(self, name):
        """Its layouts are stock — Office's `Blank`, `Title and Content` — and
        building on them is what made a client deck come out looking like plain
        PowerPoint. The design is harvested off a slide instead."""
        style, _prs, _slide = _render(name)
        assert style.chart_layout_index is None
        assert style.profile.source.startswith("slide:")

    def test_the_clients_palette_survives_either_way(self):
        """Chart colours follow the template whichever way its design is
        stated. They used to be keyed on there being a layout to build from."""
        style, _prs, _slide = _render("synsam")
        assert style.from_template
        assert style.color_for(0) == "FF5000"


class TestTheTitleFollowsTheTemplate:
    def test_a_harvested_title_sits_in_the_templates_own_box(self):
        style, _prs, slide = _render("synsam")
        title = _title_shape(slide, _LONG)
        assert title is not None
        assert title.left == style.profile.title.left
        assert title.top == style.profile.title.top
        run = title.text_frame.paragraphs[0].runs[0]
        assert run.font.name == "Avenir Next LT Pro Demi"
        assert run.font.size.pt == 22.0

    def test_a_layout_title_goes_in_the_placeholder_and_inherits(self):
        """Font, colour and position stay the customer's — the whole point of
        using their layout. The SIZE is the one thing stated, and only because
        it was fitted to their own box.

        It used to be inherited too, which sounds more faithful and is not: a
        placeholder written by python-pptx carries nothing to inherit FROM, so
        each renderer picked its own answer — LibreOffice one size, PowerPoint
        another, the preview compositor a third. Stating the fitted size is what
        makes the deck and the preview show the same slide, and what stops a
        long headline running over the subtitle at the template's own size.
        """
        _style, _prs, slide = _render("attendo")
        assert slide.shapes.title is not None
        assert slide.shapes.title.text_frame.text == _LONG
        run = slide.shapes.title.text_frame.paragraphs[0].runs[0]
        # Fitted, not invented: never larger than the template's own title size.
        assert run.font.size is not None
        assert run.font.size.pt <= 30.0
        # Still the customer's in every other respect.
        assert run.font.name is None
        assert run.font.color.type is None


class TestWhereEverythingElseGoes:
    @pytest.mark.parametrize("name", sorted(_TEMPLATES))
    def test_the_chart_starts_below_the_title(self, name):
        """Holiday Club's title box holds one line of the headline THEY wrote;
        ours ran to two and printed straight through the subtitle below it."""
        style, _prs, slide = _render(name)
        st = style.profile.title
        assert _picture(slide).top >= st.top + st.height

    def test_nothing_of_ours_lands_on_the_templates_own_footer(self):
        """Synsam's master puts its logo at 6.73in on a 7.5in slide, exactly
        where the "N = ..." line goes."""
        _style, prs, slide = _render("synsam")
        floor = content_floor(slide, int(prs.slide_width), int(prs.slide_height))
        assert floor < int(prs.slide_height)      # the logo was found
        pic = _picture(slide)
        assert pic.top + pic.height <= floor
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.startswith("N ="):
                assert shape.top < floor

    def test_the_chart_keeps_the_templates_side_margin(self):
        style, _prs, slide = _render("synsam")
        assert _picture(slide).left >= style.profile.title.left


class TestFurniture:
    def test_repeating_graphics_are_redrawn_on_a_harvested_slide(self):
        """The rule under Synsam's title is on every one of its chart slides."""
        _style, _prs, slide = _render("synsam")
        assert any(s.shape_type == MSO_SHAPE_TYPE.LINE for s in slide.shapes)

    def test_a_layout_slide_clones_nothing(self):
        """Its furniture is inherited; cloning would draw the logo twice."""
        style, _prs, _slide = _render("attendo")
        assert style.profile.furniture == []

    @pytest.mark.parametrize("name", ["synsam", "attendo"])
    def test_the_house_background_is_not_painted_over_the_clients(self, name):
        """A full-slide cream rectangle would cover whatever the customer's
        design puts behind the chart. (The agent deck ends up with one because
        ITS design is a full-slide rectangle — cloned from its own slides.)"""
        _style, prs, slide = _render(name)
        full = [s for s in slide.shapes
                if int(s.width or 0) >= int(prs.slide_width) * 0.99
                and int(s.height or 0) >= int(prs.slide_height) * 0.99]
        assert full == []


def test_a_deck_with_no_template_is_unchanged():
    """The house style is what a wizard report without a template still gets:
    cream ground, teal bar, our own title box."""
    from reportbuilder.render.base import StyleSpec

    report = dataclasses.replace(one_chart_report(), render_mode="image")
    prs = render_report(report, {"q1": known_series()}, StyleSpec(), titles={"q1": "T"})
    slide = prs.slides[0]
    bg = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(bg) == 2                                   # background + accent bar
    assert bg[0].width == int(prs.slide_width)
    assert _title_shape(slide, "T") is not None
    assert _title_shape(slide, "T").left == Inches(0.80)


# ---------------------------------------------------------------------------
# Bullet slides. A deck must not come out half in the customer's design and
# half in ours: an overview or conclusion slide in nSight cream between chart
# slides in the client's brand is exactly what a "use our template" request is
# asking us not to do.
# ---------------------------------------------------------------------------

_HEADING = "Yhteenveto tutkimuksen päätuloksista"
_BULLETS = ["Attendo tunnetaan laajasti", "  Luotettavuus korostuu"]


def _bullet_slide(name):
    path = pathlib.Path(_TEMPLATES[name]) if name else None
    if path is not None and not path.exists():
        pytest.skip(f"{path} not available locally")
    if path is None:
        from reportbuilder.render.base import StyleSpec
        style = StyleSpec()
    else:
        style = load_style_spec(str(path))
    spec = ChartSpec(
        question_ref="s1", chart_type="special_conclusion", statistic="pct",
        classifying_var=None, number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"), template_slot="s1",
        elements=ElementToggles(), slide_title=_HEADING,
        options={"bullets": _BULLETS})
    report = Report(name="s", render_mode="image", template_ref="", charts=(spec,))
    prs = render_report(report, {}, style)
    return style, prs, prs.slides[0]


class TestBulletSlidesFollowTheTemplateToo:
    @pytest.mark.parametrize("name", ["attendo", "synsam"])
    def test_no_house_ground_is_painted_over_the_clients(self, name):
        _style, prs, slide = _bullet_slide(name)
        full = [s for s in slide.shapes
                if int(s.width or 0) >= int(prs.slide_width) * 0.99
                and int(s.height or 0) >= int(prs.slide_height) * 0.99]
        assert full == []

    def test_the_heading_goes_in_the_layouts_title_placeholder(self):
        _style, _prs, slide = _bullet_slide("attendo")
        assert slide.shapes.title.text_frame.text == _HEADING

    def test_a_harvested_heading_sits_in_the_templates_box(self):
        style, _prs, slide = _bullet_slide("synsam")
        shape = _title_shape(slide, _HEADING)
        assert shape is not None
        assert shape.top == style.profile.title.top
        assert shape.text_frame.paragraphs[0].runs[0].font.size.pt == 22.0

    def test_the_bullet_glyph_is_the_clients_accent(self):
        """House teal dots down an Attendo navy deck was the giveaway."""
        _style, _prs, slide = _bullet_slide("attendo")
        box = next(s for s in slide.shapes
                   if s.has_text_frame and "Attendo tunnetaan" in s.text_frame.text)
        glyph = box.text_frame.paragraphs[0].runs[0]
        assert glyph.text.startswith("•")
        assert str(glyph.font.color.rgb) == "122D49"

    def test_the_house_style_is_unchanged_without_a_template(self):
        _style, prs, slide = _bullet_slide(None)
        grounds = [s for s in slide.shapes
                   if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
        assert len(grounds) == 2                     # cream ground + teal bar
        assert grounds[0].width == int(prs.slide_width)


class TestTheDeckIsInTheTemplatesColours:
    """The agent deck is the case that made this necessary: cream and teal on
    every slide, stock Office in its theme. Reading the theme produced a deck in
    colours that appear nowhere in the customer's own file."""

    def test_a_slide_designed_template_paints_on_its_own_ground(self):
        style, _prs, _slide = _render("agent_deck")
        assert style.background == "F7F3EC"
        assert style.accent == "13615E"
        assert style.brand_palette == []          # stock Office states no brand

    def test_a_themed_template_keeps_its_theme(self):
        style, _prs, _slide = _render("attendo")
        assert style.brand_palette[:1] == ["122D49"]
        assert style.accent == "122D49"

    def test_the_chart_is_drawn_in_the_templates_colour(self):
        """Not a colour test on pixels — the ramp the chart builder is handed."""
        from reportbuilder.render.house_style import series_colors
        from reportbuilder.render.image._mpl import chart_accent, template_palette

        class _Ctx:
            pass

        style, _prs, _slide = _render("agent_deck")
        ctx = _Ctx()
        ctx.style = style
        assert template_palette(ctx) is None      # its theme is not its brand
        assert series_colors(1, palette=template_palette(ctx),
                             accent=chart_accent(ctx)) == ["#13615E"]


class TestTheSubtitleBelongsToTheTitle:
    def test_it_sits_under_the_title_not_over_the_chart(self):
        """It used to hang off the chart, floating in a band of empty cream."""
        style, _prs, slide = _render("agent_deck", headline=_LONG)
        sub = _subtitle(slide)
        st = style.profile.title
        assert sub.top >= st.top + st.height
        assert sub.top + sub.height <= _picture(slide).top + 1

    def test_it_clears_the_graphic_the_template_rules_under_the_title(self):
        """Synsam's rule sits below its title box; placing the subtitle by the
        title's height alone landed the text across it."""
        _style, _prs, slide = _render("synsam", headline=_LONG)
        rule = next(s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE)
        assert _subtitle(slide).top >= rule.top + rule.height


class TestWhenHarvestingFails:
    """The fallback must not be the thing we set out to prevent.

    Harvesting is a heuristic over somebody else's file, so it can throw. It
    used to fall back to building on the template's layouts — which for a
    stock-Office file is a 44pt centred title placeholder on a white slide, the
    exact plain-PowerPoint deck this module exists to avoid, with nothing in the
    log to say it had happened.
    """

    def _without_harvesting(self, name):
        import reportbuilder.render.template_profile as tp

        path = pathlib.Path(_TEMPLATES[name])
        if not path.exists():
            pytest.skip(f"{path} not available locally")
        with mock.patch.object(tp, "extract_profile", side_effect=RuntimeError("boom")):
            return load_style_spec(str(path))

    def test_a_stock_office_template_falls_back_to_the_house_style(self, caplog):
        with caplog.at_level(logging.WARNING):
            style = self._without_harvesting("agent_deck")
        assert style.profile is None
        assert style.chart_layout_index is None      # not its Office layouts
        assert "could not harvest" in caplog.text    # and it is not silent

    def test_a_branded_template_still_uses_its_own_layout(self):
        """There the layouts ARE the brand, so they remain the better fallback."""
        style = self._without_harvesting("attendo")
        assert style.chart_layout_index is not None
