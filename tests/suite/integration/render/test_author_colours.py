"""Colours set in the template editor must reach the slide, not break it.

The editor lets an author override the title, subtitle and footer colour. Each
override was handed to python-pptx as the STRING "#112233" while every other
branch of the same function returns an RGBColor, and `font.color.rgb` refuses
anything else — so setting a colour did not tint the text, it raised
`ValueError: assigned value must be type RGBColor` and took the whole deck
render down on its first slide.
"""
from __future__ import annotations

import dataclasses

import pandas as pd
import pytest
from pptx.dml.color import RGBColor

from reportbuilder.export.pptx_build import build_presentation
from reportbuilder.model.question import Question, QuestionModel, ValueLabel, Variable
from reportbuilder.model.report import (
    ChartSpec, ElementToggles, NumberFormat, Report, SortSpec,
)
from reportbuilder.render.base import StyleSpec

pytestmark = pytest.mark.integration


def _model():
    var = Variable(name="q1", label="Suositteletko?", measurement="categorical",
                   value_labels=(ValueLabel(1.0, "Kyllä"), ValueLabel(2.0, "Ei")),
                   missing_values=frozenset())
    return QuestionModel(
        variables={"q1": var},
        questions=[Question(qid="q1", kind="single", variables=("q1",),
                            text="Suositteletko?")])


def _report():
    spec = ChartSpec(question_ref="q1", chart_type="horizontal_bar", statistic="pct",
                     classifying_var=None, number_format=NumberFormat(),
                     sort=SortSpec(basis="data_order"), template_slot="s1",
                     elements=ElementToggles(), slide_id="a",
                     slide_title="Otsikko", slide_description="Alaotsikko")
    return Report(name="r", render_mode="image", template_ref="", charts=(spec,))


def _styled(**colours):
    style = StyleSpec()
    for field, value in colours.items():
        try:
            style = dataclasses.replace(style, **{field: value})
        except (TypeError, ValueError):
            object.__setattr__(style, field, value)
    return style


def _render(**colours):
    return build_presentation(_report(), _model(),
                              pd.DataFrame({"q1": [1.0] * 7 + [2.0] * 3}),
                              style=_styled(**colours))


def test_a_footer_colour_does_not_take_the_deck_down():
    _render(footer_colour="112233")


def test_a_title_colour_does_not_take_the_deck_down():
    _render(title_colour="112233")


def test_a_subtitle_colour_does_not_take_the_deck_down():
    _render(subtitle_colour="112233")


def test_the_colours_actually_reach_the_text():
    """Subtitle and footer, on a house-style slide. The TITLE's colour belongs
    to the templated path (`title_colour_for`) — a house slide draws its
    headline in the house ink — so it is asserted there instead."""
    prs = _render(subtitle_colour="445566", footer_colour="778899")
    used = {run.font.color.rgb
            for sh in prs.slides[0].shapes if sh.has_text_frame
            for para in sh.text_frame.paragraphs for run in para.runs
            if run.font.color and run.font.color.type is not None}
    assert RGBColor(0x44, 0x55, 0x66) in used, used
    assert RGBColor(0x77, 0x88, 0x99) in used, used


def test_an_authored_title_colour_is_a_colour_object_not_a_string():
    """What `_textbox` needs: python-pptx accepts nothing but an RGBColor."""
    from reportbuilder.render.image.slide_chrome import title_colour_for

    class _Harvested:          # what the template's own title looks like
        colour = "000000"

    got = title_colour_for(_Harvested(), _styled(title_colour="112233"))
    assert isinstance(got, RGBColor), repr(got)
    assert got == RGBColor(0x11, 0x22, 0x33)
