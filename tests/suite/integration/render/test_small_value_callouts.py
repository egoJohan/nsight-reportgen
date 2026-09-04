"""A value too small to sit inside its own slice is drawn beside it instead.

A 100%-stacked scale puts the tail of the scale — the 1 %, the 2 % — into
slivers narrower than the number that belongs in them, and a pie does the same
with a thin wedge. The number was dropped, so the reader of the finished deck
had no way to check it: "nyt kun presiksessä ei pääse dataan tarkistamaan
kaikkia pienempiä prosenttiosuuksia".

It now goes outside the shape, close to it, with a line drawn back to the piece
it belongs to.
"""
from __future__ import annotations

import pytest

from reportbuilder.model.report import ChartSpec, ElementToggles, NumberFormat, SortSpec
from reportbuilder.render.base import RenderContext, Slot, StyleSpec
from reportbuilder.stats.series import Cell, SeriesResult

pytestmark = pytest.mark.integration


def _ctx(cats, pcts, chart_type, segments=("Total",), hide_below=0.0):
    from pptx import Presentation
    from pptx.util import Inches

    cells = {}
    for s in segments:
        for c, p in zip(cats, pcts):
            cells[(c, s)] = Cell(pct=p, count=p * 10, mean=None)
    series = SeriesResult(categories=tuple(cats), segments=tuple(segments), cells=cells,
                          base_n={s: 1000 for s in segments}, statistic="pct")
    spec = ChartSpec(question_ref="q", chart_type=chart_type, statistic="pct",
                     classifying_var=None if segments == ("Total",) else "maa",
                     number_format=NumberFormat(hide_below_pct=hide_below),
                     sort=SortSpec(basis="data_order"), template_slot="s",
                     elements=ElementToggles(title=False))
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slot = Slot(slide_index=0, left=Inches(0.7), top=Inches(1.4),
                width=Inches(12.0), height=Inches(4.6), name="s")
    return RenderContext(slide=slide, slot=slot, style=StyleSpec(), spec=spec,
                         series=series, fmt=spec.number_format)


# ── the pie ──────────────────────────────────────────────────────────────────

def _pie_annotations(pcts):
    from matplotlib.text import Annotation

    from reportbuilder.render.image.pie import _build_pie_figure

    cats = [f"c{i}" for i in range(len(pcts))]
    fig = _build_pie_figure(_ctx(cats, pcts, "pie"), donut=False)
    return [t for ax in fig.axes for t in ax.texts if isinstance(t, Annotation)]


def test_a_thin_wedge_gets_its_number_beside_the_circle():
    out = _pie_annotations([1.0, 2.0, 97.0])
    assert {a.get_text() for a in out} == {"1 %", "2 %"}


def test_the_number_sits_OUTSIDE_the_circle():
    for a in _pie_annotations([1.0, 99.0]):
        x, y = a.get_position()
        assert (x * x + y * y) ** 0.5 > 1.0, (x, y)


def test_a_line_reaches_back_to_the_wedge():
    for a in _pie_annotations([1.0, 99.0]):
        assert a.arrowprops, "no leader line drawn"
        x, y = a.xy                      # the end that touches the wedge
        assert (x * x + y * y) ** 0.5 <= 1.0, (x, y)


def test_wedges_big_enough_to_hold_their_number_keep_it_inside():
    assert _pie_annotations([30.0, 70.0]) == []


def test_nothing_is_drawn_for_a_category_nobody_chose():
    assert {a.get_text() for a in _pie_annotations([0.0, 1.0, 99.0])} == {"1 %"}


def test_a_called_out_number_never_lands_on_the_shared_legend():
    """The numbers hang BELOW the circles, which is where the legend sits on a
    row of panels. The clearance is computed from the panels' own extent, so
    whatever is drawn around them is what the legend is kept off."""
    from reportbuilder.render.image.pie import _build_pie_figure

    cats = [f"c{i}" for i in range(7)]
    pcts = [63.0, 13.0, 7.0, 5.0, 5.0, 6.0, 1.0]
    fig = _build_pie_figure(_ctx(cats, pcts, "pie",
                                 segments=("Suomi", "Ruotsi", "Saksa")), donut=False)
    from matplotlib.text import Annotation

    fig.canvas.draw()
    renderer = fig.canvas.get_renderer()
    leg = next(a for a in fig.legends).get_window_extent(renderer)
    for ax in fig.axes:
        # The panel AND what is drawn around it. A callout sits outside the
        # axes, which is why get_tightbbox does not report it — measuring the
        # panel alone is how the legend came to print over the small numbers.
        boxes = [ax.get_tightbbox(renderer)] + [
            t.get_window_extent(renderer) for t in ax.texts if isinstance(t, Annotation)]
        for box in boxes:
            overlap = min(leg.y1, box.y1) - max(leg.y0, box.y0)
            assert overlap <= 1.0, (
                f"the legend sits {overlap:.0f}px into a panel: {leg.bounds} vs {box.bounds}")


def test_the_circle_is_not_shrunk_to_make_room_for_the_numbers():
    """The callouts go in the space around the panel, not inside it: a reader
    looks at the circle, and paying for a handful of small numbers by making
    every circle smaller is the wrong trade."""
    from matplotlib.patches import Wedge

    from reportbuilder.render.image.pie import _build_pie_figure

    def radius(pcts):
        fig = _build_pie_figure(_ctx([f"c{i}" for i in range(len(pcts))], pcts, "pie"),
                                donut=False)
        ax = next(a for a in fig.axes if any(isinstance(p, Wedge) for p in a.patches))
        return next(p for p in ax.patches if isinstance(p, Wedge)).r / (
            ax.get_xlim()[1] - ax.get_xlim()[0])

    assert radius([1.0, 99.0]) == pytest.approx(radius([30.0, 70.0]), rel=0.01)


# ── the stacked bar ──────────────────────────────────────────────────────────

def _stacked_annotations(pcts, chart_type="stacked_horizontal_bar"):
    from matplotlib.text import Annotation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    import matplotlib
    seen: list = []
    orig = matplotlib.axes.Axes.annotate

    def spy(self, text, *a, **kw):
        art = orig(self, text, *a, **kw)
        seen.append((text, kw))
        return art

    matplotlib.axes.Axes.annotate = spy
    try:
        from reportbuilder.render.image import IMAGE_BUILDERS

        cats = [f"c{i}" for i in range(len(pcts))]
        IMAGE_BUILDERS[chart_type](_ctx(cats, pcts, chart_type, segments=("Suomi",)))
    finally:
        matplotlib.axes.Axes.annotate = orig
    return seen


def test_a_sliver_of_a_stacked_bar_gets_its_number_beside_the_bar():
    """Which segments count as slivers is the author's cut-off (`hide_below_pct`,
    1 % of the axis by default) — the same one that used to decide which numbers
    were dropped. Nothing is dropped now: below it, the number goes outside."""
    # 3 % holds "3 %" at this width; 2 % does not. Which is which is the
    # measurement's business — what matters is that the one that does not fit
    # is drawn rather than dropped.
    out = _stacked_annotations([2.0, 3.0, 30.0, 65.0])
    assert {t for t, _kw in out} == {"2 %"}


def test_the_stacked_callout_is_drawn_with_a_line_to_its_segment():
    for _text, kw in _stacked_annotations([2.0, 98.0]):
        assert kw.get("arrowprops"), "no leader line drawn"
        # the line ends ON the bar, the text sits off it
        assert kw["xy"][1] != kw["xytext"][1]


def test_segments_wide_enough_keep_their_number_inside():
    assert _stacked_annotations([40.0, 60.0]) == []


def test_a_stacked_column_calls_its_slivers_out_too():
    out = _stacked_annotations([2.0, 98.0], chart_type="stacked_vertical_bar")
    assert {t for t, _kw in out} == {"2 %"}


# ── where the called-out numbers actually land ───────────────────────────────

def _drawn(chart_type, pcts, segments=("Suomi",), summary=None):
    """Every callout as (text, xy on the shape, xytext where the number goes)."""
    import matplotlib

    seen: list = []
    orig = matplotlib.axes.Axes.annotate

    def spy(self, text, *a, **kw):
        # The bars are read HERE, while they exist: the builder closes its
        # figure once the PNG is made, and an axes inspected afterwards has no
        # artists left — which is how the first version of this check passed
        # against a placement that was plainly wrong.
        from matplotlib.patches import Rectangle

        cols = [(r.get_x(), r.get_x() + r.get_width())
                for c in self.containers for r in c
                if isinstance(r, Rectangle) and r.get_width()]
        seen.append((text, kw.get("xy"), kw.get("xytext"), self, cols))
        return orig(self, text, *a, **kw)

    matplotlib.axes.Axes.annotate = spy
    try:
        from reportbuilder.render.image import IMAGE_BUILDERS

        ctx = _ctx([f"c{i}" for i in range(len(pcts))], pcts, chart_type,
                   segments=segments)
        IMAGE_BUILDERS[chart_type](ctx)
    finally:
        matplotlib.axes.Axes.annotate = orig
    return seen


def test_a_column_callout_stays_clear_of_the_next_column():
    """It was placed a full column-width to the right of its own left edge, so
    the number sat on a neighbour it does not describe. Asserted against the
    drawn rectangles, not a guessed distance."""
    from matplotlib.patches import Rectangle

    out = _drawn("stacked_vertical_bar", [2.0, 2.0, 96.0], segments=("a", "b", "c"))
    assert out, "nothing was called out"
    for text, xy, xytext, _ax, cols in out:
        assert cols, "no columns seen — the check would pass vacuously"
        mine = [c for c in cols if c[0] - 1e-9 <= xy[0] <= c[1] + 1e-9]
        others = [c for c in cols if c not in mine]
        for left, right in others:
            # ha="left", so the number runs RIGHT from its anchor: landing on
            # the neighbour's left edge is landing on the neighbour.
            assert not (left - 1e-9 <= xytext[0] < right), (
                f"{text} at x={xytext[0]:.2f} sits on the column {left:.2f}..{right:.2f}")


def test_two_callouts_on_one_column_do_not_share_a_line():
    """Two slivers a couple of percent apart put their numbers a couple of
    percent apart — which on a 0-100 axis is a few pixels, so they overlapped.
    Measured in pixels, where the question actually lives."""
    out = _drawn("stacked_vertical_bar", [1.5, 1.5, 97.0], segments=("a",))
    assert len(out) >= 2
    ax = out[0][3]
    ax.figure.canvas.draw()
    ys = sorted(ax.transData.transform((0.0, t[2][1]))[1] for t in out)
    # a line of 8.5pt text is ~12px at the figure's dpi
    assert all(b - a >= 11.0 for a, b in zip(ys, ys[1:])), (
        f"callouts {ys} are closer than a line of text")


def test_callouts_stay_inside_the_plot():
    """Eight slivers in a row used to march the numbers off the right edge."""
    out = _drawn("stacked_horizontal_bar", [2.0] * 8 + [84.0])
    assert out
    for _text, _xy, xytext, _ax, _cols in out:
        assert xytext[0] <= 100.0, xytext
