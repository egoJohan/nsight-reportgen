"""Where a panel's base is written, and what it must not sit on.

Reported: "kolmen piirakan käytössä piirakkakohtaiset N-luvut ja Legend menee
päällekkäin" — with three pies, each panel's `n = …` printed inside the shared
legend at the foot of the chart.

The legend's height IS measured and the panels are lifted to clear it, but the
measurement is of the axes; a base written UNDER a circle hangs below its axes,
outside what was measured, so the lift left it where the legend would land. It
now goes on the title's second line — above the circle, where a legend at the
foot cannot reach it, and where the funnel has always put it.
"""
from __future__ import annotations

import pytest

from reportbuilder.model.report import ChartSpec, ElementToggles, NumberFormat, SortSpec
from reportbuilder.render.base import RenderContext, Slot, StyleSpec
from reportbuilder.stats.series import Cell, SeriesResult

pytestmark = pytest.mark.integration

CATS = ("Kyllä", "En")
LONG = ("Olen palkansaajana kokopäivätöissä", "Olen palkansaajana osa-aikatöissä",
        "Toimin yrittäjänä / freelancerina", "Olen maanviljelijä",
        "Olen koulussa / opiskelen ja teen silloin tällöin töitä",
        "Olen eläkkeellä mutta teen silloin tällöin töitä",
        "Olen työtön, ollut työttömänä alle vuoden")


def _ctx(segments, cats=CATS, chart_type="pie"):
    from pptx import Presentation
    from pptx.util import Inches

    cells = {(c, s): Cell(pct=100.0 / len(cats), count=100.0, mean=None)
             for s in segments for c in cats}
    series = SeriesResult(categories=tuple(cats), segments=tuple(segments), cells=cells,
                          base_n={s: 1049 + i for i, s in enumerate(segments)},
                          statistic="pct")
    spec = ChartSpec(question_ref="q", chart_type=chart_type, statistic="pct",
                     classifying_var="maa", number_format=NumberFormat(),
                     sort=SortSpec(basis="data_order"), template_slot="s",
                     elements=ElementToggles(title=False))
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slot = Slot(slide_index=0, left=Inches(0.7), top=Inches(1.4),
                width=Inches(12.0), height=Inches(4.6), name="s")
    return RenderContext(slide=slide, slot=slot, style=StyleSpec(), spec=spec,
                         series=series, fmt=spec.number_format)


def _base_of(ax) -> str:
    """A panel's base, off the quiet line under its name. It is a text of its
    own rather than a second line of the title because a title carries ONE
    style and the base is neither bold nor the same size."""
    return next((t.get_text() for t in ax.texts
                 if t.get_text().startswith("n = ")), "")


def _pie_fig(segments, cats=CATS):
    from reportbuilder.render.image.pie import _build_pie_figure

    return _build_pie_figure(_ctx(segments, cats), donut=False)


# ── every panel still says what it is drawn on ───────────────────────────────

def test_each_panel_states_its_own_base():
    axes = [ax for ax in _pie_fig(["Suomi", "Ruotsi", "Saksa"]).axes if ax.get_title()]
    assert [ax.get_title() for ax in axes] == ["Suomi", "Ruotsi", "Saksa"]
    assert [_base_of(ax) for ax in axes] == ["n = 1049", "n = 1050", "n = 1051"]


def test_the_base_is_quieter_than_the_name_it_sits_under():
    """It says what the circle is drawn on, not what it is."""
    ax = next(ax for ax in _pie_fig(["Suomi"]).axes if ax.get_title())
    base = next(t for t in ax.texts if t.get_text().startswith("n = "))
    assert base.get_fontweight() != "bold"
    assert base.get_fontsize() < ax.title.get_fontsize()


def test_the_base_is_no_longer_written_under_the_circle():
    """Under the axes it is outside what the legend's clearance measures."""
    fig = _pie_fig(["Suomi", "Ruotsi", "Saksa"])
    assert [ax.get_xlabel() for ax in fig.axes if ax.get_title()] == ["", "", ""]


def test_a_single_panel_states_its_base_too():
    titled = [ax for ax in _pie_fig(["Suomi"]).axes if ax.get_title()]
    assert titled and _base_of(titled[0]) == "n = 1049"


def test_a_funnel_row_keeps_a_base_per_panel():
    """The funnel has always titled its panels this way — pies now match it."""
    import matplotlib

    seen: list[str] = []
    orig = matplotlib.axes.Axes.set_title

    def spy(self, label, *a, **kw):
        seen.append(str(label))
        return orig(self, label, *a, **kw)

    matplotlib.axes.Axes.set_title = spy
    try:
        from reportbuilder.render.image.funnel import build_image_funnel

        build_image_funnel(_ctx(["Suomi", "Ruotsi", "Saksa"], chart_type="funnel"))
    finally:
        matplotlib.axes.Axes.set_title = orig
    assert seen and all("n = " in t for t in seen), seen


# ── and nothing on a panel may sit inside the shared legend ──────────────────

def test_the_shared_legend_clears_every_panel():
    """Asserted on each panel's FULL extent — the box including whatever is
    written above or below it — so anything put there later has to clear the
    legend too, rather than this holding only for today's layout."""
    fig = _pie_fig(["Suomi", "Ruotsi", "Saksa"], cats=LONG)
    fig.canvas.draw()
    renderer = fig.canvas.get_renderer()
    leg_box = next(a for a in fig.legends).get_window_extent(renderer)
    for ax in fig.axes:
        panel = ax.get_tightbbox(renderer)
        # A shared edge is not a collision: the panels are lifted to sit just
        # above the legend, so the two boxes touch by a fraction of a pixel.
        # What is being asserted is that neither is INSIDE the other.
        overlap = min(leg_box.y1, panel.y1) - max(leg_box.y0, panel.y0)
        assert overlap <= 1.0, (
            f"the legend sits {overlap:.0f}px into a panel: legend "
            f"{leg_box.bounds}, panel {panel.bounds}")
