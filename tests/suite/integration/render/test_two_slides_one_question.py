"""Two slides on the same question are two different charts.

A question can appear more than once in a deck on purpose: a total-level slide
and one split by region ("Extra slide" in Select), or — since a slide can be
drawn on part of the classifying variable — the same battery through design 1
and through design 2, which is made by duplicating the slide and changing the
tick.

The computed series were collected in a dict keyed by QUESTION, so the second
spec's numbers overwrote the first's and every slide of that question rendered
the last one. The editor showed it correctly, because a preview computes one
spec at a time; only the exported deck was wrong, which is the worst place for
it to be.
"""
from __future__ import annotations

import pandas as pd
import pytest

from reportbuilder.export.pptx_build import build_presentation
from reportbuilder.model.question import Question, QuestionModel, ValueLabel, Variable
from reportbuilder.model.report import (
    ChartSpec, ElementToggles, NumberFormat, Report, SortSpec,
)

pytestmark = pytest.mark.integration


def _model():
    q1 = Variable(name="q1", label="Suositteletko?", measurement="categorical",
                  value_labels=(ValueLabel(1.0, "Kyllä"), ValueLabel(2.0, "Ei")),
                  missing_values=frozenset())
    maa = Variable(name="maa", label="Maa", measurement="categorical",
                   value_labels=(ValueLabel(1.0, "Suomi"), ValueLabel(2.0, "Ruotsi")),
                   missing_values=frozenset())
    return QuestionModel(
        variables={"q1": q1, "maa": maa},
        questions=[Question(qid="q1", kind="single", variables=("q1",),
                            text="Suositteletko?")])


def _spec(slide_id, slot, **over):
    base = dict(question_ref="q1", chart_type="horizontal_bar", statistic="pct",
                classifying_var="maa", number_format=NumberFormat(),
                sort=SortSpec(basis="data_order"), template_slot=slot,
                elements=ElementToggles(), slide_id=slide_id)
    base.update(over)
    return ChartSpec(**base)


def _df():
    # Suomi says Kyllä 9 times in 10; Ruotsi says Ei 9 times in 10. Two groups
    # that could not be confused for one another.
    return pd.DataFrame({
        "q1": [1.0] * 9 + [2.0] + [1.0] + [2.0] * 9,
        "maa": [1.0] * 10 + [2.0] * 10,
    })


def _texts(prs, i):
    return " || ".join(sh.text_frame.text for sh in prs.slides[i].shapes
                       if sh.has_text_frame and sh.text_frame.text.strip())


def _picture(prs, i) -> bytes:
    """The drawn chart. Asserted on rather than on the slide's text, because the
    footer's group names come from the SPEC — they are printed whatever series
    the slide was given, so text alone cannot tell two slides apart."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    pics = [sh for sh in prs.slides[i].shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pics, f"slide {i} drew no chart"
    return pics[0].image.blob


def test_each_slide_is_drawn_on_its_own_group():
    report = Report(name="r", render_mode="image", template_ref="", charts=(
        _spec("a", "s1", classifying_values=("Suomi",)),
        _spec("b", "s2", classifying_values=("Ruotsi",)),
    ))
    prs = build_presentation(report, _model(), _df())
    # Suomi answers Kyllä 9 times in 10, Ruotsi 1 in 10: two charts that cannot
    # look alike unless they were given the same numbers.
    assert _picture(prs, 0) != _picture(prs, 1)
    assert "Suomi" in _texts(prs, 0) and "Ruotsi" in _texts(prs, 1)


def test_two_statistics_of_one_question_do_not_collapse_into_one():
    """The same collision without any group filter — a total slide and a split
    slide, which Select offers as an "Extra slide"."""
    report = Report(name="r", render_mode="image", template_ref="", charts=(
        _spec("a", "s1", classifying_var=None),
        _spec("b", "s2"),
    ))
    prs = build_presentation(report, _model(), _df())
    assert _picture(prs, 0) != _picture(prs, 1)
