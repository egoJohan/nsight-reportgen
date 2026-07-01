"""Shared, import-safe helpers for the fresh backend suite (`suite.*`).

Non-fixture utilities live here so tests can `from suite._helpers import ...`
(importing from a conftest is discouraged). Fixtures live in conftest.py.
"""
from __future__ import annotations

import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from reportbuilder.model.report import (
    ChartSpec, ElementToggles, NumberFormat, SortSpec,
)
from reportbuilder.render.base import RenderContext, Slot, StyleSpec


def make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def make_slot() -> Slot:
    return Slot(slide_index=0, left=Inches(1), top=Inches(1),
                width=Inches(8), height=Inches(5), name="slot1")


def make_spec(chart_type: str, **overrides) -> ChartSpec:
    kw = dict(
        question_ref="q1", chart_type=chart_type, statistic="pct",
        classifying_var=None, number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"), template_slot="slot1",
        elements=ElementToggles(title=True, legend=True, data_labels=True),
    )
    kw.update(overrides)
    return ChartSpec(**kw)


def make_ctx(chart_type: str, series, **spec_overrides):
    """A fully-formed image RenderContext (slide+slot+style+spec+series)."""
    prs, slide = make_slide()
    slot = make_slot()
    spec = make_spec(chart_type, **spec_overrides)
    ctx = RenderContext(slide=slide, slot=slot, style=StyleSpec(),
                        spec=spec, series=series, fmt=spec.number_format)
    return prs, slide, slot, ctx


def picture_shapes(slide):
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


def assert_single_picture(slide, slot):
    """Exactly one PICTURE, contained in the slot, undistorted PNG."""
    from PIL import Image

    pics = picture_shapes(slide)
    assert len(pics) == 1, f"expected 1 PICTURE shape, found {len(pics)}"
    pic = pics[0]
    assert 0 < pic.width <= slot.width
    assert 0 < pic.height <= slot.height
    blob = pic.image.blob
    assert blob[:4] == b"\x89PNG", "picture is not a PNG"
    px_w, px_h = Image.open(io.BytesIO(blob)).size
    assert abs(pic.width / pic.height - px_w / px_h) < 0.02, "image distorted"
    return pic


class RecordingChat:
    """A deterministic stand-in for `egohive_chat` / an injected `chat=`.

    `reply` may be a string (returned for every prompt) or a callable
    `(prompt) -> str`. All prompts seen are recorded in `.prompts`.
    """

    def __init__(self, reply="Canned reply"):
        self.reply = reply
        self.prompts: list[str] = []
        self.calls = 0

    def __call__(self, prompt: str, **_kw) -> str:
        self.calls += 1
        self.prompts.append(prompt)
        return self.reply(prompt) if callable(self.reply) else self.reply
