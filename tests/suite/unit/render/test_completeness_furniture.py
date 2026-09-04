"""The completeness check counts CHARTS, not every picture on the slide.

A template's harvested furniture is redrawn with `add_picture` — a logo, a rule,
a background band — by design. The check that every chart made it onto the deck
counted picture shapes, so one logo plus one chart read as two charts where one
was expected and the whole render was refused with a CompletenessError. The
deck was complete; the count was wrong.
"""
from __future__ import annotations

import io

import pytest
from pptx import Presentation
from pptx.util import Inches

from reportbuilder.model.report import (
    ChartSpec, ElementToggles, NumberFormat, Report, SortSpec,
)
from reportbuilder.render.deck import assert_complete


def _png() -> io.BytesIO:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (8, 8), "white").save(buf, format="PNG")
    buf.seek(0)
    return buf


def _report():
    spec = ChartSpec(question_ref="q1", chart_type="horizontal_bar", statistic="pct",
                     classifying_var=None, number_format=NumberFormat(),
                     sort=SortSpec(basis="data_order"), template_slot="s1",
                     elements=ElementToggles(), slide_id="a")
    return Report(name="r", render_mode="image", template_ref="", charts=(spec,))


def _deck_with(furniture: int):
    from reportbuilder.render.image._mpl import name_chart_picture

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for _ in range(furniture):                      # the template's own images
        slide.shapes.add_picture(_png(), 0, 0, Inches(1), Inches(1))
    chart = slide.shapes.add_picture(_png(), Inches(2), Inches(2),
                                     Inches(1), Inches(1))
    name_chart_picture(chart)                       # this one IS the chart
    return prs


def test_a_slide_with_no_furniture_still_counts_its_chart():
    assert_complete(_deck_with(0), _report())


def test_a_logo_on_the_slide_is_not_counted_as_a_chart():
    assert_complete(_deck_with(1), _report())


def test_a_missing_chart_is_still_caught():
    from reportbuilder.render.deck import CompletenessError

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])      # furniture only, no chart
    with pytest.raises(CompletenessError):
        assert_complete(prs, _report())
