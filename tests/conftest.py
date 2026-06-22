import uuid
from pathlib import Path

import pandas as pd
import pyreadstat
import pytest


@pytest.fixture
def tiny_sav(tmp_path: Path) -> Path:
    """A 6-respondent survey: one categorical awareness var, one segment var, one open-ended var."""
    df = pd.DataFrame(
        {
            "aware_attendo": [1, 1, 0, 1, 0, 1],
            "experience": [1, 2, 1, 2, 1, 2],
            "image_word": ["kallis", "Kallis ", "luotettava", "kallis", "hyvä", "luotettava"],
        }
    )
    out = tmp_path / "tiny.sav"
    pyreadstat.write_sav(
        df,
        str(out),
        column_labels=["Tunnetko Attendo", "Kokemus", "Kuvaile Attendoa"],
        variable_value_labels={
            "aware_attendo": {1.0: "Kyllä", 0.0: "Ei"},
            "experience": {1.0: "Kokemusta", 2.0: "Ei kokemusta"},
        },
    )
    return out


@pytest.fixture
def chart_pptx(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    data = CategoryChartData()
    data.categories = ["Attendo", "Esperi"]
    data.add_series("Series 1", (0.1, 0.2))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                                Inches(1), Inches(1), Inches(5), Inches(4), data)
    gf.name = "awareness_chart"
    out = tmp_path / "tmpl.pptx"
    prs.save(out)
    return out


@pytest.fixture
def dup_chart_pptx(tmp_path):
    """One slide with two bar charts, both named 'dup'."""
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    data1 = CategoryChartData()
    data1.categories = ["A", "B"]
    data1.add_series("S1", (0.1, 0.9))
    gf1 = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                                 Inches(1), Inches(1), Inches(3), Inches(3), data1)
    gf1.name = "dup"

    data2 = CategoryChartData()
    data2.categories = ["A", "B"]
    data2.add_series("S2", (0.3, 0.7))
    gf2 = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                                 Inches(5), Inches(1), Inches(3), Inches(3), data2)
    gf2.name = "dup"

    out = tmp_path / "dup.pptx"
    prs.save(out)
    return out


@pytest.fixture
def multiseries_chart_pptx(tmp_path):
    """One slide with one chart named 'multi', two series W2025 and W2024."""
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    data = CategoryChartData()
    data.categories = ["Attendo", "Esperi"]
    data.add_series("W2025", (0.8, 0.2))
    data.add_series("W2024", (0.7, 0.3))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                                Inches(1), Inches(1), Inches(5), Inches(4), data)
    gf.name = "multi"

    out = tmp_path / "multi.pptx"
    prs.save(out)
    return out
