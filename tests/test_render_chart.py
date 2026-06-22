import pytest
from nsight.render.template import Template


def test_template_finds_shape_by_slide_and_name(chart_pptx):
    tmpl = Template(chart_pptx)
    shape = tmpl.shape(slide_idx=0, name="awareness_chart")
    assert shape.has_chart


from nsight.render.fill_chart import fill_chart


def test_fill_chart_replaces_values(chart_pptx, tmp_path):
    tmpl = Template(chart_pptx)
    shape = tmpl.shape(slide_idx=0, name="awareness_chart")
    fill_chart(shape, categories=["Attendo", "Esperi"], series={"Series 1": [0.75, 0.25]})
    out = tmp_path / "out.pptx"
    tmpl.save(out)

    from pptx import Presentation
    chart = [s for s in Presentation(str(out)).slides[0].shapes if s.has_chart][0].chart
    plot = chart.plots[0]
    assert list(plot.categories) == ["Attendo", "Esperi"]
    assert list(plot.series[0].values) == [0.75, 0.25]


from nsight.render.renderer import SlideFill, ChartFill, render


def test_render_applies_chart_fill(chart_pptx, tmp_path):
    out = tmp_path / "rendered.pptx"
    render(
        template_path=chart_pptx,
        out_path=out,
        fills=[SlideFill(slide_idx=0, charts=[
            ChartFill(name="awareness_chart", categories=["Attendo", "Esperi"],
                      series={"Series 1": [0.9, 0.1]})
        ])],
    )
    from pptx import Presentation
    chart = [s for s in Presentation(str(out)).slides[0].shapes if s.has_chart][0].chart
    assert list(chart.plots[0].series[0].values) == [0.9, 0.1]


# ---------------------------------------------------------------------------
# Change 1 — nth-occurrence shape lookup
# ---------------------------------------------------------------------------

def test_shape_occurrence_0_returns_first(dup_chart_pptx):
    tmpl = Template(dup_chart_pptx)
    sh = tmpl.shape(slide_idx=0, name="dup", occurrence=0)
    # The first chart has series "S1"; verify by its series name
    assert sh.chart.plots[0].series[0].name == "S1"


def test_shape_occurrence_1_returns_second(dup_chart_pptx):
    tmpl = Template(dup_chart_pptx)
    sh = tmpl.shape(slide_idx=0, name="dup", occurrence=1)
    assert sh.chart.plots[0].series[0].name == "S2"


def test_shape_occurrence_out_of_range_raises_key_error(dup_chart_pptx):
    tmpl = Template(dup_chart_pptx)
    import pytest
    with pytest.raises(KeyError) as exc_info:
        tmpl.shape(slide_idx=0, name="dup", occurrence=2)
    msg = str(exc_info.value)
    assert "dup" in msg
    assert "2" in msg or "found" in msg.lower()


# ---------------------------------------------------------------------------
# Change 2 — read existing chart series / replace_one_series
# ---------------------------------------------------------------------------

from nsight.render.fill_chart import read_chart_series, read_chart_categories, replace_one_series


def test_read_chart_series_returns_all_series(multiseries_chart_pptx):
    tmpl = Template(multiseries_chart_pptx)
    shape = tmpl.shape(slide_idx=0, name="multi")
    series = read_chart_series(shape)
    assert set(series.keys()) == {"W2025", "W2024"}
    assert list(series["W2025"]) == [0.8, 0.2]
    assert list(series["W2024"]) == [0.7, 0.3]


def test_read_chart_categories_returns_strings(multiseries_chart_pptx):
    tmpl = Template(multiseries_chart_pptx)
    shape = tmpl.shape(slide_idx=0, name="multi")
    cats = read_chart_categories(shape)
    assert cats == ["Attendo", "Esperi"]


def test_read_chart_series_raises_on_non_chart(chart_pptx):
    """read_chart_series raises ValueError when shape has no chart."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation(str(chart_pptx))
    # Add a plain text box (no chart) to slide 0
    slide = prs.slides[0]
    # Use existing awareness_chart shape then patch has_chart manually via a non-chart shape
    # Instead: build a fake object
    class FakeShape:
        has_chart = False
        name = "notachart"
    import pytest
    with pytest.raises(ValueError):
        read_chart_series(FakeShape())


def test_replace_one_series_updates_target_preserves_others(multiseries_chart_pptx, tmp_path):
    tmpl = Template(multiseries_chart_pptx)
    shape = tmpl.shape(slide_idx=0, name="multi")
    replace_one_series(shape, series_name="W2025", values=[0.86, 0.13])
    out = tmp_path / "replaced.pptx"
    tmpl.save(out)

    from pptx import Presentation
    prs2 = Presentation(str(out))
    chart = [s for s in prs2.slides[0].shapes if s.has_chart][0].chart
    plot = chart.plots[0]
    series_map = {s.name: list(s.values) for s in plot.series}
    assert series_map["W2025"] == pytest.approx([0.86, 0.13])
    assert series_map["W2024"] == pytest.approx([0.7, 0.3])
    assert list(plot.categories) == ["Attendo", "Esperi"]


def test_replace_one_series_bad_name_raises_key_error(multiseries_chart_pptx):
    tmpl = Template(multiseries_chart_pptx)
    shape = tmpl.shape(slide_idx=0, name="multi")
    import pytest
    with pytest.raises(KeyError):
        replace_one_series(shape, series_name="W2099", values=[0.5, 0.5])


# ---------------------------------------------------------------------------
# Change 3 — ChartFill single-series-by-category mode in render()
# ---------------------------------------------------------------------------

def test_render_single_series_by_category_preserves_others(multiseries_chart_pptx, tmp_path):
    """render() with series_name replaces only that series; other series unchanged."""
    out = tmp_path / "single_series.pptx"
    render(
        template_path=multiseries_chart_pptx,
        out_path=out,
        fills=[SlideFill(slide_idx=0, charts=[
            ChartFill(
                name="multi",
                categories=[],
                series={},
                series_name="W2025",
                values_by_category={"Attendo": 0.86, "Esperi": 0.13},
            )
        ])],
    )
    from pptx import Presentation
    prs2 = Presentation(str(out))
    chart = [s for s in prs2.slides[0].shapes if s.has_chart][0].chart
    plot = chart.plots[0]
    series_map = {s.name: list(s.values) for s in plot.series}
    # W2025 aligned to category order -> [0.86, 0.13]
    assert series_map["W2025"] == pytest.approx([0.86, 0.13])
    # W2024 must be untouched
    assert series_map["W2024"] == pytest.approx([0.7, 0.3])
    # categories unchanged
    assert [str(c) for c in plot.categories] == ["Attendo", "Esperi"]


def test_render_single_series_missing_category_keeps_existing(multiseries_chart_pptx, tmp_path):
    """Categories absent from values_by_category keep their existing value."""
    out = tmp_path / "partial.pptx"
    render(
        template_path=multiseries_chart_pptx,
        out_path=out,
        fills=[SlideFill(slide_idx=0, charts=[
            ChartFill(
                name="multi",
                categories=[],
                series={},
                series_name="W2025",
                values_by_category={"Attendo": 0.9},  # Esperi absent
            )
        ])],
    )
    from pptx import Presentation
    prs2 = Presentation(str(out))
    chart = [s for s in prs2.slides[0].shapes if s.has_chart][0].chart
    plot = chart.plots[0]
    series_map = {s.name: list(s.values) for s in plot.series}
    # Attendo updated, Esperi keeps original 0.2
    assert series_map["W2025"] == pytest.approx([0.9, 0.2])
    # W2024 unchanged
    assert series_map["W2024"] == pytest.approx([0.7, 0.3])


def test_render_chartfill_occurrence_targets_nth(dup_chart_pptx, tmp_path):
    """ChartFill with occurrence=1 replaces only the second chart."""
    out = tmp_path / "occ.pptx"
    render(
        template_path=dup_chart_pptx,
        out_path=out,
        fills=[SlideFill(slide_idx=0, charts=[
            ChartFill(
                name="dup",
                categories=["A", "B"],
                series={"S2": [0.6, 0.4]},
                occurrence=1,
            )
        ])],
    )
    from pptx import Presentation
    prs2 = Presentation(str(out))
    charts = [s for s in prs2.slides[0].shapes if s.has_chart]
    first_vals = list(charts[0].chart.plots[0].series[0].values)
    second_vals = list(charts[1].chart.plots[0].series[0].values)
    # First chart (occurrence=0) must be untouched: S1 with [0.1, 0.9]
    assert first_vals == pytest.approx([0.1, 0.9])
    # Second chart (occurrence=1) replaced: S2 -> [0.6, 0.4]
    assert second_vals == pytest.approx([0.6, 0.4])
