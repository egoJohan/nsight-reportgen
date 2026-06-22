import pytest
from pathlib import Path
from nsight.render.template import Template
from nsight.render.fill_table import fill_table


@pytest.fixture
def table_pptx(tmp_path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    gf = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    gf.name = "awareness_table"
    out = tmp_path / "t.pptx"
    prs.save(out)
    return out


def test_fill_table_sets_cells(table_pptx, tmp_path):
    tmpl = Template(table_pptx)
    shape = tmpl.shape(slide_idx=0, name="awareness_table")
    fill_table(shape, {(0, 0): "Attendo", (0, 1): "92 %", (1, 1): "-1 %"})
    out = tmp_path / "o.pptx"
    tmpl.save(out)

    from pptx import Presentation
    tbl = [s for s in Presentation(str(out)).slides[0].shapes if s.has_table][0].table
    assert tbl.cell(0, 1).text == "92 %"
    assert tbl.cell(1, 1).text == "-1 %"
