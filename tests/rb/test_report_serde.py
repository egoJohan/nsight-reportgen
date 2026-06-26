"""Tests for Report JSON (de)serialization (Task 0.9)."""
import json
import pytest

from reportbuilder.model.report import (
    ChartSpec,
    ElementToggles,
    NumberFormat,
    Report,
    SortSpec,
    report_from_json,
    report_to_json,
)


def _make_report() -> Report:
    chart_a = ChartSpec(
        question_ref="q1",
        chart_type="bar",
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="pct"),
        template_slot="slot_a",
        elements=ElementToggles(),
        scatter_xy=None,
    )
    chart_b = ChartSpec(
        question_ref="q2",
        chart_type="column",
        statistic="pct",
        classifying_var="region",
        number_format=NumberFormat(pct_decimals=1),
        sort=SortSpec(basis="topbox_sum", topbox_codes=(4.0, 5.0)),
        template_slot="slot_b",
        elements=ElementToggles(legend=False),
        scatter_xy=None,
    )
    chart_c = ChartSpec(
        question_ref="q3",
        chart_type="scatter",
        statistic="mean",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"),
        template_slot="slot_c",
        elements=ElementToggles(),
        scatter_xy=("satisfaction", "loyalty"),
    )
    return Report(
        name="Test Report",
        render_mode="native",
        template_ref="tmpl_001",
        charts=(chart_a, chart_b, chart_c),
    )


def test_round_trip_equality():
    """Full round-trip: report_from_json(report_to_json(r)) == r."""
    r = _make_report()
    assert report_from_json(report_to_json(r)) == r


def test_to_json_produces_valid_json_string():
    """report_to_json returns a str that parses as valid JSON."""
    r = _make_report()
    result = report_to_json(r)
    assert isinstance(result, str)
    parsed = json.loads(result)
    assert isinstance(parsed, dict)


def test_from_json_accepts_parsed_dict():
    """report_from_json accepts an already-parsed dict."""
    r = _make_report()
    parsed_dict = json.loads(report_to_json(r))
    result = report_from_json(parsed_dict)
    assert result == r


def test_tuples_restored_after_round_trip():
    """topbox_codes and scatter_xy are tuples (not lists) after round-trip."""
    r = _make_report()
    result = report_from_json(report_to_json(r))
    assert result.charts[1].sort.topbox_codes == (4.0, 5.0)
    assert isinstance(result.charts[1].sort.topbox_codes, tuple)
    assert result.charts[2].scatter_xy == ("satisfaction", "loyalty")
    assert isinstance(result.charts[2].scatter_xy, tuple)


# ---------------------------------------------------------------------------
# slide_title / slide_description round-trip tests (REQ-C-24a, D-04)
# ---------------------------------------------------------------------------


def test_slide_title_round_trips():
    """ChartSpec.slide_title is preserved by report_to_json / report_from_json. (REQ-C-24a)"""
    chart = ChartSpec(
        question_ref="q1",
        chart_type="bar",
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="pct"),
        template_slot="s1",
        elements=ElementToggles(),
        slide_title="My Custom Title",
    )
    report = Report(
        name="R", render_mode="image", template_ref="t", charts=(chart,)
    )
    restored = report_from_json(report_to_json(report))
    assert restored.charts[0].slide_title == "My Custom Title"
    assert restored.charts[0].slide_description is None


def test_slide_description_round_trips():
    """ChartSpec.slide_description is preserved by report_to_json / report_from_json. (REQ-D-04)"""
    chart = ChartSpec(
        question_ref="q1",
        chart_type="bar",
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="pct"),
        template_slot="s1",
        elements=ElementToggles(),
        slide_title="Title",
        slide_description="A subtitle line",
    )
    report = Report(
        name="R", render_mode="image", template_ref="t", charts=(chart,)
    )
    restored = report_from_json(report_to_json(report))
    assert restored.charts[0].slide_title == "Title"
    assert restored.charts[0].slide_description == "A subtitle line"


def test_slide_title_defaults_to_none():
    """ChartSpec.slide_title/description default to None (backward-compatible). (REQ-C-24a)"""
    chart = ChartSpec(
        question_ref="q1",
        chart_type="bar",
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="pct"),
        template_slot="s1",
        elements=ElementToggles(),
    )
    assert chart.slide_title is None
    assert chart.slide_description is None
    # Round-trip through JSON preserves None.
    report = Report(name="R", render_mode="image", template_ref="t", charts=(chart,))
    restored = report_from_json(report_to_json(report))
    assert restored.charts[0].slide_title is None
    assert restored.charts[0].slide_description is None


def test_slide_chrome_uses_slide_title(monkeypatch) -> None:
    """add_image_slide_chrome uses spec.slide_title as the title when set. (REQ-C-24a)"""
    from pptx import Presentation
    from pptx.util import Inches
    from reportbuilder.render.base import RenderContext, Slot, StyleSpec
    from reportbuilder.render.image.slide_chrome import add_image_slide_chrome
    from reportbuilder.stats.series import Cell, SeriesResult

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slot = Slot(slide_index=0, left=Inches(1), top=Inches(1),
                width=Inches(8), height=Inches(5), name="s1")
    spec = ChartSpec(
        question_ref="q1",
        chart_type="bar",
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="pct"),
        template_slot="s1",
        elements=ElementToggles(),
        slide_title="Override Title",
        slide_description="Override Description",
    )
    cats = ("A",)
    segs = ("Total",)
    cells = {("A", "Total"): Cell(pct=50.0)}
    series = SeriesResult(categories=cats, segments=segs, cells=cells,
                          base_n={"Total": 10}, statistic="pct")
    ctx = RenderContext(
        slide=slide, slot=slot, style=StyleSpec(),
        spec=spec, series=series, fmt=spec.number_format,
        title="Question Text",  # should be overridden by spec.slide_title
    )

    add_image_slide_chrome(ctx)

    textboxes = [s for s in slide.shapes if s.has_text_frame]
    # "Override Title" must appear in a textbox (not "Question Text").
    override_found = any(
        "Override Title" in s.text_frame.text for s in textboxes
    )
    assert override_found, "slide_title 'Override Title' not found in chrome textboxes"
    # "Override Description" must appear as the subtitle.
    desc_found = any(
        "Override Description" in s.text_frame.text for s in textboxes
    )
    assert desc_found, "slide_description 'Override Description' not found in chrome textboxes"
