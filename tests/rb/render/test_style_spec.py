"""Tests for load_style_spec — TDD for Task 5.1 (REQ-C-25/27a)."""
from __future__ import annotations
import pytest
from pptx import Presentation
from pptx.util import Emu, Pt

from reportbuilder.render.style_spec import load_style_spec
from reportbuilder.render.base import Slot


def _blank_pptx(tmp_path):
    """Return path to a minimal 1-blank-slide presentation."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    path = tmp_path / "template.pptx"
    prs.save(str(path))
    return str(path)


def test_default_font_classes(tmp_path):
    """Minimal template with no special shapes yields documented font defaults."""
    path = _blank_pptx(tmp_path)
    spec = load_style_spec(path)

    assert spec.font_for("title") == ("Arial", 14)
    assert spec.font_for("data_labels") == ("Arial", 10)
    # Unknown class falls back to ("Arial", 10)
    assert spec.font_for("unknown_class") == ("Arial", 10)
    # Palette entries are 6-char hex strings
    assert isinstance(spec.color_for(0), str)
    assert len(spec.color_for(0)) == 6
    # Palette wraps: index 8 == index 0 (palette length is 8)
    assert spec.color_for(8) == spec.color_for(0)


def test_slots_from_placeholders(tmp_path):
    """Textbox shapes on slides are captured as Slots with correct geometry."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left_a, top_a, w_a, h_a = Emu(914400), Emu(457200), Emu(3200400), Emu(2400300)
    left_b, top_b, w_b, h_b = Emu(4114800), Emu(914400), Emu(3200400), Emu(2400300)

    tb_a = slide.shapes.add_textbox(left_a, top_a, w_a, h_a)
    tb_a.name = "chart_left"
    tb_b = slide.shapes.add_textbox(left_b, top_b, w_b, h_b)
    tb_b.name = "chart_right"

    path = tmp_path / "template.pptx"
    prs.save(str(path))

    spec = load_style_spec(str(path))

    slot_a = spec.slot("chart_left")
    assert isinstance(slot_a, Slot)
    assert slot_a.left == int(left_a)
    assert slot_a.top == int(top_a)
    assert slot_a.width == int(w_a)
    assert slot_a.height == int(h_a)
    assert slot_a.slide_index == 0

    assert set(spec.slots()) >= {"chart_left", "chart_right"}


def test_font_class_read_from_template(tmp_path):
    """A shape named 'style:<class>' causes font_for() to return that shape's font."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb = slide.shapes.add_textbox(Emu(100), Emu(100), Emu(914400), Emu(457200))
    tb.name = "style:title"
    tf = tb.text_frame
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = "Title Style Marker"
    run.font.name = "Calibri"
    run.font.size = Pt(18)

    path = tmp_path / "template.pptx"
    prs.save(str(path))

    spec = load_style_spec(str(path))
    assert spec.font_for("title") == ("Calibri", 18)
