"""Task D3: the slide's OWN text (title/subtitle/footer, special-slide surface,
demographics panel) adapts to the template's background — the piece Task D1
(test_image_dark_background.py) deliberately left out because slide_chrome.py,
special_slide.py, and demographics_grid.py were being edited elsewhere at the
time.

Same rule as the chart furniture: house_style.furniture_colors(background),
shared with contrast_ink's 0.55 luminance threshold. A light background
(including the house default) must resolve to exactly today's PX_CREAM/
PX_INK/PX_MUTED constants — byte-identical to before this task. Only a dark
background gets a derived, legible set.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from reportbuilder.model.report import ChartSpec, ElementToggles, NumberFormat, SortSpec
from reportbuilder.render.base import RenderContext, Slot, StyleSpec
from reportbuilder.render.house_style import PX_CREAM, PX_INK, PX_MUTED
from reportbuilder.stats.series import Cell, SeriesResult

# Same dark background Task D1 used — well below the 0.55 luminance threshold,
# unambiguously "dark" by the house rule.
_DARK_BG = "0B0B12"
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _slot() -> Slot:
    return Slot(slide_index=0, left=Inches(1), top=Inches(1),
                width=Inches(8), height=Inches(5), name="slot1")


def _series() -> SeriesResult:
    cats = ("Kyllä", "Ei")
    cells = {(c, "Total"): Cell(pct=v) for c, v in zip(cats, [60.0, 40.0])}
    return SeriesResult(categories=cats, segments=("Total",), cells=cells,
                        base_n={"Total": 250}, statistic="pct")


def _chart_ctx(background: str, **spec_kw):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style = StyleSpec()
    style.background = background
    spec_fields = dict(
        question_ref="q1", chart_type="horizontal_bar", statistic="pct",
        classifying_var=None, number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"), template_slot="slot1",
        elements=ElementToggles(title=True, data_labels=True, legend=True, n=True),
        slide_title="Onko tämä kysymys?",
        slide_description="Selittävä alaotsikko.",
    )
    spec_fields.update(spec_kw)
    spec = ChartSpec(**spec_fields)
    ctx = RenderContext(slide=slide, slot=_slot(), style=style, spec=spec,
                        series=_series(), fmt=spec.number_format,
                        title="Onko tämä kysymys?")
    return prs, slide, ctx


def _runs(slide):
    """(text, colour RGBColor-or-None) for every run on every textbox on *slide*."""
    out = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                col = None
                try:
                    col = run.font.color.rgb
                except Exception:
                    pass
                out.append((run.text, col))
    return out


def _colour_for(slide, needle: str) -> RGBColor:
    for text, col in _runs(slide):
        if needle in text:
            return col
    raise AssertionError(f"no run containing {needle!r} found")


def _shape_fill_rgb(slide, shape_type=MSO_SHAPE_TYPE.AUTO_SHAPE, index=0):
    shapes = [s for s in slide.shapes if s.shape_type == shape_type]
    return shapes[index].fill.fore_color.rgb


# ---------------------------------------------------------------------------
# slide_chrome.add_image_slide_chrome — title, subtitle, footer
# ---------------------------------------------------------------------------

class TestSlideChromeDarkBackground:
    def test_light_default_background_is_exactly_house_cream(self):
        from reportbuilder.render.image.slide_chrome import add_image_slide_chrome

        _prs, slide, ctx = _chart_ctx("")
        add_image_slide_chrome(ctx)
        assert _shape_fill_rgb(slide) == PX_CREAM

    def test_dark_background_matches_the_template_not_cream(self):
        from reportbuilder.render.image.slide_chrome import add_image_slide_chrome

        _prs, slide, ctx = _chart_ctx(_DARK_BG)
        add_image_slide_chrome(ctx)
        assert _shape_fill_rgb(slide) == RGBColor(0x0B, 0x0B, 0x12)

    def test_title_is_house_ink_on_light_default(self):
        from reportbuilder.render.image.slide_chrome import add_image_slide_chrome

        _prs, slide, ctx = _chart_ctx("")
        add_image_slide_chrome(ctx)
        assert _colour_for(slide, "Onko tämä kysymys?") == PX_INK

    def test_title_flips_to_white_on_dark_background(self):
        from reportbuilder.render.image.slide_chrome import add_image_slide_chrome

        _prs, slide, ctx = _chart_ctx(_DARK_BG)
        add_image_slide_chrome(ctx)
        assert _colour_for(slide, "Onko tämä kysymys?") == _WHITE

    def test_subtitle_is_house_muted_on_light_default(self):
        from reportbuilder.render.image.slide_chrome import add_image_slide_chrome

        _prs, slide, ctx = _chart_ctx("")
        add_image_slide_chrome(ctx)
        assert _colour_for(slide, "Selittävä alaotsikko.") == PX_MUTED

    def test_subtitle_flips_to_light_on_dark_background(self):
        from reportbuilder.render.image.slide_chrome import add_image_slide_chrome
        from matplotlib.colors import to_rgb

        _prs, slide, ctx = _chart_ctx(_DARK_BG)
        add_image_slide_chrome(ctx)
        col = _colour_for(slide, "Selittävä alaotsikko.")
        assert col != PX_MUTED
        r, g, b = to_rgb(f"#{col}")
        assert (r + g + b) / 3.0 > 0.5

    def test_footer_is_house_muted_on_light_default(self):
        from reportbuilder.render.image.slide_chrome import add_image_slide_chrome

        _prs, slide, ctx = _chart_ctx("")
        add_image_slide_chrome(ctx)
        assert _colour_for(slide, "N = 250") == PX_MUTED

    def test_footer_flips_to_light_on_dark_background(self):
        from reportbuilder.render.image.slide_chrome import add_image_slide_chrome
        from matplotlib.colors import to_rgb

        _prs, slide, ctx = _chart_ctx(_DARK_BG)
        add_image_slide_chrome(ctx)
        col = _colour_for(slide, "N = 250")
        assert col != PX_MUTED
        r, g, b = to_rgb(f"#{col}")
        assert (r + g + b) / 3.0 > 0.5

    def test_footer_flips_to_light_using_the_scale_caption_too(self):
        """The scale-endpoint caption (right side of the footer row) shares the
        footer's derived muted colour and must flip the same way (lighter than
        house MUTED, legible against the dark ground — not necessarily pure
        white, which is the TITLE's ink, not the muted footer text)."""
        from reportbuilder.render.image.slide_chrome import add_image_slide_chrome
        from matplotlib.colors import to_rgb

        cats = ("1", "4", "7")
        cells = {(c, "Total"): Cell(pct=v) for c, v in zip(cats, [20.0, 50.0, 30.0])}
        series = SeriesResult(categories=cats, segments=("Total",), cells=cells,
                              base_n={"Total": 40}, statistic="pct",
                              caption="1 = eri mieltä · 7 = samaa mieltä")
        _prs, slide, ctx = _chart_ctx(_DARK_BG, chart_type="stacked_horizontal_bar")
        ctx = RenderContext(slide=slide, slot=ctx.slot, style=ctx.style, spec=ctx.spec,
                            series=series, fmt=ctx.fmt, title=ctx.title)
        add_image_slide_chrome(ctx)
        col = _colour_for(slide, "eri mieltä")
        assert col != PX_MUTED
        r, g, b = to_rgb(f"#{col}")
        assert (r + g + b) / 3.0 > 0.5


# ---------------------------------------------------------------------------
# special_slide.render_special_slide — surface + heading/bullets
# ---------------------------------------------------------------------------

def _special_spec(**kw) -> ChartSpec:
    fields = dict(
        question_ref="", chart_type="special_overview", statistic="pct",
        classifying_var=None, number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"), template_slot="auto",
        elements=ElementToggles(),
        slide_title="Tutkimuksen taustaa",
        options={"bullets": ["Eka havainto", "Toka havainto"]},
    )
    fields.update(kw)
    return ChartSpec(**fields)


def _special_slide(background: str):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style = StyleSpec()
    style.background = background
    from reportbuilder.render.image.special_slide import render_special_slide
    render_special_slide(slide, _slot(), style, _special_spec())
    return slide


class TestSpecialSlideDarkBackground:
    def test_fill_is_exactly_house_cream_on_light_default(self):
        slide = _special_slide("")
        assert _shape_fill_rgb(slide) == PX_CREAM

    def test_fill_matches_the_template_background_when_dark(self):
        slide = _special_slide(_DARK_BG)
        assert _shape_fill_rgb(slide) == RGBColor(0x0B, 0x0B, 0x12)

    def test_heading_is_house_ink_on_light_default(self):
        slide = _special_slide("")
        assert _colour_for(slide, "Tutkimuksen taustaa") == PX_INK

    def test_heading_flips_to_white_on_dark_background(self):
        slide = _special_slide(_DARK_BG)
        assert _colour_for(slide, "Tutkimuksen taustaa") == _WHITE

    def test_bullet_text_is_house_ink_on_light_default(self):
        slide = _special_slide("")
        assert _colour_for(slide, "Eka havainto") == PX_INK

    def test_bullet_text_flips_to_white_on_dark_background(self):
        slide = _special_slide(_DARK_BG)
        assert _colour_for(slide, "Eka havainto") == _WHITE


# ---------------------------------------------------------------------------
# demographics_grid.render_demographics_grid — panel fill + heading
# ---------------------------------------------------------------------------

def _grid_slide(background: str):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style = StyleSpec()
    style.background = background
    spec = ChartSpec(
        question_ref="", chart_type="special_demographics", statistic="pct",
        classifying_var=None, number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"), template_slot="auto",
        elements=ElementToggles(), slide_title="Vastaajat",
        options={"charts": []},
    )
    from reportbuilder.render.image.demographics_grid import render_demographics_grid
    render_demographics_grid(slide, _slot(), style, spec, series_by_ref={}, titles={})
    return slide


class TestDemographicsGridDarkBackground:
    def test_panel_fill_is_exactly_house_cream_on_light_default(self):
        slide = _grid_slide("")
        assert _shape_fill_rgb(slide) == PX_CREAM

    def test_panel_fill_matches_the_template_background_when_dark(self):
        slide = _grid_slide(_DARK_BG)
        assert _shape_fill_rgb(slide) == RGBColor(0x0B, 0x0B, 0x12)

    def test_heading_is_house_ink_on_light_default(self):
        slide = _grid_slide("")
        assert _colour_for(slide, "Vastaajat") == PX_INK

    def test_heading_flips_to_white_on_dark_background(self):
        slide = _grid_slide(_DARK_BG)
        assert _colour_for(slide, "Vastaajat") == _WHITE
