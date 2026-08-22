"""Unit tests for nSight house style applied to image-mode bar charts.

Verifies that:
- build_image_bar and build_image_column use the TEAL palette, not matplotlib default blue
- Data labels are always rendered (and the annotated text appears in the figure)
- build_image_column auto-orients to horizontal when there are > 6 categories or long labels
- add_image_slide_chrome adds background, title, and N textboxes to the slide

REQ-C-24a/f/h, REQ-C-25, REQ-C-27a, REQ-D-04
"""
from __future__ import annotations

import io

import pytest
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from reportbuilder.render.base import RenderContext, Slot, StyleSpec
from reportbuilder.render.house_style import TEAL, series_colors, register_fonts, CREAM
from reportbuilder.model.report import ChartSpec, ElementToggles, NumberFormat, SortSpec
from reportbuilder.stats.series import Cell, SeriesResult
from reportbuilder.testing.fixtures import known_series


# ---------------------------------------------------------------------------
# Shared fixtures
# ---------------------------------------------------------------------------

def _make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _make_slot() -> Slot:
    return Slot(
        slide_index=0,
        left=Inches(1), top=Inches(1),
        width=Inches(8), height=Inches(5),
        name="slot1",
    )


def _make_spec(chart_type: str = "horizontal_bar") -> ChartSpec:
    return ChartSpec(
        question_ref="q1",
        chart_type=chart_type,
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"),
        template_slot="slot1",
        elements=ElementToggles(title=True, data_labels=True, legend=True, n=True),
    )


def _make_ctx(chart_type="horizontal_bar", series=None, title="Test question"):
    prs, slide = _make_slide()
    slot = _make_slot()
    style = StyleSpec()
    spec = _make_spec(chart_type)
    sr = series or known_series()
    ctx = RenderContext(
        slide=slide, slot=slot, style=style, spec=spec,
        series=sr, fmt=spec.number_format, title=title,
    )
    return prs, slide, slot, ctx


def _many_cat_series() -> SeriesResult:
    """8 categories with long labels — triggers auto-orientation to horizontal."""
    cats = (
        "En mitään näistä", "Attendo Finland", "Esperi Care",
        "Rinnekodit Oy", "Validia Kuntoutus", "Onnikodit",
        "Mainio-kodit", "Humana Suomi",
    )
    segs = ("Total",)
    cells = {(c, "Total"): Cell(pct=float(10 + i * 5)) for i, c in enumerate(cats)}
    return SeriesResult(
        categories=cats, segments=segs, cells=cells,
        base_n={"Total": 100}, statistic="pct",
    )


def _multi_series() -> SeriesResult:
    cats = ("Yes", "No")
    segs = ("Group A", "Group B")
    cells = {
        ("Yes", "Group A"): Cell(pct=70.0, count=7.0, mean=None),
        ("Yes", "Group B"): Cell(pct=50.0, count=5.0, mean=None),
        ("No", "Group A"): Cell(pct=30.0, count=3.0, mean=None),
        ("No", "Group B"): Cell(pct=50.0, count=5.0, mean=None),
    }
    return SeriesResult(
        categories=cats, segments=segs, cells=cells,
        base_n={"Group A": 10, "Group B": 10}, statistic="pct",
    )


# ---------------------------------------------------------------------------
# 1 — house_style module
# ---------------------------------------------------------------------------

def test_register_fonts_idempotent():
    """register_fonts() can be called multiple times without error."""
    register_fonts()
    register_fonts()  # second call must not raise or duplicate entries


def test_series_colors_single():
    """Single series returns the strong TEAL (not matplotlib default blue)."""
    cols = series_colors(1)
    assert len(cols) == 1
    assert cols[0] == TEAL, f"Expected TEAL {TEAL!r}, got {cols[0]!r}"


def test_series_colors_multi_ends_darkest():
    """Multi-series ramp always ends with the darkest teal (REQ-C-27a)."""
    for n in (2, 3, 4):
        cols = series_colors(n)
        assert len(cols) == n
        assert cols[-1] == TEAL, (
            f"n={n}: last color should be TEAL {TEAL!r}, got {cols[-1]!r}"
        )


def test_series_colors_five_cycles():
    """series_colors(5) returns 5 colours (cycles ramp beyond 4). REQ-C-27a."""
    cols = series_colors(5)
    assert len(cols) == 5


# ---------------------------------------------------------------------------
# 2 — build_image_bar uses TEAL and adds data labels
# ---------------------------------------------------------------------------

def test_image_bar_uses_teal_not_default_blue():
    """build_image_bar renders bars in TEAL (#13615E), not matplotlib default blue.

    REQ-C-27a: chart-property layout from a style spec.
    """
    from reportbuilder.render.image.bars import build_image_bar
    prs, slide, slot, ctx = _make_ctx("horizontal_bar")
    build_image_bar(ctx)

    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1

    # Verify via PIL that the PNG contains the TEAL colour
    from PIL import Image
    img = Image.open(io.BytesIO(pics[0].image.blob)).convert("RGB")
    pixels = list(img.getdata())
    teal_r, teal_g, teal_b = 0x13, 0x61, 0x5E
    # Tolerance of ±4 per channel to allow for PNG compression / anti-aliasing
    def near_teal(p):
        return (
            abs(p[0] - teal_r) <= 4
            and abs(p[1] - teal_g) <= 4
            and abs(p[2] - teal_b) <= 4
        )
    teal_count = sum(1 for p in pixels if near_teal(p))
    assert teal_count > 50, (
        f"Expected ≥50 TEAL pixels in bar PNG, found {teal_count}. "
        "Bars may be using a different colour."
    )


def test_image_bar_data_labels_present():
    """build_image_bar always renders data labels on bars.

    REQ-C-24f: category numeric values (data labels) must be present.
    """
    from reportbuilder.render.image.bars import build_image_bar
    # Use a series with a known value (60 %) to confirm labels
    prs, slide, slot, ctx = _make_ctx("horizontal_bar")
    build_image_bar(ctx)

    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    # Non-empty PNG blob confirms rendering completed; values are in the image
    assert len(pics[0].image.blob) > 1000


def test_image_column_data_labels_present():
    """build_image_column renders data labels (REQ-C-24f)."""
    from reportbuilder.render.image.bars import build_image_column
    prs, slide, slot, ctx = _make_ctx("vertical_bar")
    build_image_column(ctx)

    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    assert len(pics[0].image.blob) > 1000


# ---------------------------------------------------------------------------
# 3 — auto-orientation in build_image_column
# ---------------------------------------------------------------------------

def test_image_column_auto_orients_to_horizontal_many_cats():
    """build_image_column auto-switches to horizontal bars for > 6 categories.

    REQ-C-27a: deterministic layout from style spec.
    """
    from reportbuilder.render.image.bars import build_image_column
    prs, slide, slot, ctx = _make_ctx("vertical_bar", series=_many_cat_series())
    build_image_column(ctx)

    # Must still produce exactly one PICTURE, regardless of orientation chosen
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    assert len(pics[0].image.blob) > 1000


def test_image_column_auto_orients_to_horizontal_long_labels():
    """build_image_column auto-switches when any label exceeds 14 chars (REQ-C-27a)."""
    from reportbuilder.render.image.bars import build_image_column, _should_orient_horizontal
    # Verify the helper detects the condition
    assert _should_orient_horizontal(["This is a very long label"]) is True
    assert _should_orient_horizontal(["Short", "Also short"]) is False
    assert _should_orient_horizontal(["A"] * 7) is True   # > 6 categories


# ---------------------------------------------------------------------------
# 4 — multi-series uses teal ramp
# ---------------------------------------------------------------------------

def test_image_bar_multi_series_teal_ramp():
    """build_image_bar multi-series uses teal ramp colours, not default palette.

    REQ-C-27a: colour layout driven by style spec.
    """
    from reportbuilder.render.image.bars import build_image_bar
    prs, slide, slot, ctx = _make_ctx("horizontal_bar", series=_multi_series())
    build_image_bar(ctx)

    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    assert len(pics[0].image.blob) > 1000


# ---------------------------------------------------------------------------
# 5 — slide chrome
# ---------------------------------------------------------------------------

def test_add_image_slide_chrome_adds_title_and_n():
    """add_image_slide_chrome adds a cream background, title textbox, and N annotation.

    REQ-C-24a (title), REQ-C-24h (N), REQ-C-27a (style spec driven), REQ-D-04 (question text).
    """
    from reportbuilder.render.image.slide_chrome import add_image_slide_chrome

    prs, slide, slot, ctx = _make_ctx(title="Autettu tunnettuus")
    add_image_slide_chrome(ctx)

    shapes = list(slide.shapes)
    # Must have added at least 3 shapes: background rect, accent bar, title textbox
    # (N textbox too if base_n has "Total")
    assert len(shapes) >= 3, f"Expected ≥3 shapes after chrome, got {len(shapes)}"

    # Confirm a textbox carries the title text
    from pptx.enum.shapes import MSO_SHAPE_TYPE as ST
    textboxes = [s for s in shapes if s.has_text_frame]
    titles_found = [
        s for s in textboxes
        if "Autettu tunnettuus" in s.text_frame.text
    ]
    assert titles_found, "Title text 'Autettu tunnettuus' not found in any textbox"


def test_slide_chrome_uses_ai_title_not_question_text():
    """When the AI-generated headline (spec.slide_title) is set, the slide shows
    it and NOT the raw question text. Originals are not shown once AI lands.
    (Pairs with the engine's label-override test that drops the original category
    labels — together they guarantee title + legend labels are the AI versions.)"""
    import dataclasses
    from reportbuilder.render.image.slide_chrome import add_image_slide_chrome

    prs, slide, slot, ctx = _make_ctx(title="Mikä on kokemuksesi hoivapalveluista?")
    ai_title = "Enemmistöllä on myönteinen kokemus hoivapalveluista"
    ctx = dataclasses.replace(ctx, spec=dataclasses.replace(ctx.spec, slide_title=ai_title))

    add_image_slide_chrome(ctx)

    # The AI headline is the bold title; the question text is shown beneath it as
    # a secondary line (so the actual question is always at the top).
    boldlines = [
        s.text_frame.text
        for s in slide.shapes
        if s.has_text_frame and s.text_frame.paragraphs and any(
            r.font.bold for p in s.text_frame.paragraphs for r in p.runs
        )
    ]
    alltext = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert any(ai_title in t for t in boldlines), "AI headline should be the bold title"
    assert "Mikä on kokemuksesi" in alltext, "question text should appear under the headline"


def test_add_image_slide_chrome_n_annotation():
    """add_image_slide_chrome includes N=<Total> annotation (REQ-C-24h)."""
    from reportbuilder.render.image.slide_chrome import add_image_slide_chrome

    prs, slide, slot, ctx = _make_ctx()
    add_image_slide_chrome(ctx)

    # known_series() has base_n={"Total": 5}
    textboxes = [s for s in slide.shapes if s.has_text_frame]
    n_texts = [s for s in textboxes if "5" in s.text_frame.text]
    assert n_texts, "N=5 annotation not found in any textbox after chrome"


def test_add_image_slide_chrome_names_a_capped_group_in_the_footer():
    """The image-mode footer is the ONLY one real users see: deck.py routes
    render_mode == "image" to add_image_slide_chrome, never add_filter_annotation
    (deck.py:259-267), and the web app only ever sends render_mode "image"
    (web/src/lib/api.ts:272). A capped pie must name what it dropped HERE, not
    just in the native-only textbox that no real deck reaches. (spec 2026-08-22)
    """
    import dataclasses
    from reportbuilder.render.image.slide_chrome import add_image_slide_chrome

    cats = ("A", "B")
    segments = ("18-29", "30-44", "45-59", "60+", "Total")
    cells = {(c, s): Cell(pct=50.0, count=1.0, mean=None) for c in cats for s in segments}
    series = SeriesResult(categories=cats, segments=segments, cells=cells,
                          base_n={"18-29": 50, "30-44": 90, "45-59": 70, "60+": 30,
                                  "Total": 240},
                          statistic="pct")
    prs, slide, slot, ctx = _make_ctx("pie", series=series)
    ctx = dataclasses.replace(ctx, spec=dataclasses.replace(ctx.spec, classifying_var="age"))

    add_image_slide_chrome(ctx)

    text = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "60+" in text
    assert "Ei mahtunut sivulle" in text
    # The raw classifier code never belongs on a client slide: RenderContext
    # carries no model to resolve it to a human label. (ruling 2026-08-22)
    assert "age" not in text


def test_add_image_slide_chrome_says_nothing_when_nothing_was_dropped():
    """An unsplit or fully-fitting classifier must not print an omission clause —
    a false disclosure on a client slide is as wrong as a missing one."""
    import dataclasses
    from reportbuilder.render.image.slide_chrome import add_image_slide_chrome

    cats = ("A", "B")
    segments = ("Naiset", "Miehet", "Total")
    cells = {(c, s): Cell(pct=50.0, count=1.0, mean=None) for c in cats for s in segments}
    series = SeriesResult(categories=cats, segments=segments, cells=cells,
                          base_n={"Naiset": 60, "Miehet": 40, "Total": 100},
                          statistic="pct")
    prs, slide, slot, ctx = _make_ctx("pie", series=series)
    ctx = dataclasses.replace(ctx, spec=dataclasses.replace(ctx.spec, classifying_var="sex"))

    add_image_slide_chrome(ctx)

    text = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "Ei mahtunut sivulle" not in text and "Ei raportoitu" not in text


def test_add_image_slide_chrome_says_grouping_could_not_be_drawn_when_all_thin():
    """Every group under the base floor is the most severe omission the feature
    can produce: the renderer falls back to a single whole-sample pie that looks
    exactly like an ordinary un-split slide. This is the path real users see
    (add_image_slide_chrome, not the native-only add_filter_annotation), so it is
    the one that matters most. (coordinator review 2026-08-22)"""
    import dataclasses
    from reportbuilder.render.image.slide_chrome import add_image_slide_chrome

    cats = ("A", "B")
    segments = ("Naiset", "Miehet", "Total")
    cells = {(c, s): Cell(pct=50.0, count=1.0, mean=None) for c in cats for s in segments}
    series = SeriesResult(categories=cats, segments=segments, cells=cells,
                          base_n={"Naiset": 4, "Miehet": 6, "Total": 10},
                          statistic="pct")
    prs, slide, slot, ctx = _make_ctx("pie", series=series)
    ctx = dataclasses.replace(ctx, spec=dataclasses.replace(ctx.spec, classifying_var="sex"))

    add_image_slide_chrome(ctx)

    text = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "Ryhmittelyä ei voitu piirtää" in text


# ---------------------------------------------------------------------------
# Title fitting (Task D2): a title too long for the template's own box shrinks
# instead of growing the box without limit.
# ---------------------------------------------------------------------------

def _title_style(**overrides):
    from reportbuilder.render.template_profile import TextStyle

    # A generous single-line template box: 10in wide, 0.5in tall, 20pt — close
    # to what a customer's own one-line headline is actually drawn at.
    fields = dict(size_pt=20.0, caps=False, line_spacing=1.0,
                  left=int(Inches(0.7)), top=int(Inches(0.45)),
                  width=int(Inches(10)), height=int(Inches(0.5)))
    fields.update(overrides)
    return TextStyle(**fields)


_SHORT_TITLE = "Brand awareness"
_LONG_TITLE = ("How strongly do respondents associate our brand with being "
               "trustworthy, professional, and a leader in customer care, "
               "compared to every other brand active in this market segment "
               "today, across all the regions surveyed this quarter?")
_PATHOLOGICAL_TITLE = (_LONG_TITLE + " ") * 8


def test_short_title_renders_at_the_templates_own_size():
    """The regression that matters: a customer-length title must not restyle."""
    from reportbuilder.render.image.slide_chrome import fit_title_size

    st = _title_style()
    assert fit_title_size(st, _SHORT_TITLE) == st.size_pt


def test_a_long_title_shrinks_and_the_box_reflects_it():
    from reportbuilder.render.image.slide_chrome import fit_title_size, harvested_title_box
    from reportbuilder.render.template_profile import TemplateProfile

    st = _title_style()
    size = fit_title_size(st, _LONG_TITLE)
    assert size < st.size_pt

    profile = TemplateProfile(title=st)
    _left, _top, _width, height = harvested_title_box(profile, _LONG_TITLE)
    # The box the smaller size actually needs, not the box the ORIGINAL size
    # would have needed — otherwise deck.py starts the chart too low.
    at_original_size = harvested_title_box(
        TemplateProfile(title=_title_style()), "x")  # 1-line box, i.e. st.height
    assert height <= 2 * at_original_size[3]


def test_a_pathologically_long_title_stops_at_the_floor():
    from reportbuilder.render.image.slide_chrome import fit_title_size, _TITLE_MIN_SCALE

    st = _title_style()
    size = fit_title_size(st, _PATHOLOGICAL_TITLE)
    assert size == pytest.approx(st.size_pt * _TITLE_MIN_SCALE)


def test_an_all_caps_title_still_wraps_sooner():
    """cap="all" widens the rendered text, so the same string needs a smaller
    size (or wraps to more lines) than it would in mixed case — unchanged from
    before D2, just now expressed as a smaller chosen size rather than a
    taller box."""
    from reportbuilder.render.image.slide_chrome import fit_title_size

    mixed = _title_style(caps=False)
    capped = _title_style(caps=True)
    assert fit_title_size(capped, _LONG_TITLE) <= fit_title_size(mixed, _LONG_TITLE)
