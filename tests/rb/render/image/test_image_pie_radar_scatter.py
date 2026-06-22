"""TDD tests for image-mode pie, doughnut, radar, scatter builders (Task 5.12).

Covers: build_image_pie, build_image_doughnut, build_image_radar, build_image_scatter.

Each test asserts:
- Exactly ONE MSO_SHAPE_TYPE.PICTURE shape is on the slide.
- The picture width/height equals the slot width/height.
- The PNG blob is non-empty (or PIL can open it if available).
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from reportbuilder.render.base import RenderContext, Slot, StyleSpec
from reportbuilder.model.report import ChartSpec, ElementToggles, NumberFormat, SortSpec
from reportbuilder.stats.series import Cell, SeriesResult


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _make_slot() -> Slot:
    return Slot(
        slide_index=0,
        left=Inches(1),
        top=Inches(1),
        width=Inches(8),
        height=Inches(5),
        name="slot1",
    )


def _make_spec(chart_type: str, scatter_xy=None) -> ChartSpec:
    return ChartSpec(
        question_ref="q1",
        chart_type=chart_type,
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"),
        template_slot="slot1",
        elements=ElementToggles(
            title=True,
            legend=True,
            data_labels=True,
            axis_names=True,
        ),
        scatter_xy=scatter_xy,
    )


def _single_segment_series() -> SeriesResult:
    """Single segment (Total) with Yes/No categories — for pie/doughnut."""
    cats = ("Yes", "No")
    segs = ("Total",)
    cells = {
        ("Yes", "Total"): Cell(pct=60.0, count=60.0, mean=None),
        ("No", "Total"): Cell(pct=40.0, count=40.0, mean=None),
    }
    return SeriesResult(
        categories=cats,
        segments=segs,
        cells=cells,
        base_n={"Total": 100},
        statistic="pct",
    )


def _multi_segment_series() -> SeriesResult:
    """Multi-segment fixture for radar builder."""
    cats = ("Attr A", "Attr B", "Attr C")
    segs = ("Group X", "Group Y")
    cells = {
        ("Attr A", "Group X"): Cell(pct=70.0, count=70.0, mean=None),
        ("Attr A", "Group Y"): Cell(pct=50.0, count=50.0, mean=None),
        ("Attr B", "Group X"): Cell(pct=80.0, count=80.0, mean=None),
        ("Attr B", "Group Y"): Cell(pct=60.0, count=60.0, mean=None),
        ("Attr C", "Group X"): Cell(pct=55.0, count=55.0, mean=None),
        ("Attr C", "Group Y"): Cell(pct=45.0, count=45.0, mean=None),
    }
    return SeriesResult(
        categories=cats,
        segments=segs,
        cells=cells,
        base_n={"Group X": 100, "Group Y": 100},
        statistic="pct",
    )


def _scatter_series() -> SeriesResult:
    """Two-segment x/y series for scatter builder."""
    cats = ("Item 1", "Item 2", "Item 3", "Item 4")
    segs = ("x_score", "y_score")
    cells = {
        ("Item 1", "x_score"): Cell(pct=10.0, count=10.0, mean=None),
        ("Item 1", "y_score"): Cell(pct=20.0, count=20.0, mean=None),
        ("Item 2", "x_score"): Cell(pct=30.0, count=30.0, mean=None),
        ("Item 2", "y_score"): Cell(pct=40.0, count=40.0, mean=None),
        ("Item 3", "x_score"): Cell(pct=50.0, count=50.0, mean=None),
        ("Item 3", "y_score"): Cell(pct=60.0, count=60.0, mean=None),
        ("Item 4", "x_score"): Cell(pct=70.0, count=70.0, mean=None),
        ("Item 4", "y_score"): Cell(pct=80.0, count=80.0, mean=None),
    }
    return SeriesResult(
        categories=cats,
        segments=segs,
        cells=cells,
        base_n={"x_score": 4, "y_score": 4},
        statistic="pct",
    )


def _make_ctx(chart_type: str, series: SeriesResult, scatter_xy=None):
    prs, slide = _make_slide()
    slot = _make_slot()
    style = StyleSpec()
    spec = _make_spec(chart_type, scatter_xy=scatter_xy)
    ctx = RenderContext(
        slide=slide,
        slot=slot,
        style=style,
        spec=spec,
        series=series,
        fmt=spec.number_format,
    )
    return prs, slide, slot, ctx


def _assert_picture(slide, slot):
    """Assert exactly one PICTURE shape with correct slot dimensions."""
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1, f"Expected 1 PICTURE shape, found {len(pics)}"
    pic = pics[0]
    assert pic.width == slot.width, f"Picture width {pic.width} != slot width {slot.width}"
    assert pic.height == slot.height, f"Picture height {pic.height} != slot height {slot.height}"

    # Try PIL first; fall back to checking blob
    try:
        import io
        from PIL import Image
        img = Image.open(io.BytesIO(pic.image.blob))
        w, h = img.size
        assert w > 0 and h > 0, "PIL image size must be positive"
    except ImportError:
        assert len(pic.image.blob) > 0, "Picture blob must be non-empty"


# ---------------------------------------------------------------------------
# Registry test
# ---------------------------------------------------------------------------

def test_registry_has_pie_doughnut_radar_scatter():
    """IMAGE_BUILDERS must include the 4 task-5.12 keys."""
    from reportbuilder.render.image import IMAGE_BUILDERS
    expected = {"pie", "doughnut", "radar", "scatter"}
    assert expected.issubset(set(IMAGE_BUILDERS)), (
        f"Missing keys: {expected - set(IMAGE_BUILDERS)}"
    )


# ---------------------------------------------------------------------------
# build_image_pie
# ---------------------------------------------------------------------------

def test_image_pie_places_picture():
    """build_image_pie places exactly one slot-sized PICTURE shape."""
    from reportbuilder.render.image.pie import build_image_pie

    prs, slide, slot, ctx = _make_ctx("pie", _single_segment_series())
    result = build_image_pie(ctx)
    assert result is None, "builder should return None"
    _assert_picture(slide, slot)


# ---------------------------------------------------------------------------
# build_image_doughnut
# ---------------------------------------------------------------------------

def test_image_doughnut_places_picture():
    """build_image_doughnut places exactly one slot-sized PICTURE shape."""
    from reportbuilder.render.image.pie import build_image_doughnut

    prs, slide, slot, ctx = _make_ctx("doughnut", _single_segment_series())
    result = build_image_doughnut(ctx)
    assert result is None, "builder should return None"
    _assert_picture(slide, slot)


# ---------------------------------------------------------------------------
# build_image_radar
# ---------------------------------------------------------------------------

def test_image_radar_places_picture():
    """build_image_radar places exactly one slot-sized PICTURE shape (multi-segment)."""
    from reportbuilder.render.image.radar import build_image_radar

    prs, slide, slot, ctx = _make_ctx("radar", _multi_segment_series())
    result = build_image_radar(ctx)
    assert result is None, "builder should return None"
    _assert_picture(slide, slot)


# ---------------------------------------------------------------------------
# build_image_scatter
# ---------------------------------------------------------------------------

def test_image_scatter_places_picture():
    """build_image_scatter places exactly one slot-sized PICTURE shape."""
    from reportbuilder.render.image.scatter import build_image_scatter

    prs, slide, slot, ctx = _make_ctx(
        "scatter",
        _scatter_series(),
        scatter_xy=("x_score", "y_score"),
    )
    result = build_image_scatter(ctx)
    assert result is None, "builder should return None"
    _assert_picture(slide, slot)


def test_build_image_scatter_requires_xy():
    """build_image_scatter raises ValueError when scatter_xy is None."""
    from reportbuilder.render.image.scatter import build_image_scatter

    prs, slide, slot, ctx = _make_ctx("scatter", _scatter_series(), scatter_xy=None)
    with pytest.raises(ValueError, match="scatter"):
        build_image_scatter(ctx)
