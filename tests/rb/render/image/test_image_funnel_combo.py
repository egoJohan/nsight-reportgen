"""TDD tests for image-mode funnel and combo chart builders (Task 5.13).

Covers: build_image_funnel, build_image_combo, IMAGE_BUILDERS all-eleven check.

Each test asserts:
- Exactly ONE MSO_SHAPE_TYPE.PICTURE shape is on the slide.
- The picture width/height equals the slot width/height.
- The PNG blob is non-empty (PIL check if available).
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from reportbuilder.render.base import RenderContext, Slot, StyleSpec
from reportbuilder.model.report import ChartSpec, ElementToggles, NumberFormat, SortSpec
from reportbuilder.stats.series import Cell, SeriesResult


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _make_slot() -> Slot:
    return Slot(
        slide_index=0,
        left=Inches(1),
        top=Inches(1),
        width=Inches(8),
        height=Inches(5),
        name="slot1",
    )


def _make_spec(chart_type: str) -> ChartSpec:
    return ChartSpec(
        question_ref="q1",
        chart_type=chart_type,
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"),
        template_slot="slot1",
        elements=ElementToggles(
            title=True,
            legend=True,
            data_labels=True,
        ),
    )


def _funnel_series() -> SeriesResult:
    """4-category single-segment series for funnel."""
    cats = ("Awareness", "Interest", "Consideration", "Purchase")
    segs = ("Total",)
    cells = {
        ("Awareness", "Total"): Cell(pct=100.0, count=1000.0, mean=None),
        ("Interest", "Total"): Cell(pct=60.0, count=600.0, mean=None),
        ("Consideration", "Total"): Cell(pct=30.0, count=300.0, mean=None),
        ("Purchase", "Total"): Cell(pct=10.0, count=100.0, mean=None),
    }
    return SeriesResult(
        categories=cats,
        segments=segs,
        cells=cells,
        base_n={"Total": 1000},
        statistic="pct",
    )


def _combo_series() -> SeriesResult:
    """2-segment series for combo (bars + line)."""
    cats = ("Q1", "Q2", "Q3", "Q4")
    segs = ("Revenue", "Growth")
    cells = {
        ("Q1", "Revenue"): Cell(pct=25.0, count=250.0, mean=None),
        ("Q2", "Revenue"): Cell(pct=30.0, count=300.0, mean=None),
        ("Q3", "Revenue"): Cell(pct=28.0, count=280.0, mean=None),
        ("Q4", "Revenue"): Cell(pct=35.0, count=350.0, mean=None),
        ("Q1", "Growth"): Cell(pct=5.0, count=50.0, mean=None),
        ("Q2", "Growth"): Cell(pct=20.0, count=200.0, mean=None),
        ("Q3", "Growth"): Cell(pct=-7.0, count=-70.0, mean=None),
        ("Q4", "Growth"): Cell(pct=25.0, count=250.0, mean=None),
    }
    return SeriesResult(
        categories=cats,
        segments=segs,
        cells=cells,
        base_n={"Revenue": 1000, "Growth": 1000},
        statistic="pct",
    )


def _make_ctx(chart_type: str, series: SeriesResult) -> tuple:
    prs, slide = _make_slide()
    slot = _make_slot()
    style = StyleSpec()
    spec = _make_spec(chart_type)
    ctx = RenderContext(
        slide=slide,
        slot=slot,
        style=style,
        spec=spec,
        series=series,
        fmt=spec.number_format,
    )
    return prs, slide, slot, ctx


def _assert_picture(slide, slot):
    """Assert exactly one PICTURE shape with correct slot dimensions and valid PNG."""
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1, f"Expected 1 PICTURE shape, found {len(pics)}"
    pic = pics[0]
    assert pic.width == slot.width, f"Picture width {pic.width} != slot width {slot.width}"
    assert pic.height == slot.height, f"Picture height {pic.height} != slot height {slot.height}"

    try:
        import io
        from PIL import Image
        img = Image.open(io.BytesIO(pic.image.blob))
        w, h = img.size
        assert w > 0 and h > 0, "PIL image size must be positive"
    except ImportError:
        assert len(pic.image.blob) > 0, "Picture blob must be non-empty"


# ---------------------------------------------------------------------------
# test_image_funnel_real
# ---------------------------------------------------------------------------

def test_image_funnel_real():
    """build_image_funnel places exactly one slot-sized PICTURE with a valid PNG."""
    from reportbuilder.render.image.funnel import build_image_funnel

    prs, slide, slot, ctx = _make_ctx("funnel", _funnel_series())
    result = build_image_funnel(ctx)
    assert result is None, "builder should return None"
    _assert_picture(slide, slot)


# ---------------------------------------------------------------------------
# test_image_combo
# ---------------------------------------------------------------------------

def test_image_combo():
    """build_image_combo (2-segment) places exactly one slot-sized PICTURE with a valid PNG."""
    from reportbuilder.render.image.combo import build_image_combo

    prs, slide, slot, ctx = _make_ctx("combo", _combo_series())
    result = build_image_combo(ctx)
    assert result is None, "builder should return None"
    _assert_picture(slide, slot)


# ---------------------------------------------------------------------------
# test_image_builders_all_eleven  (REQ-C-13)
# ---------------------------------------------------------------------------

def test_image_builders_all_eleven():
    """IMAGE_BUILDERS must contain exactly the 11 canonical chart_type ids (REQ-C-13)."""
    from reportbuilder.render.image import IMAGE_BUILDERS

    expected = {
        "line",
        "pie",
        "vertical_bar",
        "stacked_vertical_bar",
        "horizontal_bar",
        "stacked_horizontal_bar",
        "radar",
        "doughnut",
        "scatter",
        "funnel",
        "combo",
    }
    assert set(IMAGE_BUILDERS) == expected, (
        f"Missing: {expected - set(IMAGE_BUILDERS)}; "
        f"Extra: {set(IMAGE_BUILDERS) - expected}"
    )
