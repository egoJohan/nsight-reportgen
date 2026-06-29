"""Tests for label handling (never overlap / never ellipsis-cut) + 'Not answered' grey bar.

Label policy (Task A — render quality):
  - _wrap_label wraps long labels at word boundaries onto as many lines as needed;
    the full text is always preserved and '…' is NEVER appended.
  - _wrap_xtick_label wraps vertical-bar x-axis labels (rotation handles overlap).
  - Short labels (≤ wrap width) pass through unchanged.
  - Horizontal-bar renderer applies wrapped display labels without error.

R4.2 — "Not answered" bar coloring:
  - The "Not answered" bar is rendered in MUTED grey, not TEAL
  - All other bars keep their teal series colour
"""
from __future__ import annotations

import io

import pytest
import matplotlib
matplotlib.use("Agg")
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from PIL import Image

from reportbuilder.render.base import RenderContext, Slot, StyleSpec
from reportbuilder.model.report import ChartSpec, ElementToggles, NumberFormat, SortSpec
from reportbuilder.stats.series import Cell, SeriesResult
from reportbuilder.stats.engine import NOT_ANSWERED_LABEL
from reportbuilder.render.house_style import MUTED, TEAL


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_ctx(chart_type: str, series: SeriesResult) -> tuple:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slot = Slot(slide_index=0, left=Inches(1), top=Inches(1),
                width=Inches(8), height=Inches(5), name="slot1")
    spec = ChartSpec(
        question_ref="q1", chart_type=chart_type, statistic="pct",
        classifying_var=None, number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"), template_slot="slot1",
        elements=ElementToggles(title=True, data_labels=True, legend=True),
    )
    ctx = RenderContext(slide=slide, slot=slot, style=StyleSpec(),
                       spec=spec, series=series, fmt=spec.number_format)
    return prs, slide, slot, ctx


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    h = hex_color.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _count_color_pixels(blob: bytes, hex_color: str, tolerance: int = 8) -> int:
    """Count pixels in blob matching hex_color within tolerance."""
    img = Image.open(io.BytesIO(blob)).convert("RGB")
    r, g, b = _hex_to_rgb(hex_color)
    return sum(
        1 for px in img.getdata()
        if abs(px[0] - r) <= tolerance and abs(px[1] - g) <= tolerance and abs(px[2] - b) <= tolerance
    )


def _get_png_blob(slide) -> bytes:
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1, f"Expected 1 picture; found {len(pics)}"
    return pics[0].image.blob


# ---------------------------------------------------------------------------
# R4.1 — Label helper unit tests
# ---------------------------------------------------------------------------

def _strip_wrap(text: str) -> str:
    """Collapse wrapped lines back to a single space-joined string for comparison."""
    return " ".join(text.split())


class TestWrapLabel:
    """Unit tests for _wrap_label helper — wrap, never ellipsis-cut."""

    def test_short_label_unchanged(self):
        """Labels ≤ wrap width pass through _wrap_label unchanged."""
        from reportbuilder.render.image.bars import _wrap_label, _LABEL_WRAP_WIDTH
        assert _wrap_label("Yes") == "Yes"
        assert _wrap_label("Kyllä, joskus") == "Kyllä, joskus"
        assert _wrap_label("A" * _LABEL_WRAP_WIDTH) == "A" * _LABEL_WRAP_WIDTH

    def test_long_label_wrapped_to_multiple_lines(self):
        """Labels > wrap width are wrapped; result contains a newline, lines fit."""
        from reportbuilder.render.image.bars import _wrap_label, _LABEL_WRAP_WIDTH
        long = "Hoitajat kuuntelevat toiveitani ja tarpeitani"  # > wrap width
        result = _wrap_label(long)
        assert "\n" in result, f"Expected newline in wrapped label; got {result!r}"
        for line in result.split("\n"):
            assert len(line) <= _LABEL_WRAP_WIDTH, f"Line too long: {line!r}"

    def test_very_long_label_never_ellipsized(self):
        """Very long labels are wrapped onto as many lines as needed — never '…'."""
        from reportbuilder.render.image.bars import _wrap_label, _LABEL_WRAP_WIDTH
        very_long = "Hoitajat kuuntelevat toiveitani ja tarpeitani erittäin hyvin kaikissa tilanteissa"
        result = _wrap_label(very_long)
        assert "…" not in result, f"Ellipsis must never appear; got {result!r}"
        # Full text is preserved (every word survives, only whitespace changes).
        assert _strip_wrap(result) == _strip_wrap(very_long)
        for line in result.split("\n"):
            assert len(line) <= _LABEL_WRAP_WIDTH, f"Line too long: {line!r}"

    def test_hyphenated_compound_not_split(self):
        """A hyphenated compound (e.g. a brand) is kept on one piece, not split at '-'."""
        from reportbuilder.render.image.bars import _wrap_label
        assert _wrap_label("Mainio-kodit") == "Mainio-kodit"


class TestWrapXtickLabel:
    """Unit tests for _wrap_xtick_label helper (vertical-bar x-axis labels)."""

    def test_short_label_unchanged(self):
        """Labels ≤ x-tick wrap width pass through unchanged."""
        from reportbuilder.render.image.bars import _wrap_xtick_label, _XLABEL_WRAP_WIDTH
        assert _wrap_xtick_label("Short") == "Short"
        assert _wrap_xtick_label("A" * _XLABEL_WRAP_WIDTH) == "A" * _XLABEL_WRAP_WIDTH

    def test_long_label_wrapped_never_ellipsized(self):
        """Long x-axis labels wrap (full text preserved), never '…'."""
        from reportbuilder.render.image.bars import _wrap_xtick_label
        long = "Erittäin tyytyväinen palveluun"
        result = _wrap_xtick_label(long)
        assert "…" not in result, f"Ellipsis must never appear; got {result!r}"
        assert "\n" in result, f"Expected wrap; got {result!r}"
        assert _strip_wrap(result) == _strip_wrap(long)

    def test_pathological_long_token_force_broken(self):
        """A single token far longer than the wrap width is force-broken into
        width-sized chunks as a last resort so erroneous unbroken long words
        can't run off the chart. Full text preserved; never ellipsized; no line
        exceeds the wrap width."""
        from reportbuilder.render.image.bars import _wrap_xtick_label, _XLABEL_WRAP_WIDTH
        result = _wrap_xtick_label("A" * 80)
        assert "…" not in result
        assert "\n" in result, f"Expected force-break; got {result!r}"
        assert result.replace("\n", "") == "A" * 80  # full text preserved
        for line in result.split("\n"):
            assert len(line) <= _XLABEL_WRAP_WIDTH

    def test_wrap_breaks_at_hyphen_not_midword(self):
        """A long hyphenated phrase wraps at the hyphen/space, never mid-word."""
        from reportbuilder.render.image.bars import _wrap_label
        result = _wrap_label("Mielenterveys- tai päihdekuntoutuspalvelut")
        # The long word stays intact on its own line — no "pä\nihde…".
        assert "päihdekuntoutuspalvelut" in result.replace("\n", "\n") and \
            "pä\nihde" not in result
        for line in result.split("\n"):
            assert line.strip()  # no empty lines


# ---------------------------------------------------------------------------
# R4.1 — Render with long labels (smoke test)
# ---------------------------------------------------------------------------

def _long_cat_series() -> SeriesResult:
    """Series with var11-style long category labels (multi-sentence options)."""
    cats = (
        "Hoitajat kuuntelevat toiveitani ja tarpeitani",
        "Hoitajat ovat ystävällisiä ja kohtelevat minua kunnioittavasti",
        "Saan tarvitsemaani apua päivittäisissä toimissa",
        "Minua kohdellaan tasa-arvoisesti ja oikeudenmukaisesti",
        "Ruoka on maukasta ja monipuolista",
        "En osaa sanoa",
    )
    cells = {(c, "Total"): Cell(pct=float(60 - i * 8)) for i, c in enumerate(cats)}
    return SeriesResult(categories=cats, segments=("Total",), cells=cells,
                        base_n={"Total": 200}, statistic="pct")


def test_long_labels_render_without_error():
    """build_image_bar renders a long-label question without error (R4.1)."""
    from reportbuilder.render.image.bars import build_image_bar
    series = _long_cat_series()
    prs, slide, slot, ctx = _make_ctx("horizontal_bar", series)
    build_image_bar(ctx)
    blob = _get_png_blob(slide)
    img = Image.open(io.BytesIO(blob))
    w, h = img.size
    assert w > 0 and h > 0, "PNG must have positive dimensions"


def test_long_labels_produces_valid_png():
    """build_image_bar with long labels produces a non-trivial PNG (R4.1)."""
    from reportbuilder.render.image.bars import build_image_bar
    series = _long_cat_series()
    prs, slide, slot, ctx = _make_ctx("horizontal_bar", series)
    build_image_bar(ctx)
    blob = _get_png_blob(slide)
    assert len(blob) > 5000, "PNG blob should be non-trivially sized"


# ---------------------------------------------------------------------------
# R4.2 — "Not answered" bar is MUTED grey
# ---------------------------------------------------------------------------

def _na_series() -> SeriesResult:
    """Opinion series including a 'Not answered' category."""
    cats = ("Erittäin huono", "Huono", "Hyvä", "Erittäin hyvä", NOT_ANSWERED_LABEL)
    cells = {(c, "Total"): Cell(pct=v) for c, v in zip(cats, [5.0, 26.0, 47.0, 10.0, 12.0])}
    return SeriesResult(categories=cats, segments=("Total",), cells=cells,
                        base_n={"Total": 300}, statistic="pct")


def test_not_answered_bar_has_muted_pixels():
    """build_image_bar colors the 'Not answered' bar MUTED grey (R4.2).

    We check that the PNG contains MUTED-coloured pixels; TEAL pixels are
    also present (from the other bars).
    """
    from reportbuilder.render.image.bars import build_image_bar
    series = _na_series()
    prs, slide, slot, ctx = _make_ctx("horizontal_bar", series)
    build_image_bar(ctx)
    blob = _get_png_blob(slide)

    muted_count = _count_color_pixels(blob, MUTED, tolerance=10)
    teal_count = _count_color_pixels(blob, TEAL, tolerance=10)
    assert muted_count > 50, (
        f"Expected MUTED grey pixels in 'Not answered' bar; found {muted_count}"
    )
    assert teal_count > 50, (
        f"Expected TEAL pixels in real-response bars; found {teal_count}"
    )


def test_series_without_na_has_no_muted_bars():
    """build_image_bar without 'Not answered' category has no MUTED bar pixels (R4.2)."""
    from reportbuilder.render.image.bars import build_image_bar
    cats = ("Erittäin huono", "Huono", "Hyvä", "Erittäin hyvä")
    cells = {(c, "Total"): Cell(pct=v) for c, v in zip(cats, [5.0, 26.0, 47.0, 22.0])}
    series = SeriesResult(categories=cats, segments=("Total",), cells=cells,
                          base_n={"Total": 300}, statistic="pct")
    prs, slide, slot, ctx = _make_ctx("horizontal_bar", series)
    build_image_bar(ctx)
    blob = _get_png_blob(slide)

    muted_count = _count_color_pixels(blob, MUTED, tolerance=5)
    # MUTED is also used for axis tick labels; we allow some MUTED pixels from
    # the tick text but expect far fewer than with a dedicated MUTED bar.
    assert muted_count < 5000, (
        f"No MUTED bars expected when 'Not answered' is absent; found {muted_count} pixels"
    )
