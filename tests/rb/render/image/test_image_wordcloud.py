"""Task J.2: image-mode word-cloud builder (build_image_wordcloud).

Asserts the builder places exactly one PICTURE shape on the slide, the saved PNG is
non-empty, has a cream-ish background, and that the layout is deterministic
(random_state=42) so two renders of the same series produce byte-identical images.
"""
from __future__ import annotations

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from reportbuilder.model.report import ChartSpec, ElementToggles, NumberFormat, SortSpec
from reportbuilder.render.base import RenderContext, Slot, StyleSpec
from reportbuilder.render.house_style import CREAM
from reportbuilder.render.image import IMAGE_BUILDERS
from reportbuilder.render.image.wordcloud import (
    build_image_wordcloud,
    _resolve_font_path,
)
from reportbuilder.render.plugins import plugin
from reportbuilder.stats.series import Cell, SeriesResult


def _slot() -> Slot:
    return Slot(slide_index=0, left=Inches(1), top=Inches(1),
                width=Inches(8), height=Inches(5), name="slot1")


def _spec() -> ChartSpec:
    return ChartSpec(
        question_ref="q", chart_type="wordcloud", statistic="count",
        classifying_var=None, number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"), template_slot="slot1",
        elements=ElementToggles(),
    )


def _series() -> SeriesResult:
    cats = ("kallis", "huono", "luotettava", "hyvä", "ahne")
    counts = {"kallis": 185.0, "huono": 75.0, "luotettava": 61.0,
              "hyvä": 60.0, "ahne": 57.0}
    cells = {(c, "Total"): Cell(pct=None, count=counts[c], mean=None) for c in cats}
    return SeriesResult(categories=cats, segments=("Total",), cells=cells,
                        base_n={"Total": 817}, statistic="count")


def _ctx():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spec = _spec()
    ctx = RenderContext(slide=slide, slot=_slot(), style=StyleSpec(),
                        spec=spec, series=_series(), fmt=spec.number_format)
    return prs, slide, ctx


def test_wordcloud_registered_in_image_builders_and_plugin():
    assert "wordcloud" in IMAGE_BUILDERS
    p = plugin("wordcloud")
    assert p.label == "Word Cloud"
    assert p.image_build is build_image_wordcloud
    # Never auto-suggested for normal questions.
    assert p.suitability(None, _series()) is None


def test_build_image_wordcloud_places_one_picture():
    _prs, slide, ctx = _ctx()
    build_image_wordcloud(ctx)
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    assert pics[0].width > 0 and pics[0].height > 0


def test_resolve_font_path_returns_existing_ttf():
    import os
    fp = _resolve_font_path()
    assert fp and os.path.exists(fp)
    assert fp.lower().endswith(".ttf")


def test_wordcloud_png_is_cream_background(tmp_path):
    from reportbuilder.render.image.wordcloud import _TEAL_CLOUD
    from wordcloud import WordCloud

    sr = _series()
    freqs = {c: float(sr.cell(c, "Total").count) for c in sr.categories}
    ranked = sorted(freqs, key=lambda w: -freqs[w])
    rank = {w: i for i, w in enumerate(ranked)}
    n = len(ranked)

    def cf(word, **kw):
        r = rank.get(word, 0)
        idx = 0 if n <= 1 else int(round(r / (n - 1) * (len(_TEAL_CLOUD) - 1)))
        return _TEAL_CLOUD[idx]

    wc = WordCloud(background_color=CREAM, color_func=cf,
                   font_path=_resolve_font_path(), random_state=42,
                   prefer_horizontal=0.9, max_words=len(freqs),
                   width=800, height=500)
    wc.generate_from_frequencies(freqs)
    out = tmp_path / "wc.png"
    wc.to_file(str(out))
    assert out.stat().st_size > 0
    with Image.open(out) as im:
        im = im.convert("RGB")
        # A corner pixel should be the cream background colour (#F7F3EC).
        assert im.getpixel((1, 1)) == (0xF7, 0xF3, 0xEC)


def test_wordcloud_layout_is_deterministic(tmp_path):
    from reportbuilder.render.image.wordcloud import _TEAL_CLOUD
    from wordcloud import WordCloud

    sr = _series()
    freqs = {c: float(sr.cell(c, "Total").count) for c in sr.categories}

    def build():
        wc = WordCloud(background_color=CREAM, color_func=lambda *a, **k: _TEAL_CLOUD[0],
                       font_path=_resolve_font_path(), random_state=42,
                       prefer_horizontal=0.9, max_words=len(freqs),
                       width=800, height=500)
        wc.generate_from_frequencies(freqs)
        return wc.to_array().tobytes()

    assert build() == build()
