"""Task D1: chart furniture adapts to the template's background colour.

Covers:
- house_style.furniture_colors — the derived-colour formula itself.
- Image-mode builders (pie legend, bar axis/legend) actually USE the derived
  colours, and a light background still resolves to today's exact constants.
- The native (python-pptx) chart path (render/elements.py) sets legend/
  data-label/axis-tick-label COLOUR from the same derived ink, not just fonts.
"""
from __future__ import annotations

import io

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

from reportbuilder.model.report import ChartSpec, ElementToggles, NumberFormat, SortSpec
from reportbuilder.render.base import RenderContext, Slot, StyleSpec
from reportbuilder.render.house_style import (
    CREAM, GRIDC, INK, MUTED, contrast_ink, furniture_colors,
)
from reportbuilder.stats.series import Cell, SeriesResult
from reportbuilder.testing.fixtures import known_series

# A dark template background — chosen well below the 0.55 luminance threshold
# contrast_ink already uses, so it is unambiguously "dark" by the house rule.
_DARK_BG = "0B0B12"
# A light, non-CREAM template background (e.g. a plain-white template) — the
# constraint is "any light background", not "only CREAM".
_LIGHT_BG = "FFFFFF"


# ---------------------------------------------------------------------------
# 1 — house_style.furniture_colors
# ---------------------------------------------------------------------------

class TestFurnitureColors:
    def test_no_background_matches_house_constants_exactly(self):
        """No template opinion (the base StyleSpec case) → today's palette, unchanged."""
        assert furniture_colors("") == (INK, MUTED, GRIDC)

    def test_house_cream_matches_house_constants_exactly(self):
        assert furniture_colors(CREAM.lstrip("#")) == (INK, MUTED, GRIDC)

    def test_any_light_background_matches_house_constants_exactly(self):
        """Not just CREAM — a light TEMPLATE background (e.g. plain white) must
        resolve to exactly today's colours too (REQ: light templates unchanged)."""
        assert furniture_colors(_LIGHT_BG) == (INK, MUTED, GRIDC)
        assert furniture_colors("#" + _LIGHT_BG) == (INK, MUTED, GRIDC)

    def test_dark_background_flips_ink_to_white(self):
        ink, _muted, _grid = furniture_colors(_DARK_BG)
        assert ink == "#FFFFFF"

    def test_dark_background_grid_and_muted_differ_from_house_defaults(self):
        ink, muted, grid = furniture_colors(_DARK_BG)
        assert muted != MUTED
        assert grid != GRIDC

    def test_dark_background_grid_is_lighter_than_the_ground(self):
        """Brief: 'on a dark ground the grid must be lighter than the ground'."""
        from matplotlib.colors import to_rgb
        bg_lum = sum(to_rgb(f"#{_DARK_BG}")) / 3.0
        _ink, _muted, grid = furniture_colors(_DARK_BG)
        grid_lum = sum(to_rgb(grid)) / 3.0
        assert grid_lum > bg_lum

    def test_light_background_grid_is_darker_than_the_ground(self):
        """Mirror of the dark case: GRIDC sits darker than CREAM (unchanged today)."""
        from matplotlib.colors import to_rgb
        bg_lum = sum(to_rgb(CREAM)) / 3.0
        grid_lum = sum(to_rgb(GRIDC)) / 3.0
        assert grid_lum < bg_lum

    def test_ink_muted_grid_contrast_ordering_holds_on_dark(self):
        """ink has the strongest contrast against the ground, grid the weakest —
        same ordering the house palette already has on light (INK > MUTED > GRIDC
        distance from CREAM)."""
        from matplotlib.colors import to_rgb
        bg_lum = sum(to_rgb(f"#{_DARK_BG}")) / 3.0
        ink, muted, grid = furniture_colors(_DARK_BG)
        d_ink = abs(sum(to_rgb(ink)) / 3.0 - bg_lum)
        d_muted = abs(sum(to_rgb(muted)) / 3.0 - bg_lum)
        d_grid = abs(sum(to_rgb(grid)) / 3.0 - bg_lum)
        assert d_ink > d_muted > d_grid

    def test_agrees_with_contrast_ink_on_what_dark_means(self):
        """furniture_colors and contrast_ink must never disagree about a colour's
        darkness — same 0.55 luminance threshold."""
        assert contrast_ink(f"#{_DARK_BG}") == "#FFFFFF"
        assert furniture_colors(_DARK_BG)[0] == contrast_ink(f"#{_DARK_BG}")
        assert contrast_ink(f"#{_LIGHT_BG}") == INK
        assert furniture_colors(_LIGHT_BG)[0] == INK

    def test_accepts_hex_with_or_without_leading_hash(self):
        assert furniture_colors(_DARK_BG) == furniture_colors(f"#{_DARK_BG}")


# ---------------------------------------------------------------------------
# 2 — image-mode builders: pie legend paints the slide background, not a
# hardcoded CREAM rectangle, and its text/edge follow furniture_colors.
# ---------------------------------------------------------------------------

def _slot() -> Slot:
    return Slot(slide_index=0, left=Inches(1), top=Inches(1),
                width=Inches(8), height=Inches(5), name="slot1")


def _pie_series() -> SeriesResult:
    cats = ("Kyllä", "Ei", "En osaa sanoa")
    cells = {(c, "Total"): Cell(pct=v) for c, v in zip(cats, [50.0, 30.0, 20.0])}
    return SeriesResult(categories=cats, segments=("Total",), cells=cells,
                        base_n={"Total": 100}, statistic="pct")


def _pie_ctx(background: str):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style = StyleSpec()
    style.background = background  # base StyleSpec has no slots; this is fine
    spec = ChartSpec(
        question_ref="q1", chart_type="pie", statistic="pct", classifying_var=None,
        number_format=NumberFormat(), sort=SortSpec(basis="data_order"),
        template_slot="slot1",
        elements=ElementToggles(title=True, data_labels=True, legend=True, n=True),
    )
    sr = _pie_series()
    ctx = RenderContext(slide=slide, slot=_slot(), style=style, spec=spec,
                        series=sr, fmt=spec.number_format, title="t")
    return prs, slide, ctx


def _opaque_pixels(png_blob: bytes) -> list[tuple[int, int, int]]:
    from PIL import Image
    img = Image.open(io.BytesIO(png_blob))
    w, h = img.size
    out = []
    for y in range(0, h, 3):
        for x in range(0, w, 3):
            p = img.getpixel((x, y))
            if len(p) == 4 and p[3] > 200:
                out.append(p[:3])
    return out


def test_pie_legend_stays_cream_on_a_light_default_template():
    """Regression: no template opinion → the legend box is exactly house CREAM,
    byte for byte, as it always has been."""
    from reportbuilder.render.image.pie import build_image_pie
    _prs, slide, ctx = _pie_ctx("")
    build_image_pie(ctx)

    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    pixels = _opaque_pixels(pics[0].image.blob)
    cream_rgb = (0xF7, 0xF3, 0xEC)
    assert any(p == cream_rgb for p in pixels), (
        "Expected the legend box to be exactly house CREAM on a light/default template"
    )


def test_pie_legend_matches_dark_template_background_not_cream():
    """A dark template background: the legend box paints the SLIDE's own dark
    background (not a stray cream rectangle), and its text is white, not INK."""
    from reportbuilder.render.image.pie import build_image_pie
    _prs, slide, ctx = _pie_ctx(_DARK_BG)
    build_image_pie(ctx)

    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    pixels = _opaque_pixels(pics[0].image.blob)
    cream_rgb = (0xF7, 0xF3, 0xEC)
    dark_rgb = (0x0B, 0x0B, 0x12)
    assert not any(p == cream_rgb for p in pixels), (
        "A dark template must never paint a cream rectangle behind the legend"
    )
    assert any(p == dark_rgb for p in pixels), (
        "Expected the legend box to match the dark slide background"
    )
    # White (ink-on-dark) legend text/edges should be present too.
    assert any(p == (255, 255, 255) for p in pixels), (
        "Expected white legend text/edges on a dark template"
    )


# ---------------------------------------------------------------------------
# 3 — bar axis/tick furniture: _apply_bar_style / _apply_column_style pick
# up the derived ink/muted/grid rather than the bare house constants.
# ---------------------------------------------------------------------------

def _bar_ctx(background: str):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style = StyleSpec()
    style.background = background
    spec = ChartSpec(
        question_ref="q1", chart_type="horizontal_bar", statistic="pct",
        classifying_var=None, number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"), template_slot="slot1",
        elements=ElementToggles(title=True, data_labels=True, legend=True, n=True),
    )
    sr = known_series()
    ctx = RenderContext(slide=slide, slot=_slot(), style=style, spec=spec,
                        series=sr, fmt=spec.number_format, title="t")
    return prs, slide, ctx


def test_bar_axis_ticks_are_house_muted_on_light_default():
    import matplotlib
    matplotlib.use("Agg")
    from reportbuilder.render.image._mpl import new_figure
    from reportbuilder.render.image.bars import _apply_bar_style

    _prs, _slide, ctx = _bar_ctx("")
    fig, ax = new_figure(ctx)
    _apply_bar_style(ax, ctx, 100.0, "pct")
    tick = ax.get_xticklabels()[0]
    assert tick.get_color() == MUTED
    import matplotlib.pyplot as plt
    plt.close(fig)


def test_bar_axis_ticks_flip_to_light_on_dark_background():
    import matplotlib
    matplotlib.use("Agg")
    from reportbuilder.render.image._mpl import new_figure
    from reportbuilder.render.image.bars import _apply_bar_style

    _prs, _slide, ctx = _bar_ctx(_DARK_BG)
    fig, ax = new_figure(ctx)
    _apply_bar_style(ax, ctx, 100.0, "pct")
    tick = ax.get_xticklabels()[0]
    assert tick.get_color() != MUTED
    # Should read as light against the dark ground.
    from matplotlib.colors import to_rgb
    r, g, b = to_rgb(tick.get_color())
    assert (r + g + b) / 3.0 > 0.5
    import matplotlib.pyplot as plt
    plt.close(fig)


# ---------------------------------------------------------------------------
# 4 — native (python-pptx) chart path: apply_elements sets legend/data-label/
# axis-tick-label COLOUR, not just fonts, from the same derived ink.
# ---------------------------------------------------------------------------

def _make_column_chart(slide):
    chart_data = CategoryChartData()
    chart_data.categories = ["Cat A", "Cat B"]
    chart_data.add_series("Series 1", (30.0, 70.0))
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1), Inches(6), Inches(4),
        chart_data,
    )
    return gf.chart


def _native_ctx(background: str) -> tuple:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart = _make_column_chart(slide)
    style = StyleSpec()
    style.background = background
    spec = ChartSpec(
        question_ref="q1", chart_type="vertical_bar", statistic="pct",
        classifying_var=None, number_format=NumberFormat(pct_decimals=0),
        sort=SortSpec(basis="data_order"), template_slot="slot1",
        elements=ElementToggles(title=True, legend=True, n=True, axis_names=True,
                                filter_var=True, data_labels=True),
    )
    ctx = RenderContext(slide=slide, slot=_slot(), style=style, spec=spec,
                        series=known_series(), fmt=spec.number_format)
    return chart, ctx


def test_native_chart_text_is_house_ink_on_light_default():
    from reportbuilder.render.elements import apply_elements

    chart, ctx = _native_ctx("")
    apply_elements(chart, ctx, title="Q")

    expected = RGBColor.from_string(INK.lstrip("#"))
    assert chart.legend.font.color.rgb == expected
    assert chart.plots[0].data_labels.font.color.rgb == expected
    assert chart.value_axis.tick_labels.font.color.rgb == expected
    assert chart.category_axis.tick_labels.font.color.rgb == expected


def test_native_chart_text_flips_to_white_on_dark_background():
    from reportbuilder.render.elements import apply_elements

    chart, ctx = _native_ctx(_DARK_BG)
    apply_elements(chart, ctx, title="Q")

    white = RGBColor.from_string("FFFFFF")
    assert chart.legend.font.color.rgb == white
    assert chart.plots[0].data_labels.font.color.rgb == white
    assert chart.value_axis.tick_labels.font.color.rgb == white
    assert chart.category_axis.tick_labels.font.color.rgb == white
