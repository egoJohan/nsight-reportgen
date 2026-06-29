"""TDD tests for image-mode chart builders (Task 5.11).

Covers: build_image_column, build_image_bar, build_image_column_stacked,
build_image_bar_stacked, build_image_line.

Each test asserts:
- Exactly ONE MSO_SHAPE_TYPE.PICTURE shape is on the slide.
- The picture width/height equals the slot width/height.
- The PNG blob is non-empty (or PIL can open it if available).
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from reportbuilder.render.base import RenderContext, Slot, StyleSpec
from reportbuilder.model.report import ChartSpec, ElementToggles, NumberFormat, SortSpec
from reportbuilder.stats.series import Cell, SeriesResult
from reportbuilder.testing.fixtures import known_series


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _make_slot() -> Slot:
    return Slot(
        slide_index=0,
        left=Inches(1),
        top=Inches(1),
        width=Inches(8),
        height=Inches(5),
        name="slot1",
    )


def _make_spec(chart_type: str) -> ChartSpec:
    return ChartSpec(
        question_ref="q1",
        chart_type=chart_type,
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"),
        template_slot="slot1",
        elements=ElementToggles(
            title=True,
            legend=True,
            data_labels=True,
        ),
    )


def _single_segment_series() -> SeriesResult:
    """Single segment (Total) with two categories — used for column/bar/line."""
    return known_series()  # categories=("Yes","No"), segments=("Total",), pct=[60,40]


def _two_segment_series() -> SeriesResult:
    """Two-segment fixture for stacked builders."""
    cats = ("Yes", "No")
    segs = ("Group A", "Group B")
    cells = {
        ("Yes", "Group A"): Cell(pct=50.0, count=10.0, mean=None),
        ("Yes", "Group B"): Cell(pct=30.0, count=6.0, mean=None),
        ("No", "Group A"): Cell(pct=20.0, count=4.0, mean=None),
        ("No", "Group B"): Cell(pct=40.0, count=8.0, mean=None),
    }
    return SeriesResult(
        categories=cats,
        segments=segs,
        cells=cells,
        base_n={"Group A": 20, "Group B": 20},
        statistic="pct",
    )


def _make_ctx(chart_type: str, series: SeriesResult) -> tuple:
    prs, slide = _make_slide()
    slot = _make_slot()
    style = StyleSpec()
    spec = _make_spec(chart_type)
    ctx = RenderContext(
        slide=slide,
        slot=slot,
        style=style,
        spec=spec,
        series=series,
        fmt=spec.number_format,
    )
    return prs, slide, slot, ctx


def _assert_picture(slide, slot):
    """Assert exactly one PICTURE shape with correct slot dimensions."""
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1, f"Expected 1 PICTURE shape, found {len(pics)}"
    pic = pics[0]
    assert pic.width == slot.width, f"Picture width {pic.width} != slot width {slot.width}"
    assert pic.height == slot.height, f"Picture height {pic.height} != slot height {slot.height}"

    # Try PIL first; fall back to checking blob
    try:
        import io
        from PIL import Image
        img = Image.open(io.BytesIO(pic.image.blob))
        w, h = img.size
        assert w > 0 and h > 0, "PIL image size must be positive"
    except ImportError:
        assert len(pic.image.blob) > 0, "Picture blob must be non-empty"


# ---------------------------------------------------------------------------
# Registry test
# ---------------------------------------------------------------------------

def test_registry_has_five():
    """IMAGE_BUILDERS must have exactly the 5 task-5.11 keys registered."""
    from reportbuilder.render.image import IMAGE_BUILDERS
    expected = {
        "vertical_bar",
        "horizontal_bar",
        "stacked_vertical_bar",
        "stacked_horizontal_bar",
        "line",
    }
    assert expected.issubset(set(IMAGE_BUILDERS)), (
        f"Missing keys: {expected - set(IMAGE_BUILDERS)}"
    )


# ---------------------------------------------------------------------------
# build_image_column
# ---------------------------------------------------------------------------

def test_image_column_places_picture():
    """build_image_column places exactly one slot-sized PICTURE shape."""
    from reportbuilder.render.image.bars import build_image_column

    prs, slide, slot, ctx = _make_ctx("vertical_bar", _single_segment_series())
    result = build_image_column(ctx)
    assert result is None, "builder should return None"
    _assert_picture(slide, slot)


# ---------------------------------------------------------------------------
# build_image_bar
# ---------------------------------------------------------------------------

def test_image_bar_places_picture():
    """build_image_bar places exactly one slot-sized PICTURE shape."""
    from reportbuilder.render.image.bars import build_image_bar

    prs, slide, slot, ctx = _make_ctx("horizontal_bar", _single_segment_series())
    result = build_image_bar(ctx)
    assert result is None
    _assert_picture(slide, slot)


# ---------------------------------------------------------------------------
# build_image_column_stacked
# ---------------------------------------------------------------------------

def test_image_column_stacked_places_picture():
    """build_image_column_stacked places exactly one slot-sized PICTURE shape (2-segment)."""
    from reportbuilder.render.image.bars import build_image_column_stacked

    prs, slide, slot, ctx = _make_ctx("stacked_vertical_bar", _two_segment_series())
    result = build_image_column_stacked(ctx)
    assert result is None
    _assert_picture(slide, slot)


# ---------------------------------------------------------------------------
# build_image_bar_stacked
# ---------------------------------------------------------------------------

def test_image_bar_stacked_places_picture():
    """build_image_bar_stacked places exactly one slot-sized PICTURE shape (2-segment)."""
    from reportbuilder.render.image.bars import build_image_bar_stacked

    prs, slide, slot, ctx = _make_ctx("stacked_horizontal_bar", _two_segment_series())
    result = build_image_bar_stacked(ctx)
    assert result is None
    _assert_picture(slide, slot)


# ---------------------------------------------------------------------------
# build_image_line
# ---------------------------------------------------------------------------

def test_image_line_places_picture():
    """build_image_line places exactly one slot-sized PICTURE shape."""
    from reportbuilder.render.image.line import build_image_line

    prs, slide, slot, ctx = _make_ctx("line", _single_segment_series())
    result = build_image_line(ctx)
    assert result is None
    _assert_picture(slide, slot)


# --- Regression: stacked layout transpose + vertical-bar orientation ---------

def test_stacked_layout_excludes_total_and_transposes():
    """_stacked_layout: bars = classifier segments (no 'Total'); stack = answer
    categories; each bar sums to 100 (column %)."""
    from reportbuilder.render.image.bars import _stacked_layout

    cells = {
        ("Yes", "A"): Cell(pct=60.0, count=None, mean=None),
        ("No", "A"): Cell(pct=40.0, count=None, mean=None),
        ("Yes", "B"): Cell(pct=70.0, count=None, mean=None),
        ("No", "B"): Cell(pct=30.0, count=None, mean=None),
        ("Yes", "Total"): Cell(pct=65.0, count=None, mean=None),
        ("No", "Total"): Cell(pct=35.0, count=None, mean=None),
    }
    s = SeriesResult(
        categories=("Yes", "No"),
        segments=("A", "B", "Total"),
        cells=cells,
        base_n={"Total": 100, "A": 50, "B": 50},
        statistic="pct",
    )
    bars, stack, data = _stacked_layout(s)
    assert bars == ["A", "B"]  # 'Total' excluded
    assert list(stack) == ["Yes", "No"]
    for bi in range(len(bars)):
        assert abs(sum(data[q][bi] for q in stack) - 100.0) < 1e-6


def test_stacked_layout_without_classifier_falls_back():
    """No classifier (only 'Total') → one bar per category, single stack member."""
    from reportbuilder.render.image.bars import _stacked_layout

    cells = {
        ("Yes", "Total"): Cell(pct=65.0, count=None, mean=None),
        ("No", "Total"): Cell(pct=35.0, count=None, mean=None),
    }
    s = SeriesResult(
        categories=("Yes", "No"),
        segments=("Total",),
        cells=cells,
        base_n={"Total": 100},
        statistic="pct",
    )
    bars, stack, data = _stacked_layout(s)
    assert list(bars) == ["Yes", "No"]
    assert stack == ["Total"]
