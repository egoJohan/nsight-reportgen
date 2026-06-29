"""Special (non-chart) slides: model helper, serde round-trip, and render.

These cover the Overview/Conclusion/Demographics slides that ride inside
Report.charts as ChartSpecs with question_ref="" and options["bullets"].
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from reportbuilder.model.report import (
    ChartSpec,
    ElementToggles,
    NumberFormat,
    Report,
    SortSpec,
    is_special_slide,
    report_from_json,
    report_to_json,
)
from reportbuilder.export.pptx_build import build_pptx
from reportbuilder.testing.fixtures import tiny_model_and_data


def _spec(**kw) -> ChartSpec:
    d = dict(
        question_ref="",
        chart_type="",
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"),
        template_slot="auto",
        elements=ElementToggles(),
    )
    d.update(kw)
    return ChartSpec(**d)


def _special(**kw) -> ChartSpec:
    return _spec(
        chart_type="special_overview",
        slide_title="Tutkimuksen taustaa",
        options={"bullets": ["Eka havainto", "Toka havainto"]},
        **kw,
    )


def test_is_special_slide_true_for_special_types():
    for t in ("special_overview", "special_conclusion", "special_demographics"):
        assert is_special_slide(_spec(chart_type=t)) is True


def test_is_special_slide_false_for_charts():
    assert is_special_slide(_spec(chart_type="horizontal_bar")) is False


def test_special_spec_round_trips():
    rep = Report(name="r", render_mode="image", template_ref="", charts=(_special(),))
    restored = report_from_json(report_to_json(rep))
    spec = restored.charts[0]
    assert is_special_slide(spec)
    assert spec.slide_title == "Tutkimuksen taustaa"
    assert spec.options["bullets"] == ["Eka havainto", "Toka havainto"]


def _count_pictures(path: str) -> int:
    prs = Presentation(path)
    return sum(
        1
        for slide in prs.slides
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    )


def test_build_pptx_mixed_deck_passes_completeness(tmp_path):
    """A special slide + a normal chart render without tripping assert_complete;
    the special slide adds NO picture (it's text)."""
    model, df = tiny_model_and_data()
    normal = _spec(question_ref="q1", chart_type="vertical_bar", statistic="pct")
    rep = Report(
        name="r", render_mode="image", template_ref="", charts=(_special(), normal)
    )
    out = str(tmp_path / "deck.pptx")
    build_pptx(rep, model, df, out)
    prs = Presentation(out)
    assert len(prs.slides) == 2
    # Exactly one picture (the normal chart); the special slide is text only.
    assert _count_pictures(out) == 1


def test_build_pptx_all_special_deck(tmp_path):
    model, df = tiny_model_and_data()
    rep = Report(
        name="r", render_mode="image", template_ref="", charts=(_special(),)
    )
    out = str(tmp_path / "deck.pptx")
    build_pptx(rep, model, df, out)
    assert _count_pictures(out) == 0  # no charts → no pictures, still valid


def test_special_slide_tolerates_string_bullets(tmp_path):
    """A bare string in options['bullets'] renders as ONE bullet, not per-char."""
    model, df = tiny_model_and_data()
    spec = _spec(
        chart_type="special_overview",
        slide_title="H",
        options={"bullets": "a single string"},
    )
    rep = Report(name="r", render_mode="image", template_ref="", charts=(spec,))
    out = str(tmp_path / "deck.pptx")
    build_pptx(rep, model, df, out)
    prs = Presentation(out)
    alltext = " ".join(
        s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame
    )
    assert "a single string" in alltext
    # Would-be per-char rendering produces standalone bullet glyphs; one bullet only.
    assert alltext.count("•") == 1


def test_special_slide_renders_heading_and_bullets(tmp_path):
    model, df = tiny_model_and_data()
    rep = Report(
        name="r", render_mode="image", template_ref="", charts=(_special(),)
    )
    out = str(tmp_path / "deck.pptx")
    build_pptx(rep, model, df, out)
    prs = Presentation(out)
    alltext = " ".join(
        s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame
    )
    assert "Tutkimuksen taustaa" in alltext
    assert "Eka havainto" in alltext and "Toka havainto" in alltext


def test_md_runs_parses_bold_italic():
    from reportbuilder.render.image.special_slide import _md_runs
    assert _md_runs("a **b** c *d*") == [
        ("a ", False, False),
        ("b", True, False),
        (" c ", False, False),
        ("d", False, True),
    ]
    assert _md_runs("plain") == [("plain", False, False)]


def test_bullet_markdown_renders_bold(tmp_path):
    """A **bold** bullet produces a bold run in the rendered slide."""
    model, df = tiny_model_and_data()
    spec = _spec(
        chart_type="special_conclusion",
        slide_title="C",
        options={"bullets": ["Tämä on **tärkeä** kohta"]},
    )
    rep = Report(name="r", render_mode="image", template_ref="", charts=(spec,))
    out = str(tmp_path / "d.pptx")
    build_pptx(rep, model, df, out)
    prs = Presentation(out)
    runs = [
        r
        for s in prs.slides[0].shapes
        if s.has_text_frame
        for p in s.text_frame.paragraphs
        for r in p.runs
    ]
    bold = [r.text for r in runs if r.font.bold]
    assert "tärkeä" in bold
