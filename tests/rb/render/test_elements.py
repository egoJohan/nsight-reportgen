"""TDD tests for render/elements.py — REQ-C-24a..i, C-25 (Task 5.3)."""
from __future__ import annotations
import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from reportbuilder.model.report import (
    ChartSpec, NumberFormat, ElementToggles, SortSpec,
)
from reportbuilder.render.base import RenderContext, Slot, StyleSpec
from reportbuilder.testing.fixtures import known_series, one_chart_report


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_prs_and_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _make_column_chart(slide):
    chart_data = CategoryChartData()
    chart_data.categories = ["Cat A", "Cat B"]
    chart_data.add_series("Series 1", (30.0, 70.0))
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1), Inches(6), Inches(4),
        chart_data,
    )
    return gf.chart


def _make_slot():
    return Slot(
        slide_index=0,
        left=int(Inches(1)),
        top=int(Inches(1)),
        width=int(Inches(6)),
        height=int(Inches(4)),
        name="slot1",
    )


def _make_chart_spec(elements: ElementToggles, classifying_var: str | None = None) -> ChartSpec:
    """Build a ChartSpec with given toggles and optional classifying_var."""
    return ChartSpec(
        question_ref="q1",
        chart_type="vertical_bar",
        statistic="pct",
        classifying_var=classifying_var,
        number_format=NumberFormat(pct_decimals=0),
        sort=SortSpec(basis="data_order"),
        template_slot="slot1",
        elements=elements,
    )


def _make_ctx(slide, elements: ElementToggles, classifying_var: str | None = None) -> RenderContext:
    spec = _make_chart_spec(elements, classifying_var=classifying_var)
    return RenderContext(
        slide=slide,
        slot=_make_slot(),
        style=StyleSpec(),
        spec=spec,
        series=known_series(),
        fmt=spec.number_format,
    )


# ---------------------------------------------------------------------------
# Test 1: number_format_code
# ---------------------------------------------------------------------------

class TestNumberFormatCode:
    def test_pct_no_decimals(self):
        from reportbuilder.render.elements import number_format_code
        fmt = NumberFormat(pct_decimals=0)
        assert number_format_code(fmt, "pct") == '0"%"'

    def test_pct_one_decimal(self):
        from reportbuilder.render.elements import number_format_code
        fmt = NumberFormat(pct_decimals=1)
        assert number_format_code(fmt, "pct") == '0.0"%"'

    def test_mean_two_decimals(self):
        from reportbuilder.render.elements import number_format_code
        fmt = NumberFormat(mean_decimals=2)
        assert number_format_code(fmt, "mean") == "0.00"

    def test_count_returns_zero(self):
        from reportbuilder.render.elements import number_format_code
        fmt = NumberFormat()
        assert number_format_code(fmt, "count") == "0"


# ---------------------------------------------------------------------------
# Test 2: title applied
# ---------------------------------------------------------------------------

class TestTitleApplied:
    def test_title_applied(self):
        from reportbuilder.render.elements import apply_elements
        _, slide = _make_prs_and_slide()
        chart = _make_column_chart(slide)
        elements = ElementToggles(title=True, legend=False, n=False, axis_names=False,
                                  filter_var=False, data_labels=False)
        ctx = _make_ctx(slide, elements)

        apply_elements(chart, ctx, title="My Question")

        assert chart.has_title
        assert chart.chart_title.text_frame.text == "My Question"
        # Check the font was applied
        name, size = ctx.style.font_for("title")
        run = chart.chart_title.text_frame.paragraphs[0].runs[0]
        assert run.font.name == name
        assert run.font.size == Pt(size)

    def test_title_not_applied_when_toggle_off(self):
        from reportbuilder.render.elements import apply_elements
        _, slide = _make_prs_and_slide()
        chart = _make_column_chart(slide)
        elements = ElementToggles(title=False, legend=False, n=False, axis_names=False,
                                  filter_var=False, data_labels=False)
        ctx = _make_ctx(slide, elements)

        apply_elements(chart, ctx, title="Should Not Appear")

        assert not chart.has_title


# ---------------------------------------------------------------------------
# Test 3: data_labels and legend toggle
# ---------------------------------------------------------------------------

class TestDataLabelsAndLegendToggle:
    def test_all_toggles_on(self):
        from reportbuilder.render.elements import apply_elements
        _, slide = _make_prs_and_slide()
        chart = _make_column_chart(slide)
        elements = ElementToggles(title=True, legend=True, n=True, axis_names=True,
                                  filter_var=True, data_labels=True)
        ctx = _make_ctx(slide, elements)

        apply_elements(chart, ctx, title="Q")

        # Data labels
        assert chart.plots[0].has_data_labels is True
        dl = chart.plots[0].data_labels
        assert dl.number_format == '0"%"'
        assert dl.number_format_is_linked is False

        # Legend
        assert chart.has_legend is True

    def test_all_toggles_off(self):
        from reportbuilder.render.elements import apply_elements
        _, slide = _make_prs_and_slide()
        chart = _make_column_chart(slide)
        elements = ElementToggles(title=False, legend=False, n=False, axis_names=False,
                                  filter_var=False, data_labels=False)
        ctx = _make_ctx(slide, elements)

        apply_elements(chart, ctx, title="Q")

        assert chart.has_legend is False
        # data_labels should be falsy (False or None)
        assert not chart.plots[0].has_data_labels


# ---------------------------------------------------------------------------
# Test 4: N and filter annotations
# ---------------------------------------------------------------------------

class TestNAndFilterAnnotations:
    def _count_textboxes_with(self, slide, text_fragment: str) -> int:
        """Count textboxes whose text contains the given fragment."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        count = 0
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # Textboxes typically lack a chart inside
            if shape.has_text_frame and text_fragment in shape.text_frame.text:
                count += 1
        return count

    def test_n_annotation_added_when_toggle_on(self):
        from reportbuilder.render.elements import add_n_annotation
        _, slide = _make_prs_and_slide()
        elements = ElementToggles(n=True, title=False, legend=False, axis_names=False,
                                  filter_var=False, data_labels=False)
        ctx = _make_ctx(slide, elements)
        initial_count = len(slide.shapes)

        add_n_annotation(ctx)

        assert len(slide.shapes) == initial_count + 1
        base_n = ctx.series.base_n["Total"]
        assert self._count_textboxes_with(slide, str(base_n)) == 1

    def test_n_annotation_not_added_when_toggle_off(self):
        from reportbuilder.render.elements import add_n_annotation
        _, slide = _make_prs_and_slide()
        elements = ElementToggles(n=False, title=False, legend=False, axis_names=False,
                                  filter_var=False, data_labels=False)
        ctx = _make_ctx(slide, elements)
        initial_count = len(slide.shapes)

        add_n_annotation(ctx)

        assert len(slide.shapes) == initial_count

    def test_filter_annotation_added_when_toggle_and_classifying_var_on(self):
        from reportbuilder.render.elements import add_filter_annotation
        _, slide = _make_prs_and_slide()
        elements = ElementToggles(filter_var=True, title=False, legend=False, axis_names=False,
                                  n=False, data_labels=False)
        ctx = _make_ctx(slide, elements, classifying_var="region")
        initial_count = len(slide.shapes)

        add_filter_annotation(ctx)

        assert len(slide.shapes) == initial_count + 1
        assert self._count_textboxes_with(slide, "region") == 1

    def test_filter_annotation_not_added_when_toggle_off(self):
        from reportbuilder.render.elements import add_filter_annotation
        _, slide = _make_prs_and_slide()
        elements = ElementToggles(filter_var=False, title=False, legend=False, axis_names=False,
                                  n=False, data_labels=False)
        ctx = _make_ctx(slide, elements, classifying_var="region")
        initial_count = len(slide.shapes)

        add_filter_annotation(ctx)

        assert len(slide.shapes) == initial_count

    def test_filter_annotation_not_added_when_no_classifying_var(self):
        from reportbuilder.render.elements import add_filter_annotation
        _, slide = _make_prs_and_slide()
        elements = ElementToggles(filter_var=True, title=False, legend=False, axis_names=False,
                                  n=False, data_labels=False)
        ctx = _make_ctx(slide, elements, classifying_var=None)
        initial_count = len(slide.shapes)

        add_filter_annotation(ctx)

        assert len(slide.shapes) == initial_count
