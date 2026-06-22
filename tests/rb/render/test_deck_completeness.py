"""TDD tests for deck completeness + native-mode no-picture gate (Task 5.15).

REQ-C-18: assert_complete — deck has exactly one rendered chart object per ChartSpec.
REQ-C-23a: assert_no_pictures_in_chart_slots — native-mode reports have zero picture shapes.
"""
from __future__ import annotations

import dataclasses
import io

import pytest
from PIL import Image as PILImage
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from reportbuilder.model.report import Report, ChartSpec, SortSpec, NumberFormat, ElementToggles
from reportbuilder.render.base import StyleSpec
from reportbuilder.render.deck import (
    render_report,
    assert_complete,
    assert_no_pictures_in_chart_slots,
    CompletenessError,
    NativePurityError,
)
from reportbuilder.testing.fixtures import known_series


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_chart(question_ref: str, slot: str, chart_type: str = "vertical_bar") -> ChartSpec:
    return ChartSpec(
        question_ref=question_ref,
        chart_type=chart_type,
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"),
        template_slot=slot,
        elements=ElementToggles(),
    )


def _three_chart_report(render_mode: str = "native") -> Report:
    """A report with 3 charts, each on a distinct slot."""
    return Report(
        name="three-chart",
        render_mode=render_mode,
        template_ref="t.pptx",
        charts=(
            _make_chart("q1", "slot1"),
            _make_chart("q2", "slot2"),
            _make_chart("q3", "slot3"),
        ),
    )


def _series_by_ref(report: Report) -> dict:
    return {spec.question_ref: known_series() for spec in report.charts}


def _tiny_png_bytes() -> bytes:
    """Return a minimal valid PNG (1x1 white pixel) as bytes."""
    img = PILImage.new("RGB", (1, 1), color=(255, 255, 255))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _inject_extra_chart(prs: Presentation) -> None:
    """Add an extra native chart shape onto the first slide."""
    slide = prs.slides[0]
    cd = CategoryChartData()
    cd.categories = ["A", "B"]
    cd.add_series("Extra", (10.0, 20.0))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1), Inches(4), Inches(3),
        cd,
    )


def _inject_picture(prs: Presentation, tmp_path) -> None:
    """Save a tiny PNG to tmp_path and add it as a PICTURE shape on the first slide."""
    png_path = tmp_path / "tiny.png"
    png_path.write_bytes(_tiny_png_bytes())
    slide = prs.slides[0]
    slide.shapes.add_picture(
        str(png_path),
        Inches(0), Inches(0),
        Inches(1), Inches(1),
    )


# ---------------------------------------------------------------------------
# Test 1: completeness pass — 3-chart native report
# ---------------------------------------------------------------------------

def test_completeness_native_passes():
    """render_report with a 3-chart native report returns without raising; assert_complete passes."""
    report = _three_chart_report("native")
    series_by_ref = _series_by_ref(report)
    style = StyleSpec()

    prs = render_report(report, series_by_ref, style)

    # render_report itself must not raise (asserts already wired in)
    # Additionally confirm the explicit guard is also satisfied
    assert_complete(prs, report)


# ---------------------------------------------------------------------------
# Test 2: extra chart injected → CompletenessError
# ---------------------------------------------------------------------------

def test_completeness_extra_chart_fails():
    """Injecting an extra chart into a native deck raises CompletenessError."""
    report = _three_chart_report("native")
    series_by_ref = _series_by_ref(report)
    style = StyleSpec()

    prs = render_report(report, series_by_ref, style)

    _inject_extra_chart(prs)

    with pytest.raises(CompletenessError, match="expected 3"):
        assert_complete(prs, report)


# ---------------------------------------------------------------------------
# Test 3: native no-picture gate passes including funnel chart
# ---------------------------------------------------------------------------

def test_native_no_picture_gate_passes_incl_funnel():
    """Native report with a funnel chart passes assert_no_pictures_in_chart_slots.

    The funnel is implemented as a stacked bar (c:barChart), so it produces a
    native chart shape, NOT a picture shape.
    """
    report = Report(
        name="funnel-report",
        render_mode="native",
        template_ref="t.pptx",
        charts=(
            _make_chart("q1", "slot1", chart_type="funnel"),
        ),
    )
    series_by_ref = {"q1": known_series()}
    style = StyleSpec()

    prs = render_report(report, series_by_ref, style)

    # Must not raise
    assert_no_pictures_in_chart_slots(prs, report)


# ---------------------------------------------------------------------------
# Test 4: injected picture in native deck → NativePurityError
# ---------------------------------------------------------------------------

def test_native_picture_injected_fails(tmp_path):
    """Injecting a PICTURE shape into a native deck raises NativePurityError."""
    report = _three_chart_report("native")
    series_by_ref = _series_by_ref(report)
    style = StyleSpec()

    prs = render_report(report, series_by_ref, style)

    _inject_picture(prs, tmp_path)

    with pytest.raises(NativePurityError, match="picture shape"):
        assert_no_pictures_in_chart_slots(prs, report)


# ---------------------------------------------------------------------------
# Test 5: image mode allows pictures — gate is a no-op; completeness counts pictures
# ---------------------------------------------------------------------------

def test_image_mode_allows_pictures():
    """Image mode: assert_no_pictures_in_chart_slots is a no-op; assert_complete passes."""
    report = _three_chart_report("image")
    series_by_ref = _series_by_ref(report)
    style = StyleSpec()

    prs = render_report(report, series_by_ref, style)

    # Gate must be a no-op (must not raise)
    assert_no_pictures_in_chart_slots(prs, report)

    # Completeness must pass (image mode counts pictures == 3 chart specs)
    assert_complete(prs, report)
