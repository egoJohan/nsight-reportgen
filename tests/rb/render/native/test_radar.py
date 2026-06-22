"""TDD tests for native radar chart builder (Task 5.7).

Covers: build_radar (RADAR_MARKERS, one series per segment, colors via line stroke).
REQ-C-13, C-24g.
"""
from __future__ import annotations
import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

from reportbuilder.render.base import Slot, RenderContext, StyleSpec
from reportbuilder.model.report import ChartSpec, SortSpec, NumberFormat, ElementToggles
from reportbuilder.stats.series import Cell, SeriesResult


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _make_slot() -> Slot:
    return Slot(
        slide_index=0,
        left=Inches(1),
        top=Inches(1),
        width=Inches(8),
        height=Inches(5),
        name="slot1",
    )


def _make_spec(chart_type: str) -> ChartSpec:
    return ChartSpec(
        question_ref="q1",
        chart_type=chart_type,
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"),
        template_slot="slot1",
        elements=ElementToggles(),
    )


def _make_ctx(series: SeriesResult, chart_type: str) -> tuple[Presentation, RenderContext]:
    prs, slide = _make_slide()
    slot = _make_slot()
    style = StyleSpec()
    spec = _make_spec(chart_type)
    ctx = RenderContext(
        slide=slide,
        slot=slot,
        style=style,
        spec=spec,
        series=series,
        fmt=spec.number_format,
    )
    return prs, ctx


def _make_radar_series() -> SeriesResult:
    """Multi-segment radar series: 3 categories (axes), 3 segments (lines)."""
    categories = ("Quality", "Price", "Service")
    segments = ("Brand A", "Brand B", "Total")
    cells = {
        ("Quality", "Brand A"): Cell(pct=70.0, count=70.0, mean=None),
        ("Price",   "Brand A"): Cell(pct=55.0, count=55.0, mean=None),
        ("Service", "Brand A"): Cell(pct=80.0, count=80.0, mean=None),
        ("Quality", "Brand B"): Cell(pct=60.0, count=60.0, mean=None),
        ("Price",   "Brand B"): Cell(pct=75.0, count=75.0, mean=None),
        ("Service", "Brand B"): Cell(pct=65.0, count=65.0, mean=None),
        ("Quality", "Total"):   Cell(pct=65.0, count=65.0, mean=None),
        ("Price",   "Total"):   Cell(pct=65.0, count=65.0, mean=None),
        ("Service", "Total"):   Cell(pct=72.0, count=72.0, mean=None),
    }
    return SeriesResult(
        categories=categories,
        segments=segments,
        cells=cells,
        base_n={seg: 100 for seg in segments},
        statistic="pct",
    )


# ---------------------------------------------------------------------------
# Test 1: build_radar
# ---------------------------------------------------------------------------

def test_build_radar():
    """build_radar returns a chart GF with c:radarChart XML, correct series, no PICTUREs."""
    from reportbuilder.render.native.radar import build_radar

    series = _make_radar_series()
    prs, ctx = _make_ctx(series, "radar")
    gf = build_radar(ctx)

    assert gf.has_chart, "shape must be a chart graphic frame"

    chart_space = gf.chart._chartSpace

    # radarChart element present
    assert chart_space.find(".//" + qn("c:radarChart")) is not None, (
        "XML must contain c:radarChart element"
    )

    # One series per segment (Brand A, Brand B, Total)
    plot = gf.chart.plots[0]
    assert len(plot.series) == 3, f"expected 3 series (one per segment), got {len(plot.series)}"

    # Values match the fixture for first series (Brand A: Quality=70, Price=55, Service=80)
    values_a = list(plot.series[0].values)
    assert values_a == [70.0, 55.0, 80.0], (
        f"expected [70.0, 55.0, 80.0] for Brand A, got {values_a}"
    )

    # Values match the fixture for second series (Brand B: Quality=60, Price=75, Service=65)
    values_b = list(plot.series[1].values)
    assert values_b == [60.0, 75.0, 65.0], (
        f"expected [60.0, 75.0, 65.0] for Brand B, got {values_b}"
    )

    # Zero PICTURE shapes — chart must be native
    pictures = [s for s in ctx.slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures == [], f"Found {len(pictures)} PICTURE shapes — chart must be native"


# ---------------------------------------------------------------------------
# Test 2: registry
# ---------------------------------------------------------------------------

def test_registry_radar():
    """NATIVE_BUILDERS['radar'] must be build_radar."""
    from reportbuilder.render.native import NATIVE_BUILDERS
    from reportbuilder.render.native.radar import build_radar

    assert NATIVE_BUILDERS["radar"] is build_radar, (
        "NATIVE_BUILDERS['radar'] must be build_radar"
    )
