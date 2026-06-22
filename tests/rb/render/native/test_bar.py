"""TDD tests for native bar chart builders (Task 5.5).

Covers: build_horizontal_bar, build_stacked_vertical_bar, build_stacked_horizontal_bar.
REQ-C-13, C-24b/g.
"""
from __future__ import annotations
import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

from reportbuilder.render.base import Slot, RenderContext, StyleSpec
from reportbuilder.stats.series import Cell, SeriesResult
from reportbuilder.testing.fixtures import known_series, one_chart_report

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

CATEGORIES = ("A", "B", "C")
SEGMENTS_STACKED = ("Seg1", "Seg2", "Total")

# PCT values per category per segment for a stacked fixture
_PCT = {
    ("A", "Seg1"): 30.0, ("B", "Seg1"): 20.0, ("C", "Seg1"): 50.0,
    ("A", "Seg2"): 25.0, ("B", "Seg2"): 35.0, ("C", "Seg2"): 40.0,
    ("A", "Total"): 55.0, ("B", "Total"): 55.0, ("C", "Total"): 90.0,
}


def _make_stacked_series() -> SeriesResult:
    cells = {
        (cat, seg): Cell(pct=_PCT[(cat, seg)], count=float(_PCT[(cat, seg)]), mean=None)
        for cat in CATEGORIES
        for seg in SEGMENTS_STACKED
    }
    return SeriesResult(
        categories=CATEGORIES,
        segments=SEGMENTS_STACKED,
        cells=cells,
        base_n={seg: 100 for seg in SEGMENTS_STACKED},
        statistic="pct",
    )


def _make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _make_slot() -> Slot:
    return Slot(
        slide_index=0,
        left=Inches(1),
        top=Inches(1),
        width=Inches(8),
        height=Inches(5),
        name="slot1",
    )


def _make_ctx(series: SeriesResult) -> tuple[Presentation, RenderContext]:
    prs, slide = _make_slide()
    slot = _make_slot()
    style = StyleSpec()
    spec = one_chart_report().charts[0]
    ctx = RenderContext(
        slide=slide,
        slot=slot,
        style=style,
        spec=spec,
        series=series,
        fmt=spec.number_format,
    )
    return prs, ctx


# ---------------------------------------------------------------------------
# Test 1: build_horizontal_bar
# ---------------------------------------------------------------------------

def test_build_horizontal_bar():
    """build_horizontal_bar returns a chart GF with barDir=bar, correct values, no PICTUREs."""
    from reportbuilder.render.native.bar import build_horizontal_bar

    prs, ctx = _make_ctx(known_series())
    gf = build_horizontal_bar(ctx)

    assert gf.has_chart, "shape must be a chart graphic frame"

    xml = gf.chart._chartSpace.xml

    # barChart element present
    assert "c:barChart" in xml, "XML must contain c:barChart element"

    # barDir val="bar" means horizontal
    bar_chart_el = gf.chart._chartSpace.find(".//" + qn("c:barChart"))
    assert bar_chart_el is not None, "c:barChart element must exist"
    bar_dir_el = bar_chart_el.find(qn("c:barDir"))
    assert bar_dir_el is not None, "c:barDir element must exist inside c:barChart"
    assert bar_dir_el.get("val") == "bar", (
        f"c:barDir val must be 'bar' for horizontal chart, got {bar_dir_el.get('val')!r}"
    )

    # Series values match known_series (pct: Yes=60.0, No=40.0)
    plot = gf.chart.plots[0]
    assert len(plot.series) == 1
    values = list(plot.series[0].values)
    assert values == [60.0, 40.0], f"expected [60.0, 40.0], got {values}"

    # Zero PICTURE shapes
    pictures = [s for s in ctx.slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures == [], f"Found {len(pictures)} PICTURE shapes — chart must be native"


# ---------------------------------------------------------------------------
# Test 2: build_stacked_vertical_bar
# ---------------------------------------------------------------------------

def test_build_stacked_vertical_bar():
    """build_stacked_vertical_bar: stacked column; grouping=stacked, overlap=100, 3 series, no PICTUREs."""
    from reportbuilder.render.native.bar import build_stacked_vertical_bar

    prs, ctx = _make_ctx(_make_stacked_series())
    gf = build_stacked_vertical_bar(ctx)

    assert gf.has_chart, "shape must be a chart graphic frame"

    chart_space = gf.chart._chartSpace

    # Must use barChart with barDir="col" (vertical)
    bar_chart_el = chart_space.find(".//" + qn("c:barChart"))
    assert bar_chart_el is not None, "c:barChart element must exist"
    bar_dir_el = bar_chart_el.find(qn("c:barDir"))
    assert bar_dir_el is not None, "c:barDir must be present"
    assert bar_dir_el.get("val") == "col", (
        f"c:barDir val must be 'col' for vertical/column chart, got {bar_dir_el.get('val')!r}"
    )

    # grouping="stacked"
    grouping_el = bar_chart_el.find(qn("c:grouping"))
    assert grouping_el is not None, "c:grouping must be present"
    assert grouping_el.get("val") == "stacked", (
        f"c:grouping val must be 'stacked', got {grouping_el.get('val')!r}"
    )

    # overlap=100
    plot = gf.chart.plots[0]
    assert plot.overlap == 100, f"overlap must be 100, got {plot.overlap}"

    # 3 series present (one per segment)
    assert len(plot.series) == 3, f"expected 3 series, got {len(plot.series)}"

    # Zero PICTURE shapes
    pictures = [s for s in ctx.slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures == [], f"Found {len(pictures)} PICTURE shapes — chart must be native"


# ---------------------------------------------------------------------------
# Test 3: build_stacked_horizontal_bar
# ---------------------------------------------------------------------------

def test_build_stacked_horizontal_bar():
    """build_stacked_horizontal_bar: barDir=bar, grouping=stacked, overlap=100, 3 series, no PICTUREs."""
    from reportbuilder.render.native.bar import build_stacked_horizontal_bar

    prs, ctx = _make_ctx(_make_stacked_series())
    gf = build_stacked_horizontal_bar(ctx)

    assert gf.has_chart, "shape must be a chart graphic frame"

    chart_space = gf.chart._chartSpace

    # barChart with barDir="bar" (horizontal)
    bar_chart_el = chart_space.find(".//" + qn("c:barChart"))
    assert bar_chart_el is not None, "c:barChart element must exist"
    bar_dir_el = bar_chart_el.find(qn("c:barDir"))
    assert bar_dir_el is not None, "c:barDir must be present"
    assert bar_dir_el.get("val") == "bar", (
        f"c:barDir val must be 'bar' for horizontal chart, got {bar_dir_el.get('val')!r}"
    )

    # grouping="stacked"
    grouping_el = bar_chart_el.find(qn("c:grouping"))
    assert grouping_el is not None, "c:grouping must be present"
    assert grouping_el.get("val") == "stacked", (
        f"c:grouping val must be 'stacked', got {grouping_el.get('val')!r}"
    )

    # overlap=100
    plot = gf.chart.plots[0]
    assert plot.overlap == 100, f"overlap must be 100, got {plot.overlap}"

    # 3 series present
    assert len(plot.series) == 3, f"expected 3 series, got {len(plot.series)}"

    # Zero PICTURE shapes
    pictures = [s for s in ctx.slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures == [], f"Found {len(pictures)} PICTURE shapes — chart must be native"


# ---------------------------------------------------------------------------
# Test 4: registry
# ---------------------------------------------------------------------------

def test_registry_has_three_bar_builders():
    """NATIVE_BUILDERS must map horizontal_bar, stacked_vertical_bar, stacked_horizontal_bar."""
    from reportbuilder.render.native import NATIVE_BUILDERS
    from reportbuilder.render.native.bar import (
        build_horizontal_bar,
        build_stacked_vertical_bar,
        build_stacked_horizontal_bar,
    )

    assert NATIVE_BUILDERS["horizontal_bar"] is build_horizontal_bar
    assert NATIVE_BUILDERS["stacked_vertical_bar"] is build_stacked_vertical_bar
    assert NATIVE_BUILDERS["stacked_horizontal_bar"] is build_stacked_horizontal_bar
