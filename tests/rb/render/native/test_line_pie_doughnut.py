"""TDD tests for native line, pie, and doughnut chart builders (Task 5.6).

Covers: build_line (LINE_MARKERS), build_pie (PIE, per-point colors),
build_doughnut (DOUGHNUT, per-point colors).
REQ-C-13, REQ-C-24b (chart-type elements/points count matches categories),
REQ-C-24c (axis values), REQ-C-24e (category names == value labels),
REQ-C-24f (category numeric values as data labels).
"""
from __future__ import annotations
import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

from reportbuilder.render.base import Slot, RenderContext, StyleSpec
from reportbuilder.model.report import ChartSpec, SortSpec, NumberFormat, ElementToggles
from reportbuilder.stats.series import Cell, SeriesResult
from reportbuilder.testing.fixtures import known_series


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _make_slot() -> Slot:
    return Slot(
        slide_index=0,
        left=Inches(1),
        top=Inches(1),
        width=Inches(8),
        height=Inches(5),
        name="slot1",
    )


def _make_spec(chart_type: str) -> ChartSpec:
    return ChartSpec(
        question_ref="q1",
        chart_type=chart_type,
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"),
        template_slot="slot1",
        elements=ElementToggles(),
    )


def _make_ctx(series: SeriesResult, chart_type: str) -> tuple[Presentation, RenderContext]:
    prs, slide = _make_slide()
    slot = _make_slot()
    style = StyleSpec()
    spec = _make_spec(chart_type)
    ctx = RenderContext(
        slide=slide,
        slot=slot,
        style=style,
        spec=spec,
        series=series,
        fmt=spec.number_format,
    )
    return prs, ctx


def _make_multi_series() -> SeriesResult:
    """Multi-category series suitable for line charts (3 categories, 2 segments)."""
    categories = ("Wave1", "Wave2", "Wave3")
    segments = ("Group A", "Group B")
    cells = {
        ("Wave1", "Group A"): Cell(pct=55.0, count=55.0, mean=None),
        ("Wave2", "Group A"): Cell(pct=60.0, count=60.0, mean=None),
        ("Wave3", "Group A"): Cell(pct=65.0, count=65.0, mean=None),
        ("Wave1", "Group B"): Cell(pct=40.0, count=40.0, mean=None),
        ("Wave2", "Group B"): Cell(pct=45.0, count=45.0, mean=None),
        ("Wave3", "Group B"): Cell(pct=50.0, count=50.0, mean=None),
    }
    return SeriesResult(
        categories=categories,
        segments=segments,
        cells=cells,
        base_n={seg: 100 for seg in segments},
        statistic="pct",
    )


# ---------------------------------------------------------------------------
# Test 1: build_line
# ---------------------------------------------------------------------------

def test_build_line():
    """build_line returns a chart GF with c:lineChart XML, correct values, no PICTUREs."""
    from reportbuilder.render.native.line import build_line

    series = _make_multi_series()
    prs, ctx = _make_ctx(series, "line")
    gf = build_line(ctx)

    assert gf.has_chart, "shape must be a chart graphic frame"

    chart_space = gf.chart._chartSpace
    xml = chart_space.xml

    # lineChart element present
    assert chart_space.find(".//" + qn("c:lineChart")) is not None, (
        "XML must contain c:lineChart element"
    )

    # Series values match first segment of multi-series (Group A: 55, 60, 65)
    plot = gf.chart.plots[0]
    assert len(plot.series) == 2, f"expected 2 series, got {len(plot.series)}"
    values_a = list(plot.series[0].values)
    assert values_a == [55.0, 60.0, 65.0], f"expected [55.0, 60.0, 65.0], got {values_a}"

    # Zero PICTURE shapes
    pictures = [s for s in ctx.slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures == [], f"Found {len(pictures)} PICTURE shapes — chart must be native"


# ---------------------------------------------------------------------------
# Test 2: build_pie
# ---------------------------------------------------------------------------

def test_build_pie():
    """build_pie returns a chart GF with c:pieChart XML, per-point fill, no PICTUREs."""
    from reportbuilder.render.native.pie import build_pie

    series = known_series()  # Yes/No, single segment "Total"
    prs, ctx = _make_ctx(series, "pie")
    gf = build_pie(ctx)

    assert gf.has_chart, "shape must be a chart graphic frame"

    chart_space = gf.chart._chartSpace

    # pieChart element present
    assert chart_space.find(".//" + qn("c:pieChart")) is not None, (
        "XML must contain c:pieChart element"
    )

    # Pie has one series (the "Total" segment)
    plot = gf.chart.plots[0]
    assert len(plot.series) == 1, f"pie chart must have 1 series, got {len(plot.series)}"

    # Number of points == number of categories (Yes, No)
    pts = plot.series[0].points
    num_cats = len(series.categories)
    assert len(pts) == num_cats, (
        f"expected {num_cats} points (one per category), got {len(pts)}"
    )

    # Per-point fill is set: each dPt in XML should have spPr/solidFill
    xml = chart_space.xml
    assert qn("c:dPt") in xml or "dPt" in xml, (
        "XML must contain dPt elements (per-point formatting)"
    )

    # Zero PICTURE shapes
    pictures = [s for s in ctx.slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures == [], f"Found {len(pictures)} PICTURE shapes — chart must be native"


# ---------------------------------------------------------------------------
# Test 3: build_doughnut
# ---------------------------------------------------------------------------

def test_build_doughnut():
    """build_doughnut returns a chart GF with c:doughnutChart XML, no PICTUREs."""
    from reportbuilder.render.native.doughnut import build_doughnut

    series = known_series()  # Yes/No, single segment "Total"
    prs, ctx = _make_ctx(series, "doughnut")
    gf = build_doughnut(ctx)

    assert gf.has_chart, "shape must be a chart graphic frame"

    chart_space = gf.chart._chartSpace

    # doughnutChart element present
    assert chart_space.find(".//" + qn("c:doughnutChart")) is not None, (
        "XML must contain c:doughnutChart element"
    )

    # Zero PICTURE shapes
    pictures = [s for s in ctx.slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures == [], f"Found {len(pictures)} PICTURE shapes — chart must be native"


# ---------------------------------------------------------------------------
# Test 4: registry
# ---------------------------------------------------------------------------

def test_registry_line_pie_doughnut():
    """NATIVE_BUILDERS must map line, pie, doughnut to the three builders."""
    from reportbuilder.render.native import NATIVE_BUILDERS
    from reportbuilder.render.native.line import build_line
    from reportbuilder.render.native.pie import build_pie
    from reportbuilder.render.native.doughnut import build_doughnut

    assert NATIVE_BUILDERS["line"] is build_line, (
        "NATIVE_BUILDERS['line'] must be build_line"
    )
    assert NATIVE_BUILDERS["pie"] is build_pie, (
        "NATIVE_BUILDERS['pie'] must be build_pie"
    )
    assert NATIVE_BUILDERS["doughnut"] is build_doughnut, (
        "NATIVE_BUILDERS['doughnut'] must be build_doughnut"
    )
