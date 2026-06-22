"""End-to-end rendering spike: native column chart → PDF → number-survival + judge (Task 1.4)."""
import os
import shutil
import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from reportbuilder.render.base import Slot
from reportbuilder.stats.series import Cell, SeriesResult
from reportbuilder.render.native.column import build_column_chart
from reportbuilder.export.pdf_convert import pptx_to_pdf
from reportbuilder.export.preview import pdf_page_to_png
from reportbuilder.testing.fidelity import numbers_from_pptx, numbers_from_pdf, assert_series_match
from reportbuilder.testing.judge import judge_image
from reportbuilder.testing.rubrics import rubric_for

_needs_soffice = pytest.mark.skipif(shutil.which("soffice") is None, reason="LibreOffice required")


def _spike_series() -> SeriesResult:
    cats = ("Strongly agree", "Agree", "Neither", "Disagree", "Strongly disagree")
    pct = (28.0, 34.0, 18.0, 13.0, 7.0)
    return SeriesResult(
        categories=cats, segments=("Total",),
        cells={(c, "Total"): Cell(p, None, None) for c, p in zip(cats, pct)},
        base_n={"Total": 412}, statistic="pct",
    )


def _build_deck(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slot = Slot(slide_index=0, left=Inches(0.8), top=Inches(0.8),
                width=Inches(8.5), height=Inches(5.2), name="slot1")
    build_column_chart(slide, slot, _spike_series())
    p = tmp_path / "spike.pptx"
    prs.save(str(p))
    return str(p), slide


def test_native_no_pictures_and_series_match(tmp_path):
    # GO signal 1: native chart, zero pictures, layer-1 numbers correct.
    pptx, slide = _build_deck(tmp_path)
    assert [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE] == []
    assert_series_match(numbers_from_pptx(pptx), _spike_series())


@_needs_soffice
def test_spike_numbers_survive_pdf(tmp_path):
    # GO signal 2 (objective, no API key): the data-label numbers survive LibreOffice conversion.
    pptx, _ = _build_deck(tmp_path)
    pdf = pptx_to_pdf(pptx, str(tmp_path))
    nums = numbers_from_pdf(pdf)
    for expected in (28.0, 34.0, 18.0, 13.0, 7.0):
        assert any(abs(n - expected) <= 0.5 for n in nums), f"{expected} missing from {nums}"


@_needs_soffice
@pytest.mark.judge
def test_spike_layout_judged_clean(tmp_path):
    # GO signal 3 (deferred without ANTHROPIC_API_KEY): Claude judges layout cleanliness.
    if not os.environ.get("ANTHROPIC_API_KEY"):
        pytest.skip("ANTHROPIC_API_KEY not set")
    pptx, _ = _build_deck(tmp_path)
    pdf = pptx_to_pdf(pptx, str(tmp_path))
    png = pdf_page_to_png(pdf, 0, str(tmp_path / "spike.png"), resolution=150)
    verdict = judge_image(png, rubric_for("R3-LAYOUT"), requirement_id="R3-LAYOUT")
    assert verdict.passed is True, verdict.reasoning
