"""TDD tests for native column chart builder (Task 1.2 + Task 5.4)."""
from __future__ import annotations
import io
import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

from reportbuilder.render.base import Slot, RenderContext, StyleSpec
from reportbuilder.stats.series import Cell, SeriesResult
from reportbuilder.testing.fixtures import known_series, one_chart_report

# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

CATEGORIES = ("Strongly agree", "Agree", "Neutral")
SEGMENT = "Total"
PCT_VALUES = (45.0, 30.0, 25.0)


def _make_series() -> SeriesResult:
    cells = {
        (cat, SEGMENT): Cell(pct=pct, count=float(pct), mean=None)
        for cat, pct in zip(CATEGORIES, PCT_VALUES)
    }
    return SeriesResult(
        categories=CATEGORIES,
        segments=(SEGMENT,),
        cells=cells,
        base_n={SEGMENT: 100},
        statistic="pct",
    )


def _make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _make_slot() -> Slot:
    return Slot(
        slide_index=0,
        left=Inches(1),
        top=Inches(1),
        width=Inches(8),
        height=Inches(5),
        name="slot1",
    )


# ---------------------------------------------------------------------------
# Test 1: registry key (updated: NATIVE_BUILDERS["vertical_bar"] = build_vertical_bar)
# ---------------------------------------------------------------------------

def test_registry_vertical_bar_is_ctx_builder():
    """NATIVE_BUILDERS['vertical_bar'] must be build_vertical_bar (ctx-based builder)."""
    from reportbuilder.render.native import NATIVE_BUILDERS
    from reportbuilder.render.native.column import build_vertical_bar

    assert NATIVE_BUILDERS["vertical_bar"] is build_vertical_bar


# ---------------------------------------------------------------------------
# Test 2: returns a graphic frame with chart + correct series values
# ---------------------------------------------------------------------------

def test_build_column_chart_returns_graphic_frame_with_correct_values():
    """build_column_chart returns a graphic frame whose chart has the fixture pct values."""
    from reportbuilder.render.native.column import build_column_chart

    prs, slide = _make_slide()
    slot = _make_slot()
    series = _make_series()

    gf = build_column_chart(slide, slot, series)

    assert gf.has_chart, "shape must be a chart graphic frame"
    chart = gf.chart
    assert len(chart.plots) > 0
    plot = chart.plots[0]
    assert len(plot.series) == 1
    values = tuple(plot.series[0].values)
    assert values == PCT_VALUES, f"expected {PCT_VALUES}, got {values}"


# ---------------------------------------------------------------------------
# Test 3: zero PICTURE shapes (editability gate)
# ---------------------------------------------------------------------------

def test_no_picture_shapes_on_slide():
    """Slide must contain ZERO MSO_SHAPE_TYPE.PICTURE shapes (chart must be native)."""
    from reportbuilder.render.native.column import build_column_chart

    prs, slide = _make_slide()
    slot = _make_slot()
    series = _make_series()

    build_column_chart(slide, slot, series)

    picture_shapes = [
        s for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert picture_shapes == [], (
        f"Found {len(picture_shapes)} PICTURE shapes — chart must be a native c:chart object"
    )


# ---------------------------------------------------------------------------
# Test 4: c:manualLayout element present in chartSpace
# ---------------------------------------------------------------------------

def test_manual_layout_element_present():
    """c:manualLayout must be present inside gf.chart._chartSpace after injection."""
    from reportbuilder.render.native.column import build_column_chart

    prs, slide = _make_slide()
    slot = _make_slot()
    series = _make_series()

    gf = build_column_chart(slide, slot, series)
    cs = gf.chart._chartSpace
    manual_layout = cs.find(".//" + qn("c:manualLayout"))
    assert manual_layout is not None, "c:manualLayout must be injected into chartSpace"


# ---------------------------------------------------------------------------
# Test 5: c:dLblPos present with val="outEnd"
# ---------------------------------------------------------------------------

def test_dlbl_pos_out_end():
    """c:dLblPos must be present in chartSpace with val='outEnd'."""
    from reportbuilder.render.native.column import build_column_chart

    prs, slide = _make_slide()
    slot = _make_slot()
    series = _make_series()

    gf = build_column_chart(slide, slot, series)
    cs = gf.chart._chartSpace
    dlbl_pos = cs.find(".//" + qn("c:dLblPos"))
    assert dlbl_pos is not None, "c:dLblPos element must be injected"
    assert dlbl_pos.get("val") == "outEnd", (
        f"c:dLblPos val must be 'outEnd', got {dlbl_pos.get('val')!r}"
    )


# ---------------------------------------------------------------------------
# Test 6: save + reopen without error (REQ-C-29a, the real feasibility gate)
# ---------------------------------------------------------------------------

def test_save_and_reopen_without_error():
    """Deck must save and reopen cleanly — guards REQ-C-29a (valid OOXML)."""
    from reportbuilder.render.native.column import build_column_chart

    prs, slide = _make_slide()
    slot = _make_slot()
    series = _make_series()

    build_column_chart(slide, slot, series)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    # Must not raise
    prs2 = Presentation(buf)
    assert len(prs2.slides) == 1, "reopened presentation must have 1 slide"


# ---------------------------------------------------------------------------
# Task 5.4 — New tests
# ---------------------------------------------------------------------------

# Test 5.4.1: series_chart_data maps SeriesResult
def test_series_chart_data_maps_seriesresult():
    """series_chart_data with known_series() → categories=('Yes','No'), Total values=[60.0,40.0]."""
    from reportbuilder.render.native.column import series_chart_data

    s = known_series()
    cd = series_chart_data(s, s.statistic)

    cat_labels = tuple(c.label for c in cd.categories)
    assert cat_labels == ("Yes", "No"), f"got categories: {cat_labels}"
    # cd is CategoryChartData; iterate series to find 'Total'
    series_list = list(cd)
    assert len(series_list) == 1, f"expected 1 series, got {len(series_list)}"
    total_series = series_list[0]
    assert total_series.name == "Total"
    assert list(total_series.values) == [60.0, 40.0], f"got {list(total_series.values)}"


# Test 5.4.2: _value_for None guard
def test_value_for_none_guard():
    """_value_for returns 0.0 when the statistic value on a Cell is None."""
    from reportbuilder.render.native.column import _value_for

    series = SeriesResult(
        categories=("A",),
        segments=("Total",),
        cells={("A", "Total"): Cell(pct=None, count=None, mean=None)},
        base_n={"Total": 0},
        statistic="pct",
    )
    result = _value_for(series, "A", "Total", "pct")
    assert result == 0.0, f"expected 0.0 for None cell, got {result}"


# Test 5.4.3: build_vertical_bar — native chart, colored series, zero pictures
def test_build_vertical_bar_native_and_colored():
    """build_vertical_bar(ctx) returns gf with correct values, colored series, no PICTURE shapes."""
    from reportbuilder.render.native.column import build_vertical_bar

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slot = Slot(
        slide_index=0,
        left=Inches(1),
        top=Inches(1),
        width=Inches(8),
        height=Inches(5),
        name="slot1",
    )
    style = StyleSpec()
    spec = one_chart_report().charts[0]
    s = known_series()
    ctx = RenderContext(
        slide=slide,
        slot=slot,
        style=style,
        spec=spec,
        series=s,
        fmt=spec.number_format,
    )

    gf = build_vertical_bar(ctx)

    # Graphic frame is a chart
    assert gf.has_chart, "shape must be a chart graphic frame"

    chart = gf.chart
    plot = chart.plots[0]

    # Correct series values
    assert list(plot.series[0].values) == [60.0, 40.0], (
        f"expected [60.0, 40.0], got {list(plot.series[0].values)}"
    )

    # Series fill color matches style.color_for(0)
    expected_rgb = RGBColor.from_string(style.color_for(0))
    actual_rgb = plot.series[0].format.fill.fore_color.rgb
    assert actual_rgb == expected_rgb, (
        f"expected RGB {expected_rgb}, got {actual_rgb}"
    )

    # Zero PICTURE shapes on the slide
    picture_shapes = [
        sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert picture_shapes == [], (
        f"Found {len(picture_shapes)} PICTURE shapes — chart must be native"
    )


# Test 5.4.4: series_chart_data honors an explicitly-passed statistic
def test_series_chart_data_honors_explicit_statistic():
    """series_chart_data uses the passed statistic, not series.statistic.

    Build a SeriesResult where pct and count differ, call series_chart_data
    with statistic='count', and assert the series values equal the COUNT values.
    """
    from reportbuilder.render.native.column import series_chart_data

    categories = ("Yes", "No")
    segment = "Total"
    cells = {
        ("Yes", segment): Cell(pct=60.0, count=120.0, mean=None),
        ("No", segment): Cell(pct=40.0, count=80.0, mean=None),
    }
    # series.statistic is "pct" — if the parameter were ignored, values would be pct
    s = SeriesResult(
        categories=categories,
        segments=(segment,),
        cells=cells,
        base_n={segment: 200},
        statistic="pct",
    )

    cd = series_chart_data(s, "count")

    series_list = list(cd)
    assert len(series_list) == 1
    values = list(series_list[0].values)
    assert values == [120.0, 80.0], (
        f"expected count values [120.0, 80.0], got {values} — "
        "statistic parameter must override series.statistic"
    )
