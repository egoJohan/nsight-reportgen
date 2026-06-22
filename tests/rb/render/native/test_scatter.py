"""TDD tests for native scatter chart builder (Task 5.8).

Covers: xy_chart_data (pairs mapped), build_scatter (XY_SCATTER / c:scatterChart,
zero pictures), ValueError when scatter_xy is None, registry entry.
REQ-C-13, REQ-C-24c (axis values), REQ-C-24d (axis names).
"""
from __future__ import annotations
import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

from reportbuilder.render.base import Slot, RenderContext, StyleSpec
from reportbuilder.model.report import ChartSpec, SortSpec, NumberFormat, ElementToggles
from reportbuilder.stats.series import Cell, SeriesResult


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _make_slot() -> Slot:
    return Slot(
        slide_index=0,
        left=Inches(1),
        top=Inches(1),
        width=Inches(8),
        height=Inches(5),
        name="slot1",
    )


def _make_spec(scatter_xy: tuple[str, str] | None = ("satisfaction", "loyalty")) -> ChartSpec:
    return ChartSpec(
        question_ref="q1",
        chart_type="scatter",
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"),
        template_slot="slot1",
        elements=ElementToggles(),
        scatter_xy=scatter_xy,
    )


def _make_series() -> SeriesResult:
    """Scatter series: 3 categories (points), 2 segments (x/y axes)."""
    categories = ("P1", "P2", "P3")
    segments = ("satisfaction", "loyalty")
    cells = {
        ("P1", "satisfaction"): Cell(pct=80.0, count=80.0, mean=None),
        ("P1", "loyalty"):      Cell(pct=60.0, count=60.0, mean=None),
        ("P2", "satisfaction"): Cell(pct=50.0, count=50.0, mean=None),
        ("P2", "loyalty"):      Cell(pct=70.0, count=70.0, mean=None),
        ("P3", "satisfaction"): Cell(pct=30.0, count=30.0, mean=None),
        ("P3", "loyalty"):      Cell(pct=40.0, count=40.0, mean=None),
    }
    return SeriesResult(
        categories=categories,
        segments=segments,
        cells=cells,
        base_n={"satisfaction": 100, "loyalty": 100},
        statistic="pct",
    )


def _make_ctx(scatter_xy: tuple[str, str] | None = ("satisfaction", "loyalty")) -> tuple[Presentation, RenderContext]:
    prs, slide = _make_slide()
    slot = _make_slot()
    style = StyleSpec()
    spec = _make_spec(scatter_xy)
    series = _make_series()
    ctx = RenderContext(
        slide=slide,
        slot=slot,
        style=style,
        spec=spec,
        series=series,
        fmt=spec.number_format,
    )
    return prs, ctx


# ---------------------------------------------------------------------------
# Test 1: xy_chart_data maps category rows to (x, y) pairs
# ---------------------------------------------------------------------------

def test_xy_chart_data_pairs():
    """xy_chart_data produces one series with 3 points matching (sat, loyalty) values."""
    from reportbuilder.render.native.scatter import xy_chart_data

    series = _make_series()
    xd = xy_chart_data(series, ("satisfaction", "loyalty"))

    xd_series = list(xd)
    assert len(xd_series) == 1, f"expected 1 series, got {len(xd_series)}"

    pts = list(xd_series[0])
    assert len(pts) == 3, f"expected 3 data points, got {len(pts)}"

    expected = [(80.0, 60.0), (50.0, 70.0), (30.0, 40.0)]
    actual = [(pt.x, pt.y) for pt in pts]
    assert actual == expected, f"expected {expected}, got {actual}"


# ---------------------------------------------------------------------------
# Test 2: build_scatter produces a native chart with scatterChart XML
# ---------------------------------------------------------------------------

def test_build_scatter_native():
    """build_scatter returns a graphic frame with c:scatterChart in XML, no PICTURE shapes."""
    from reportbuilder.render.native.scatter import build_scatter

    prs, ctx = _make_ctx()
    gf = build_scatter(ctx)

    assert gf.has_chart, "returned shape must be a chart graphic frame"
    assert gf.shape_type == MSO_SHAPE_TYPE.CHART, (
        f"expected MSO_SHAPE_TYPE.CHART, got {gf.shape_type}"
    )

    chart_space = gf.chart._chartSpace
    assert chart_space.find(".//" + qn("c:scatterChart")) is not None, (
        "chart XML must contain c:scatterChart element"
    )

    pictures = [s for s in ctx.slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures == [], f"Found {len(pictures)} PICTURE shape(s) — chart must be native"


# ---------------------------------------------------------------------------
# Test 3: build_scatter raises ValueError when scatter_xy is None
# ---------------------------------------------------------------------------

def test_build_scatter_requires_xy():
    """build_scatter raises ValueError mentioning 'scatter' when scatter_xy is None."""
    from reportbuilder.render.native.scatter import build_scatter

    prs, ctx = _make_ctx(scatter_xy=None)
    with pytest.raises(ValueError, match="scatter"):
        build_scatter(ctx)


# ---------------------------------------------------------------------------
# Test 4: NATIVE_BUILDERS registry
# ---------------------------------------------------------------------------

def test_registry_scatter():
    """NATIVE_BUILDERS['scatter'] must be build_scatter."""
    from reportbuilder.render.native import NATIVE_BUILDERS
    from reportbuilder.render.native.scatter import build_scatter

    assert NATIVE_BUILDERS["scatter"] is build_scatter, (
        "NATIVE_BUILDERS['scatter'] must be build_scatter"
    )
