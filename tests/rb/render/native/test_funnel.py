"""TDD tests for native funnel chart builder (Task 5.9).

Covers: funnel_spacer_values (centering math), build_funnel (BAR_STACKED with
leading transparent spacer series, two series total, overlap=100, value series
values correct, zero PICTURE shapes), registry entry.
REQ-C-13, C-23a.
"""
from __future__ import annotations
import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL

from reportbuilder.render.base import Slot, RenderContext, StyleSpec
from reportbuilder.model.report import ChartSpec, SortSpec, NumberFormat, ElementToggles
from reportbuilder.stats.series import Cell, SeriesResult


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

CATEGORIES = ("Stage1", "Stage2", "Stage3")
VALUES = (80.0, 60.0, 40.0)


def _make_funnel_series() -> SeriesResult:
    """3-category single-'Total'-segment SeriesResult for funnel testing."""
    cells = {
        ("Stage1", "Total"): Cell(pct=80.0, count=80.0, mean=None),
        ("Stage2", "Total"): Cell(pct=60.0, count=60.0, mean=None),
        ("Stage3", "Total"): Cell(pct=40.0, count=40.0, mean=None),
    }
    return SeriesResult(
        categories=CATEGORIES,
        segments=("Total",),
        cells=cells,
        base_n={"Total": 100},
        statistic="pct",
    )


def _make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _make_slot() -> Slot:
    return Slot(
        slide_index=0,
        left=Inches(1),
        top=Inches(1),
        width=Inches(8),
        height=Inches(5),
        name="slot1",
    )


def _make_spec() -> ChartSpec:
    return ChartSpec(
        question_ref="q1",
        chart_type="funnel",
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"),
        template_slot="slot1",
        elements=ElementToggles(),
    )


def _make_ctx() -> tuple[Presentation, RenderContext]:
    prs, slide = _make_slide()
    slot = _make_slot()
    style = StyleSpec()
    spec = _make_spec()
    series = _make_funnel_series()
    ctx = RenderContext(
        slide=slide,
        slot=slot,
        style=style,
        spec=spec,
        series=series,
        fmt=spec.number_format,
    )
    return prs, ctx


# ---------------------------------------------------------------------------
# Test 1: funnel_spacer_values centers bars correctly
# ---------------------------------------------------------------------------

def test_funnel_spacer_values_center_bars():
    """funnel_spacer_values((80,60,40)) == (0.0, 10.0, 20.0) — each = (max-v)/2."""
    from reportbuilder.render.native.funnel import funnel_spacer_values

    result = funnel_spacer_values((80.0, 60.0, 40.0))
    assert result == (0.0, 10.0, 20.0), (
        f"expected (0.0, 10.0, 20.0), got {result}"
    )


# ---------------------------------------------------------------------------
# Test 2: build_funnel produces a native BAR_STACKED chart (not a picture)
# ---------------------------------------------------------------------------

def test_build_funnel_native_no_picture():
    """build_funnel returns a native c:barChart/stacked with spacer+value series; zero PICTUREs."""
    from reportbuilder.render.native.funnel import build_funnel

    prs, ctx = _make_ctx()
    gf = build_funnel(ctx)

    # Must be a chart graphic frame
    assert gf.has_chart, "returned shape must be a chart graphic frame"

    chart_space = gf.chart._chartSpace
    xml = chart_space.xml

    # barChart element present
    bar_chart_el = chart_space.find(".//" + qn("c:barChart"))
    assert bar_chart_el is not None, "XML must contain c:barChart element"

    # grouping="stacked"
    grouping_el = bar_chart_el.find(qn("c:grouping"))
    assert grouping_el is not None, "c:grouping must be present"
    assert grouping_el.get("val") == "stacked", (
        f"c:grouping val must be 'stacked', got {grouping_el.get('val')!r}"
    )

    # overlap=100
    plot = gf.chart.plots[0]
    assert plot.overlap == 100, f"plot.overlap must be 100, got {plot.overlap}"

    # Exactly two series (spacer + value)
    assert len(plot.series) == 2, f"expected 2 series (spacer + value), got {len(plot.series)}"

    # Spacer series fill is transparent (background/noFill)
    spacer_fill_type = plot.series[0].format.fill.type
    assert spacer_fill_type == MSO_FILL.BACKGROUND, (
        f"spacer series fill must be MSO_FILL.BACKGROUND (noFill), got {spacer_fill_type}"
    )
    # Also verify via XML
    assert "a:noFill" in plot.series[0]._element.xml, (
        "spacer series element XML must contain a:noFill"
    )

    # Value series values match original data
    value_series_values = tuple(plot.series[1].values)
    assert value_series_values == VALUES, (
        f"value series values must be {VALUES}, got {value_series_values}"
    )

    # REQ-C-23a: Zero PICTURE shapes — funnel must be native, never an image
    pictures = [s for s in ctx.slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures == [], (
        f"Found {len(pictures)} PICTURE shape(s) — funnel must be a native chart, not a picture"
    )


# ---------------------------------------------------------------------------
# Test 3: NATIVE_BUILDERS registry
# ---------------------------------------------------------------------------

def test_registry_funnel():
    """NATIVE_BUILDERS['funnel'] must be build_funnel."""
    from reportbuilder.render.native import NATIVE_BUILDERS
    from reportbuilder.render.native.funnel import build_funnel

    assert NATIVE_BUILDERS["funnel"] is build_funnel, (
        "NATIVE_BUILDERS['funnel'] must be build_funnel"
    )
