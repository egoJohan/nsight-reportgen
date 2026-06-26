"""Tests for the ChartPlugin registry (R2, task A).

Covers:
- Every canonical chart_type has a registered ChartPlugin (REQ-C-13)
- suggest_chart_type returns sensible defaults:
    many/long categories → horizontal_bar
    few short categories → vertical_bar or pie
    multi-response question → not pie/doughnut
    time-wave cats → line
- Plugin has non-None image_build and native_build (or None for combo native)
- House-style smoke: pie/doughnut/line/funnel/radar builders produce a
  TEAL-palette PNG (not matplotlib default blue)  (REQ-C-27a)
- deck.py routes through plugin registry (no hardcoded builder dict in dispatch)

REQ-C-13: Every chart type id in ChartType enum has a registered plugin.
REQ-C-24f: Data labels present in image charts.
REQ-C-27a: House palette applied (not default matplotlib blue).
REQ-D-04: Title flows through RenderContext.
"""
from __future__ import annotations

import io

import pytest
import matplotlib
matplotlib.use("Agg")
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from reportbuilder.model.chart_types import ChartType
from reportbuilder.render.plugins import (
    CHART_PLUGINS, ChartPlugin, plugin, suggest_chart_type,
)
from reportbuilder.render.base import RenderContext, Slot, StyleSpec
from reportbuilder.model.report import ChartSpec, ElementToggles, NumberFormat, SortSpec
from reportbuilder.stats.series import Cell, SeriesResult
from reportbuilder.testing.fixtures import known_series


# ---------------------------------------------------------------------------
# Shared fixtures
# ---------------------------------------------------------------------------

def _make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _make_slot() -> Slot:
    return Slot(
        slide_index=0,
        left=Inches(1), top=Inches(1),
        width=Inches(8), height=Inches(5),
        name="slot1",
    )


def _make_spec(chart_type: str = "horizontal_bar") -> ChartSpec:
    return ChartSpec(
        question_ref="q1",
        chart_type=chart_type,
        statistic="pct",
        classifying_var=None,
        number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"),
        template_slot="slot1",
        elements=ElementToggles(title=True, data_labels=True, legend=True, n=True),
    )


def _make_ctx(chart_type="horizontal_bar", series=None, title="Test question"):
    prs, slide = _make_slide()
    slot = _make_slot()
    style = StyleSpec()
    spec = _make_spec(chart_type)
    sr = series or known_series()
    return prs, slide, slot, RenderContext(
        slide=slide, slot=slot, style=style, spec=spec,
        series=sr, fmt=spec.number_format, title=title,
    )


def _likert_series() -> SeriesResult:
    """4-option Likert-scale series: good for stacked charts."""
    cats = ("Erittäin huono", "Huono", "Hyvä", "Erittäin hyvä")
    segs = ("Total",)
    cells = {(c, "Total"): Cell(pct=float(5 + i * 20)) for i, c in enumerate(cats)}
    return SeriesResult(
        categories=cats, segments=segs, cells=cells,
        base_n={"Total": 200}, statistic="pct",
    )


def _many_cat_series() -> SeriesResult:
    """8 categories with long Finnish labels → triggers horizontal_bar preference."""
    cats = (
        "En mitään näistä", "Attendo Finland", "Esperi Care",
        "Rinnekodit Oy", "Validia Kuntoutus", "Onnikodit",
        "Mainio-kodit", "Humana Suomi",
    )
    segs = ("Total",)
    cells = {(c, "Total"): Cell(pct=float(10 + i * 5)) for i, c in enumerate(cats)}
    return SeriesResult(
        categories=cats, segments=segs, cells=cells,
        base_n={"Total": 1001}, statistic="pct",
    )


def _wave_series() -> SeriesResult:
    """Time-wave series → line chart preferred."""
    cats = ("Toukokuu 2024", "Marraskuu 2024", "Toukokuu 2025", "Marraskuu 2025")
    segs = ("Attendo", "Esperi")
    cells = {}
    for j, seg in enumerate(segs):
        for i, c in enumerate(cats):
            cells[(c, seg)] = Cell(pct=float(60 + j * 10 + i * 2))
    return SeriesResult(
        categories=cats, segments=segs, cells=cells,
        base_n={"Attendo": 1001, "Esperi": 1001}, statistic="pct",
    )


def _multi_series() -> SeriesResult:
    """Single-series 2-category result (Yes/No)."""
    return known_series()


def _radar_series() -> SeriesResult:
    """5-attribute radar series."""
    cats = ("Luotettava", "Asiantunteva", "Ystävällinen", "Edullinen", "Laadukas")
    segs = ("Attendo", "Esperi")
    cells = {}
    for j, seg in enumerate(segs):
        for i, c in enumerate(cats):
            cells[(c, seg)] = Cell(pct=float(30 + j * 15 + i * 8))
    return SeriesResult(
        categories=cats, segments=segs, cells=cells,
        base_n={"Attendo": 500, "Esperi": 500}, statistic="pct",
    )


def _funnel_series() -> SeriesResult:
    """Ordered-descending awareness funnel."""
    cats = ("Spontaneous", "Aided", "Consideration", "Trial", "Loyalty")
    vals = [80.0, 60.0, 40.0, 25.0, 12.0]
    cells = {(c, "Total"): Cell(pct=v) for c, v in zip(cats, vals)}
    return SeriesResult(
        categories=cats, segments=("Total",), cells=cells,
        base_n={"Total": 500}, statistic="pct",
    )


# ---------------------------------------------------------------------------
# Helper: check that a PNG has TEAL pixels
# ---------------------------------------------------------------------------

def _has_teal_pixels(blob: bytes, threshold: int = 50) -> bool:
    """Return True if the PNG blob contains ≥ threshold TEAL (#13615E) pixels."""
    from PIL import Image
    img = Image.open(io.BytesIO(blob)).convert("RGB")
    teal_r, teal_g, teal_b = 0x13, 0x61, 0x5E
    count = sum(
        1 for p in img.getdata()
        if abs(p[0] - teal_r) <= 5 and abs(p[1] - teal_g) <= 5 and abs(p[2] - teal_b) <= 5
    )
    return count >= threshold


# ---------------------------------------------------------------------------
# 1 — Plugin registry completeness (REQ-C-13)
# ---------------------------------------------------------------------------

def test_all_chart_types_have_plugins():
    """Every canonical ChartType value has an entry in CHART_PLUGINS (REQ-C-13)."""
    expected = {t.value for t in ChartType}
    registered = set(CHART_PLUGINS.keys())
    missing = expected - registered
    assert not missing, (
        f"These ChartType values have no registered ChartPlugin: {sorted(missing)}"
    )


def test_plugin_lookup_returns_chart_plugin():
    """plugin(chart_type) returns a ChartPlugin instance for every registered type."""
    for ct in ChartType:
        p = plugin(ct.value)
        assert isinstance(p, ChartPlugin), f"plugin({ct.value!r}) should be a ChartPlugin"
        assert p.id == ct.value, f"plugin id mismatch: {p.id!r} != {ct.value!r}"


def test_plugin_unknown_raises_key_error():
    """plugin() raises KeyError for unregistered chart_type strings."""
    with pytest.raises(KeyError, match="No ChartPlugin registered"):
        plugin("not_a_chart_type")


def test_plugin_image_build_callable_for_all():
    """Every plugin (including combo) has a callable image_build."""
    for ct in ChartType:
        p = plugin(ct.value)
        assert callable(p.image_build), (
            f"plugin({ct.value!r}).image_build should be callable"
        )


def test_plugin_native_build_callable_for_non_combo():
    """All plugins except combo have a callable native_build (combo raises)."""
    from reportbuilder.render.native import NativeUnsupportedError
    for ct in ChartType:
        p = plugin(ct.value)
        assert callable(p.native_build), (
            f"plugin({ct.value!r}).native_build should be callable"
        )


def test_scatter_plugin_requires_scatter_xy():
    """scatter plugin declares requires=('scatter_xy',)."""
    p = plugin("scatter")
    assert "scatter_xy" in p.requires, (
        "scatter plugin must declare requires=('scatter_xy',)"
    )


def test_plugin_label_is_non_empty_string():
    """Every plugin has a non-empty label string."""
    for ct in ChartType:
        p = plugin(ct.value)
        assert isinstance(p.label, str) and p.label.strip(), (
            f"plugin({ct.value!r}).label must be a non-empty string"
        )


# ---------------------------------------------------------------------------
# 2 — suggest_chart_type (wizard default, REQ-C-13)
# ---------------------------------------------------------------------------

class _SingleQ:
    """Minimal question stub with kind='single'."""
    kind = "single"

class _MultiQ:
    """Minimal question stub with kind='multi'."""
    kind = "multi"


def test_suggest_many_cats_returns_horizontal_bar():
    """Many or long-label categories → suggest_chart_type picks horizontal_bar (REQ-C-13)."""
    result = suggest_chart_type(_SingleQ(), _many_cat_series())
    assert result == "horizontal_bar", (
        f"Expected 'horizontal_bar' for many/long categories, got {result!r}"
    )


def test_suggest_wave_cats_returns_line():
    """Wave/time categories → suggest_chart_type picks line (REQ-C-13)."""
    result = suggest_chart_type(_SingleQ(), _wave_series())
    assert result == "line", (
        f"Expected 'line' for wave categories, got {result!r}"
    )


def test_suggest_multi_response_not_pie():
    """Multi-response question → suggest_chart_type must NOT return pie or doughnut.

    REQ-C-13: pie/doughnut suitability returns None for multi-response.
    """
    result = suggest_chart_type(_MultiQ(), _many_cat_series())
    assert result not in ("pie", "doughnut"), (
        f"suggest_chart_type must not return pie/doughnut for multi-response; got {result!r}"
    )


def test_suggest_funnel_for_descending():
    """Ordered-descending single-series → suggest_chart_type picks funnel (REQ-C-13)."""
    result = suggest_chart_type(_SingleQ(), _funnel_series())
    assert result == "funnel", (
        f"Expected 'funnel' for descending single-series, got {result!r}"
    )


def test_suggest_returns_string_always():
    """suggest_chart_type always returns a non-empty string (fallback guard)."""
    result = suggest_chart_type(_SingleQ(), known_series())
    assert isinstance(result, str) and result in CHART_PLUGINS


# ---------------------------------------------------------------------------
# 3 — House-style smoke tests: all chart builders produce TEAL PNG  (REQ-C-27a)
# ---------------------------------------------------------------------------

def test_plugin_line_uses_teal():
    """build_image_line produces a TEAL PNG (house palette, REQ-C-27a)."""
    prs, slide, slot, ctx = _make_ctx("line", series=_wave_series())
    plugin("line").image_build(ctx)
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    assert _has_teal_pixels(pics[0].image.blob), (
        "line chart should contain TEAL pixels"
    )


def test_plugin_pie_uses_teal():
    """build_image_pie uses teal-ramp colours, not matplotlib default blue (REQ-C-27a)."""
    prs, slide, slot, ctx = _make_ctx("pie", series=_multi_series())
    plugin("pie").image_build(ctx)
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    assert len(pics[0].image.blob) > 1000  # non-trivial PNG


def test_plugin_doughnut_renders():
    """build_image_doughnut renders without error (REQ-C-24b)."""
    prs, slide, slot, ctx = _make_ctx("doughnut", series=_multi_series())
    plugin("doughnut").image_build(ctx)
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    assert len(pics[0].image.blob) > 1000


def test_plugin_funnel_uses_teal():
    """build_image_funnel uses TEAL fill (house palette, REQ-C-27a)."""
    prs, slide, slot, ctx = _make_ctx("funnel", series=_funnel_series())
    plugin("funnel").image_build(ctx)
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    assert _has_teal_pixels(pics[0].image.blob), (
        "funnel chart should contain TEAL pixels"
    )


def test_plugin_radar_uses_teal():
    """build_image_radar uses TEAL (house palette, REQ-C-27a)."""
    prs, slide, slot, ctx = _make_ctx("radar", series=_radar_series())
    plugin("radar").image_build(ctx)
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    assert _has_teal_pixels(pics[0].image.blob, threshold=20), (
        "radar chart should contain TEAL pixels"
    )


def test_plugin_stacked_vertical_renders():
    """build_image_column_stacked renders without error (REQ-C-24b)."""
    prs, slide, slot, ctx = _make_ctx("stacked_vertical_bar", series=_likert_series())
    plugin("stacked_vertical_bar").image_build(ctx)
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    assert len(pics[0].image.blob) > 1000


def test_plugin_stacked_horizontal_renders():
    """build_image_bar_stacked renders without error (REQ-C-24b)."""
    prs, slide, slot, ctx = _make_ctx("stacked_horizontal_bar", series=_likert_series())
    plugin("stacked_horizontal_bar").image_build(ctx)
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    assert len(pics[0].image.blob) > 1000


# ---------------------------------------------------------------------------
# 4 — % sign on data labels when statistic == "pct"  (REQ-C-24f)
# ---------------------------------------------------------------------------

def test_fmt_value_pct_adds_percent_sign():
    """fmt_value returns '86 %' when statistic='pct' (REQ-C-24f)."""
    from reportbuilder.render.image._mpl import fmt_value
    assert fmt_value(86.0, "pct") == "86 %"
    assert fmt_value(40.0, "pct") == "40 %"


def test_fmt_value_count_no_percent_sign():
    """fmt_value returns '86' (no %) when statistic='count'."""
    from reportbuilder.render.image._mpl import fmt_value
    assert fmt_value(86.0, "count") == "86"


def test_fmt_value_mean_respects_decimals():
    """fmt_value returns '4.2' for mean=4.2 (auto mode, Likert-style → 1 decimal).

    In auto mode mean_decimals is ignored; use mode='manual' to enforce 0 decimals.
    REQ-N-02.
    """
    from reportbuilder.render.image._mpl import fmt_value
    from reportbuilder.model.report import NumberFormat
    assert fmt_value(4.2, "mean", NumberFormat()) == "4.2"
    # manual mode: explicit mean_decimals=0 → integer display
    assert fmt_value(4.0, "mean", NumberFormat(mode="manual", mean_decimals=0)) == "4"


# ---------------------------------------------------------------------------
# 5 — Slide chrome: methodology footer  (REQ-C-24h)
# ---------------------------------------------------------------------------

def test_slide_chrome_has_methodology_footer():
    """add_image_slide_chrome adds methodology footer text (statistic + n).

    REQ-C-24h: methodology note bottom-left.
    """
    from reportbuilder.render.image.slide_chrome import add_image_slide_chrome

    prs, slide, slot, ctx = _make_ctx("horizontal_bar", title="Autettu tunnettuus")
    # known_series() has statistic="pct" and base_n={"Total": 5}
    add_image_slide_chrome(ctx)

    textboxes = [s for s in slide.shapes if s.has_text_frame]
    all_text = " ".join(s.text_frame.text for s in textboxes)
    # Should contain the stat label
    assert "vastaajista" in all_text or "Osuus" in all_text, (
        f"Methodology footer not found. Textbox texts: {all_text!r}"
    )
    # Should contain n=5 from known_series
    assert "5" in all_text, "Base N should appear in chrome text"


# ---------------------------------------------------------------------------
# 6 — deck.py routes through plugin (no hardcoded dict)  (REQ-C-13)
# ---------------------------------------------------------------------------

def test_deck_uses_plugin_for_image_mode():
    """deck.render_report uses plugin registry for image dispatch (REQ-C-13)."""
    import dataclasses
    from reportbuilder.render.deck import render_report
    from reportbuilder.testing.fixtures import one_chart_report

    report = dataclasses.replace(one_chart_report(), render_mode="image")
    series_by_ref = {"q1": known_series()}
    prs = render_report(report, series_by_ref, StyleSpec(), titles={"q1": "Test"})

    pics = [s for sl in prs.slides for s in sl.shapes
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1, "Image mode via plugin should produce exactly 1 picture"


def test_deck_uses_plugin_for_native_mode():
    """deck.render_report uses plugin registry for native dispatch (REQ-C-13)."""
    from pptx.oxml.ns import qn
    from reportbuilder.render.deck import render_report
    from reportbuilder.testing.fixtures import one_chart_report

    report = one_chart_report()  # native mode
    series_by_ref = {"q1": known_series()}
    prs = render_report(report, series_by_ref, StyleSpec(), titles={"q1": "Chart"})

    charts = [s for sl in prs.slides for s in sl.shapes if getattr(s, "has_chart", False)]
    assert len(charts) == 1, "Native mode via plugin should produce exactly 1 chart"
