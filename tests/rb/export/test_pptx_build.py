from pptx import Presentation
from reportbuilder.export.pptx_build import build_pptx
from reportbuilder.testing.fixtures import tiny_model_and_data, one_chart_report
from reportbuilder.model.report import ChartSpec, Report, SortSpec, NumberFormat, ElementToggles


def test_build_pptx_writes_valid_deck(tmp_path):
    model, data = tiny_model_and_data()
    report = one_chart_report()                       # native, 1 chart, question_ref "q1"
    out = build_pptx(report, model, data, str(tmp_path / "r.pptx"))
    assert out.endswith("r.pptx")
    prs = Presentation(out)                            # reopens without exception (REQ-C-29a)
    chart_slides = [s for s in prs.slides if any(getattr(sh, "has_chart", False) for sh in s.shapes)]
    assert len(chart_slides) == len(report.charts)    # completeness (REQ-C-18)


def test_build_pptx_two_charts(tmp_path):
    """Two charts over distinct categorical question_refs — REQ-C-18 completeness."""
    import pandas as pd
    from reportbuilder.model.question import Variable, ValueLabel, Question, QuestionModel

    # Build a two-categorical-question model (age is scale/no labels so not usable as-is;
    # replace with a second categorical question q2 so both chart refs compute cleanly).
    q1_var = Variable(name="q1", label="Satisfaction", measurement="categorical",
                      value_labels=(ValueLabel(1.0, "Yes"), ValueLabel(2.0, "No")),
                      missing_values=frozenset())
    q2_var = Variable(name="q2", label="Likelihood", measurement="categorical",
                      value_labels=(ValueLabel(1.0, "High"), ValueLabel(2.0, "Low")),
                      missing_values=frozenset())
    model = QuestionModel(
        variables={"q1": q1_var, "q2": q2_var},
        questions=[
            Question(qid="q1", kind="single", variables=("q1",), text="Satisfaction"),
            Question(qid="q2", kind="single", variables=("q2",), text="Likelihood"),
        ],
    )
    data = pd.DataFrame({"q1": [1.0, 1.0, 2.0, 2.0, 1.0], "q2": [1.0, 2.0, 1.0, 2.0, 1.0]})

    def _spec(question_ref: str, slot: str) -> ChartSpec:
        return ChartSpec(
            question_ref=question_ref,
            chart_type="vertical_bar",
            statistic="pct",
            classifying_var=None,
            number_format=NumberFormat(),
            sort=SortSpec(basis="data_order"),
            template_slot=slot,
            elements=ElementToggles(),
        )

    report = Report(
        name="R2",
        render_mode="native",
        template_ref="t.pptx",
        charts=(_spec("q1", "slot1"), _spec("q2", "slot2")),
    )
    out = build_pptx(report, model, data, str(tmp_path / "two.pptx"))
    assert out.endswith("two.pptx")
    prs = Presentation(out)                            # reopens without exception (REQ-C-29a)
    chart_slides = [s for s in prs.slides if any(getattr(sh, "has_chart", False) for sh in s.shapes)]
    assert len(chart_slides) == len(report.charts)    # 2 chart slides (REQ-C-18)


def test_build_pptx_cancels_when_signalled(tmp_path):
    """A cancel_check that returns True aborts the build with RenderCancelled — so a
    mistakenly-started long run stops promptly instead of grinding to the end."""
    import pytest
    from reportbuilder.render.deck import RenderCancelled
    model, data = tiny_model_and_data()
    report = one_chart_report()
    with pytest.raises(RenderCancelled):
        build_pptx(report, model, data, str(tmp_path / "x.pptx"),
                   cancel_check=lambda: True)
    # A falsey cancel_check builds normally.
    out = build_pptx(report, model, data, str(tmp_path / "ok.pptx"),
                     cancel_check=lambda: False)
    assert out.endswith("ok.pptx")
