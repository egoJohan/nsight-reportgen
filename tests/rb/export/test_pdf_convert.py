"""TDD tests for LibreOffice PPTX->PDF conversion and PDF page rasterization (Task 1.3)."""
from __future__ import annotations
import os
import shutil
import subprocess

import pytest
from pptx import Presentation
from pptx.util import Inches

from reportbuilder.render.base import Slot
from reportbuilder.render.native.column import build_column_chart
from reportbuilder.testing.fixtures import known_series

pytestmark = pytest.mark.skipif(
    shutil.which("soffice") is None,
    reason="LibreOffice not installed",
)


def _make_deck(tmp_path) -> str:
    """Build a one-chart PPTX using build_column_chart and save to tmp_path."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slot = Slot(
        slide_index=0,
        left=Inches(1),
        top=Inches(1),
        width=Inches(8),
        height=Inches(5),
        name="slot1",
    )
    series = known_series()
    build_column_chart(slide, slot, series)
    deck_path = str(tmp_path / "deck.pptx")
    prs.save(deck_path)
    return deck_path


def test_pptx_to_pdf_produces_pdf(tmp_path):
    """pptx_to_pdf converts a real chart deck to a non-empty PDF."""
    from reportbuilder.export.pdf_convert import pptx_to_pdf

    deck = _make_deck(tmp_path)
    pdf = pptx_to_pdf(deck, str(tmp_path))

    assert pdf.endswith(".pdf"), f"expected .pdf extension, got {pdf!r}"
    assert os.path.getsize(pdf) > 0, "PDF must be non-empty"


def test_pdf_page_to_png(tmp_path):
    """pdf_page_to_png rasterizes first page of the converted PDF to a non-empty PNG."""
    from reportbuilder.export.pdf_convert import pptx_to_pdf
    from reportbuilder.export.preview import pdf_page_to_png

    deck = _make_deck(tmp_path)
    pdf = pptx_to_pdf(deck, str(tmp_path))
    png = pdf_page_to_png(pdf, 0, str(tmp_path / "p0.png"))

    assert os.path.getsize(png) > 0, "PNG must be non-empty"


def test_a_long_deck_converts_in_parallel_and_keeps_its_order(tmp_path):
    """Slices go to separate soffice processes and poppler stitches them back.

    The page count must match the deck and the pages must stay in order — a
    merge that reordered or dropped slides would be invisible in a thumbnail
    strip and catastrophic in front of a client.
    """
    import shutil as _shutil

    from pptx import Presentation
    from pptx.util import Inches, Pt

    from reportbuilder.export.pdf_convert import (
        pdf_page_count, pptx_to_pdf_parallel,
    )

    if _shutil.which("pdfunite") is None:
        pytest.skip("poppler pdfunite required")

    # 20 numbered slides: enough to cross the threshold where slicing pays off.
    prs = Presentation()
    for i in range(20):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = f"SLIDE {i + 1}"
        run.font.size = Pt(54)
    src = str(tmp_path / "long.pptx")
    prs.save(src)

    pdf = pptx_to_pdf_parallel(src, str(tmp_path / "out"))
    assert pdf_page_count(pdf) == 20

    text = subprocess.run(["pdftotext", pdf, "-"], capture_output=True,
                          text=True, check=True).stdout
    found = [line.strip() for line in text.splitlines() if line.strip().startswith("SLIDE")]
    assert found == [f"SLIDE {i + 1}" for i in range(20)]


def test_a_short_deck_is_not_sliced(tmp_path, monkeypatch):
    """Splitting costs a LibreOffice startup per slice; below the threshold that
    is slower than just converting the deck."""
    from pptx import Presentation

    import reportbuilder.export.pdf_convert as pc

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    src = str(tmp_path / "short.pptx")
    prs.save(src)

    calls = []
    monkeypatch.setattr(pc, "pptx_to_pdf",
                        lambda path, out, priority=False: calls.append(path) or "x.pdf")
    assert pc.pptx_to_pdf_parallel(src, str(tmp_path)) == "x.pdf"
    assert calls == [src]                      # the whole deck, in one go


def test_parallelism_follows_the_cores_this_process_may_use(monkeypatch):
    """Not the machine's cores — the ones we are scheduled on.

    A render pinned to one core (taskset, or a container given one) would
    otherwise slice its deck six ways and pay six LibreOffice startups for
    parallelism that cannot happen.
    """
    from reportbuilder.export import _cpu

    monkeypatch.setattr(_cpu.os, "process_cpu_count", lambda: 1, raising=False)
    assert _cpu.usable_cores() == 1
    assert _cpu.workers_for(6) == 1              # one core -> one worker

    monkeypatch.setattr(_cpu.os, "process_cpu_count", lambda: 16, raising=False)
    assert _cpu.workers_for(6) == 6              # capped: RAM runs out first
    assert _cpu.workers_for(32) == 15            # one core left for the server
