import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


@pytest.fixture
def tmp_native_pptx(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = ["Yes", "No"]
    cd.add_series("Total", (60.0, 40.0))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(6), Inches(4), cd)
    out = tmp_path / "deck.pptx"
    prs.save(str(out))
    return str(out)
