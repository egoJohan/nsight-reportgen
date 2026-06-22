import pytest
from pathlib import Path
from nsight.render.template import Template
from nsight.render.fill_text import set_lines, set_text


@pytest.fixture
def text_pptx(tmp_path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.name = "key_message"
    tb.text_frame.text = "PLACEHOLDER"
    out = tmp_path / "x.pptx"
    prs.save(out)
    return out


def test_set_text_replaces_first_paragraph_keeping_box(text_pptx, tmp_path):
    tmpl = Template(text_pptx)
    shape = tmpl.shape(slide_idx=0, name="key_message")
    set_text(shape, "Attendo on tunnetuin.")
    out = tmp_path / "o.pptx"
    tmpl.save(out)
    from pptx import Presentation
    sh = [s for s in Presentation(str(out)).slides[0].shapes if s.name == "key_message"][0]
    assert sh.text_frame.text == "Attendo on tunnetuin."


@pytest.fixture
def header_lines_pptx(tmp_path: Path) -> Path:
    """A box with a header paragraph + several word paragraphs (deck-like)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(4))
    tb.name = "top10"
    tf = tb.text_frame
    tf.paragraphs[0].text = "TOP 10"
    p_label = tf.add_paragraph()
    p_label.text = "Marras 25"
    tf.add_paragraph().text = ""
    # word paragraphs (the slots that should be overwritten); first carries the font.
    p0 = tf.add_paragraph()
    r0 = p0.add_run()
    r0.text = "Old (1)"
    r0.font.size = Pt(12)
    for w in ("Old (2)", "Old (3)"):
        tf.add_paragraph().text = w
    out = tmp_path / "h.pptx"
    prs.save(out)
    return out


def test_set_lines_preserves_header_and_word_font(header_lines_pptx, tmp_path):
    from pptx.util import Pt
    tmpl = Template(header_lines_pptx)
    shape = tmpl.shape(slide_idx=0, name="top10")
    set_lines(shape, ["Kallis (185)", "Huono (75)", "Kiire (66)", "Hyvä (60)"], start=3)
    out = tmp_path / "o.pptx"
    tmpl.save(out)
    from pptx import Presentation
    sh = [s for s in Presentation(str(out)).slides[0].shapes if s.name == "top10"][0]
    paras = sh.text_frame.paragraphs
    # headers preserved
    assert [p.text for p in paras[:3]] == ["TOP 10", "Marras 25", ""]
    # exactly the 4 new word lines, no leftover "Old (3)"
    assert [p.text for p in paras[3:]] == ["Kallis (185)", "Huono (75)", "Kiire (66)", "Hyvä (60)"]
    # font of the original first word paragraph is carried onto every new line
    for p in paras[3:]:
        assert p.runs[0].font.size == Pt(12)
