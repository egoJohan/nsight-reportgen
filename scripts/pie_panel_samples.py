"""Visual samples for the split pie/doughnut/funnel panels (spec 2026-08-22).

Two levels, because they catch different things:

  C. REPORT — one openable PPTX/PDF whose slides walk the SAME question through
     1, 2, 3 and 4 background-variable groups, so the cap and the footer can be
     compared side by side. This is the one to send to a human.

  A. FIGURE level — synthetic series straight into `_build_pie_figure`, so every
     edge case is reachable on demand: the degraded split, the single surviving
     group, the cap, a count statistic, six categories with long labels. Fast, no
     LibreOffice.

  B. SLIDE level — real survey data through `build_pptx` → PDF → PNG, so what you
     look at is the actual deliverable: chart in its slot, title, and the
     methodology footer that names any omitted group.

Usage:
    uv run python scripts/pie_panel_samples.py            # both levels
    uv run python scripts/pie_panel_samples.py figure     # A only (no soffice)
    uv run python scripts/pie_panel_samples.py slide      # B only
    uv run python scripts/pie_panel_samples.py report     # C only — the deck to review

Output lands in `work/pie_panels/` — `work/` is git-ignored, and NOT /tmp, which
is a ramfs on the dev machine.
"""
from __future__ import annotations

import os
import shutil
import sys
import tempfile

from pptx import Presentation
from pptx.util import Inches

from reportbuilder.model.report import (
    ChartSpec, ElementToggles, NumberFormat, Report, SortSpec,
)
from reportbuilder.render.base import RenderContext, Slot, StyleSpec
from reportbuilder.render.image._mpl import render_png
from reportbuilder.stats.series import Cell, SeriesResult

OUT = os.path.join("work", "pie_panels")

# A five-point opinion scale with real Finnish labels — long enough to push the
# shared legend, which is where the panel layout is most likely to break.
LONG_CATS = ("Erittäin tyytyväinen", "Melko tyytyväinen", "Ei kumpaakaan",
             "Melko tyytymätön", "Erittäin tyytymätön", "En osaa sanoa")


def _ctx(series, chart_type="pie", **spec_kw):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slot = Slot(slide_index=0, left=Inches(1), top=Inches(1),
                width=Inches(8), height=Inches(5), name="slot1")
    kw = dict(question_ref="q1", chart_type=chart_type, statistic="pct",
              classifying_var=None, number_format=NumberFormat(),
              sort=SortSpec(basis="data_order"), template_slot="slot1",
              elements=ElementToggles(title=True, legend=True, data_labels=True))
    kw.update(spec_kw)
    spec = ChartSpec(**kw)
    return RenderContext(slide=slide, slot=slot, style=StyleSpec(),
                         spec=spec, series=series, fmt=spec.number_format)


def _series(cats, segments, bases, shares, statistic="pct"):
    """`shares` maps segment -> per-category values (already within-group %)."""
    cells = {}
    for seg in segments:
        vals = shares[seg]
        for c, v in zip(cats, vals):
            cells[(c, seg)] = Cell(pct=float(v), count=float(v), mean=None)
    return SeriesResult(categories=tuple(cats), segments=tuple(segments),
                        cells=cells, base_n=dict(bases), statistic=statistic)


def _save(fig, name):
    # `render_png` writes into the system temp dir, which is a separate filesystem
    # here (/tmp is a ramfs), so os.replace would raise EXDEV. Copy, then drop the
    # ramfs original rather than leaving it resident in RAM.
    src = render_png(fig)
    dest = os.path.join(OUT, name + ".png")
    shutil.copyfile(src, dest)
    try:
        os.remove(src)
    except OSError:
        pass
    print("  ok:", dest)


def figure_samples():
    """Level A — every branch of the panel renderer, including the ones real data
    rarely produces."""
    from reportbuilder.render.image.pie import _build_pie_figure

    yes_no = ("En voi", "Kyllä voin")

    def case(name, series, *, donut=False, **spec_kw):
        try:
            _save(_build_pie_figure(_ctx(series, donut and "doughnut" or "pie",
                                         **spec_kw), donut=donut), name)
        except Exception as e:                       # a broken case must be VISIBLE
            print("  FAIL:", name, repr(e)[:160])

    # The regression baseline: no classifier at all. Must look like it always has —
    # one circle, no panel title, legend on the RIGHT.
    case("a01_unsplit", _series(yes_no, ("Total",), {"Total": 1023},
                                {"Total": (54, 46)}))

    # The approved layout: two and three groups.
    case("a02_two_groups",
         _series(yes_no, ("Naiset", "Miehet", "Total"),
                 {"Naiset": 512, "Miehet": 486, "Total": 998},
                 {"Naiset": (54, 46), "Miehet": (41, 59), "Total": (48, 52)}),
         classifying_var="sex")
    case("a03_three_groups",
         _series(yes_no, ("Naiset", "Miehet", "Muut", "Total"),
                 {"Naiset": 512, "Miehet": 486, "Muut": 25, "Total": 1023},
                 {"Naiset": (54, 46), "Miehet": (41, 59), "Muut": (60, 40),
                  "Total": (48, 52)}),
         classifying_var="sex")

    # Five groups: only the three largest are drawn, in the variable's own order.
    ages = ("18-29", "30-44", "45-59", "60-74", "Yli 74")
    case("a04_five_groups_capped",
         _series(yes_no, (*ages, "Total"),
                 {"18-29": 150, "30-44": 290, "45-59": 270, "60-74": 130,
                  "Yli 74": 90, "Total": 930},
                 {a: (50 + i * 3, 50 - i * 3) for i, a in enumerate(ages)}
                 | {"Total": (52, 48)}),
         classifying_var="age")

    # One group survives the base floor: still TITLED, so the reader knows which.
    case("a05_one_surviving_group",
         _series(yes_no, ("Naiset", "Miehet", "Total"),
                 {"Naiset": 512, "Miehet": 4, "Total": 516},
                 {"Naiset": (54, 46), "Miehet": (25, 75), "Total": (54, 46)}),
         classifying_var="sex")

    # Every group too thin: degrades to ONE whole-sample pie, untitled.
    case("a06_degraded_all_thin",
         _series(yes_no, ("Naiset", "Miehet", "Total"),
                 {"Naiset": 4, "Miehet": 6, "Total": 10},
                 {"Naiset": (50, 50), "Miehet": (50, 50), "Total": (50, 50)}),
         classifying_var="sex")

    # Six long-labelled categories — the shared legend's worst case.
    six = {"Naiset": (30, 25, 15, 12, 10, 8), "Miehet": (22, 28, 18, 14, 10, 8),
           "Muut": (18, 22, 20, 18, 12, 10), "Total": (25, 25, 17, 15, 11, 7)}
    case("a07_six_long_categories",
         _series(LONG_CATS, ("Naiset", "Miehet", "Muut", "Total"),
                 {"Naiset": 512, "Miehet": 486, "Muut": 120, "Total": 1118}, six),
         classifying_var="sex")

    # Long GROUP names, against a panel a third of the slot wide.
    long_groups = ("Naiseksi itsensä identifioivat", "Mieheksi itsensä identifioivat",
                   "Muuksi tai ei halua sanoa")
    case("a08_long_group_names",
         _series(yes_no, (*long_groups, "Total"),
                 {g: 300 for g in long_groups} | {"Total": 900},
                 {g: (55 - i * 5, 45 + i * 5) for i, g in enumerate(long_groups)}
                 | {"Total": (50, 50)}),
         classifying_var="sex")

    # A count statistic keeps the engine's "Total" segment alive — it must still
    # never become a fourth circle.
    case("a09_count_statistic",
         _series(yes_no, ("Naiset", "Miehet", "Total"),
                 {"Naiset": 512, "Miehet": 486, "Total": 998},
                 {"Naiset": (276, 236), "Miehet": (199, 287),
                  "Total": (475, 523)}, statistic="count"),
         classifying_var="sex", statistic="count")

    # Doughnut shares the renderer: the hole must survive the panel split.
    case("a10_doughnut_three_groups",
         _series(yes_no, ("Naiset", "Miehet", "Muut", "Total"),
                 {"Naiset": 512, "Miehet": 486, "Muut": 25, "Total": 1023},
                 {"Naiset": (54, 46), "Miehet": (41, 59), "Muut": (60, 40),
                  "Total": (48, 52)}),
         donut=True, classifying_var="sex")

    # Legend switched off: panels keep their titles and bases, no legend band.
    case("a11_legend_off",
         _series(yes_no, ("Naiset", "Miehet", "Total"),
                 {"Naiset": 512, "Miehet": 486, "Total": 998},
                 {"Naiset": (54, 46), "Miehet": (41, 59), "Total": (48, 52)}),
         classifying_var="sex",
         elements=ElementToggles(title=True, legend=False, data_labels=True,
                                 axis_names=False))

    # A near-zero slice: its % label is suppressed on the wedge but the category
    # still belongs to the legend.
    case("a12_tiny_slice",
         _series(("Kyllä", "Ei", "En osaa sanoa"),
                 ("Naiset", "Miehet", "Total"),
                 {"Naiset": 512, "Miehet": 486, "Total": 998},
                 {"Naiset": (62, 37, 1), "Miehet": (55, 44, 1),
                  "Total": (58, 41, 1)}),
         classifying_var="sex")

    # Funnel gets the same split (spec 2026-08-22): three groups, each
    # descending, side by side. `build_image_funnel` places its own picture onto
    # the slide rather than returning a Figure, so pull the PNG back out of the
    # slide's picture shape instead of going through `_save`.
    def case_funnel(name, series, **spec_kw):
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        from reportbuilder.render.image.funnel import build_image_funnel

        try:
            ctx = _ctx(series, "funnel", **spec_kw)
            build_image_funnel(ctx)
            pics = [s for s in ctx.slide.shapes
                    if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
            dest = os.path.join(OUT, name + ".png")
            with open(dest, "wb") as f:
                f.write(pics[0].image.blob)
            print("  ok:", dest)
        except Exception as e:                       # a broken case must be VISIBLE
            print("  FAIL:", name, repr(e)[:160])

    funnel_cats = ("Tuntee", "Harkitsee", "Ostanut")
    # TWO groups — the likeliest split of all (by gender), and an EVEN panel
    # count: `invert_yaxis()` once per panel on a SHARED y-axis used to cancel
    # itself here and draw the funnel upside down, while the three-group case
    # below came out right by parity. Both must point the same way.
    case_funnel(
        "a13b_funnel_two_groups",
        _series(funnel_cats, ("Naiset", "Miehet", "Total"),
                {"Naiset": 512, "Miehet": 486, "Total": 998},
                {"Naiset": (85, 55, 20), "Miehet": (42, 27, 10),
                 "Total": (64, 41, 15)}),
        classifying_var="sex")
    case_funnel(
        "a13_funnel_three_groups",
        _series(funnel_cats, ("Naiset", "Miehet", "Muut", "Total"),
                {"Naiset": 512, "Miehet": 486, "Muut": 120, "Total": 1118},
                {"Naiset": (85, 55, 20), "Miehet": (78, 48, 15),
                 "Muut": (90, 60, 25), "Total": (83, 53, 20)}),
        classifying_var="sex")


def slide_samples():
    """Level B — real survey data, rendered as whole slides through PowerPoint."""
    from reportbuilder.export.pdf_convert import pptx_to_pdf
    from reportbuilder.export.preview import rasterize_pages
    from reportbuilder.export.pptx_build import build_pptx
    from reportbuilder.ingest.multi_group import enrich_model
    from reportbuilder.ingest.sav_reader import read_sav

    sav = "input/spss AttendoSuomi-Brandiseuranta_112025.sav"
    if not os.path.exists(sav):
        print("  skip: no SAV at", sav)
        return
    df, model = read_sav(sav)
    model = enrich_model(model)

    # var20: a five-point opinion scale (a real pie subject).
    # var9 = gender (4 groups, two of them small), var7 = age (7 groups),
    # var10 = region (5 groups). Between them: the base floor, the cap, and both.
    cases = [
        ("b01_opinion_unsplit", "var20", None, "pie"),
        ("b02_opinion_by_gender", "var20", "var9", "pie"),
        ("b03_opinion_by_age_capped", "var20", "var7", "pie"),
        ("b04_opinion_by_region_capped", "var20", "var10", "pie"),
        ("b05_opinion_by_gender_doughnut", "var20", "var9", "doughnut"),
    ]
    for name, ref, clf, ctype in cases:
        sp = ChartSpec(question_ref=ref, chart_type=ctype, statistic="pct",
                       classifying_var=clf, number_format=NumberFormat(),
                       sort=SortSpec(basis="data_order"), template_slot="s1",
                       elements=ElementToggles(), options={})
        d = tempfile.mkdtemp()
        try:
            pptx = os.path.join(d, "x.pptx")
            build_pptx(Report(name="pie-panels", render_mode="image",
                              template_ref="", charts=(sp,)), model, df, pptx)
            pdf = pptx_to_pdf(pptx, d)
            pages = rasterize_pages(pdf, os.path.join(d, "pg"), dpi=110)
            dest = os.path.join(OUT, name + ".png")
            shutil.copyfile(pages[0], dest)
            print("  ok:", dest)
        except Exception as e:
            print("  FAIL:", name, repr(e)[:200])
        finally:
            shutil.rmtree(d, ignore_errors=True)


# --------------------------------------------------------------------------- #
# Level C — the reviewable deck: one question, 1..4 groups.
# --------------------------------------------------------------------------- #
REPORT_OUT = os.path.join("work", "pie_panels_report")

# Age codes present in the Attendo data (var7). Real respondents, so every base
# below is a real base — only the GROUPING is synthesised, to hit exactly 1, 2, 3
# and 4 groups on one question.
_A = {"25_34": 10004.0, "35_44": 10005.0, "45_54": 10006.0,
      "55_64": 10007.0, "65_74": 10008.0}

# name -> (label, {new_code: (new_label, [source age codes])})
_DERIVED = {
    "demo1": ("Kaikki vastaajat", {
        1.0: ("Kaikki vastaajat", list(_A.values())),
    }),
    "demo2": ("Ikä, 2 ryhmää", {
        1.0: ("Alle 45 vuotta", [_A["25_34"], _A["35_44"]]),
        2.0: ("45 vuotta tai yli", [_A["45_54"], _A["55_64"], _A["65_74"]]),
    }),
    "demo3": ("Ikä, 3 ryhmää", {
        1.0: ("Alle 45 vuotta", [_A["25_34"], _A["35_44"]]),
        2.0: ("45–54 vuotta", [_A["45_54"]]),
        3.0: ("55 vuotta tai yli", [_A["55_64"], _A["65_74"]]),
    }),
    "demo4": ("Ikä, 4 ryhmää", {
        1.0: ("25–34 vuotta", [_A["25_34"]]),
        2.0: ("35–44 vuotta", [_A["35_44"]]),
        3.0: ("45–54 vuotta", [_A["45_54"]]),
        4.0: ("55 vuotta tai yli", [_A["55_64"], _A["65_74"]]),
    }),
}


def _add_derived_classifiers(df, model):
    """Add demo1..demo4 — the same respondents grouped 1, 2, 3 and 4 ways.

    Derived rather than picked from the file because no single real variable
    gives a clean 1/2/3/4 progression on one question, and the point of the deck
    is to vary ONLY the group count.
    """
    import pandas as pd
    from reportbuilder.model.question import ValueLabel, Variable

    age = pd.to_numeric(df["var7"], errors="coerce")
    for name, (label, groups) in _DERIVED.items():
        col = pd.Series(float("nan"), index=df.index)
        for code, (_lbl, sources) in groups.items():
            col = col.where(~age.isin(sources), code)
        df[name] = col
        model.variables[name] = Variable(
            name=name, label=label, measurement="categorical",
            value_labels=tuple(ValueLabel(value=c, label=l)
                                for c, (l, _s) in sorted(groups.items())),
            missing_values=frozenset(),
        )
    return df, model


def group_count_report():
    """Level C — one sample per group count, for a human to look at.

    NOTE ON WHY EACH CASE IS ITS OWN REPORT: `_build` caches computed series in
    `series_by_ref` keyed by QUESTION REF alone (`export/pptx_build.py:54,87`) and
    `deck.py:245` reads them back the same way. Put two charts of the SAME question
    in one Report and the last one computed wins — every slide then renders that
    series. So a single 7-slide deck of one question silently showed the same split
    seven times. One Report per case sidesteps it; the underlying defect is
    reported separately, not worked around in shipped code.
    """
    from PIL import Image
    from reportbuilder.export.pdf_convert import pptx_to_pdf
    from reportbuilder.export.preview import rasterize_pages
    from reportbuilder.export.pptx_build import build_pptx
    from reportbuilder.ingest.multi_group import enrich_model
    from reportbuilder.ingest.sav_reader import read_sav

    sav = "input/spss AttendoSuomi-Brandiseuranta_112025.sav"
    if not os.path.exists(sav):
        print("  skip: no SAV at", sav)
        return
    os.makedirs(REPORT_OUT, exist_ok=True)
    df, model = read_sav(sav)
    model = enrich_model(model)
    df, model = _add_derived_classifiers(df, model)

    # (file stem, caption, classifier, chart type) — the SAME question throughout,
    # so the only thing that changes is the number of groups.
    cases = [
        ("0-groups-baseline", "0 ryhmää — ei taustamuuttujaa (nykyinen piirakka)", None, "pie"),
        ("1-group", "1 ryhmä — yksi paneeli, otsikoituna", "demo1", "pie"),
        ("2-groups", "2 ryhmää", "demo2", "pie"),
        ("3-groups", "3 ryhmää — suurin sallittu", "demo3", "pie"),
        ("4-groups-capped", "4 ryhmää — kolme suurinta piirretään, loput alaviitteessä",
         "demo4", "pie"),
        ("3-groups-doughnut", "3 ryhmää — donitsi", "demo3", "doughnut"),
        ("3-groups-funnel", "3 ryhmää — suppilo", "demo3", "funnel"),
    ]

    pngs = []
    for stem, cap, clf, ctype in cases:
        sp = ChartSpec(question_ref="var20", chart_type=ctype, statistic="pct",
                       classifying_var=clf, number_format=NumberFormat(),
                       sort=SortSpec(basis="data_order"), template_slot="s1",
                       elements=ElementToggles(), options={})
        d = tempfile.mkdtemp()
        try:
            pptx = os.path.join(REPORT_OUT, f"{stem}.pptx")
            build_pptx(Report(name=cap, render_mode="image", template_ref="",
                              charts=(sp,)), model, df, pptx)
            pdf = pptx_to_pdf(pptx, d)
            page = rasterize_pages(pdf, os.path.join(d, "pg"), dpi=110)[0]
            dest = os.path.join(REPORT_OUT, f"{stem}.png")
            shutil.copyfile(page, dest)
            pngs.append(dest)
            print(f"  ok: {dest}  — {cap}")
        except Exception as e:
            print("  FAIL:", stem, repr(e)[:200])
        finally:
            shutil.rmtree(d, ignore_errors=True)

    # One flip-through PDF of every case, assembled from the PNGs (no PDF library
    # in this environment; PIL writes multi-page PDFs natively).
    if pngs:
        combined = os.path.join(REPORT_OUT, "ALL-CASES.pdf")
        first, *rest = [Image.open(p).convert("RGB") for p in pngs]
        first.save(combined, save_all=True, append_images=rest)
        print("  ok:", combined)

    with open(os.path.join(REPORT_OUT, "README.md"), "w") as fh:
        fh.write(_REPORT_README)
    print("  ok:", os.path.join(REPORT_OUT, "README.md"))


_REPORT_README = """# Useampi piirakkakuvio samalla sivulla — samples

Flip through `ALL-CASES.pdf`, or open any single case's `.pptx` / `.png`. Every case
charts the SAME question —
*"Mikä on yleinen käsityksesi ... yksityisistä yrityksistä?"*, a five-point scale —
so the only thing that changes down the deck is how many background-variable groups
it is split by.

The classifiers `demo1`..`demo4` are real respondents regrouped, so every base is a
real base; only the grouping is synthesised, because no single variable in the file
gives a clean 1-2-3-4 progression on one question.

| File | Groups | What to look at |
|---|---|---|
| `0-groups-baseline` | none | The current pie, unchanged. Your baseline — it must look exactly as it does today. |
| `1-group` | 1 | One circle, but TITLED with its group and its own base, so a reader knows which group it is. |
| `2-groups` | 2 | Two pies, shared legend beneath, each its own 100%. |
| `3-groups` | 3 | Three pies — the most that fit. |
| `4-groups-capped` | 4 | Only the three largest are drawn. The footer names the one left out: *Ei mahtunut sivulle: …* |
| `3-groups-doughnut` | 3 | The same split as a doughnut. |
| `3-groups-funnel` | 3 | The same split as a funnel. |

Worth checking on `4-groups-capped`: the omitted group is named in the footer and
nowhere else. That footer line is the only record of the omission that travels with
the deck — the editor's warning stays in the editor.

**Why seven separate files rather than one seven-slide deck.** A Report caches each
computed series by question ref alone, so two charts of the SAME question in one deck
both render the last one's data. The first attempt at this deck showed the identical
3-group split on all five pie slides. That is a pre-existing defect in the deck
builder, unrelated to this feature but made much easier to hit by it — see the branch
notes.

Regenerate with:

    uv run python scripts/pie_panel_samples.py report
"""


def main():
    os.makedirs(OUT, exist_ok=True)
    which = sys.argv[1] if len(sys.argv) > 1 else "both"
    if which in ("both", "figure"):
        print("Level A — figure samples")
        figure_samples()
    if which in ("both", "slide"):
        print("Level B — whole-slide samples (real data)")
        slide_samples()
    if which in ("both", "report"):
        print("Level C — reviewable deck (1..4 groups)")
        group_count_report()


if __name__ == "__main__":
    main()
