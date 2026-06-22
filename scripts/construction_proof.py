"""CONSTRUCTION PROOF.

Answers the question "did the engine actually build the content, or just copy the
original deck?" by:

  1. Blanking (zeroing) every chart series the engine claims to produce, in a copy of
     the ORIGINAL Attendo template -> work/attendo_blanked.pptx
  2. Regenerating from the SPSS file into the blanked deck -> work/attendo_from_blanked.pptx
  3. Showing the regenerated values reconstruct the ORIGINAL deck's numbers.

If the engine were merely copying the template, the blanked charts would stay empty.
They don't: the numbers come back, computed from the .sav.

Run: uv run python scripts/construction_proof.py
"""
from __future__ import annotations

from pptx import Presentation

from nsight import config
from nsight.agent.workflow import generate_fills
from nsight.attendo_bindings import AIDED_AWARENESS_VARS, WEIGHT_VAR
from nsight.brief import parse_brief
from nsight.fidelity.extract import extract_deck
from nsight.render.fill_chart import read_chart_series, replace_one_series
from nsight.render.renderer import render
from nsight.store.survey_store import SurveyStore


def _shape_on(slide, name, occ):
    return [s for s in slide.shapes if s.name == name][occ]


def _series_vals(deck, idx, name, occ, series):
    return [c for c in deck.slides[idx].charts if c.name == name][occ].series.get(series)


def main() -> None:
    store = SurveyStore(db_path=config.SURVEY_DB, codebook_path=config.CODEBOOK_JSON)
    if store.frame().empty:
        store.ingest(config.ATTENDO_SAV)
    frame = store.frame()
    brief = parse_brief(config.BRIEFS_DIR / "attendo.md")
    fills = generate_fills(brief, frame, brand_vars=AIDED_AWARENESS_VARS, weight=WEIGHT_VAR,
                           narrator=lambda j, n: j.title or j.id)

    targets = [(sf.slide_idx, cf.name, cf.occurrence, cf.series_name)
               for sf in fills for cf in sf.charts if cf.series_name]

    # 1) blank those series in a copy of the original
    prs = Presentation(str(config.ATTENDO_TEMPLATE))
    for idx, name, occ, series in targets:
        sh = _shape_on(prs.slides[idx], name, occ)
        n = len(read_chart_series(sh)[series])
        replace_one_series(sh, series_name=series, values=[0.0] * n)
    blanked = config.WORK_DIR / "attendo_blanked.pptx"
    prs.save(str(blanked))

    # 2) regenerate from the blanked template
    regen = config.WORK_DIR / "attendo_from_blanked.pptx"
    render(template_path=blanked, out_path=regen, fills=fills)

    # 3) compare
    orig, bl, rg = (extract_deck(config.ATTENDO_TEMPLATE), extract_deck(blanked),
                    extract_deck(regen))
    print(f"Blanked deck:      {blanked}")
    print(f"Regenerated deck:  {regen}\n")
    exact = within1 = total = 0
    for idx, name, occ, series in targets:
        o = _series_vals(orig, idx, name, occ, series)
        b = _series_vals(bl, idx, name, occ, series)
        r = _series_vals(rg, idx, name, occ, series)
        for ov, bv, rv in zip(o, b, r):
            total += 1
            op = round(ov * 100) if ov is not None else None
            rp = round(rv * 100) if rv is not None else None
            if op == rp:
                exact += 1
            if op is not None and rp is not None and abs(op - rp) <= 1:
                within1 += 1
        print(f"slide {idx} [{name}] '{series}':")
        print(f"  blanked : {[round(x, 3) for x in b]}")
        print(f"  original: {[round(x, 3) if x is not None else None for x in o]}")
        print(f"  REGEN   : {[round(x, 3) if x is not None else None for x in r]}\n")
    print(f"Reconstructed from blank: {exact}/{total} exact, {within1}/{total} within +/-1pp")


if __name__ == "__main__":
    main()
