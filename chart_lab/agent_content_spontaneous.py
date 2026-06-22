#!/usr/bin/env python3
"""Build the spontaneous brand-awareness (top-of-mind) slide — nSight house style.

Horizontal STACKED bar per care provider: a darker "top-of-mind" segment
(mentioned first) + a lighter "muut spontaanit maininnat" segment, summing to
the total spontaneous awareness ("Kaikki"). Providers sorted by total
descending; totals labelled at bar end, top-of-mind share labelled inside the
dark segment. Current wave Marraskuu 2025, n=1001.

All numbers are read verbatim from chart_lab/spontaneous.json. The two aggregate
buckets "Muu" and "En osaa sanoa" are NOT real providers and do not follow the
top-of-mind/muut split (muut is null), so they are excluded from the stacked
chart; the "En osaa sanoa" share is surfaced as a footnote reference.
"""
import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
from matplotlib.patches import Patch

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ---------------------------------------------------------------- paths
ROOT = Path("/home/johan/Projects/nsight/proto")
DATA = ROOT / "chart_lab" / "spontaneous.json"
PNG = ROOT / "chart_lab" / "agent_content_spontaneous_chart.png"
OUT = ROOT / "work" / "agent_content_spontaneous.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------- fonts
for f in ["/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
          "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"]:
    fm.fontManager.addfont(f)
plt.rcParams["font.family"] = "Liberation Sans"

# ---------------------------------------------------------------- data
d = json.loads(DATA.read_text())
cats = d["categories"]
s = d["series"]
TOM = "Top of mind"
MUUT = "Muut spontaanit maininnat"
KAIKKI = "Kaikki"

# Aggregate / non-provider buckets — excluded from the stacked chart.
EXCLUDE = {"Muu", "En osaa sanoa"}

# Build provider records, asserting the split sums to the total (data integrity).
recs = []
for i, name in enumerate(cats):
    if name in EXCLUDE:
        continue
    tom = s[TOM][i]
    muut = s[MUUT][i]
    tot = s[KAIKKI][i]
    assert muut is not None, f"{name}: muut is null"
    assert tom + muut == tot, f"{name}: {tom}+{muut} != {tot}"
    recs.append({"name": name, "tom": tom, "muut": muut, "tot": tot})

# Sort by total spontaneous awareness, descending.
recs.sort(key=lambda r: r["tot"], reverse=True)

# "En osaa sanoa" reference (share who could not name any provider).
eos_idx = cats.index("En osaa sanoa")
EOS = s[KAIKKI][eos_idx]

# ---------------------------------------------------------------- colours
CREAM = "#F7F3EC"
INK = "#2B2B2B"
MUTED = "#6E6A63"
TOM_COLOR = "#13615E"      # dark teal — mentioned first (top of mind)
MUUT_COLOR = "#8FBFBC"     # light teal — other spontaneous mentions
GRID = "#DAD3C7"

# ---------------------------------------------------------------- figure
FIG_W, FIG_H = 11.6, 6.05
fig = plt.figure(figsize=(FIG_W, FIG_H), dpi=200)
fig.patch.set_facecolor(CREAM)
ax = fig.add_axes([0.175, 0.075, 0.795, 0.885])
ax.set_facecolor(CREAM)

n = len(recs)
ypos = list(range(n))[::-1]      # first record at top
bar_h = 0.62
XMAX = 42

for r, y in zip(recs, ypos):
    tom, muut, tot = r["tom"], r["muut"], r["tot"]
    # dark top-of-mind segment
    ax.barh(y, tom, height=bar_h, color=TOM_COLOR, edgecolor="none", zorder=3)
    # light "muut maininnat" segment, stacked
    ax.barh(y, muut, left=tom, height=bar_h, color=MUUT_COLOR,
            edgecolor="none", zorder=3)
    # total label at bar end (bold)
    ax.text(tot + 0.7, y, f"{tot}", va="center", ha="left",
            fontsize=10.5, fontweight="bold", color=INK, zorder=5)
    # top-of-mind value inside the dark segment when it is wide enough,
    # otherwise just above the bar so small bars stay legible.
    if tom >= 5:
        ax.text(tom / 2, y, f"{tom}", va="center", ha="center",
                fontsize=9.5, fontweight="bold", color="#FFFFFF", zorder=6)
    elif tom > 0:
        ax.text(tom / 2, y + bar_h * 0.62, f"{tom}", va="bottom", ha="center",
                fontsize=7.6, fontweight="bold", color=TOM_COLOR, zorder=6)

# category labels
ax.set_yticks(ypos)
ax.set_yticklabels([r["name"] for r in recs], fontsize=11.5, color=INK)
for lbl in ax.get_yticklabels():
    if lbl.get_text() == "Attendo":
        lbl.set_fontweight("bold")

# x axis
ax.set_xlim(0, XMAX)
ticks = [0, 10, 20, 30, 40]
ax.set_xticks(ticks)
ax.set_xticklabels(["0", "10", "20", "30", "40 %"], fontsize=9.5, color=MUTED)
ax.tick_params(axis="x", length=0)
ax.tick_params(axis="y", length=0)
ax.set_ylim(-0.7, n - 0.3)

for x in ticks[1:]:
    ax.axvline(x, color=GRID, lw=0.8, zorder=1)
for sp in ax.spines.values():
    sp.set_visible(False)
ax.spines["left"].set_visible(True)
ax.spines["left"].set_color("#C9C1B4")
ax.spines["left"].set_linewidth(1.0)

# legend
handles = [
    Patch(facecolor=TOM_COLOR, label="Ensimmäisenä mainittu (top of mind)"),
    Patch(facecolor=MUUT_COLOR, label="Muut spontaanit maininnat"),
]
leg = ax.legend(handles=handles, loc="lower right",
                bbox_to_anchor=(1.0, 0.015),
                frameon=True, fontsize=9.5, ncol=1,
                handlelength=1.1, handleheight=1.0,
                labelspacing=0.5, borderpad=0.85)
leg.get_frame().set_facecolor("#FFFFFF")
leg.get_frame().set_edgecolor(GRID)
leg.get_frame().set_linewidth(0.8)

fig.savefig(PNG, dpi=200, facecolor=CREAM, bbox_inches="tight", pad_inches=0.04)
plt.close(fig)
print("chart ->", PNG)

# ---------------------------------------------------------------- pptx
PX_CREAM = RGBColor(0xF7, 0xF3, 0xEC)
PX_INK = RGBColor(0x2B, 0x2B, 0x2B)
PX_TEAL = RGBColor(0x13, 0x61, 0x5E)
PX_MUTED = RGBColor(0x6E, 0x6A, 0x63)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = PX_CREAM
bg.line.fill.background()
bg.shadow.inherit = False
slide.shapes._spTree.remove(bg._element)
slide.shapes._spTree.insert(2, bg._element)


def textbox(l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for txt, sz, col, bold in line:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col
            r.font.name = "Liberation Sans"
    return tb


# Accent bar
acc = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42),
                             Inches(0.10), Inches(0.92))
acc.fill.solid(); acc.fill.fore_color.rgb = PX_TEAL
acc.line.fill.background(); acc.shadow.inherit = False

# Title (Finnish key message) + sub-line
attendo = next(r for r in recs if r["name"] == "Attendo")
runner = recs[1]
textbox(Inches(0.80), Inches(0.40), Inches(11.9), Inches(1.0),
        [[("Attendo on tunnetuin yksityinen hoivapalveluiden tarjoaja – "
           "ja usein ensimmäisenä mieleen tuleva", 22, PX_INK, True)],
         [(f"Spontaani tunnettuus {attendo['tot']} %, joista {attendo['tom']} % "
           f"mainitsi Attendon ensimmäisenä – selvästi {runner['name']}n "
           f"({runner['tot']} %) edellä", 13.5, PX_MUTED, False)]],
        space_after=3)

# Section label
textbox(Inches(0.80), Inches(1.55), Inches(11.9), Inches(0.32),
        [[("SPONTAANI TUNNETTUUS  ·  osuus vastaajista (%)", 11, PX_TEAL, True)]])

# Chart image
iw, ih = Image.open(PNG).size
ar = iw / ih
pic_w = Inches(12.15)
pic_h = Emu(int(pic_w / ar))
max_h = Inches(4.95)
if pic_h > max_h:
    pic_h = max_h
    pic_w = Emu(int(pic_h * ar))
pic_l = Inches(0.62)
pic_t = Inches(1.92)
slide.shapes.add_picture(str(PNG), pic_l, pic_t, width=pic_w, height=pic_h)

# Footer: question caption (+ EOS reference) and base n
textbox(Inches(0.80), Inches(6.98), Inches(10.0), Inches(0.5),
        [[("Kysymys: ", 9.5, PX_MUTED, True),
          ("”Mitä hoivapalveluita tarjoavia yrityksiä tunnet vähintään nimeltä. "
           "Listaa kaikki, jotka muistat.”", 9.5, PX_MUTED, False)],
         [(f"Vastaajista {EOS} % ei osannut nimetä yhtäkään tarjoajaa.",
           9.5, PX_MUTED, False)]],
        space_after=1)
textbox(Inches(10.9), Inches(6.98), Inches(1.9), Inches(0.45),
        [[("Marraskuu 2025", 9.5, PX_MUTED, False)],
         [("Kaikki vastaajat, n = 1001", 9.5, PX_MUTED, True)]],
        align=PP_ALIGN.RIGHT)

prs.save(OUT)
print("slide ->", OUT)
