#!/usr/bin/env python3
"""
Attendo brand-image profile, MULTI-WAVE dense slide.

14 image attributes x 4 survey waves (Toukokuu 24 -> Marraskuu 25), top-2-box
% agreement. Grouped horizontal bars: 4 waves per attribute in a light->dark
ramp with the current wave (Marraskuu 25) emphasised. Sorted by current-wave
agreement. Negatives ("Ahne","Valinpitamaton", low=good) distinguished in a
warm ramp and tagged. Per-row change vs previous wave shown on the right.

Numbers come verbatim from chart_lab/brandimage_attendo_waves.json.
Output: work/agent_dense_brandimage.pptx (1 slide, 16:9).
"""
import json
import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

# ---- house style / fonts (no Arial) -------------------------------------
plt.rcParams["font.family"] = "Liberation Sans"
plt.rcParams["font.sans-serif"] = ["Liberation Sans", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False

HERE = os.path.dirname(os.path.abspath(__file__))
PROTO = os.path.dirname(HERE)
DATA = os.path.join(HERE, "brandimage_attendo_waves.json")
WORKDIR = os.path.join(PROTO, "work")
PNG = os.path.join(WORKDIR, "agent_dense_brandimage.png")
PPTX = os.path.join(WORKDIR, "agent_dense_brandimage.pptx")

os.makedirs(WORKDIR, exist_ok=True)

# ---- load data ----------------------------------------------------------
with open(DATA) as fh:
    d = json.load(fh)

attrs = d["attribute_categories"]
waves = list(d["series_by_wave"].keys())          # 4 waves, chronological
series = d["series_by_wave"]
assert len(waves) == 4, f"expected 4 waves, got {waves}"
CURRENT = waves[-1]                                # "Marraskuu 25"
PREV = waves[-2]                                   # "Toukokuu 25"
N_BASE = 863

NEGATIVES = {"Ahne", "Valinpitamaton", "Välinpitämätön"}

rows = []
for i, a in enumerate(attrs):
    vals = {w: series[w][i] for w in waves}
    rows.append({
        "attr": a,
        "vals": vals,
        "cur": vals[CURRENT],
        "prev": vals[PREV],
        "neg": a in NEGATIVES,
    })

# sort by current-wave agreement, descending
rows.sort(key=lambda r: r["cur"], reverse=True)

# ---- colours ------------------------------------------------------------
POS_RAMP = ["#cfe3ec", "#9cc4d6", "#5e95b0", "#1f6b8c"]   # current = darkest
NEG_RAMP = ["#f3d9c9", "#e7b394", "#d68a5e", "#b85c2a"]
INK = "#243845"
GRID = "#d9e0e4"
RISE = "#1f8a52"
FALL = "#c0392b"
FLAT = "#8a97a0"

# ---- figure geometry (16:9) --------------------------------------------
FIG_W, FIG_H = 13.333, 7.5
fig = plt.figure(figsize=(FIG_W, FIG_H), dpi=200)
fig.patch.set_facecolor("white")

ax = fig.add_axes([0.255, 0.085, 0.575, 0.70])
ax.set_facecolor("white")

n = len(rows)
group_h = 1.0
bar_h = group_h / 5.6
y_centers = [n - 1 - i for i in range(n)]   # top row highest y

xmax = max(max(r["vals"].values()) for r in rows)
ax.set_xlim(0, xmax * 1.18)
ax.set_ylim(-0.6, n - 1 + 0.6)

for gx in range(0, int(xmax) + 10, 10):
    ax.axvline(gx, color=GRID, lw=0.8, zorder=0)

offsets = [1.5, 0.5, -0.5, -1.5]  # wave0 top ... wave3 bottom within group
for r, yc in zip(rows, y_centers):
    ramp = NEG_RAMP if r["neg"] else POS_RAMP
    for wi, w in enumerate(waves):
        v = r["vals"][w]
        is_cur = (w == CURRENT)
        y = yc + offsets[wi] * bar_h
        ax.barh(
            y, v, height=bar_h * 0.92,
            color=ramp[wi],
            edgecolor=(INK if is_cur else "none"),
            linewidth=1.1 if is_cur else 0,
            zorder=3,
        )
        if is_cur:
            ax.text(v + xmax * 0.012, y, f"{v}", va="center", ha="left",
                    fontsize=9.5, fontweight="bold", color=INK, zorder=5)

for r, yc in zip(rows, y_centers):
    ax.text(-xmax * 0.022, yc + (0.16 if r["neg"] else 0.0), r["attr"],
            va="center", ha="right",
            fontsize=10.2, color=INK,
            fontweight=("bold" if not r["neg"] else "normal"),
            transform=ax.transData)
    if r["neg"]:
        ax.text(-xmax * 0.022, yc - 0.30, "(negatiivinen, matala = hyvä)",
                va="center", ha="right", fontsize=7.0,
                style="italic", color="#b85c2a", transform=ax.transData)

for spine in ["top", "right", "left"]:
    ax.spines[spine].set_visible(False)
ax.spines["bottom"].set_color(GRID)
ax.set_yticks([])
ax.tick_params(axis="x", colors="#5c6b73", labelsize=8.5, length=0)
ax.set_xlabel("Samaa mieltä (%), top-2-box (arvot 4–5)",
              fontsize=8.8, color="#5c6b73", labelpad=4)

# ---- change-vs-previous column (right of bars) --------------------------
chg_x = 0.860
fig.text(chg_x, 0.792, "Muutos vs.\n" + PREV, ha="center", va="bottom",
         fontsize=8.4, color="#5c6b73", fontweight="bold")

ax_pos = ax.get_position()
y0, y1 = ax.get_ylim()

def data_y_to_fig(yc):
    frac = (yc - y0) / (y1 - y0)
    return ax_pos.y0 + frac * ax_pos.height

for r, yc in zip(rows, y_centers):
    delta = r["cur"] - r["prev"]
    fy = data_y_to_fig(yc)
    if delta > 0:
        col, arrow, txt = RISE, "▲", f"+{delta}"
    elif delta < 0:
        col, arrow, txt = FALL, "▼", f"{delta}"
    else:
        col, arrow, txt = FLAT, "▬", "0"
    fig.text(chg_x - 0.020, fy, arrow, ha="center", va="center",
             fontsize=8.2, color=col)
    fig.text(chg_x + 0.000, fy, txt, ha="left", va="center",
             fontsize=9.4, color=col, fontweight="bold")

fig.add_artist(plt.Line2D([0.835, 0.835],
                          [ax_pos.y0, ax_pos.y0 + ax_pos.height],
                          color=GRID, lw=1.0))

# ---- title block --------------------------------------------------------
fig.text(0.038, 0.952,
         "Attendon positiiviset mielikuvat ovat vahvistuneet hieman",
         fontsize=20, fontweight="bold", color=INK, ha="left", va="top")
fig.text(0.038, 0.905,
         "Ammattitaitoisuus, hyvän arjen mahdollistaminen ja luotettavuus "
         "nousivat eniten viimeisellä aallolla",
         fontsize=11.5, color="#5c6b73", ha="left", va="top")

# ---- legend: wave ramp --------------------------------------------------
leg_y = 0.852
lx = 0.255
fig.text(lx, leg_y + 0.011, "Aalto:", fontsize=8.6, color="#5c6b73",
         ha="left", va="center", fontweight="bold")
swx = lx + 0.035
for wi, w in enumerate(waves):
    is_cur = (w == CURRENT)
    fig.patches.append(
        plt.Rectangle((swx, leg_y), 0.018, 0.020, transform=fig.transFigure,
                      facecolor=POS_RAMP[wi],
                      edgecolor=(INK if is_cur else "none"),
                      linewidth=1.0 if is_cur else 0, zorder=5)
    )
    fig.text(swx + 0.021, leg_y + 0.010,
             w + ("  (nyk.)" if is_cur else ""),
             fontsize=8.4, color=(INK if is_cur else "#5c6b73"),
             fontweight=("bold" if is_cur else "normal"),
             ha="left", va="center")
    swx += 0.021 + (0.092 if is_cur else 0.072)

fig.patches.append(
    plt.Rectangle((swx + 0.01, leg_y), 0.018, 0.020, transform=fig.transFigure,
                  facecolor=NEG_RAMP[3], edgecolor="none", zorder=5))
fig.text(swx + 0.033, leg_y + 0.010, "negatiiviset attribuutit",
         fontsize=8.4, color="#b85c2a", ha="left", va="center")

# ---- footer caption -----------------------------------------------------
fig.text(0.038, 0.028,
         "Mittari: osuus vastaajista, jotka ovat samaa mieltä (top-2-box, "
         "arvot 4–5 asteikolla 1–5).  Attendo, 4 mittausaaltoa.",
         fontsize=8.2, color="#8a97a0", ha="left", va="bottom")
fig.text(0.962, 0.028, f"n = {N_BASE}", fontsize=8.6, color="#8a97a0",
         ha="right", va="bottom", fontweight="bold")

fig.savefig(PNG, dpi=200, facecolor="white")
plt.close(fig)
print("PNG written:", PNG)

# ---- assemble PPTX ------------------------------------------------------
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
slide.shapes.add_picture(PNG, 0, 0, width=prs.slide_width,
                         height=prs.slide_height)
prs.save(PPTX)
print("PPTX written:", PPTX)
