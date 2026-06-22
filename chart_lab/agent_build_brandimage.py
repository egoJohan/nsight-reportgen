#!/usr/bin/env python3
"""Build per-provider brand-image profile slides (nSight house style).

ONE slide PER private care provider (8 providers -> 8 slides) in a single
PowerPoint, produced by a single reusable generator that LOOPS the providers
found in the data file.

Each slide shows how strongly respondents associate that provider with each of
14 image attributes (% who agree the trait fits), current wave Marraskuu 2025,
as a SORTED horizontal bar chart (strongest -> weakest association), every bar
%-labelled. Twelve attributes are positive traits; two are NEGATIVE traits
("Ahne", "Välinpitämätön") for which a LOW value is good. The negatives are
visually distinguished (warm-red bars + a "matala = hyvä" marker) so a low
score is read as a positive, never confused with a weak positive trait.

All numbers are read VERBATIM from the JSON data file; nothing is invented.
Each provider's base n is parsed from its data key ("Attendo, n=863" -> 863).

Data shape (see radar_0.json):
    {
      "categories": [<14 attribute names>],   # includes the 2 negatives
      "series": {
        "Attendo, n=863": [14 ints],          # % agreement, in categories order
        ...
      }
    }
"""
import json
import re
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
from matplotlib.patches import Patch

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ----------------------------------------------------------------- config
ROOT = Path("/home/johan/Projects/nsight/proto")

# Attributes for which a LOW value is good.
NEGATIVE_ATTRS = ("Ahne", "Välinpitämätön")

# nSight house palette
CREAM = "#F7F3EC"
INK = "#2B2B2B"
MUTED = "#6E6A63"
TEAL = "#13615E"          # strong positive accent
TEAL_LT = "#7DB8A6"       # lighter positive (weaker associations)
RED = "#B23A2E"           # negative-trait bars
RED_LT = "#E08C82"
GRIDC = "#DAD3C7"

PX_CREAM = RGBColor(0xF7, 0xF3, 0xEC)
PX_INK = RGBColor(0x2B, 0x2B, 0x2B)
PX_TEAL = RGBColor(0x13, 0x61, 0x5E)
PX_RED = RGBColor(0xB2, 0x3A, 0x2E)
PX_MUTED = RGBColor(0x6E, 0x6A, 0x63)


def _register_fonts():
    for f in [
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
    ]:
        try:
            fm.fontManager.addfont(f)
        except Exception:
            pass
    plt.rcParams["font.family"] = "Liberation Sans"


def parse_provider(key):
    """'Attendo, n=863' -> ('Attendo', 863)."""
    m = re.match(r"\s*(.+?)\s*,\s*n\s*=\s*(\d+)\s*$", key)
    if not m:
        return key.strip(), None
    return m.group(1).strip(), int(m.group(2))


# ----------------------------------------------------------------- title text
def make_title(name, ranked):
    """Finnish key-message title naming the provider + strongest associations.

    `ranked` is the list of (attr, value, is_neg) sorted strongest->weakest by
    *favourability* (positives high first; negatives are placed by their raw %
    among the others but flagged). For the headline we take the top positive
    associations and call out the negatives separately. All strings precomputed.
    """
    positives = [(a, v) for (a, v, neg) in ranked if not neg]
    negs = {a: v for (a, v, neg) in ranked if neg}
    top = positives[:2]
    top_str = " ja ".join(f"{a.lower()} ({v} %)" for a, v in top)
    title = f"{name} liitetään vahvimmin attribuutteihin {top_str}"

    neg_bits = ", ".join(f"{a.lower()} {negs[a]} %" for a in NEGATIVE_ATTRS if a in negs)
    if neg_bits:
        subtitle = (
            f"14 attribuutin mielikuvaprofiili. Kielteiset mielikuvat matalalla "
            f"({neg_bits}) – matala arvo on myönteinen."
        )
    else:
        subtitle = "14 attribuutin mielikuvaprofiili."
    return title, subtitle


# ----------------------------------------------------------------- chart
def render_chart(name, attrs, values, neg_flags, png_path):
    """Sorted horizontal bar chart for one provider.

    attrs/values/neg_flags are already SORTED for display (top row first).
    """
    n = len(attrs)
    FIG_W, FIG_H = 11.6, 6.15
    fig = plt.figure(figsize=(FIG_W, FIG_H), dpi=200)
    fig.patch.set_facecolor(CREAM)
    ax = fig.add_axes([0.30, 0.075, 0.66, 0.88])
    ax.set_facecolor(CREAM)

    ypos = list(range(n))[::-1]   # first item at top

    for y, a, v, neg in zip(ypos, attrs, values, neg_flags):
        if neg:
            color = RED
        else:
            color = TEAL
        ax.barh(y, v, height=0.62, color=color, edgecolor="none", zorder=3)
        # value label at bar end
        ax.text(v + 1.0, y, f"{v} %", va="center", ha="left",
                fontsize=10.5, fontweight="bold",
                color=RED if neg else INK, zorder=5)

    # y labels (attribute names); negatives marked with a low-is-good tag
    ylabels = []
    for a, neg in zip(attrs, neg_flags):
        ylabels.append((a + "  (kielteinen)") if neg else a)
    ax.set_yticks(ypos)
    ax.set_yticklabels(ylabels, fontsize=11.5, color=INK)
    for lbl, neg in zip(ax.get_yticklabels(), neg_flags):
        if neg:
            lbl.set_color(RED)
            lbl.set_fontweight("bold")

    # x axis
    ax.set_xlim(0, 100)
    ax.set_xticks([0, 20, 40, 60, 80, 100])
    ax.set_xticklabels(["0", "20", "40", "60", "80", "100 %"],
                       fontsize=9.5, color=MUTED)
    ax.tick_params(axis="x", length=0)
    ax.tick_params(axis="y", length=0)
    ax.set_ylim(-0.7, n - 0.3)

    for x in [20, 40, 60, 80, 100]:
        ax.axvline(x, color=GRIDC, lw=0.8, zorder=1)
    for s in ax.spines.values():
        s.set_visible(False)
    ax.spines["left"].set_visible(True)
    ax.spines["left"].set_color("#C9C1B4")
    ax.spines["left"].set_linewidth(1.0)

    # legend explaining the two bar meanings
    handles = [
        Patch(facecolor=TEAL, label="Myönteinen mielikuva (korkea = hyvä)"),
        Patch(facecolor=RED, label="Kielteinen mielikuva (matala = hyvä)"),
    ]
    leg = ax.legend(handles=handles, loc="lower right",
                    bbox_to_anchor=(1.0, 0.012), frameon=True,
                    fontsize=9.5, handlelength=1.1, handleheight=1.0,
                    labelspacing=0.45, borderpad=0.8)
    leg.get_frame().set_facecolor("#FFFFFF")
    leg.get_frame().set_edgecolor(GRIDC)
    leg.get_frame().set_linewidth(0.8)
    for t in leg.get_texts():
        t.set_color(INK)

    fig.savefig(png_path, dpi=200, facecolor=CREAM,
                bbox_inches="tight", pad_inches=0.04)
    plt.close(fig)


# ----------------------------------------------------------------- slide
def _textbox(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, space_after=0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for txt, sz, col, bold in line:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col
            r.font.name = "Liberation Sans"
    return tb


def add_slide(prs, name, base_n, title, subtitle, png_path, wave, question):
    SW, SH = prs.slide_width, prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = PX_CREAM
    bg.line.fill.background(); bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # accent bar
    acc = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42),
                                 Inches(0.10), Inches(0.92))
    acc.fill.solid(); acc.fill.fore_color.rgb = PX_TEAL
    acc.line.fill.background(); acc.shadow.inherit = False

    # title + subtitle
    _textbox(slide, Inches(0.80), Inches(0.38), Inches(12.0), Inches(1.1),
             [[(title, 21, PX_INK, True)],
              [(subtitle, 12.5, PX_MUTED, False)]],
             space_after=3)

    # section label
    _textbox(slide, Inches(0.80), Inches(1.62), Inches(12.0), Inches(0.32),
             [[(f"MIELIKUVAPROFIILI · {name.upper()} · osuus vastaajista (%)",
                11, PX_TEAL, True)]])

    # chart
    iw, ih = Image.open(png_path).size
    ar = iw / ih
    pic_w = Inches(12.15)
    pic_h = Emu(int(pic_w / ar))
    max_h = Inches(5.0)
    if pic_h > max_h:
        pic_h = max_h
        pic_w = Emu(int(pic_h * ar))
    pic_l = Inches(0.62)
    pic_t = Inches(1.92)
    slide.shapes.add_picture(str(png_path), pic_l, pic_t,
                             width=pic_w, height=pic_h)

    # footer: question + base n
    _textbox(slide, Inches(0.80), Inches(7.00), Inches(9.6), Inches(0.45),
             [[("Kysymys: ", 9.5, PX_MUTED, True),
               ("”" + question + "”", 9.5, PX_MUTED, False)]])
    _textbox(slide, Inches(10.2), Inches(7.00), Inches(2.6), Inches(0.45),
             [[(f"{wave} · {name}, n = {base_n}", 9.5, PX_MUTED, True)]],
             align=PP_ALIGN.RIGHT)


# ----------------------------------------------------------------- build
def build(
    data_path=ROOT / "chart_lab" / "radar_0.json",
    png_dir=ROOT / "chart_lab",
    out_path=ROOT / "work" / "agent_brand_images.pptx",
    *,
    wave="Marraskuu 2025",
    question=("Missä määrin seuraavat ominaisuudet mielestäsi sopivat tähän "
              "hoivapalveluiden tarjoajaan?"),
):
    _register_fonts()
    d = json.loads(Path(data_path).read_text())
    cats = d["categories"]
    series = d["series"]
    neg_idx = {i for i, c in enumerate(cats) if c in NEGATIVE_ATTRS}

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    first_png = None
    for key in series:                       # LOOP providers, one slide each
        name, base_n = parse_provider(key)
        raw = series[key]
        assert len(raw) == len(cats), f"{key}: expected {len(cats)} values"

        # Display order: the 12 POSITIVE attributes sorted strongest -> weakest
        # by raw % (strongest association at the top), then the 2 NEGATIVE
        # attributes grouped together at the bottom (also sorted by raw %),
        # colour-distinguished. This keeps the positive ranking immediately
        # legible while isolating the negatives as a separate, clearly flagged
        # block where a LOW value is the good outcome.
        pos = [(cats[i], raw[i], False) for i in range(len(cats)) if i not in neg_idx]
        neg = [(cats[i], raw[i], True) for i in range(len(cats)) if i in neg_idx]
        pos.sort(key=lambda r: r[1], reverse=True)
        neg.sort(key=lambda r: r[1], reverse=True)
        rows = pos + neg          # positives (top) then negatives (bottom)

        attrs = [r[0] for r in rows]
        values = [r[1] for r in rows]
        negflags = [r[2] for r in rows]

        # headline uses the favourability-ranked view (top positives + negs)
        ranked = [(r[0], r[1], r[2]) for r in rows]
        title, subtitle = make_title(name, ranked)

        png_path = png_dir / f"agent_brandimage_{name.replace(' ', '_')}.png"
        render_chart(name, attrs, values, negflags, png_path)
        if first_png is None:
            first_png = png_path

        add_slide(prs, name, base_n, title, subtitle, png_path, wave, question)
        print(f"slide -> {name} (n={base_n})")

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"\n{len(prs.slides._sldIdLst)} slides -> {out_path}")
    return out_path


if __name__ == "__main__":
    build()
