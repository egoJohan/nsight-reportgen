#!/usr/bin/env python3
"""Build the brand-image profile slide (nSight house style).

Eight private-care brands are profiled across 14 image attributes (% prompted
agreement, wave Marraskuu 2025). Eight lines on one radar overlap into an
unreadable tangle, so this uses SMALL MULTIPLES: one mini-radar per brand,
each overlaid on the dashed all-brand average as a fixed reference, so every
profile's strengths/weaknesses read at a glance and brands are comparable.

The 14 long attribute names are handled with numbered spokes + a numbered key.
The two negative traits (Ahne, Välinpitämätön) are placed together and marked,
because for them a *small* profile is good.

All numbers are taken verbatim from radar_0.json (averages computed
deterministically from those exact values).
Run:  uv run python chart_lab/agent_build_radar.py
"""
import json
from pathlib import Path

import numpy as np

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ---------------------------------------------------------------- paths
ROOT = Path("/home/johan/Projects/nsight/proto")
DATA = ROOT / "chart_lab" / "radar_0.json"
PNG = ROOT / "chart_lab" / "agent_radar_chart.png"
OUT = ROOT / "work" / "agent_slide_radar.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------- fonts
for f in ["/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
          "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"]:
    fm.fontManager.addfont(f)
plt.rcParams["font.family"] = "Liberation Sans"

# ---------------------------------------------------------------- data
d = json.loads(DATA.read_text())
attrs_raw = d["categories"]                       # survey order, 14 names
# brand -> values (in survey order); keep the n from the key separately
series = {k.split(",")[0].strip(): v for k, v in d["series"].items()}
bases = {k.split(",")[0].strip(): int(k.split("n=")[1]) for k in d["series"]}
brands = list(series.keys())

# Indices of the two negative attributes (low = good for these).
NEG_NAMES = {"Ahne", "Välinpitämätön"}
neg_idx = [i for i, a in enumerate(attrs_raw) if a in NEG_NAMES]
pos_idx = [i for i, a in enumerate(attrs_raw) if a not in NEG_NAMES]
# Display order: negatives first (grouped, marked), then positives in survey order.
order = neg_idx + pos_idx
attrs = [attrs_raw[i] for i in order]
ser = {b: [series[b][i] for i in order] for b in brands}
NEG_SET = set(range(len(neg_idx)))                # positions of negatives in new order

# Average across all 8 brands, per attribute (deterministic from exact data).
N = len(attrs)
avg = [float(np.mean([ser[b][p] for b in brands])) for p in range(N)]

# ---------------------------------------------------------------- colours
CREAM = "#F4EFE6"; INK = "#2B2B2B"; MUTED = "#6E685C"; GREY = "#B7AFA0"
GRID = "#DED6C7"; SPINE = "#CFC8BA"
# Onnikodit is the headline brand -> strong teal; others muted-but-distinct.
TEAL = "#13615E"
PALETTE = {
    "Attendo": "#A4554B", "Esperi": "#B8823A", "Mainio-kodit": "#5E8C57",
    "Onnikodit": TEAL, "Ykköskodit": "#6E7BA6", "Humana": "#7E5BA6",
    "Rinnekodit": "#3F8C84", "Validia": "#8A6D3B",
}

# ---------------------------------------------------------------- figure
ang = np.linspace(0, 2 * np.pi, N, endpoint=False).tolist()
ang_c = ang + ang[:1]
avg_c = avg + avg[:1]
RMAX = 60

# Order panels: headline brand (Onnikodit) first, then by overall profile size.
panel_order = sorted(brands, key=lambda b: np.mean(ser[b]), reverse=True)
panel_order = ["Onnikodit"] + [b for b in panel_order if b != "Onnikodit"]

fig = plt.figure(figsize=(12.8, 6.05), dpi=220)
fig.patch.set_facecolor(CREAM)

for k, brand in enumerate(panel_order):
    ax = fig.add_subplot(2, 4, k + 1, polar=True)
    ax.set_facecolor(CREAM)
    ax.set_theta_offset(np.pi / 2)
    ax.set_theta_direction(-1)
    col = PALETTE[brand]
    is_lead = (brand == "Onnikodit")

    v = ser[brand] + ser[brand][:1]
    # average reference (dashed grey) under the brand profile
    ax.plot(ang_c, avg_c, color=GREY, lw=1.0, ls=(0, (3, 2)), zorder=2)
    ax.fill(ang_c, v, color=col, alpha=0.32 if is_lead else 0.24, zorder=3)
    ax.plot(ang_c, v, color=col, lw=2.4 if is_lead else 1.7, zorder=4)

    ax.set_xticks(ang)
    labels = [str(p + 1) for p in range(N)]
    ax.set_xticklabels(labels, fontsize=7.5, color=MUTED)
    # mark the negative-trait spoke numbers in red
    for tl, p in zip(ax.get_xticklabels(), range(N)):
        if p in NEG_SET:
            tl.set_color("#B23A3A"); tl.set_fontweight("bold")
    ax.set_ylim(0, RMAX)
    ax.set_yticks([20, 40])
    ax.set_yticklabels([])
    ax.grid(color=GRID, lw=0.7)
    ax.spines["polar"].set_color(SPINE)
    ax.set_title(brand, fontsize=11.5,
                 fontweight="bold" if is_lead else "bold",
                 color=col if is_lead else INK, pad=9)

fig.subplots_adjust(left=0.015, right=0.985, top=0.96, bottom=0.04,
                    wspace=0.38, hspace=0.40)
fig.savefig(PNG, dpi=220, facecolor=CREAM, bbox_inches="tight", pad_inches=0.06)
plt.close(fig)
print("chart ->", PNG)

# ---------------------------------------------------------------- pptx
PX_CREAM = RGBColor(0xF4, 0xEF, 0xE6)
PX_INK = RGBColor(0x2B, 0x2B, 0x2B)
PX_TEAL = RGBColor(0x13, 0x61, 0x5E)
PX_MUTED = RGBColor(0x6E, 0x68, 0x5C)
PX_RED = RGBColor(0xB2, 0x3A, 0x3A)
PX_GREY = RGBColor(0xB7, 0xAF, 0xA0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
slide = prs.slides.add_slide(prs.slide_layouts[6])

# background
bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = PX_CREAM
bg.line.fill.background(); bg.shadow.inherit = False
slide.shapes._spTree.remove(bg._element)
slide.shapes._spTree.insert(2, bg._element)


def textbox(l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=0, line_spacing=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if line_spacing:
            p.line_spacing = line_spacing
        for txt, sz, col, bold in line:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col
            r.font.name = "Liberation Sans"
    return tb


# Accent bar top-left
acc = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42),
                             Inches(0.10), Inches(0.92))
acc.fill.solid(); acc.fill.fore_color.rgb = PX_TEAL
acc.line.fill.background(); acc.shadow.inherit = False

# Title (Finnish key message) + sub
textbox(Inches(0.80), Inches(0.40), Inches(11.9), Inches(1.0),
        [[("Brändiprofiilit ovat muodoltaan samankaltaisia – Onnikodit "
           "erottuu vahvimpana lähes kaikilla mielikuvilla",
           21, PX_INK, True)],
         [("Onnikodit yltää keskiarvon yläpuolelle kaikilla "
           "myönteisillä attribuuteilla; Attendo ja Esperi profiloituvat "
           "vahvimmin kielteisiin mielikuviin",
           13, PX_MUTED, False)]],
        space_after=3)

# Section label above chart
textbox(Inches(0.80), Inches(1.58), Inches(11.9), Inches(0.32),
        [[("BRÄNDIMIELIKUVAT  ·  osuus, joka liittää attribuutin brändiin "
           "(autettu, % vastaajista)  ·  Marraskuu 2025",
           11, PX_TEAL, True)]])

# Chart image
iw, ih = Image.open(PNG).size
ar = iw / ih
pic_w = Inches(9.55)
pic_h = Emu(int(pic_w / ar))
max_h = Inches(4.45)
if pic_h > max_h:
    pic_h = max_h
    pic_w = Emu(int(pic_h * ar))
pic_l = Inches(0.55)
pic_t = Inches(1.95)
slide.shapes.add_picture(str(PNG), pic_l, pic_t, width=pic_w, height=pic_h)

# ---- Right-hand legend panel: numbered attribute key -------------------
KEY_L = Inches(10.30)
KEY_T = Inches(2.05)
KEY_W = Inches(2.78)
panel = slide.shapes.add_shape(1, KEY_L, KEY_T, KEY_W, Inches(4.55))
panel.fill.solid(); panel.fill.fore_color.rgb = RGBColor(0xFA, 0xF7, 0xF1)
panel.line.color.rgb = RGBColor(0xDE, 0xD6, 0xC7); panel.line.width = Pt(0.75)
panel.shadow.inherit = False

# key header
textbox(KEY_L + Inches(0.16), KEY_T + Inches(0.12), KEY_W - Inches(0.30),
        Inches(0.30),
        [[("ATTRIBUUTIT", 10.5, PX_TEAL, True)]])

# numbered list; negatives (1-2) in red and tagged "(kielteinen)"
key_lines = []
for p, a in enumerate(attrs):
    is_neg = p in NEG_SET
    num_col = PX_RED if is_neg else PX_INK
    txt_col = PX_INK if not is_neg else PX_RED
    tag = "  (kielteinen)" if is_neg else ""
    key_lines.append([(f"{p+1}.  ", 9.5, num_col, True),
                      (a + tag, 9.5, txt_col, is_neg)])
textbox(KEY_L + Inches(0.16), KEY_T + Inches(0.46), KEY_W - Inches(0.28),
        Inches(3.55), key_lines, space_after=2.0, line_spacing=1.0)

# average reference note inside the panel (bottom)
textbox(KEY_L + Inches(0.16), KEY_T + Inches(4.06), KEY_W - Inches(0.28),
        Inches(0.42),
        [[("– – –  ", 10, PX_GREY, True),
          ("keskiarvo, kaikki brändit", 9.5, PX_MUTED, False)],
         [("Punainen = kielteinen mielikuva (pieni = hyvä)",
           8.3, PX_RED, False)]],
        space_after=1.0, line_spacing=1.05)

# Footer: question caption + base note (n varies by brand)
textbox(Inches(0.80), Inches(6.98), Inches(8.2), Inches(0.50),
        [[("Kysymys: ", 9.5, PX_MUTED, True),
          ("”Missä määrin seuraavat ominaisuudet kuvaavat kutakin "
           "hoivapalveluiden tarjoajaa?” · samaa mieltä -osuus",
           9.5, PX_MUTED, False)]])

base_txt = "Pohja vaihtelee brändeittäin: " + ", ".join(
    f"{b} n={bases[b]}" for b in brands)
textbox(Inches(0.80), Inches(7.20), Inches(11.9), Inches(0.28),
        [[(base_txt, 8.3, PX_MUTED, False)]])

prs.save(OUT)
print("slide ->", OUT)
