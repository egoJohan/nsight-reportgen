#!/usr/bin/env python3
"""Build the MULTI-WAVE spontaneous brand-image (words) slide — nSight style.

Shows the TOP-10 spontaneously-mentioned words describing Attendo for ALL FOUR
survey waves side by side, as four ranked columns. Each cell is a word + its
EXACT mention count, coloured by sentiment (kielteinen / myönteinen /
neutraali). Thin connector lines link the same word across adjacent waves so
rank movement is visible; words new to a wave's TOP-10 carry a small "uusi"
flag. The current wave (Marraskuu 2025) is emphasised with a highlighted
header band and a heavier panel border.

All words, counts and sentiment values are read VERBATIM from words_waves.json.
Nothing is hard-coded or invented.

Run: uv run python chart_lab/agent_dense_words.py
"""
import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
from matplotlib.patches import FancyBboxPatch, Patch
from matplotlib.lines import Line2D

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ---------------------------------------------------------------- paths
ROOT = Path("/home/johan/Projects/nsight/proto")
DATA = ROOT / "chart_lab" / "words_waves.json"
PNG = ROOT / "chart_lab" / "agent_dense_words.png"
OUT = ROOT / "work" / "agent_dense_words.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------- fonts
FONT_REG = "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
FONT_BOLD = "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"
for f in (FONT_REG, FONT_BOLD):
    fm.fontManager.addfont(f)
plt.rcParams["font.family"] = "Liberation Sans"

# ---------------------------------------------------------------- data
d = json.loads(DATA.read_text())
QUESTION = d["question"]
BASE = d["base"]                       # "kaikki vastaajat, n=863"
WAVES = list(d["waves"].keys())        # 4 wave names, chronological order
waves = {w: [(word, cnt) for word, cnt in d["waves"][w]] for w in WAVES}
SENT = d["sentiment"]                   # word -> sentiment (verbatim)
CURRENT = WAVES[-1]                     # current wave = last (Marraskuu 2025)

# ---------------------------------------------------------------- brand colours
CREAM = "#F7F3EC"
PANEL = "#FFFFFF"
INK = "#2B2B2B"
TEAL = "#13615E"
MUTED = "#6E6A63"
GRID = "#DAD3C7"
HILITE = "#E9F0EE"          # soft teal tint for current-wave header band

SENT_COLOR = {
    "kielteinen": "#C0473A",   # negative — warm red
    "myönteinen": "#3F7D4E",   # positive — green
    "neutraali":  "#B3A98F",   # neutral — warm grey
}
SENT_LABEL = {
    "kielteinen": "Kielteinen",
    "myönteinen": "Myönteinen",
    "neutraali":  "Neutraali",
}
SENT_ORDER = ["kielteinen", "myönteinen", "neutraali"]
present = [s for s in SENT_ORDER
           if any(SENT[w] == s for col in waves.values() for w, _ in col)]

# ---------------------------------------------------------------- layout grid
N_WAVES = len(WAVES)
N_ROWS = 10                                   # TOP-10

# canvas in data coords (0..100 x, 0..100 y)
COL_X0, COL_X1 = 4.0, 99.0
col_span = (COL_X1 - COL_X0) / N_WAVES
# centre x of each wave column
col_cx = [COL_X0 + col_span * (i + 0.5) for i in range(N_WAVES)]
# pill (cell) geometry
PILL_W = col_span * 0.82
HDR_Y = 94.0                                   # header band centre
ROW_TOP = 86.5
ROW_BOT = 6.5
row_h = (ROW_TOP - ROW_BOT) / (N_ROWS - 1)
row_y = [ROW_TOP - row_h * r for r in range(N_ROWS)]   # rank 1 at top

# rank lookup for each wave: word -> row index (0-based)
rank_of = {w: {word: i for i, (word, _) in enumerate(waves[w])} for w in WAVES}

# ---------------------------------------------------------------- figure
FIG_W, FIG_H = 12.6, 6.25
fig = plt.figure(figsize=(FIG_W, FIG_H), dpi=200)
fig.patch.set_facecolor(CREAM)
ax = fig.add_axes([0.0, 0.0, 1.0, 1.0])
ax.set_xlim(0, 100)
ax.set_ylim(-4, 102)
ax.axis("off")

PILL_H = row_h * 0.80

def pill_xy(col_i, row_i):
    cx = col_cx[col_i]
    cy = row_y[row_i]
    return cx, cy


# ---- connector lines between adjacent waves (drawn first, behind pills) ----
for ci in range(N_WAVES - 1):
    left_words = {word for word, _ in waves[WAVES[ci]]}
    for word, _ in waves[WAVES[ci + 1]]:
        if word in left_words:
            r0 = rank_of[WAVES[ci]][word]
            r1 = rank_of[WAVES[ci + 1]][word]
            x0 = col_cx[ci] + PILL_W / 2.0
            x1 = col_cx[ci + 1] - PILL_W / 2.0
            y0 = row_y[r0]
            y1 = row_y[r1]
            col = SENT_COLOR[SENT[word]]
            # emphasise lines feeding into the current wave
            lw = 1.7 if WAVES[ci + 1] == CURRENT else 1.1
            alpha = 0.55 if WAVES[ci + 1] == CURRENT else 0.32
            ax.plot([x0, x1], [y0, y1], color=col, lw=lw, alpha=alpha,
                    solid_capstyle="round", zorder=2)

# ---- rank gutter numbers (left of first column) ----
for r in range(N_ROWS):
    ax.text(COL_X0 - 2.6, row_y[r], f"{r + 1}", va="center", ha="right",
            fontsize=11, color=MUTED, fontweight="bold", zorder=4)

# ---- column header bands + wave labels ----
for ci, wv in enumerate(WAVES):
    is_cur = (wv == CURRENT)
    cx = col_cx[ci]
    band_w = col_span * 0.94
    # header band
    ax.add_patch(FancyBboxPatch(
        (cx - band_w / 2.0, HDR_Y - 3.4), band_w, 6.8,
        boxstyle="round,pad=0,rounding_size=1.6",
        facecolor=(HILITE if is_cur else "none"),
        edgecolor=(TEAL if is_cur else GRID),
        linewidth=(1.8 if is_cur else 1.0), zorder=3))
    # split "Toukokuu 2024" -> month / year for a tidy 2-line header
    parts = wv.split(" ")
    month, year = parts[0], parts[-1]
    ax.text(cx, HDR_Y + 1.0, month, va="center", ha="center",
            fontsize=13.5, fontweight="bold",
            color=(TEAL if is_cur else INK), zorder=5)
    ax.text(cx, HDR_Y - 2.2, year, va="center", ha="center",
            fontsize=10.5, color=(TEAL if is_cur else MUTED),
            fontweight=("bold" if is_cur else "normal"), zorder=5)
    if is_cur:
        ax.text(cx, HDR_Y + 6.4, "NYKYINEN AALTO", va="center", ha="center",
                fontsize=8.5, fontweight="bold", color=TEAL, zorder=5)

# ---- emphasis border around the current-wave column ----
cur_i = WAVES.index(CURRENT)
ax.add_patch(FancyBboxPatch(
    (col_cx[cur_i] - col_span * 0.49, ROW_BOT - row_h * 0.62),
    col_span * 0.98, (ROW_TOP - ROW_BOT) + row_h * 1.24,
    boxstyle="round,pad=0,rounding_size=1.8",
    facecolor="none", edgecolor=TEAL, linewidth=1.6, alpha=0.55, zorder=1))

# ---- the cells (word + count pills) ----
for ci, wv in enumerate(WAVES):
    is_cur = (wv == CURRENT)
    prev_words = (set(w for w, _ in waves[WAVES[ci - 1]]) if ci > 0 else None)
    cx = col_cx[ci]
    for r, (word, cnt) in enumerate(waves[wv]):
        cy = row_y[r]
        col = SENT_COLOR[SENT[word]]
        # pill: light tinted background with a coloured left sentiment chip
        pw, ph = PILL_W, PILL_H
        ax.add_patch(FancyBboxPatch(
            (cx - pw / 2.0, cy - ph / 2.0), pw, ph,
            boxstyle="round,pad=0,rounding_size=1.3",
            facecolor=PANEL, edgecolor=GRID,
            linewidth=(1.3 if is_cur else 0.8), zorder=4))
        # sentiment chip (left edge)
        chip_w = pw * 0.07
        ax.add_patch(FancyBboxPatch(
            (cx - pw / 2.0, cy - ph / 2.0), chip_w, ph,
            boxstyle="round,pad=0,rounding_size=1.3",
            facecolor=col, edgecolor="none", zorder=5))
        # word (left) + count (right), inside the pill
        tx = cx - pw / 2.0 + chip_w + pw * 0.045
        ax.text(tx, cy, word, va="center", ha="left",
                fontsize=(12.0 if is_cur else 11.0),
                fontweight="bold", color=INK, zorder=6)
        ax.text(cx + pw / 2.0 - pw * 0.05, cy, f"{cnt}", va="center",
                ha="right", fontsize=(12.0 if is_cur else 11.0),
                fontweight="bold", color=col, zorder=6)
        # "uusi" flag for words not in the previous wave's TOP-10
        if prev_words is not None and word not in prev_words:
            ax.text(cx - pw / 2.0 + chip_w + pw * 0.045, cy - ph * 0.42,
                    "uusi", va="center", ha="left", fontsize=6.8,
                    style="italic", color=TEAL, zorder=6)

# ---- sentiment legend (bottom) ----
handles = [Patch(facecolor=SENT_COLOR[s], label=SENT_LABEL[s]) for s in present]
handles.append(Line2D([0], [0], color=MUTED, lw=1.4, alpha=0.5,
                      label="Sama sana edellisessä aallossa"))
leg = ax.legend(handles=handles, loc="lower center",
                bbox_to_anchor=(0.5, -0.015), ncol=len(handles),
                frameon=False, fontsize=10.0,
                handlelength=1.2, handleheight=1.0,
                columnspacing=1.8, handletextpad=0.55)
for t in leg.get_texts():
    t.set_color(INK)
leg.set_title("Sävy", prop={"size": 10.0, "weight": "bold"})
leg.get_title().set_color(INK)

fig.savefig(PNG, dpi=200, facecolor=CREAM, bbox_inches="tight", pad_inches=0.06)
plt.close(fig)
print("chart ->", PNG)

# ================================================================ PPTX
PX_CREAM = RGBColor(0xF7, 0xF3, 0xEC)
PX_INK = RGBColor(0x2B, 0x2B, 0x2B)
PX_TEAL = RGBColor(0x13, 0x61, 0x5E)
PX_RED = RGBColor(0xC0, 0x47, 0x3A)
PX_MUTED = RGBColor(0x6E, 0x6A, 0x63)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = PX_CREAM
bg.line.fill.background(); bg.shadow.inherit = False
slide.shapes._spTree.remove(bg._element)
slide.shapes._spTree.insert(2, bg._element)


def textbox(l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for txt, sz, col, bold in line:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col
            r.font.name = "Liberation Sans"
    return tb


# accent bar (red — flags the negative-dominated headline)
acc = slide.shapes.add_shape(1, Inches(0.55), Inches(0.40),
                             Inches(0.10), Inches(0.92))
acc.fill.solid(); acc.fill.fore_color.rgb = PX_RED
acc.line.fill.background(); acc.shadow.inherit = False

# Title (Finnish key message). Derive the two leading negative words of the
# current wave from the data so the headline stays exact.
cur_ranked = waves[CURRENT]
cur_neg = [w for w, _ in cur_ranked if SENT[w] == "kielteinen"]
neg_str = ", ".join(w.lower() for w in cur_neg[:2])
textbox(Inches(0.80), Inches(0.38), Inches(12.0), Inches(1.0),
        [[("Attendon mielikuvaa hallitsevat yhä kielteiset sanat — "
           "kuva pysyy vakaana", 21, PX_INK, True)],
         [(f"Kärjessä {neg_str}; sama kärki toistuu kaikissa neljässä "
           "mittausaallossa", 13, PX_MUTED, False)]],
        space_after=3)

# section label
textbox(Inches(0.80), Inches(1.52), Inches(12.0), Inches(0.32),
        [[("SPONTAANI MIELIKUVA  ·  TOP 10 mainituinta sanaa  ·  "
           "neljä mittausaaltoa", 11, PX_TEAL, True)]])

# chart image (high-res embed, aspect preserved)
iw, ih = Image.open(PNG).size
ar = iw / ih
pic_w = Inches(12.55)
pic_h = Emu(int(pic_w / ar))
max_h = Inches(5.02)
if pic_h > max_h:
    pic_h = max_h
    pic_w = Emu(int(pic_h * ar))
pic_l = Emu(int((SW - pic_w) / 2))
pic_t = Inches(1.92)
slide.shapes.add_picture(str(PNG), pic_l, pic_t, width=pic_w, height=pic_h)

# footer: question + base n
textbox(Inches(0.80), Inches(7.04), Inches(9.6), Inches(0.42),
        [[("Kysymys: ", 9.5, PX_MUTED, True),
          ("”" + QUESTION + "”", 9.5, PX_MUTED, False)]])
base_disp = BASE[0].upper() + BASE[1:]
base_disp = base_disp.replace("n=", "n = ")
textbox(Inches(9.2), Inches(7.04), Inches(3.55), Inches(0.42),
        [[(base_disp, 9.5, PX_MUTED, True)]], align=PP_ALIGN.RIGHT)

prs.save(OUT)
print("slide ->", OUT)
