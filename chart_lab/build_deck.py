"""Assemble the approved, data-driven charts into an actual generated PowerPoint deck.
This is the deliverable: a presentation produced automatically from the charts (which are
produced from the verified survey numbers). Run: uv run python chart_lab/build_deck.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

LAB = Path(__file__).parent
OUT = LAB.parent / "work" / "attendo_nsight_generated.pptx"
CREAM = RGBColor(0xF4, 0xEF, 0xE6)
INK = RGBColor(0x2B, 0x2B, 0x2B)
GREY = RGBColor(0x8A, 0x85, 0x7B)

# slide title (key message) + chart image, in deck order
SLIDES = [
    ("Attendo on edelleen selvästi tunnetuin yksityinen hoivapalvelujen tarjoaja",
     "v4.png"),
    ("Yleinen käsitys yksityisistä on kohentunut ja on nyt yhtä myönteinen kuin julkisista",
     "stacked_s3.png"),
    ("Brändien mielikuvaprofiilit ovat samankaltaisia; Onnikodit yhdistetään lähes kaikkiin",
     "radar_sm1.png"),
    ("Attendon ja Esperin tunnettuus on pysynyt vakaana aaltojen yli",
     "trend_t1.png"),
    ("Attendon spontaani mielikuva painottuu yhä kielteisiin sanoihin (kallis, huono)",
     "words_w2.png"),
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def cream_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM


def textbox(slide, text, left, top, width, height, size, color, bold=True, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = "Liberation Sans"
    return tb


# --- title slide ---
s = prs.slides.add_slide(blank); cream_bg(s)
textbox(s, "Attendo — Bränditutkimus", 0.9, 2.4, 11.5, 1.0, 40, INK)
textbox(s, "Marraskuu 2025", 0.9, 3.5, 11.5, 0.7, 24, GREY, bold=False)
textbox(s, "Automaattisesti tuotettu · data: SPSS (n=1001) · graafit: data-ajettu",
        0.9, 6.6, 11.5, 0.5, 13, GREY, bold=False)

# --- content slides: key message + chart image ---
for title, img in SLIDES:
    s = prs.slides.add_slide(blank); cream_bg(s)
    textbox(s, title, 0.6, 0.4, 12.1, 1.1, 20, INK)
    pic_path = LAB / img
    # place chart centred below the title, max area ~ 12.1 x 5.4 in
    from PIL import Image  # noqa: PLC0415
    iw, ih = Image.open(pic_path).size
    maxw, maxh = 12.1, 5.5
    scale = min(maxw / (iw / 96), maxh / (ih / 96))
    w = (iw / 96) * scale; h = (ih / 96) * scale
    left = (13.333 - w) / 2; top = 1.6 + (5.5 - h) / 2
    s.shapes.add_picture(str(pic_path), Inches(left), Inches(top), Inches(w), Inches(h))

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"wrote {OUT} ({OUT.stat().st_size} bytes, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
