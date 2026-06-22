#!/usr/bin/env python3
"""Build the spontaneous brand-image (words) slide — nSight house style.

Ranked horizontal bar chart of the TOP-10 spontaneous words used to describe
Attendo, coloured by sentiment (kielteinen / myönteinen / neutraali). All
numbers and sentiment values are read VERBATIM from words_data.json — nothing
is hard-coded or invented.

Run: uv run python chart_lab/agent_build_words.py
"""
import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
from matplotlib.patches import Patch

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ---------------------------------------------------------------- paths
ROOT = Path("/home/johan/Projects/nsight/proto")
DATA = ROOT / "chart_lab" / "words_data.json"
PNG = ROOT / "chart_lab" / "agent_words_chart.png"
OUT = ROOT / "work" / "agent_slide_words.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------- fonts
for f in ["/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
          "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"]:
    fm.fontManager.addfont(f)
plt.rcParams["font.family"] = "Liberation Sans"

# ---------------------------------------------------------------- data
d = json.loads(DATA.read_text())
meta = d["meta"]
words = d["words"]                                   # verbatim from file

QUESTION = meta["question"]
BASE = meta["base"]                                  # "kaikki vastaajat, n=863"
WAVE = meta["wave"]

# Rank by count, highest at top. Sort is stable so ties keep file order.
words_ranked = sorted(words, key=lambda w: w["count"], reverse=True)

# Sentiment values are taken verbatim from the data file.
SENT_COLOR = {
    "kielteinen": "#C0473A",   # negative — warm red
    "myönteinen": "#3F7D4E",   # positive — green
    "neutraali":  "#B3A98F",   # neutral — warm grey
}
SENT_LABEL = {
    "kielteinen": "Kielteinen",
    "myönteinen": "Myönteinen",
    "neutraali":  "Neutraali",
}
# Legend order: present only sentiments that actually occur, fixed order.
SENT_ORDER = ["kielteinen", "myönteinen", "neutraali"]
present = [s for s in SENT_ORDER if any(w["sentiment"] == s for w in words_ranked)]

# ---------------------------------------------------------------- colours
CREAM = "#F7F3EC"
INK = "#2B2B2B"
MUTED = "#6E6A63"
GRID = "#DAD3C7"

# ---------------------------------------------------------------- figure
FIG_W, FIG_H = 11.6, 5.55
fig = plt.figure(figsize=(FIG_W, FIG_H), dpi=200)
fig.patch.set_facecolor(CREAM)
ax = fig.add_axes([0.175, 0.115, 0.80, 0.86])
ax.set_facecolor(CREAM)

n = len(words_ranked)
ypos = list(range(n))[::-1]          # first (largest) word at top
maxc = max(w["count"] for w in words_ranked)
xmax = maxc * 1.14                    # headroom for value labels

for y, w in zip(ypos, words_ranked):
    c = w["count"]
    col = SENT_COLOR[w["sentiment"]]
    ax.barh(y, c, height=0.62, color=col, edgecolor="none", zorder=3)
    ax.text(c + maxc * 0.013, y, f"{c}", va="center", ha="left",
            fontsize=11.5, fontweight="bold", color=INK, zorder=5)

# word labels (left)
ax.set_yticks(ypos)
ax.set_yticklabels([w["word"] for w in words_ranked],
                   fontsize=13, color=INK, fontweight="bold")
ax.set_ylim(-0.7, n - 0.3)

# x axis (counts -> respondents who mentioned the word)
ax.set_xlim(0, xmax)
xticks = [0, 50, 100, 150]
ax.set_xticks(xticks)
ax.set_xticklabels([str(t) for t in xticks], fontsize=9.5, color=MUTED)
ax.tick_params(axis="x", length=0)
ax.tick_params(axis="y", length=0)

for x in xticks[1:]:
    ax.axvline(x, color=GRID, lw=0.8, zorder=1)
for s in ax.spines.values():
    s.set_visible(False)
ax.spines["left"].set_visible(True)
ax.spines["left"].set_color("#C9C1B4")
ax.spines["left"].set_linewidth(1.0)

ax.set_xlabel("Mainintojen lukumäärä", fontsize=10, color=MUTED, labelpad=6)

# legend — sentiment
handles = [Patch(facecolor=SENT_COLOR[s], label=SENT_LABEL[s]) for s in present]
leg = ax.legend(handles=handles, loc="lower right",
                bbox_to_anchor=(1.0, 0.015),
                frameon=True, fontsize=10.5, ncol=1,
                title="Sävy", title_fontsize=10.5,
                handlelength=1.1, handleheight=1.0,
                labelspacing=0.5, borderpad=0.8)
leg.get_frame().set_facecolor("#FFFFFF")
leg.get_frame().set_edgecolor(GRID)
leg.get_frame().set_linewidth(0.8)
leg.get_title().set_color(INK)
leg.get_title().set_fontweight("bold")

fig.savefig(PNG, dpi=200, facecolor=CREAM, bbox_inches="tight", pad_inches=0.04)
plt.close(fig)
print("chart ->", PNG)

# ---------------------------------------------------------------- pptx
PX_CREAM = RGBColor(0xF7, 0xF3, 0xEC)
PX_INK = RGBColor(0x2B, 0x2B, 0x2B)
PX_TEAL = RGBColor(0x13, 0x61, 0x5E)
PX_RED = RGBColor(0xC0, 0x47, 0x3A)
PX_MUTED = RGBColor(0x6E, 0x6A, 0x63)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
slide = prs.slides.add_slide(prs.slide_layouts[6])

# background
bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = PX_CREAM
bg.line.fill.background()
bg.shadow.inherit = False
slide.shapes._spTree.remove(bg._element)
slide.shapes._spTree.insert(2, bg._element)


def textbox(l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for txt, sz, col, bold in line:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col
            r.font.name = "Liberation Sans"
    return tb


# Accent bar (red — flags the negative-dominated headline)
acc = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42),
                             Inches(0.10), Inches(0.92))
acc.fill.solid(); acc.fill.fore_color.rgb = PX_RED
acc.line.fill.background(); acc.shadow.inherit = False

# Title (Finnish key message) — derive top-2 negative words from data
neg_words = [w["word"].lower() for w in words_ranked
             if w["sentiment"] == "kielteinen"][:2]
neg_str = ", ".join(neg_words)
textbox(Inches(0.80), Inches(0.40), Inches(11.9), Inches(1.0),
        [[("Attendon spontaani mielikuva on yhä kielteisten sanojen hallitsema",
           23, PX_INK, True)],
         [(f"Kärjessä {neg_str} – kolme yleisintä mainintaa ovat kielteisiä",
           13.5, PX_MUTED, False)]],
        space_after=3)

# Section label above chart
textbox(Inches(0.80), Inches(1.52), Inches(11.9), Inches(0.32),
        [[("SPONTAANI MIELIKUVA  ·  TOP 10 mainituinta sanaa  ·  "
           f"{WAVE}", 11, PX_TEAL, True)]])

# Chart image
iw, ih = Image.open(PNG).size
ar = iw / ih
pic_w = Inches(12.15)
pic_h = Emu(int(pic_w / ar))
max_h = Inches(5.05)
if pic_h > max_h:
    pic_h = max_h
    pic_w = Emu(int(pic_h * ar))
pic_l = Inches(0.62)
pic_t = Inches(1.92)
slide.shapes.add_picture(str(PNG), pic_l, pic_t, width=pic_w, height=pic_h)

# Footer: question caption + base n
textbox(Inches(0.80), Inches(7.02), Inches(9.6), Inches(0.45),
        [[("Kysymys: ", 9.5, PX_MUTED, True),
          (f"”{QUESTION}”", 9.5, PX_MUTED, False)]])
# base "kaikki vastaajat, n=863" -> "Kaikki vastaajat, n = 863"
base_disp = BASE[0].upper() + BASE[1:]
base_disp = base_disp.replace("n=", "n = ")
textbox(Inches(10.5), Inches(7.02), Inches(2.3), Inches(0.45),
        [[(base_disp, 9.5, PX_MUTED, True)]],
        align=PP_ALIGN.RIGHT)

prs.save(OUT)
print("slide ->", OUT)
