#!/usr/bin/env python3
"""Build the aided-awareness TREND slide (nSight house style).

Multi-line time-series: 4 survey waves on the x-axis, one line per provider.
The two leaders (Attendo, Esperi) are emphasised with strong colours, thick
lines, value labels and bold direct end labels; the remaining providers are
muted (grey) with light lines and small direct end labels — no legend.

The first wave (Toukokuu 2024) has a null value for Validia (the provider was
not measured that wave); its line simply begins at the first wave it appears.

All numbers are taken verbatim from original_data.json.

Run:  uv run python chart_lab/agent_build_trend.py
"""
import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ---------------------------------------------------------------- paths
ROOT = Path("/home/johan/Projects/nsight/proto")
DATA = ROOT / "chart_lab" / "original_data.json"
PNG = ROOT / "chart_lab" / "agent_trend_chart.png"
OUT = ROOT / "work" / "agent_slide_trend.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------- fonts
for f in ["/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
          "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"]:
    fm.fontManager.addfont(f)
plt.rcParams["font.family"] = "Liberation Sans"

# ---------------------------------------------------------------- data
d = json.loads(DATA.read_text())
cats = d["categories"]
series = d["series"]

WAVES = ["Toukokuu 2024", "Marraskuu 2024", "Toukokuu 2025", "Marraskuu 2025"]
SHORT = ["Touko\n2024", "Marras\n2024", "Touko\n2025", "Marras\n2025"]

# providers (exclude the "none of these" row), keep survey order
brands = [b for b in cats if b != "En mitään näistä"]
bidx = {b: cats.index(b) for b in brands}

# per-brand value list across waves (None preserved -> line gap)
vals = {b: [series[w][bidx[b]] for w in WAVES] for b in brands}

# ---------------------------------------------------------------- colours
CREAM = "#F4EFE6"     # nSight cream
INK = "#2B2B2B"
MUTED = "#8A857B"
GRID = "#E2DBCD"

LEAD = {"Attendo": "#13615E", "Esperi": "#C46A1E"}   # teal + warm ochre
MUTE = "#B7AFA0"

# ---------------------------------------------------------------- figure
FIG_W, FIG_H = 11.7, 5.55
fig = plt.figure(figsize=(FIG_W, FIG_H), dpi=200)
fig.patch.set_facecolor(CREAM)
# leave generous right margin for the direct end labels
ax = fig.add_axes([0.045, 0.105, 0.78, 0.875])
ax.set_facecolor(CREAM)

x = list(range(len(WAVES)))

# --- draw muted providers first (behind) -------------------------------
for b in brands:
    if b in LEAD:
        continue
    ys = vals[b]
    xs = [xi for xi, y in zip(x, ys) if y is not None]
    yy = [y for y in ys if y is not None]
    ax.plot(xs, yy, color=MUTE, lw=1.6, zorder=2, solid_capstyle="round")
    ax.plot(xs, yy, "o", color=MUTE, ms=3.4, zorder=2)

# --- draw the two leaders on top ---------------------------------------
for b, c in LEAD.items():
    ys = vals[b]
    xs = [xi for xi, y in zip(x, ys) if y is not None]
    yy = [y for y in ys if y is not None]
    ax.plot(xs, yy, color=c, lw=3.4, zorder=5, solid_capstyle="round")
    ax.plot(xs, yy, "o", color=c, ms=7.0, zorder=6,
            markeredgecolor=CREAM, markeredgewidth=1.4)
    # value labels above each point for the leaders
    for xi, y in zip(xs, yy):
        ax.annotate(f"{y}", (xi, y), textcoords="offset points",
                    xytext=(0, 10), ha="center", va="bottom",
                    fontsize=10.5, fontweight="bold", color=c, zorder=7)

# ---------------------------------------------------------------- axes
ax.set_xlim(-0.18, len(WAVES) - 1 + 0.06)
ax.set_ylim(0, 95)
ax.set_xticks(x)
ax.set_xticklabels(SHORT, fontsize=11, color=INK)
ax.set_yticks([0, 20, 40, 60, 80])
ax.set_yticklabels(["0", "20", "40", "60", "80"], fontsize=9.5, color=MUTED)
ax.tick_params(axis="both", length=0)

for yv in [20, 40, 60, 80]:
    ax.axhline(yv, color=GRID, lw=0.9, zorder=1)
ax.axhline(0, color="#C9C1B4", lw=1.0, zorder=1)

for s in ax.spines.values():
    s.set_visible(False)

# ---------------------------------------------------------------- end labels
# Build (last value, brand, colour, is_lead) for every line, then
# de-collide the y positions so nothing overlaps.
end_x = len(WAVES) - 1
labels = []
for b in brands:
    # last non-null value defines the label anchor
    last_y = None
    for y in vals[b]:
        if y is not None:
            last_y = y
    labels.append([last_y, b, LEAD.get(b, MUTE), b in LEAD])

# sort by y descending and push apart to a minimum gap
labels.sort(key=lambda r: r[0], reverse=True)
MIN_GAP = 4.4
placed_y = []
for i, r in enumerate(labels):
    ty = r[0]
    if placed_y and placed_y[-1] - ty < MIN_GAP:
        ty = placed_y[-1] - MIN_GAP
    placed_y.append(ty)
    r.append(ty)  # r = [val, brand, colour, lead, target_y]

for val, b, col, lead, ty in labels:
    txt = f"{b}  {val}" if lead else f"{b}  {val}"
    ax.annotate(
        txt,
        xy=(end_x, val), xycoords="data",
        xytext=(end_x + 0.10, ty), textcoords="data",
        ha="left", va="center",
        fontsize=11.5 if lead else 9.6,
        fontweight="bold" if lead else "normal",
        color=col if lead else MUTED,
        zorder=8,
        annotation_clip=False,
    )

fig.savefig(PNG, dpi=200, facecolor=CREAM, bbox_inches="tight", pad_inches=0.05)
plt.close(fig)
print("chart ->", PNG)

# ================================================================ pptx
PX_CREAM = RGBColor(0xF4, 0xEF, 0xE6)
PX_INK = RGBColor(0x2B, 0x2B, 0x2B)
PX_TEAL = RGBColor(0x13, 0x61, 0x5E)
PX_OCHRE = RGBColor(0xC4, 0x6A, 0x1E)
PX_MUTED = RGBColor(0x6E, 0x6A, 0x63)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
slide = prs.slides.add_slide(prs.slide_layouts[6])

# background
bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = PX_CREAM
bg.line.fill.background()
bg.shadow.inherit = False
slide.shapes._spTree.remove(bg._element)
slide.shapes._spTree.insert(2, bg._element)


def textbox(l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for txt, sz, col, bold in line:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col
            r.font.name = "Liberation Sans"
    return tb


# Accent bar at top-left
acc = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42),
                             Inches(0.10), Inches(0.92))
acc.fill.solid(); acc.fill.fore_color.rgb = PX_TEAL
acc.line.fill.background(); acc.shadow.inherit = False

# Title (Finnish key message) — exact leader figures from data
textbox(Inches(0.80), Inches(0.40), Inches(12.0), Inches(1.0),
        [[("Attendon ja Esperin tunnettuus on pysynyt vakaana – ja selvästi muita edellä",
           23, PX_INK, True)],
         [("Attendo 84 → 86 % ja Esperi 72 → 75 % aaltojen yli; muut tarjoajat 42 % ja alle",
           13.5, PX_MUTED, False)]],
        space_after=3)

# Section label above chart
textbox(Inches(0.80), Inches(1.52), Inches(12.0), Inches(0.32),
        [[("AUTETTU TUNNETTUUS  ·  kehitys mittausaalloittain (%)",
           11, PX_TEAL, True)]])

# Chart image
iw, ih = Image.open(PNG).size
ar = iw / ih
pic_w = Inches(12.15)
pic_h = Emu(int(pic_w / ar))
max_h = Inches(5.00)
if pic_h > max_h:
    pic_h = max_h
    pic_w = Emu(int(pic_h * ar))
pic_l = Inches((13.333 - pic_w / Inches(1)) / 2) if False else Inches(0.62)
pic_t = Inches(1.92)
slide.shapes.add_picture(str(PNG), pic_l, pic_t, width=pic_w, height=pic_h)

# Footer: question caption + base n
textbox(Inches(0.80), Inches(7.02), Inches(9.8), Inches(0.45),
        [[("Kysymys: ", 9.5, PX_MUTED, True),
          ("”Mitä seuraavista hoivapalveluiden tarjoajista tunnet "
           "vähintään nimeltä?”  ·  Validia mitattu "
           "ensimmäisen kerran 11/2024", 9.5, PX_MUTED, False)]])
textbox(Inches(10.6), Inches(7.02), Inches(2.2), Inches(0.45),
        [[("Kaikki vastaajat, n ≈ 1001 / aalto", 9.5, PX_MUTED, True)]],
        align=PP_ALIGN.RIGHT)

prs.save(OUT)
print("slide ->", OUT)
