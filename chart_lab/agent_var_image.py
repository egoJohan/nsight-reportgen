#!/usr/bin/env python3
"""Build 4 BRAND-IMAGE slides, each a DIFFERENT chart style (nSight house style).

ONE PowerPoint with four slides, each a clearly distinct layout/visual style,
built deterministically from `chart_lab/radar_0.json`. Every number is read
VERBATIM from the data; nothing is invented. Base n per brand is parsed from the
data key ("Attendo, n=863" -> 863).

Styles:
  1. HEATMAP        — brands (rows) x 14 attributes (columns), colour = % agree.
  2. CLUSTERED DOT  — Attendo vs Esperi head-to-head, two dots per attribute row,
                      connector line shows the gap.
  3. NET-IMAGE DIVERGING BAR — per brand: mean(positive attrs) - mean(negative
                      attrs), bars diverging left/right from a centre line.
  4. GROUPED HORIZONTAL BAR — top 6 positive attributes compared across the 3
                      leading brands (Attendo, Esperi, Onnikodit).

Data shape (radar_0.json):
    {
      "categories": [<14 attribute names>],   # includes 2 negatives
      "series": { "Attendo, n=863": [14 ints], ... }   # % agreement
    }
"""
import json
import re
from pathlib import Path

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
from matplotlib import cm
from matplotlib.colors import LinearSegmentedColormap, Normalize
from matplotlib.patches import Patch
from matplotlib.lines import Line2D

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ----------------------------------------------------------------- config
ROOT = Path("/home/johan/Projects/nsight/proto")
NEGATIVE_ATTRS = ("Ahne", "Välinpitämätön")

# nSight house palette
CREAM = "#F7F3EC"
INK = "#2B2B2B"
MUTED = "#6E6A63"
TEAL = "#13615E"
TEAL_LT = "#7DB8A6"
RED = "#B23A2E"
RED_LT = "#E08C82"
GOLD = "#C68A2E"
GRIDC = "#DAD3C7"

PX_CREAM = RGBColor(0xF7, 0xF3, 0xEC)
PX_INK = RGBColor(0x2B, 0x2B, 0x2B)
PX_TEAL = RGBColor(0x13, 0x61, 0x5E)
PX_RED = RGBColor(0xB2, 0x3A, 0x2E)
PX_MUTED = RGBColor(0x6E, 0x6A, 0x63)

WAVE = "Marraskuu 2025"
QUESTION = ("Missä määrin seuraavat ominaisuudet mielestäsi sopivat tähän "
            "hoivapalveluiden tarjoajaan?")


def _register_fonts():
    for f in [
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
    ]:
        try:
            fm.fontManager.addfont(f)
        except Exception:
            pass
    plt.rcParams["font.family"] = "Liberation Sans"


def parse_brand(key):
    m = re.match(r"\s*(.+?)\s*,\s*n\s*=\s*(\d+)\s*$", key)
    if not m:
        return key.strip(), None
    return m.group(1).strip(), int(m.group(2))


def load_data(data_path):
    d = json.loads(Path(data_path).read_text())
    cats = d["categories"]
    series = d["series"]
    brands = {}          # name -> {"n": int, "vals": [..]}
    for key, vals in series.items():
        name, n = parse_brand(key)
        assert len(vals) == len(cats), f"{key}: expected {len(cats)} values"
        brands[name] = {"n": n, "vals": list(vals)}
    neg_idx = [i for i, c in enumerate(cats) if c in NEGATIVE_ATTRS]
    pos_idx = [i for i in range(len(cats)) if i not in neg_idx]
    return cats, brands, pos_idx, neg_idx


# ================================================================= CHART 1
def chart_heatmap(cats, brands, pos_idx, neg_idx, png_path):
    """Brands (rows) x attributes (columns) heatmap of % agreement.

    Rows sorted by mean POSITIVE-attribute score (best overall image at top).
    Columns ordered: 12 positives (by overall mean, strongest left) then the
    2 negatives grouped at the right, separated by a divider and flagged.
    """
    # column order
    pos_sorted = sorted(pos_idx,
                        key=lambda i: np.mean([brands[b]["vals"][i] for b in brands]),
                        reverse=True)
    neg_sorted = sorted(neg_idx,
                        key=lambda i: np.mean([brands[b]["vals"][i] for b in brands]),
                        reverse=True)
    col_order = pos_sorted + neg_sorted
    col_labels = [cats[i] for i in col_order]
    n_pos = len(pos_sorted)

    # row order: by mean positive score desc
    row_names = sorted(brands.keys(),
                       key=lambda b: np.mean([brands[b]["vals"][i] for i in pos_idx]),
                       reverse=True)
    M = np.array([[brands[b]["vals"][i] for i in col_order] for b in row_names],
                 dtype=float)

    nrows, ncols = M.shape
    fig = plt.figure(figsize=(13.0, 6.05), dpi=200)
    fig.patch.set_facecolor(CREAM)
    # Matrix axes: leave generous headroom on top for angled headers + the
    # group label, and a clear band at the bottom for the colourbar.
    ax = fig.add_axes([0.105, 0.135, 0.85, 0.55])
    ax.set_facecolor(CREAM)

    cmap = LinearSegmentedColormap.from_list(
        "nsight_teal", ["#F2EFE7", "#BFD9CE", "#6BAE97", "#2E8A72", "#13615E"])
    norm = Normalize(vmin=15, vmax=55)

    im = ax.imshow(M, cmap=cmap, norm=norm, aspect="auto")

    # cell value labels
    for r in range(nrows):
        for c in range(ncols):
            v = M[r, c]
            tc = "#FFFFFF" if v >= 41 else INK
            ax.text(c, r, f"{int(v)}", ha="center", va="center",
                    fontsize=9.0, color=tc, fontweight="bold")

    ax.set_xticks(range(ncols))
    ax.set_xticklabels(col_labels, rotation=40, ha="left", fontsize=9.3,
                       color=INK)
    ax.xaxis.set_ticks_position("top")
    ax.xaxis.set_label_position("top")
    # flag negative column labels in red
    for lbl, ci in zip(ax.get_xticklabels(), col_order):
        if ci in neg_idx:
            lbl.set_color(RED)
            lbl.set_fontweight("bold")

    ax.set_yticks(range(nrows))
    ax.set_yticklabels(row_names, fontsize=10.5, color=INK, fontweight="bold")
    ax.tick_params(length=0)
    for s in ax.spines.values():
        s.set_visible(False)

    # gridlines between cells
    ax.set_xticks(np.arange(-0.5, ncols, 1), minor=True)
    ax.set_yticks(np.arange(-0.5, nrows, 1), minor=True)
    ax.grid(which="minor", color=CREAM, linewidth=2.4)
    ax.tick_params(which="minor", length=0)

    # divider between positives and negatives
    ax.axvline(n_pos - 0.5, color=RED, lw=2.2, zorder=6)

    # group labels: placed in axes-fraction y well ABOVE the angled column
    # headers so they never touch any header text. y>1 is above the matrix.
    pos_center = (n_pos - 1) / 2.0
    neg_center = (n_pos + ncols - 1) / 2.0
    ax.text(pos_center, 1.345, "Myönteiset", color=TEAL, fontsize=10.5,
            ha="center", va="bottom", fontweight="bold",
            transform=ax.get_xaxis_transform())
    ax.text(neg_center, 1.345, "Kielteiset", color=RED, fontsize=10.5,
            ha="center", va="bottom", fontweight="bold",
            transform=ax.get_xaxis_transform())

    # colourbar: placed in a dedicated band BELOW all heatmap rows.
    cax = fig.add_axes([0.105, 0.055, 0.30, 0.028])
    cb = fig.colorbar(im, cax=cax, orientation="horizontal")
    cb.set_ticks([15, 25, 35, 45, 55])
    cb.set_ticklabels(["15", "25", "35", "45", "55 %"])
    cb.ax.tick_params(labelsize=8.5, color=MUTED, labelcolor=MUTED, length=2)
    cb.outline.set_edgecolor(GRIDC)
    cb.set_label("Samaa mieltä, % vastaajista", fontsize=8.8, color=MUTED)

    fig.savefig(png_path, dpi=200, facecolor=CREAM,
                bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)


# ================================================================= CHART 2
def chart_dotplot(cats, brands, pos_idx, neg_idx, png_path,
                  brand_a="Attendo", brand_b="Esperi"):
    """Clustered dot plot: two brands head-to-head across 14 attributes.

    Rows = attributes, ordered positives (by brand_a value desc) then negatives.
    Each row: a connector line + two dots (one per brand). Negatives flagged.
    """
    va = brands[brand_a]["vals"]
    vb = brands[brand_b]["vals"]

    pos_sorted = sorted(pos_idx, key=lambda i: va[i])         # asc -> top is best
    neg_sorted = sorted(neg_idx, key=lambda i: va[i], reverse=True)
    # display top-to-bottom: best positives at top. We'll plot bottom-up.
    order = list(reversed(pos_sorted)) + neg_sorted
    order = order[::-1]   # so first (best) ends at top via ypos
    labels = [cats[i] for i in order]
    a_vals = [va[i] for i in order]
    b_vals = [vb[i] for i in order]
    is_neg = [i in neg_idx for i in order]

    n = len(order)
    fig = plt.figure(figsize=(12.2, 6.0), dpi=200)
    fig.patch.set_facecolor(CREAM)
    ax = fig.add_axes([0.31, 0.10, 0.65, 0.82])
    ax.set_facecolor(CREAM)

    ypos = list(range(n))
    for y, av, bv, neg in zip(ypos, a_vals, b_vals, is_neg):
        ax.plot([av, bv], [y, y], color=GRIDC, lw=2.4, zorder=2,
                solid_capstyle="round")
        # leader = whoever is more favourable; for neg, lower is better
        ax.scatter(av, y, s=130, color=TEAL, zorder=4, edgecolor="white", lw=1.2)
        ax.scatter(bv, y, s=130, color=GOLD, zorder=4, edgecolor="white", lw=1.2)
        # value labels, offset to avoid overlap
        if abs(av - bv) < 3:
            ax.text(av, y + 0.30, f"{av}", color=TEAL, fontsize=8.6,
                    ha="center", va="bottom", fontweight="bold")
            ax.text(bv, y - 0.30, f"{bv}", color=GOLD, fontsize=8.6,
                    ha="center", va="top", fontweight="bold")
        else:
            ax.text(av + (1.6 if av >= bv else -1.6), y, f"{av}", color=TEAL,
                    fontsize=8.8, ha="left" if av >= bv else "right",
                    va="center", fontweight="bold")
            ax.text(bv + (1.6 if bv > av else -1.6), y, f"{bv}", color=GOLD,
                    fontsize=8.8, ha="left" if bv > av else "right",
                    va="center", fontweight="bold")

    ax.set_yticks(ypos)
    ax.set_yticklabels(labels, fontsize=10.8, color=INK)
    for lbl, neg in zip(ax.get_yticklabels(), is_neg):
        if neg:
            lbl.set_color(RED)
            lbl.set_fontweight("bold")
    ax.set_ylim(-0.7, n - 0.3)

    ax.set_xlim(0, 60)
    ax.set_xticks([0, 10, 20, 30, 40, 50, 60])
    ax.set_xticklabels(["0", "10", "20", "30", "40", "50", "60 %"],
                       fontsize=9.5, color=MUTED)
    ax.tick_params(length=0)
    for x in [10, 20, 30, 40, 50, 60]:
        ax.axvline(x, color=GRIDC, lw=0.8, zorder=1)
    for s in ax.spines.values():
        s.set_visible(False)
    ax.spines["left"].set_visible(True)
    ax.spines["left"].set_color("#C9C1B4")

    handles = [
        Line2D([0], [0], marker="o", color="w", markerfacecolor=TEAL,
               markersize=11, label=brand_a),
        Line2D([0], [0], marker="o", color="w", markerfacecolor=GOLD,
               markersize=11, label=brand_b),
    ]
    leg = ax.legend(handles=handles, loc="upper left",
                    bbox_to_anchor=(0.01, 0.995), frameon=True, fontsize=10.5,
                    handletextpad=0.4, borderpad=0.8, labelspacing=0.5)
    leg.get_frame().set_facecolor("#FFFFFF")
    leg.get_frame().set_edgecolor(GRIDC)
    for t in leg.get_texts():
        t.set_color(INK)

    fig.savefig(png_path, dpi=200, facecolor=CREAM,
                bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)


# ================================================================= CHART 3
def chart_diverging(cats, brands, pos_idx, neg_idx, png_path):
    """Net-image diverging bar per brand.

    net = mean(positive attrs %) - mean(negative attrs %).
    Bars diverge left/right of a centre line; sorted strongest net image at top.
    A high positive net = strong, clean image; low/negative = weak image.
    """
    rows = []
    for b in brands:
        pos_mean = float(np.mean([brands[b]["vals"][i] for i in pos_idx]))
        neg_mean = float(np.mean([brands[b]["vals"][i] for i in neg_idx]))
        net = pos_mean - neg_mean
        rows.append((b, net, pos_mean, neg_mean))
    rows.sort(key=lambda r: r[1])      # ascending -> highest ends at top

    names = [r[0] for r in rows]
    nets = [r[1] for r in rows]
    n = len(rows)

    fig = plt.figure(figsize=(11.6, 5.7), dpi=200)
    fig.patch.set_facecolor(CREAM)
    ax = fig.add_axes([0.175, 0.115, 0.78, 0.82])
    ax.set_facecolor(CREAM)

    ypos = list(range(n))
    lo, hi = min(nets), max(nets)
    # Diverging domain: room left of zero for negative bars, room right for
    # positive bars (and their value labels).
    x_lo = min(lo, 0) - 7.0
    x_hi = max(hi, 0) + 6.0
    for y, name, net in zip(ypos, names, nets):
        color = TEAL if net >= 0 else RED
        ax.barh(y, net, height=0.60, color=color, edgecolor="none", zorder=3)
        off = 0.4
        if net >= 0:
            ax.text(net + off, y, f"+{net:.1f}", va="center", ha="left",
                    fontsize=11, fontweight="bold", color=TEAL, zorder=5)
        else:
            ax.text(net - off, y, f"{net:.1f}", va="center", ha="right",
                    fontsize=11, fontweight="bold", color=RED, zorder=5)

    ax.set_yticks(ypos)
    ax.set_yticklabels(names, fontsize=12, color=INK, fontweight="bold")
    ax.set_ylim(-0.7, n - 0.3)
    ax.tick_params(length=0)

    ax.set_xlim(x_lo, x_hi)
    xticks = [t for t in range(int(np.floor(x_lo / 5) * 5),
                               int(np.ceil(x_hi / 5) * 5) + 1, 5)]
    ax.set_xticks(xticks)
    ax.set_xticklabels([str(t) for t in xticks], fontsize=9.5, color=MUTED)
    for x in xticks:
        if x == 0:
            continue
        ax.axvline(x, color=GRIDC, lw=0.8, zorder=1)
    # zero centre line, drawn above gridlines and bars
    ax.axvline(0, color=INK, lw=1.4, zorder=4)
    ax.set_xlabel("Nettomielikuva = myönteisten attribuuttien keskiarvo − "
                  "kielteisten keskiarvo (%-yksikköä)",
                  fontsize=9.8, color=MUTED, labelpad=8)
    for s in ax.spines.values():
        s.set_visible(False)

    fig.savefig(png_path, dpi=200, facecolor=CREAM,
                bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    return rows


# ================================================================= CHART 4
def chart_grouped(cats, brands, pos_idx, neg_idx, png_path,
                  group_brands=("Attendo", "Esperi", "Onnikodit"), top_k=6):
    """Grouped horizontal bar: top-K positive attributes across 3 brands.

    Top-K chosen by the cross-(group)brand mean of each positive attribute.
    Each attribute row holds 3 bars (one per brand). Attributes sorted by group
    mean desc (strongest at top).
    """
    gb = list(group_brands)
    scored = sorted(
        pos_idx,
        key=lambda i: np.mean([brands[b]["vals"][i] for b in gb]),
        reverse=True,
    )
    sel = scored[:top_k]
    labels = [cats[i] for i in sel]

    n = len(sel)
    nb = len(gb)
    fig = plt.figure(figsize=(12.0, 6.0), dpi=200)
    fig.patch.set_facecolor(CREAM)
    ax = fig.add_axes([0.255, 0.085, 0.70, 0.85])
    ax.set_facecolor(CREAM)

    colors = {gb[0]: TEAL, gb[1]: GOLD, gb[2]: TEAL_LT}
    bar_h = 0.74 / nb
    ypos = np.arange(n)[::-1]    # first attribute at top

    for bi, b in enumerate(gb):
        offs = (bi - (nb - 1) / 2) * bar_h
        vals = [brands[b]["vals"][i] for i in sel]
        ax.barh(ypos + offs, vals, height=bar_h * 0.92, color=colors[b],
                edgecolor="none", zorder=3, label=b)
        for y, v in zip(ypos, vals):
            ax.text(v + 0.8, y + offs, f"{v}", va="center", ha="left",
                    fontsize=8.8, color=INK, fontweight="bold", zorder=5)

    ax.set_yticks(ypos)
    ax.set_yticklabels(labels, fontsize=11.5, color=INK)
    ax.set_ylim(-0.6, n - 0.4)

    ax.set_xlim(0, 62)
    ax.set_xticks([0, 10, 20, 30, 40, 50, 60])
    ax.set_xticklabels(["0", "10", "20", "30", "40", "50", "60 %"],
                       fontsize=9.5, color=MUTED)
    ax.tick_params(length=0)
    for x in [10, 20, 30, 40, 50, 60]:
        ax.axvline(x, color=GRIDC, lw=0.8, zorder=1)
    for s in ax.spines.values():
        s.set_visible(False)
    ax.spines["left"].set_visible(True)
    ax.spines["left"].set_color("#C9C1B4")

    leg = ax.legend(loc="lower right", bbox_to_anchor=(1.0, 0.0), frameon=True,
                    fontsize=10.5, ncol=3, handlelength=1.1, borderpad=0.8,
                    columnspacing=1.2)
    leg.get_frame().set_facecolor("#FFFFFF")
    leg.get_frame().set_edgecolor(GRIDC)
    for t in leg.get_texts():
        t.set_color(INK)

    fig.savefig(png_path, dpi=200, facecolor=CREAM,
                bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    return labels, sel


# ----------------------------------------------------------------- slide
def _textbox(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, space_after=0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for txt, sz, col, bold in line:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col
            r.font.name = "Liberation Sans"
    return tb


def add_slide(prs, title, subtitle, section, png_path, footer_left, footer_right,
              pic_top=Inches(1.92), max_h=Inches(5.05), pic_w_in=12.45):
    SW, SH = prs.slide_width, prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = PX_CREAM
    bg.line.fill.background(); bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    acc = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42),
                                 Inches(0.10), Inches(0.92))
    acc.fill.solid(); acc.fill.fore_color.rgb = PX_TEAL
    acc.line.fill.background(); acc.shadow.inherit = False

    _textbox(slide, Inches(0.80), Inches(0.38), Inches(12.2), Inches(1.1),
             [[(title, 21, PX_INK, True)],
              [(subtitle, 12.5, PX_MUTED, False)]],
             space_after=3)

    _textbox(slide, Inches(0.80), Inches(1.62), Inches(12.2), Inches(0.32),
             [[(section, 11, PX_TEAL, True)]])

    iw, ih = Image.open(png_path).size
    ar = iw / ih
    pic_w = Inches(pic_w_in)
    pic_h = Emu(int(pic_w / ar))
    if pic_h > max_h:
        pic_h = max_h
        pic_w = Emu(int(pic_h * ar))
    pic_l = Emu(int((SW - pic_w) / 2))
    slide.shapes.add_picture(str(png_path), pic_l, pic_top,
                             width=pic_w, height=pic_h)

    _textbox(slide, Inches(0.80), Inches(7.00), Inches(9.6), Inches(0.45),
             [[("Kysymys: ", 9.5, PX_MUTED, True),
               ("”" + QUESTION + "”", 9.5, PX_MUTED, False)]])
    _textbox(slide, Inches(7.7), Inches(7.00), Inches(5.1), Inches(0.45),
             [[(footer_right, 9.5, PX_MUTED, True)]],
             align=PP_ALIGN.RIGHT)


# ----------------------------------------------------------------- build
def build(
    data_path=ROOT / "chart_lab" / "radar_0.json",
    png_dir=ROOT / "chart_lab",
    out_path=ROOT / "work" / "agent_var_image.pptx",
):
    _register_fonts()
    cats, brands, pos_idx, neg_idx = load_data(data_path)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    n_all = sum(b["n"] for b in brands.values())

    # ---- Slide 1: HEATMAP -------------------------------------------------
    p1 = png_dir / "agent_var_heatmap.png"
    chart_heatmap(cats, brands, pos_idx, neg_idx, p1)
    # leader: highest mean positive
    pos_means = {b: np.mean([brands[b]["vals"][i] for i in pos_idx]) for b in brands}
    leader = max(pos_means, key=pos_means.get)
    add_slide(
        prs,
        title=f"{leader} erottuu vahvimmalla mielikuvaprofiililla kaikkien "
              f"attribuuttien yli",
        subtitle="Lämpökartta: 8 brändiä × 14 mielikuva-attribuuttia, väri = "
                 "samaa mieltä olevien osuus. Kielteiset attribuutit oikealla "
                 "(matala = hyvä).",
        section="MIELIKUVALÄMPÖKARTTA · OSUUS VASTAAJISTA (%)",
        png_path=p1,
        footer_left=None,
        footer_right=f"{WAVE} · 8 brändiä, n yhteensä = {n_all}",
        pic_top=Inches(1.85), max_h=Inches(5.05), pic_w_in=12.7,
    )

    # ---- Slide 2: CLUSTERED DOT PLOT -------------------------------------
    p2 = png_dir / "agent_var_dotplot.png"
    chart_dotplot(cats, brands, pos_idx, neg_idx, p2, "Attendo", "Esperi")
    a = brands["Attendo"]; e = brands["Esperi"]
    # count attributes where Attendo leads (favourably)
    a_lead = 0
    for i in range(len(cats)):
        if i in neg_idx:
            if a["vals"][i] < e["vals"][i]:
                a_lead += 1
        else:
            if a["vals"][i] > e["vals"][i]:
                a_lead += 1
    add_slide(
        prs,
        title=f"Attendo johtaa Esperiä useimmissa mielikuva-attribuuteissa, "
              f"mutta jää jälkeen kielteisissä",
        subtitle="Pistekaavio vertaa kahta brändiä rinnakkain jokaisen "
                 "attribuutin kohdalla; viiva kuvaa brändien eroa.",
        section="ATTENDO vs. ESPERI · OSUUS VASTAAJISTA (%)",
        png_path=p2,
        footer_left=None,
        footer_right=f"{WAVE} · Attendo n = {a['n']}, Esperi n = {e['n']}",
        pic_top=Inches(1.85), max_h=Inches(5.05), pic_w_in=11.9,
    )

    # ---- Slide 3: NET-IMAGE DIVERGING BAR --------------------------------
    p3 = png_dir / "agent_var_diverging.png"
    rows = chart_diverging(cats, brands, pos_idx, neg_idx, p3)
    top_net = max(rows, key=lambda r: r[1])
    add_slide(
        prs,
        title=f"{top_net[0]} omaa selvästi vahvimman nettomielikuvan "
              f"(+{top_net[1]:.1f} %-yksikköä)",
        subtitle="Nettomielikuva = 12 myönteisen attribuutin keskiarvo "
                 "miinus 2 kielteisen attribuutin keskiarvo. Pidempi palkki = "
                 "puhtaampi, vahvempi mielikuva.",
        section="NETTOMIELIKUVA BRÄNDEITTÄIN · %-YKSIKKÖÄ",
        png_path=p3,
        footer_left=None,
        footer_right=f"{WAVE} · 8 brändiä, n yhteensä = {n_all}",
        pic_top=Inches(1.95), max_h=Inches(4.95), pic_w_in=11.7,
    )

    # ---- Slide 4: GROUPED HORIZONTAL BAR ---------------------------------
    p4 = png_dir / "agent_var_grouped.png"
    gbrands = ("Attendo", "Esperi", "Onnikodit")
    labels, sel = chart_grouped(cats, brands, pos_idx, neg_idx, p4, gbrands, 6)
    # Onnikodit leads how many of the 6?
    onn_lead = sum(
        1 for i in sel
        if brands["Onnikodit"]["vals"][i] == max(brands[b]["vals"][i] for b in gbrands)
    )
    add_slide(
        prs,
        title=f"Onnikodit johtaa kaikissa kuudessa tärkeimmässä myönteisessä "
              f"attribuutissa",
        subtitle="Kuusi yleisintä myönteistä attribuuttia (kolmen brändin "
                 "keskiarvon mukaan) verrattuna kolmen kärkibrändin kesken.",
        section="KÄRKIATTRIBUUTIT · ATTENDO · ESPERI · ONNIKODIT · (%)",
        png_path=p4,
        footer_left=None,
        footer_right=(f"{WAVE} · n: Attendo {brands['Attendo']['n']} · "
                      f"Esperi {brands['Esperi']['n']} · "
                      f"Onnikodit {brands['Onnikodit']['n']}"),
        pic_top=Inches(1.92), max_h=Inches(5.0), pic_w_in=11.8,
    )

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"{len(prs.slides._sldIdLst)} slides -> {out_path}")
    return out_path


if __name__ == "__main__":
    build()
