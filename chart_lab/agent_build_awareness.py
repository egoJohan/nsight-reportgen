#!/usr/bin/env python3
"""Build the aided brand-awareness slide (nSight house style).

Horizontal grouped bar chart, providers ordered best->least known,
four survey waves grouped per provider, current wave (Marraskuu 2025)
emphasised. Numbers taken verbatim from original_data.json.
"""
import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
from matplotlib.patches import Patch

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------- paths
ROOT = Path("/home/johan/Projects/nsight/proto")
DATA = ROOT / "chart_lab" / "original_data.json"
PNG = ROOT / "chart_lab" / "agent_awareness_chart.png"
OUT = ROOT / "work" / "agent_slide_awareness.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------- fonts
for f in ["/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
          "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"]:
    fm.fontManager.addfont(f)
plt.rcParams["font.family"] = "Liberation Sans"

# ---------------------------------------------------------------- data
d = json.loads(DATA.read_text())
cats = d["categories"]                       # survey order
series = d["series"]

# Wave order: oldest -> newest so newest renders last/emphasised.
WAVES = ["Toukokuu 2024", "Marraskuu 2024", "Toukokuu 2025", "Marraskuu 2025"]
CURRENT = "Marraskuu 2025"

# Order providers best->least known by the CURRENT wave.
order = sorted(range(len(cats)), key=lambda i: series[CURRENT][i], reverse=True)
cats_o = [cats[i] for i in order]
vals = {w: [series[w][i] for i in order] for w in WAVES}

# ---------------------------------------------------------------- colours
CREAM = "#F7F3EC"          # light cream background
INK = "#2B2B2B"            # primary text
MUTED = "#6E6A63"          # captions
# Teal family: prior waves light -> current wave dark/strong.
WAVE_COLORS = {
    "Toukokuu 2024":  "#CFE3E2",
    "Marraskuu 2024": "#9CC6C4",
    "Toukokuu 2025":  "#5E9C9A",
    "Marraskuu 2025": "#13615E",   # current wave, strong teal
}

# ---------------------------------------------------------------- figure
# Sized to fill the slide content area at high res.
FIG_W, FIG_H = 11.6, 6.05
fig = plt.figure(figsize=(FIG_W, FIG_H), dpi=200)
fig.patch.set_facecolor(CREAM)
ax = fig.add_axes([0.165, 0.085, 0.805, 0.875])
ax.set_facecolor(CREAM)

n = len(cats_o)
ypos = list(range(n))[::-1]      # first category at top
group_h = 0.78
bar_h = group_h / len(WAVES)

for wi, w in enumerate(WAVES):
    # offset so the group is vertically centred on the category row
    off = group_h / 2 - bar_h * (wi + 0.5)
    ys = [y + off for y in ypos]
    xs = vals[w]
    is_cur = (w == CURRENT)
    for y, x in zip(ys, xs):
        if x is None:
            continue
        ax.barh(y, x, height=bar_h * 0.94, color=WAVE_COLORS[w],
                edgecolor="none", zorder=3)
        # value label at bar end
        ax.text(x + 0.8, y, f"{x}", va="center", ha="left",
                fontsize=7.6 if not is_cur else 8.6,
                fontweight="bold" if is_cur else "normal",
                color=INK if is_cur else MUTED, zorder=5)

# category labels
ax.set_yticks(ypos)
ax.set_yticklabels(cats_o, fontsize=11.5, color=INK)
# bold the leader
for lbl, idx in zip(ax.get_yticklabels(), range(n)):
    if lbl.get_text() == "Attendo":
        lbl.set_fontweight("bold")

# x axis
ax.set_xlim(0, 100)
ax.set_xticks([0, 20, 40, 60, 80, 100])
ax.set_xticklabels(["0", "20", "40", "60", "80", "100 %"],
                   fontsize=9.5, color=MUTED)
ax.tick_params(axis="x", length=0)
ax.tick_params(axis="y", length=0)
ax.set_ylim(-0.7, n - 0.3)

# grid + spines
for x in [20, 40, 60, 80, 100]:
    ax.axvline(x, color="#DAD3C7", lw=0.8, zorder=1)
for s in ax.spines.values():
    s.set_visible(False)
ax.spines["left"].set_visible(True)
ax.spines["left"].set_color("#C9C1B4")
ax.spines["left"].set_linewidth(1.0)

# legend (waves, oldest -> newest)
handles = [Patch(facecolor=WAVE_COLORS[w], label=w) for w in WAVES]
leg = ax.legend(handles=handles, loc="lower right",
                bbox_to_anchor=(1.0, 0.012),
                frameon=True, fontsize=9.5, ncol=2,
                title="Mittausajankohta", title_fontsize=9.5,
                handlelength=1.1, handleheight=1.0,
                columnspacing=1.3, labelspacing=0.45,
                borderpad=0.8)
leg.get_frame().set_facecolor("#FFFFFF")
leg.get_frame().set_edgecolor("#DAD3C7")
leg.get_frame().set_linewidth(0.8)
leg.get_title().set_color(INK)
leg.get_title().set_fontweight("bold")

fig.savefig(PNG, dpi=200, facecolor=CREAM, bbox_inches="tight", pad_inches=0.04)
plt.close(fig)
print("chart ->", PNG)

# ---------------------------------------------------------------- pptx
PX_CREAM = RGBColor(0xF7, 0xF3, 0xEC)
PX_INK = RGBColor(0x2B, 0x2B, 0x2B)
PX_TEAL = RGBColor(0x13, 0x61, 0x5E)
PX_MUTED = RGBColor(0x6E, 0x6A, 0x63)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
slide = prs.slides.add_slide(prs.slide_layouts[6])

# background
bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = PX_CREAM
bg.line.fill.background()
bg.shadow.inherit = False
slide.shapes._spTree.remove(bg._element)
slide.shapes._spTree.insert(2, bg._element)

def textbox(l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for txt, sz, col, bold in line:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col
            r.font.name = "Liberation Sans"
    return tb

# Accent bar at top-left
acc = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42),
                             Inches(0.10), Inches(0.92))
acc.fill.solid(); acc.fill.fore_color.rgb = PX_TEAL
acc.line.fill.background(); acc.shadow.inherit = False

# Title (Finnish key message)
textbox(Inches(0.80), Inches(0.40), Inches(11.9), Inches(1.0),
        [[("Attendo on selvästi tunnetuin yksityinen hoivapalveluiden tarjoaja",
           23, PX_INK, True)],
         [("Tunnettuus 86 % – selvästi Esperin (75 %) edellä ja vakaa edelliseen aaltoon nähden",
           13.5, PX_MUTED, False)]],
        space_after=3)

# Section label above chart
textbox(Inches(0.80), Inches(1.52), Inches(11.9), Inches(0.32),
        [[("AUTETTU TUNNETTUUS  ·  osuus vastaajista (%)", 11, PX_TEAL, True)]])

# Chart image
from PIL import Image
iw, ih = Image.open(PNG).size
ar = iw / ih
pic_w = Inches(12.15)
pic_h = Emu(int(pic_w / ar))
max_h = Inches(5.05)
if pic_h > max_h:
    pic_h = max_h
    pic_w = Emu(int(pic_h * ar))
pic_l = Inches(0.62)
pic_t = Inches(1.88)
slide.shapes.add_picture(str(PNG), pic_l, pic_t, width=pic_w, height=pic_h)

# Footer: question caption + base n
textbox(Inches(0.80), Inches(7.00), Inches(9.6), Inches(0.45),
        [[("Kysymys: ", 9.5, PX_MUTED, True),
          ("”Mitä seuraavista hoivapalveluiden tarjoajista tunnet "
           "vähintään nimeltä?”", 9.5, PX_MUTED, False)]])
textbox(Inches(10.5), Inches(7.00), Inches(2.3), Inches(0.45),
        [[("Kaikki vastaajat, n = 1001", 9.5, PX_MUTED, True)]],
        align=PP_ALIGN.RIGHT)

prs.save(OUT)
print("slide ->", OUT)
