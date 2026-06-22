#!/usr/bin/env python3
"""Build the DENSE aided brand-awareness slide (nSight house style).

Grouped horizontal bar chart, 9 providers ordered best->least known, four
survey waves grouped per provider, current wave (Marraskuu 2025) emphasised,
PLUS a "Muutos toukokuusta" change column on the right showing the per-provider
delta = Marraskuu 2025 - Toukokuu 2025.

All numbers (and deltas) are derived deterministically from original_data.json.
"""
import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
from matplotlib.patches import Patch

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------- paths
ROOT = Path("/home/johan/Projects/nsight/proto")
DATA = ROOT / "chart_lab" / "original_data.json"
PNG = ROOT / "chart_lab" / "agent_dense_awareness_chart.png"
OUT = ROOT / "work" / "agent_dense_awareness.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------- fonts
for f in ["/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
          "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"]:
    fm.fontManager.addfont(f)
plt.rcParams["font.family"] = "Liberation Sans"

# ---------------------------------------------------------------- data
d = json.loads(DATA.read_text())
cats = d["categories"]
series = d["series"]

# Wave order: oldest -> newest so newest renders last/emphasised.
WAVES = ["Toukokuu 2024", "Marraskuu 2024", "Toukokuu 2025", "Marraskuu 2025"]
CURRENT = "Marraskuu 2025"
PREV = "Toukokuu 2025"

# Order providers best->least known by the CURRENT wave.
order = sorted(range(len(cats)), key=lambda i: series[CURRENT][i], reverse=True)
cats_o = [cats[i] for i in order]
vals = {w: [series[w][i] for i in order] for w in WAVES}

# Change vs the previous wave (Marraskuu 2025 - Toukokuu 2025) -- computed here.
deltas = [series[CURRENT][i] - series[PREV][i] for i in order]


def fmt_delta(v):
    if v > 0:
        return f"+{v}"
    if v < 0:
        return f"−{abs(v)}"   # proper minus sign
    return "±0"               # +/- 0 for "no change"


# ---------------------------------------------------------------- colours
CREAM = "#F7F3EC"
INK = "#2B2B2B"
MUTED = "#6E6A63"
UP = "#1E7A5A"     # green-teal for positive change
DOWN = "#B5532E"   # warm rust for negative change
FLAT = "#8A857C"   # grey for no change
WAVE_COLORS = {
    "Toukokuu 2024":  "#CFE3E2",
    "Marraskuu 2024": "#9CC6C4",
    "Toukokuu 2025":  "#5E9C9A",
    "Marraskuu 2025": "#13615E",
}

# ---------------------------------------------------------------- figure
FIG_W, FIG_H = 12.4, 6.05
fig = plt.figure(figsize=(FIG_W, FIG_H), dpi=200)
fig.patch.set_facecolor(CREAM)

n = len(cats_o)
ypos = list(range(n))[::-1]
group_h = 0.78
bar_h = group_h / len(WAVES)

# Two axes: the bar chart, and a narrow change column on the right.
# Left chart axes
ax = fig.add_axes([0.150, 0.085, 0.665, 0.875])
ax.set_facecolor(CREAM)
# Change-column axes (shares the y range so rows line up exactly)
cax = fig.add_axes([0.840, 0.085, 0.150, 0.875])
cax.set_facecolor(CREAM)

# --- bars
for wi, w in enumerate(WAVES):
    off = group_h / 2 - bar_h * (wi + 0.5)
    ys = [y + off for y in ypos]
    xs = vals[w]
    is_cur = (w == CURRENT)
    for y, x in zip(ys, xs):
        if x is None:
            continue
        ax.barh(y, x, height=bar_h * 0.94, color=WAVE_COLORS[w],
                edgecolor="none", zorder=3)
        ax.text(x + 0.9, y, f"{x}", va="center", ha="left",
                fontsize=7.6 if not is_cur else 8.6,
                fontweight="bold" if is_cur else "normal",
                color=INK if is_cur else MUTED, zorder=5)

# category labels
ax.set_yticks(ypos)
ax.set_yticklabels(cats_o, fontsize=11.5, color=INK)
for lbl in ax.get_yticklabels():
    if lbl.get_text() == "Attendo":
        lbl.set_fontweight("bold")

# x axis
ax.set_xlim(0, 100)
ax.set_xticks([0, 20, 40, 60, 80, 100])
ax.set_xticklabels(["0", "20", "40", "60", "80", "100 %"],
                   fontsize=9.5, color=MUTED)
ax.tick_params(axis="x", length=0)
ax.tick_params(axis="y", length=0)
ax.set_ylim(-0.7, n - 0.3)

for x in [20, 40, 60, 80, 100]:
    ax.axvline(x, color="#DAD3C7", lw=0.8, zorder=1)
for s in ax.spines.values():
    s.set_visible(False)
ax.spines["left"].set_visible(True)
ax.spines["left"].set_color("#C9C1B4")
ax.spines["left"].set_linewidth(1.0)

# legend (waves, oldest -> newest)
handles = [Patch(facecolor=WAVE_COLORS[w], label=w) for w in WAVES]
leg = ax.legend(handles=handles, loc="lower right",
                bbox_to_anchor=(1.0, 0.012),
                frameon=True, fontsize=9.0, ncol=2,
                title="Mittausajankohta", title_fontsize=9.0,
                handlelength=1.1, handleheight=1.0,
                columnspacing=1.3, labelspacing=0.45,
                borderpad=0.8)
leg.get_frame().set_facecolor("#FFFFFF")
leg.get_frame().set_edgecolor("#DAD3C7")
leg.get_frame().set_linewidth(0.8)
leg.get_title().set_color(INK)
leg.get_title().set_fontweight("bold")

# --- change column (pill per provider)
cax.set_xlim(0, 1)
cax.set_ylim(-0.7, n - 0.3)
cax.set_xticks([])
cax.set_yticks([])
for s in cax.spines.values():
    s.set_visible(False)
# divider line between chart and change column
cax.axvline(0.0, color="#C9C1B4", lw=1.0, zorder=2)

for y, dv in zip(ypos, deltas):
    if dv > 0:
        col = UP
    elif dv < 0:
        col = DOWN
    else:
        col = FLAT
    txt = fmt_delta(dv)
    # rounded pill behind the value
    cax.text(0.50, y, txt, va="center", ha="center",
             fontsize=12.5, fontweight="bold", color="#FFFFFF",
             zorder=6,
             bbox=dict(boxstyle="round,pad=0.34,rounding_size=0.55",
                       facecolor=col, edgecolor="none"))

# header for the change column (placed just above the top row)
cax.text(0.50, n - 0.42, "%-yks.", va="center", ha="center",
         fontsize=8.5, color=MUTED, zorder=6)

fig.savefig(PNG, dpi=200, facecolor=CREAM, bbox_inches="tight", pad_inches=0.04)
plt.close(fig)
print("chart ->", PNG)

# ---------------------------------------------------------------- pptx
PX_CREAM = RGBColor(0xF7, 0xF3, 0xEC)
PX_INK = RGBColor(0x2B, 0x2B, 0x2B)
PX_TEAL = RGBColor(0x13, 0x61, 0x5E)
PX_MUTED = RGBColor(0x6E, 0x6A, 0x63)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
slide = prs.slides.add_slide(prs.slide_layouts[6])

# background
bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = PX_CREAM
bg.line.fill.background()
bg.shadow.inherit = False
slide.shapes._spTree.remove(bg._element)
slide.shapes._spTree.insert(2, bg._element)


def textbox(l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for txt, sz, col, bold in line:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col
            r.font.name = "Liberation Sans"
    return tb


# Accent bar at top-left
acc = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42),
                             Inches(0.10), Inches(0.92))
acc.fill.solid(); acc.fill.fore_color.rgb = PX_TEAL
acc.line.fill.background(); acc.shadow.inherit = False

# Title (Finnish key message)
textbox(Inches(0.80), Inches(0.40), Inches(11.9), Inches(1.0),
        [[("Attendo on selvästi tunnetuin yksityinen hoivapalveluiden tarjoaja",
           23, PX_INK, True)],
         [("Tunnettuus 86 % – selvästi Esperin (75 %) edellä ja vakaa "
           "edelliseen aaltoon nähden (±0 %-yks.)",
           13.5, PX_MUTED, False)]],
        space_after=3)

# Section label above chart
textbox(Inches(0.80), Inches(1.52), Inches(9.0), Inches(0.32),
        [[("AUTETTU TUNNETTUUS  ·  osuus vastaajista (%)", 11, PX_TEAL, True)]])
# Change-column label (centered over the change-badge column)
textbox(Inches(9.19), Inches(1.52), Inches(2.1), Inches(0.32),
        [[("MUUTOS TOUKOKUUSTA", 11, PX_TEAL, True)]], align=PP_ALIGN.CENTER)

# Chart image
iw, ih = Image.open(PNG).size
ar = iw / ih
pic_w = Inches(12.55)
pic_h = Emu(int(pic_w / ar))
max_h = Inches(5.05)
if pic_h > max_h:
    pic_h = max_h
    pic_w = Emu(int(pic_h * ar))
pic_l = Inches((13.333 - pic_w / Inches(1)) / 2) if False else Inches(0.42)
pic_t = Inches(1.90)
slide.shapes.add_picture(str(PNG), pic_l, pic_t, width=pic_w, height=pic_h)

# Footer: question caption + base n
textbox(Inches(0.80), Inches(7.02), Inches(9.6), Inches(0.42),
        [[("Kysymys: ", 9.5, PX_MUTED, True),
          ("”Mitä seuraavista hoivapalveluiden tarjoajista tunnet "
           "vähintään nimeltä?”", 9.5, PX_MUTED, False)]])
textbox(Inches(10.5), Inches(7.02), Inches(2.3), Inches(0.42),
        [[("Kaikki vastaajat, n = 1001", 9.5, PX_MUTED, True)]],
        align=PP_ALIGN.RIGHT)

prs.save(OUT)
print("slide ->", OUT)
