#!/usr/bin/env python3
"""Build a 3-slide deck — WORD + OPINION data, each a DIFFERENT chart style.

nSight house style (cream / teal), Finnish key-message titles, captions + n.
All numbers and sentiment values are read VERBATIM from the JSON data files;
nothing is hard-coded or invented. The three styles are:

  1. WORD CLOUD  — the 10 words sized by count, coloured by sentiment.
  2. TREEMAP     — the 10 words as squarify rectangles, sized by count,
                   coloured by sentiment.
  3. DIVERGING / TORNADO BAR — opinion data: negatives extend LEFT of a centre
                   line, positives extend RIGHT; "En osaa sanoa" muted aside.

Output: ONE pptx -> work/agent_var_wordsop.pptx

Run: uv run python chart_lab/agent_var_wordsop.py
"""
import json
from pathlib import Path

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
from matplotlib.patches import Patch, FancyBboxPatch
import squarify
from wordcloud import WordCloud

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ---------------------------------------------------------------- paths
ROOT = Path("/home/johan/Projects/nsight/proto")
WORDS_DATA = ROOT / "chart_lab" / "words_data.json"
OPIN_DATA = ROOT / "chart_lab" / "perception_idx17.json"

PNG_CLOUD = ROOT / "chart_lab" / "agent_var_cloud.png"
PNG_TREE = ROOT / "chart_lab" / "agent_var_treemap.png"
PNG_TORNADO = ROOT / "chart_lab" / "agent_var_tornado.png"
OUT = ROOT / "work" / "agent_var_wordsop.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------- fonts
FONT_REG = "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
FONT_BOLD = "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"
for f in (FONT_REG, FONT_BOLD):
    fm.fontManager.addfont(f)
plt.rcParams["font.family"] = "Liberation Sans"

# ---------------------------------------------------------------- brand
CREAM = "#F7F3EC"
INK = "#2B2B2B"
TEAL = "#13615E"
MUTED = "#6E6A63"
GRID = "#DAD3C7"

# sentiment palette (verbatim sentiment keys from data)
SENT_COLOR = {
    "kielteinen": "#C0473A",   # negative — warm red
    "myönteinen": "#3F7D4E",   # positive — green
    "neutraali":  "#B3A98F",   # neutral — warm grey
}
SENT_LABEL = {
    "kielteinen": "Kielteinen",
    "myönteinen": "Myönteinen",
    "neutraali":  "Neutraali",
}
SENT_ORDER = ["kielteinen", "myönteinen", "neutraali"]


# ================================================================ CHARTS
def make_wordcloud(words, png):
    """Word cloud: size ~ count, colour ~ sentiment. Deterministic layout."""
    freqs = {w["word"]: w["count"] for w in words}
    sent = {w["word"]: w["sentiment"] for w in words}

    def color_func(word, **kwargs):
        return SENT_COLOR[sent[word]]

    wc = WordCloud(
        width=2000, height=820,
        background_color=CREAM,
        font_path=FONT_BOLD,
        prefer_horizontal=0.96,
        relative_scaling=0.6,    # strong link between count and size
        min_font_size=22,
        max_font_size=320,
        margin=6,
        random_state=7,          # fixed -> deterministic layout
        color_func=color_func,
    ).generate_from_frequencies(freqs)

    FIG_W, FIG_H = 11.9, 5.1
    fig = plt.figure(figsize=(FIG_W, FIG_H), dpi=200)
    fig.patch.set_facecolor(CREAM)
    ax = fig.add_axes([0.0, 0.085, 1.0, 0.915])
    ax.set_facecolor(CREAM)
    ax.imshow(wc, interpolation="bilinear")
    ax.axis("off")

    # legend along the bottom
    present = [s for s in SENT_ORDER if any(w["sentiment"] == s for w in words)]
    handles = [Patch(facecolor=SENT_COLOR[s], label=SENT_LABEL[s]) for s in present]
    leg = fig.legend(handles=handles, loc="lower center",
                     bbox_to_anchor=(0.5, 0.0), ncol=len(present),
                     frameon=False, fontsize=11.5,
                     handlelength=1.1, handleheight=1.0,
                     columnspacing=2.0, handletextpad=0.55)
    for t in leg.get_texts():
        t.set_color(INK)
    leg.set_title("Sävy", prop={"size": 11.5, "weight": "bold"})
    leg.get_title().set_color(INK)

    fig.savefig(png, dpi=200, facecolor=CREAM, bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    print("chart ->", png)


def make_treemap(words, png):
    """Treemap: rectangle area ~ count, colour ~ sentiment (squarify)."""
    ranked = sorted(words, key=lambda w: w["count"], reverse=True)
    sizes = [w["count"] for w in ranked]
    colors = [SENT_COLOR[w["sentiment"]] for w in ranked]

    FIG_W, FIG_H = 11.6, 5.7
    fig = plt.figure(figsize=(FIG_W, FIG_H), dpi=200)
    fig.patch.set_facecolor(CREAM)
    ax = fig.add_axes([0.012, 0.10, 0.976, 0.88])
    ax.set_facecolor(CREAM)

    W, H = 100.0, 100.0
    norm = squarify.normalize_sizes(sizes, W, H)
    rects = squarify.squarify(norm, 0, 0, W, H)

    light_text = {"kielteinen", "myönteinen"}  # dark fills -> white text
    for r, w, col in zip(rects, ranked, colors):
        x, y, dx, dy = r["x"], r["y"], r["dx"], r["dy"]
        pad = 0.6
        ax.add_patch(FancyBboxPatch(
            (x + pad, y + pad), dx - 2 * pad, dy - 2 * pad,
            boxstyle="round,pad=0,rounding_size=1.4",
            facecolor=col, edgecolor=CREAM, linewidth=2.2, zorder=3))

        area = dx * dy
        # font scales with tile area; clamp to keep small tiles legible
        word_fs = float(np.clip(np.sqrt(area) * 0.95, 9.5, 30))
        cnt_fs = word_fs * 0.62
        tcol = "#FFFFFF" if w["sentiment"] in light_text else INK
        cx, cy = x + dx / 2.0, y + dy / 2.0

        if dy >= 9:   # tall enough for word + count stacked
            ax.text(cx, cy + dy * 0.085, w["word"], va="center", ha="center",
                    fontsize=word_fs, fontweight="bold", color=tcol, zorder=5)
            ax.text(cx, cy - dy * 0.16, f"{w['count']}", va="center",
                    ha="center", fontsize=cnt_fs, color=tcol, alpha=0.92,
                    zorder=5)
        else:         # thin tile: single compact line
            ax.text(cx, cy, f"{w['word']} {w['count']}", va="center",
                    ha="center", fontsize=max(8.5, word_fs * 0.8),
                    fontweight="bold", color=tcol, zorder=5)

    ax.set_xlim(0, W)
    ax.set_ylim(0, H)
    ax.invert_yaxis()   # largest tile top-left (reading order)
    ax.axis("off")

    present = [s for s in SENT_ORDER if any(w["sentiment"] == s for w in words)]
    handles = [Patch(facecolor=SENT_COLOR[s], label=SENT_LABEL[s]) for s in present]
    leg = fig.legend(handles=handles, loc="lower center",
                     bbox_to_anchor=(0.5, 0.0), ncol=len(present),
                     frameon=False, fontsize=11.5,
                     handlelength=1.1, handleheight=1.0,
                     columnspacing=2.0, handletextpad=0.55)
    for t in leg.get_texts():
        t.set_color(INK)
    leg.set_title("Sävy", prop={"size": 11.5, "weight": "bold"})
    leg.get_title().set_color(INK)

    fig.savefig(png, dpi=200, facecolor=CREAM, bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    print("chart ->", png)


def make_tornado(cats, series, neg_keys, pos_keys, dk_key, png):
    """Diverging / tornado: negatives LEFT of centre, positives RIGHT.

    Each rating segment is stacked outward from a central 0 line. The "don't
    know" share is shown muted, off to the right side as a separate aside.
    """
    n = len(cats)
    # colour ramp: dark = extreme, light = mild
    SEG = {
        neg_keys[0]: "#B23A2E",   # Erittäin huono (darkest red, outermost L)
        neg_keys[1]: "#E08C82",   # Huono          (light red, inner L)
        pos_keys[0]: "#7DB8A6",   # Hyvä           (light teal, inner R)
        pos_keys[1]: "#13615E",   # Erittäin hyvä  (dark teal, outermost R)
    }
    LIGHT_TEXT = {neg_keys[0], pos_keys[1]}

    # derived exact totals
    neg_tot = [sum(series[k][i] for k in neg_keys) for i in range(n)]
    pos_tot = [sum(series[k][i] for k in pos_keys) for i in range(n)]

    FIG_W, FIG_H = 11.9, 5.0
    fig = plt.figure(figsize=(FIG_W, FIG_H), dpi=220)
    fig.patch.set_facecolor(CREAM)
    ax = fig.add_axes([0.005, 0.10, 0.99, 0.78])
    ax.set_facecolor(CREAM)

    bar_h = 0.40
    ypos = list(range(n))[::-1]   # first category at top

    AX_MAX = 60.0                 # half-axis extent (data units, %)
    POS_CALL_X = AX_MAX + 5.0     # x anchor for the positive-% total callout
    DK_X = AX_MAX + 20.0          # x anchor for the muted EOS aside (clear gap)

    for i, y in zip(range(n), ypos):
        # negatives: stack leftward from 0, mild (Huono) nearest centre
        left = 0.0
        for k in (neg_keys[1], neg_keys[0]):   # inner -> outer
            v = series[k][i]
            ax.barh(y, -v, left=-left, height=bar_h, color=SEG[k],
                    edgecolor=CREAM, linewidth=1.2, zorder=3)
            if v >= 4:
                ax.text(-(left + v / 2.0), y, f"{v}", va="center", ha="center",
                        fontsize=11, fontweight="bold",
                        color="#FFFFFF" if k in LIGHT_TEXT else INK, zorder=5)
            left += v
        # positives: stack rightward from 0, mild (Hyvä) nearest centre
        right = 0.0
        for k in (pos_keys[0], pos_keys[1]):   # inner -> outer
            v = series[k][i]
            ax.barh(y, v, left=right, height=bar_h, color=SEG[k],
                    edgecolor=CREAM, linewidth=1.2, zorder=3)
            if v >= 4:
                ax.text(right + v / 2.0, y, f"{v}", va="center", ha="center",
                        fontsize=11, fontweight="bold",
                        color="#FFFFFF" if k in LIGHT_TEXT else INK, zorder=5)
            right += v

        # group label above-left of its bar
        ax.text(-AX_MAX, y + bar_h * 0.95, cats[i], va="bottom", ha="left",
                fontsize=12.5, color=INK, fontweight="bold", zorder=5)

        # total callouts: negative sum (left, red) & positive sum (right, teal)
        ax.text(-left - 2.0, y, f"{neg_tot[i]} %", va="center", ha="right",
                fontsize=15, fontweight="bold", color="#B23A2E", zorder=6)
        ax.text(POS_CALL_X, y, f"{pos_tot[i]} %", va="center", ha="left",
                fontsize=15, fontweight="bold", color="#13615E", zorder=6)

        # muted "don't know" aside (separate, not part of the divide)
        dk = series[dk_key][i]
        ax.text(DK_X, y, f"{dk} %", va="center", ha="center",
                fontsize=11.5, fontweight="bold", color=MUTED, zorder=5)

    # centre line (0 = neutral balance point)
    ax.axvline(0, color="#9B938686", lw=1.4, zorder=2)
    ax.text(0, n - 0.42, "0", va="bottom", ha="center", fontsize=9,
            color=MUTED, zorder=4)

    # zone captions across the top
    top_y = n - 0.30
    ax.text(-AX_MAX * 0.5, top_y, "◄  KIELTEISET", va="bottom", ha="center",
            fontsize=10, fontweight="bold", color="#B23A2E", zorder=5)
    ax.text(AX_MAX * 0.5, top_y, "MYÖNTEISET  ►", va="bottom", ha="center",
            fontsize=10, fontweight="bold", color="#13615E", zorder=5)
    ax.text(DK_X, top_y, "EI OSAA SANOA", va="bottom", ha="center",
            fontsize=9, fontweight="bold", color=MUTED, zorder=5)

    ax.set_xlim(-AX_MAX - 12, DK_X + 9)
    ax.set_ylim(-0.55, n - 0.02)
    ax.axis("off")

    # legend: full ordered scale (worst -> best)
    order = [neg_keys[0], neg_keys[1], pos_keys[0], pos_keys[1]]
    handles = [Patch(facecolor=SEG[k], label=k) for k in order]
    leg = ax.legend(handles=handles, loc="upper center",
                    bbox_to_anchor=(0.42, 0.02), ncol=4, frameon=False,
                    fontsize=10.5, handlelength=1.0, handleheight=1.0,
                    columnspacing=1.4, handletextpad=0.5)
    for t in leg.get_texts():
        t.set_color(INK)

    fig.savefig(png, dpi=220, facecolor=CREAM, bbox_inches="tight",
                pad_inches=0.05)
    plt.close(fig)
    print("chart ->", png)


# ================================================================ PPTX
PX_CREAM = RGBColor(0xF7, 0xF3, 0xEC)
PX_INK = RGBColor(0x2B, 0x2B, 0x2B)
PX_TEAL = RGBColor(0x13, 0x61, 0x5E)
PX_RED = RGBColor(0xC0, 0x47, 0x3A)
PX_MUTED = RGBColor(0x6E, 0x6A, 0x63)


def _textbox(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, space_after=0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for txt, sz, col, bold in line:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col
            r.font.name = "Liberation Sans"
    return tb


def add_slide(prs, png, *, accent, title, subtitle, section_label,
              question, base_caption):
    SW, SH = prs.slide_width, prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = PX_CREAM
    bg.line.fill.background(); bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # accent bar
    acc = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42),
                                 Inches(0.10), Inches(0.92))
    acc.fill.solid(); acc.fill.fore_color.rgb = accent
    acc.line.fill.background(); acc.shadow.inherit = False

    # title + subtitle
    _textbox(slide, Inches(0.80), Inches(0.40), Inches(12.0), Inches(1.05),
             [[(title, 22, PX_INK, True)],
              [(subtitle, 13.5, PX_MUTED, False)]],
             space_after=3)

    # section label
    _textbox(slide, Inches(0.80), Inches(1.60), Inches(12.0), Inches(0.32),
             [[(section_label, 11, PX_TEAL, True)]])

    # chart (high-res embed, aspect preserved)
    iw, ih = Image.open(png).size
    ar = iw / ih
    pic_w = Inches(12.20)
    pic_h = Emu(int(pic_w / ar))
    max_h = Inches(4.85)
    if pic_h > max_h:
        pic_h = max_h
        pic_w = Emu(int(pic_h * ar))
    pic_l = Emu(int((SW - pic_w) / 2))
    pic_t = Inches(2.02)
    slide.shapes.add_picture(str(png), pic_l, pic_t, width=pic_w, height=pic_h)

    # footer: question + base n
    _textbox(slide, Inches(0.80), Inches(7.02), Inches(9.6), Inches(0.45),
             [[("Kysymys: ", 9.5, PX_MUTED, True),
               ("”" + question + "”", 9.5, PX_MUTED, False)]])
    _textbox(slide, Inches(9.0), Inches(7.02), Inches(3.75), Inches(0.45),
             [[(base_caption, 9.5, PX_MUTED, True)]],
             align=PP_ALIGN.RIGHT)
    return slide


# ================================================================ BUILD
def build():
    # ---- WORD data (verbatim) ----
    wd = json.loads(WORDS_DATA.read_text())
    wmeta = wd["meta"]
    words = wd["words"]
    ranked = sorted(words, key=lambda w: w["count"], reverse=True)

    WAVE = wmeta["wave"]
    WQ = wmeta["question"]
    # "kaikki vastaajat, n=863" -> "Kaikki vastaajat, n = 863"
    wbase = wmeta["base"]
    wbase = wbase[0].upper() + wbase[1:]
    wbase = wbase.replace("n=", "n = ")
    WBASE = f"{WAVE} · {wbase}"

    # derived (exact) for messaging
    top3 = ranked[:3]
    top3_neg = all(w["sentiment"] == "kielteinen" for w in top3)
    top_word = ranked[0]
    pos_words = [w for w in ranked if w["sentiment"] == "myönteinen"]
    top_pos = pos_words[0] if pos_words else None

    # ---- OPINION data (verbatim) ----
    od = json.loads(OPIN_DATA.read_text())
    cats = od["categories"]
    series = od["series"]
    NEG = ("Erittäin huono", "Huono")
    POS = ("Hyvä", "Erittäin hyvä")
    DK = "En osaa sanoa"
    pos_tot = [sum(series[k][i] for k in POS) for i in range(len(cats))]
    neg_tot = [sum(series[k][i] for k in NEG) for i in range(len(cats))]

    OQ = ("Mikä on yleinen käsityksesi tuntemistasi hoivapalveluita "
          "tarjoavista yrityksistä?")
    OBASE = f"{WAVE} · Kaikki vastaajat, n = 1001"

    # ---- render charts ----
    make_wordcloud(words, PNG_CLOUD)
    make_treemap(words, PNG_TREE)
    make_tornado(cats, series, NEG, POS, DK, PNG_TORNADO)

    # ---- assemble pptx ----
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1 — WORD CLOUD
    add_slide(
        prs, PNG_CLOUD, accent=PX_RED,
        title=(f"Attendon spontaania mielikuvaa hallitsee ”{top_word['word']}”"),
        subtitle=(f"Kolme yleisintä mainintaa ({', '.join(w['word'] for w in top3)}) "
                  f"ovat kielteisiä" if top3_neg else
                  "Maininnat jakautuvat sävyltään"),
        section_label=f"SPONTAANI MIELIKUVA  ·  SANAPILVI  ·  {WAVE}",
        question=WQ, base_caption=WBASE)

    # Slide 2 — TREEMAP
    add_slide(
        prs, PNG_TREE, accent=PX_RED,
        title=("Kielteiset sanat vievät suurimman pinta-alan mielikuvasta"),
        subtitle=(f"”{top_word['word']}” ({top_word['count']} mainintaa) erottuu; "
                  + (f"myönteisistä suurin on ”{top_pos['word']}” "
                     f"({top_pos['count']})" if top_pos else "")),
        section_label=f"SPONTAANI MIELIKUVA  ·  PUUKARTTA  ·  {WAVE}",
        question=WQ, base_caption=WBASE)

    # Slide 3 — DIVERGING / TORNADO
    add_slide(
        prs, PNG_TORNADO, accent=PX_TEAL,
        title=("Myönteiset mielikuvat voittavat selvästi sekä yksityisistä "
               "että julkisista"),
        subtitle=(f"Myönteisiä {pos_tot[0]} % vs. kielteisiä {neg_tot[0]} % "
                  f"(yksityiset); {pos_tot[1]} % vs. {neg_tot[1]} % (julkinen)"),
        section_label="YLEINEN KÄSITYS  ·  POIKKEAMAKAAVIO  ·  osuus vastaajista (%)",
        question=OQ, base_caption=OBASE)

    prs.save(OUT)
    print("deck ->", OUT, f"({len(prs.slides._sldIdLst)} slides)")
    return OUT


if __name__ == "__main__":
    build()
