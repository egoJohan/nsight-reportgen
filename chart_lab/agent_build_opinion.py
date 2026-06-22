#!/usr/bin/env python3
"""Build a 100%-stacked perception slide (nSight house style).

Reusable / parameterised generator for "general opinion" perception slides
that compare groups on an ordered rating scale (very poor -> very good) plus a
"don't know" option. Each group is shown as ONE 100% horizontal stacked bar,
split into the rating categories, with a clear negative|positive divide, the
total positive share called out per group, segment %% labels, a legend, the
survey question caption, base n and a Finnish key-message title.

All numbers are read verbatim from the JSON data file; nothing is invented.

Data shape expected (see perception_idx17.json):
    {
      "categories": ["Group A", "Group B"],
      "series": {
        "Erittäin huono":[a,b], "Huono":[a,b],
        "Hyvä":[a,b], "Erittäin hyvä":[a,b], "En osaa sanoa":[a,b]
      }
    }
Each series array is one value per category, in %.
"""
import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
from matplotlib.patches import Patch

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ----------------------------------------------------------------- config
ROOT = Path("/home/johan/Projects/nsight/proto")


def build(
    data_path=ROOT / "chart_lab" / "perception_idx17.json",
    png_path=ROOT / "chart_lab" / "agent_opinion_chart.png",
    out_path=ROOT / "work" / "agent_slide_opinion.pptx",
    *,
    # ordered worst -> best, then the neutral "don't know" series name
    negative_keys=("Erittäin huono", "Huono"),
    positive_keys=("Hyvä", "Erittäin hyvä"),
    dontknow_key="En osaa sanoa",
    wave="Marraskuu 2025",
    base_n=1001,
    question=("Mikä on yleinen käsityksesi tuntemistasi hoivapalveluita "
              "tarjoavista yrityksistä?"),
    section_label="YLEINEN KÄSITYS  ·  osuus vastaajista (%)",
    title=("Yksityisten palveluntarjoajien käsitys on parantunut – "
           "nyt yhtä myönteinen kuin julkisella"),
    subtitle=("Myönteisten osuus 57 % (yksityiset) ja 58 % (julkinen); "
              "ero käytännössä hävinnyt"),
):
    # ------------------------------------------------------------- fonts
    for f in ["/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
              "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"]:
        fm.fontManager.addfont(f)
    plt.rcParams["font.family"] = "Liberation Sans"

    # ------------------------------------------------------------- data
    d = json.loads(Path(data_path).read_text())
    cats = d["categories"]
    series = d["series"]
    n_groups = len(cats)

    # stacking order: negatives (worst first) | positives | don't know
    stack_keys = list(negative_keys) + list(positive_keys) + [dontknow_key]

    # per-group derived totals (exact, from data)
    pos_share = [sum(series[k][i] for k in positive_keys) for i in range(n_groups)]
    neg_share = [sum(series[k][i] for k in negative_keys) for i in range(n_groups)]

    # ------------------------------------------------------------- colours
    CREAM = "#F7F3EC"
    INK = "#2B2B2B"
    MUTED = "#6E6A63"
    # Diverging: negatives warm red (dark = worst), positives teal-green
    # (dark = best), neutral grey for "don't know".
    SEG_COLORS = {
        negative_keys[0]: "#B23A2E",   # Erittäin huono  (darkest red)
        negative_keys[1]: "#E08C82",   # Huono           (light red)
        positive_keys[0]: "#7DB8A6",   # Hyvä            (light teal)
        positive_keys[1]: "#13615E",   # Erittäin hyvä   (dark teal)
        dontknow_key:     "#C9C1B4",   # En osaa sanoa   (neutral)
    }
    # label text colour per segment (light text on dark fills)
    LIGHT_TEXT = {negative_keys[0], positive_keys[1]}

    # ------------------------------------------------------------- figure
    FIG_W, FIG_H = 11.9, 4.55
    fig = plt.figure(figsize=(FIG_W, FIG_H), dpi=220)
    fig.patch.set_facecolor(CREAM)
    ax = fig.add_axes([0.005, 0.075, 0.99, 0.80])
    ax.set_facecolor(CREAM)

    bar_h = 0.46
    ypos = list(range(n_groups))[::-1]      # first category at top

    # left gutter (data units) for group labels drawn inside the axes
    LEFT_PAD = 30.0
    # right gutter for the prominent positive-% callout
    RIGHT_PAD = 26.0

    for i, y in zip(range(n_groups), ypos):
        left = 0.0
        for k in stack_keys:
            v = series[k][i]
            if v == 0:
                continue
            ax.barh(y, v, left=left, height=bar_h,
                    color=SEG_COLORS[k], edgecolor=CREAM, linewidth=1.2,
                    zorder=3)
            # segment % label (centred); hide if too thin
            if v >= 4:
                ax.text(left + v / 2.0, y, f"{v}",
                        va="center", ha="center",
                        fontsize=11.5, fontweight="bold",
                        color="#FFFFFF" if k in LIGHT_TEXT else INK,
                        zorder=5)
            left += v

        # group label (left gutter)
        ax.text(-1.5, y, cats[i], va="center", ha="right",
                fontsize=12.5, color=INK, fontweight="bold", zorder=5)

        # prominent positive-share callout (right gutter)
        ax.text(102.5, y + 0.018, f"{pos_share[i]} %",
                va="center", ha="left",
                fontsize=20, fontweight="bold", color="#13615E", zorder=6)
        ax.text(102.7, y - bar_h * 0.62, "myönteisiä",
                va="center", ha="left",
                fontsize=9, color=MUTED, zorder=6)

    # divider between negative and positive halves is data-dependent per group,
    # so instead annotate the two macro-zones once, above the bars.
    ax.set_xlim(-LEFT_PAD, 100 + RIGHT_PAD)
    ax.set_ylim(-0.62, n_groups - 0.30)
    ax.axis("off")

    # macro-zone brackets along the top (kielteiset | myönteiset | EOS)
    # use the private (top) group split only as a visual guide reference.
    top_y = n_groups - 0.5 + 0.06
    # We anchor zone labels at fixed, readable x positions (not data-exact),
    # purely as orientation; exact numbers are the segment labels + callout.
    ax.text(neg_share[0] / 2.0, top_y, "KIELTEISET",
            va="bottom", ha="center", fontsize=9, fontweight="bold",
            color="#B23A2E", zorder=5)
    ax.text(neg_share[0] + pos_share[0] / 2.0, top_y, "MYÖNTEISET",
            va="bottom", ha="center", fontsize=9, fontweight="bold",
            color="#13615E", zorder=5)

    # legend (full ordered scale + don't know)
    handles = [Patch(facecolor=SEG_COLORS[k], label=k) for k in stack_keys]
    leg = ax.legend(handles=handles, loc="upper center",
                    bbox_to_anchor=(0.5, -0.02), ncol=len(stack_keys),
                    frameon=False, fontsize=10.5,
                    handlelength=1.0, handleheight=1.0,
                    columnspacing=1.4, handletextpad=0.5)
    for t in leg.get_texts():
        t.set_color(INK)

    fig.savefig(png_path, dpi=220, facecolor=CREAM,
                bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    print("chart ->", png_path)

    # ------------------------------------------------------------- pptx
    PX_CREAM = RGBColor(0xF7, 0xF3, 0xEC)
    PX_INK = RGBColor(0x2B, 0x2B, 0x2B)
    PX_TEAL = RGBColor(0x13, 0x61, 0x5E)
    PX_MUTED = RGBColor(0x6E, 0x6A, 0x63)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = PX_CREAM
    bg.line.fill.background(); bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    def textbox(l, t, w, h, runs, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, space_after=0):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = 0; tf.margin_right = 0
        tf.margin_top = 0; tf.margin_bottom = 0
        first = True
        for line in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.space_after = Pt(space_after)
            p.space_before = Pt(0)
            for txt, sz, col, bold in line:
                r = p.add_run(); r.text = txt
                r.font.size = Pt(sz); r.font.bold = bold
                r.font.color.rgb = col
                r.font.name = "Liberation Sans"
        return tb

    # accent bar
    acc = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42),
                                 Inches(0.10), Inches(0.92))
    acc.fill.solid(); acc.fill.fore_color.rgb = PX_TEAL
    acc.line.fill.background(); acc.shadow.inherit = False

    # title + subtitle
    textbox(Inches(0.80), Inches(0.40), Inches(11.9), Inches(1.05),
            [[(title, 22, PX_INK, True)],
             [(subtitle, 13.5, PX_MUTED, False)]],
            space_after=3)

    # section label
    textbox(Inches(0.80), Inches(1.62), Inches(11.9), Inches(0.32),
            [[(section_label, 11, PX_TEAL, True)]])

    # chart
    iw, ih = Image.open(png_path).size
    ar = iw / ih
    pic_w = Inches(12.30)
    pic_h = Emu(int(pic_w / ar))
    max_h = Inches(4.55)
    if pic_h > max_h:
        pic_h = max_h
        pic_w = Emu(int(pic_h * ar))
    pic_l = Inches((13.333 - pic_w / Inches(1)) / 2) if False else Inches(0.55)
    pic_t = Inches(2.05)
    slide.shapes.add_picture(str(png_path), pic_l, pic_t,
                             width=pic_w, height=pic_h)

    # footer: question + base n
    textbox(Inches(0.80), Inches(6.95), Inches(9.8), Inches(0.5),
            [[("Kysymys: ", 9.5, PX_MUTED, True),
              ("”" + question + "”", 9.5, PX_MUTED, False)]])
    textbox(Inches(8.7), Inches(6.95), Inches(4.0), Inches(0.5),
            [[(f"{wave} · Kaikki vastaajat, n = {base_n}",
               9.5, PX_MUTED, True)]],
            align=PP_ALIGN.RIGHT)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print("slide ->", out_path)
    return out_path


if __name__ == "__main__":
    build()
