#!/usr/bin/env python3
"""Build the respondent-demographics slide (nSight house style).

ONE 16:9 slide, THREE panels for all respondents (n=1001):
  - Gender  -> DONUT (two slices: Mies / Nainen)
  - Age     -> vertical COLUMNS (25-34 ... 65-74)
  - Region  -> horizontal BAR (Pääkaupunkiseutu ... Itä-Suomi)

Every value is %-labelled. Numbers are read verbatim from
demographics_overview.json — no value is typed or invented.

Run:  uv run python chart_lab/agent_content_demographics.py
"""
import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ---------------------------------------------------------------- paths
ROOT = Path("/home/johan/Projects/nsight/proto")
DATA = ROOT / "chart_lab" / "demographics_overview.json"
PNG = ROOT / "chart_lab" / "agent_demographics_chart.png"
OUT = ROOT / "work" / "agent_content_demographics.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------- fonts
for f in ["/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
          "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"]:
    fm.fontManager.addfont(f)
plt.rcParams["font.family"] = "Liberation Sans"

# ---------------------------------------------------------------- data
d = json.loads(DATA.read_text())
base = d["base"]
gender = d["gender"]                 # {"Mies":48,"Nainen":52}
age = d["age"]                       # band -> %
region = d["region"]                 # area -> %

# Region: order largest -> smallest by current value.
region_items = sorted(region.items(), key=lambda kv: kv[1], reverse=True)
# Short, readable region labels (display only; values untouched).
REGION_LABEL = {
    "Pääkaupunkiseudulla": "Pääkaupunkiseutu",
    "Muualla Etelä-Suomessa": "Muu Etelä-Suomi",
    "Länsi-Suomessa": "Länsi-Suomi",
    "Pohjois-Suomessa": "Pohjois-Suomi",
    "Itä-Suomessa": "Itä-Suomi",
}

# ---------------------------------------------------------------- colours
CREAM = "#F4EFE6"
INK = "#2B2B2B"
MUTED = "#6E685C"
TEAL = "#13615E"          # strong teal
TEAL2 = "#2F6F8F"         # secondary teal/blue
TEAL_LT = "#9CC6C4"       # light teal
GRID = "#DED6C7"
SPINE = "#CFC8BA"

# ---------------------------------------------------------------- figure
# Three panels side by side; generous gaps so nothing crowds.
fig = plt.figure(figsize=(13.0, 4.55), dpi=220)
fig.patch.set_facecolor(CREAM)

# Panel axes positions [left, bottom, width, height] in figure coords.
ax_g = fig.add_axes([0.045, 0.090, 0.230, 0.720])   # donut
ax_a = fig.add_axes([0.370, 0.150, 0.255, 0.680])   # columns
ax_r = fig.add_axes([0.760, 0.140, 0.225, 0.700])   # horizontal bar
for ax in (ax_g, ax_a, ax_r):
    ax.set_facecolor(CREAM)

# ---- Panel 1: GENDER donut --------------------------------------------
g_labels = ["Mies", "Nainen"]
g_vals = [gender["Mies"], gender["Nainen"]]
g_cols = [TEAL2, TEAL]
wedges, _ = ax_g.pie(
    g_vals, colors=g_cols, startangle=90, counterclock=False,
    wedgeprops=dict(width=0.40, edgecolor=CREAM, linewidth=2.5),
)
ax_g.set_aspect("equal")
# Direct % labels at the middle of each wedge ring.
import numpy as np
cum = 0.0
total_g = sum(g_vals)
for v, lab, col in zip(g_vals, g_labels, g_cols):
    mid = 90 - (cum + v / 2) / total_g * 360
    cum += v
    rad = np.deg2rad(mid)
    x = np.cos(rad) * 0.80
    y = np.sin(rad) * 0.80
    ax_g.text(x, y, f"{v} %", ha="center", va="center",
              fontsize=14, fontweight="bold", color="#FFFFFF")
# Category labels outside the ring, on the side each slice occupies.
# startangle=90, clockwise: Mies = right half, Nainen = left half.
ax_g.text(1.36, 0.0, g_labels[0], ha="left", va="center",
          fontsize=12.5, fontweight="bold", color=INK)
ax_g.text(-1.36, 0.0, g_labels[1], ha="right", va="center",
          fontsize=12.5, fontweight="bold", color=INK)
ax_g.set_xlim(-1.95, 1.95)
ax_g.set_ylim(-1.45, 1.45)

# ---- Panel 2: AGE columns ---------------------------------------------
a_labels = list(age.keys())          # survey order (ascending bands)
a_vals = [age[k] for k in a_labels]
xpos = list(range(len(a_labels)))
ax_a.bar(xpos, a_vals, width=0.66, color=TEAL, edgecolor="none", zorder=3)
for x, v in zip(xpos, a_vals):
    ax_a.text(x, v + 0.8, f"{v} %", ha="center", va="bottom",
              fontsize=11.5, fontweight="bold", color=INK, zorder=5)
ax_a.set_xticks(xpos)
ax_a.set_xticklabels(a_labels, fontsize=11, color=INK)
ax_a.set_ylim(0, max(a_vals) + 7)
ax_a.set_yticks([])
ax_a.tick_params(axis="x", length=0)
for s in ax_a.spines.values():
    s.set_visible(False)
ax_a.spines["bottom"].set_visible(True)
ax_a.spines["bottom"].set_color(SPINE)
ax_a.spines["bottom"].set_linewidth(1.1)

# ---- Panel 3: REGION horizontal bar -----------------------------------
r_labels = [REGION_LABEL[k] for k, _ in region_items]
r_vals = [v for _, v in region_items]
ypos = list(range(len(r_labels)))[::-1]   # largest at top
ax_r.barh(ypos, r_vals, height=0.62, color=TEAL2, edgecolor="none", zorder=3)
for y, v in zip(ypos, r_vals):
    ax_r.text(v + 0.7, y, f"{v} %", ha="left", va="center",
              fontsize=11, fontweight="bold", color=INK, zorder=5)
ax_r.set_yticks(ypos)
ax_r.set_yticklabels(r_labels, fontsize=10.8, color=INK)
ax_r.set_xlim(0, max(r_vals) + 9)
ax_r.set_xticks([])
ax_r.tick_params(axis="y", length=0)
for s in ax_r.spines.values():
    s.set_visible(False)
ax_r.spines["left"].set_visible(True)
ax_r.spines["left"].set_color(SPINE)
ax_r.spines["left"].set_linewidth(1.1)
ax_r.set_ylim(-0.6, len(r_labels) - 0.4)

# ---- Panel sub-titles --------------------------------------------------
def panel_title(ax, txt):
    ax.text(0.5, 1.14, txt, transform=ax.transAxes, ha="center", va="bottom",
            fontsize=12.5, fontweight="bold", color=TEAL)

# Gender title placed via figure text (donut ax has equal aspect).
fig.text(0.155, 0.90, "SUKUPUOLI", ha="center", va="bottom",
         fontsize=12.5, fontweight="bold", color=TEAL)
fig.text(0.4975, 0.90, "IKÄRYHMÄ", ha="center", va="bottom",
         fontsize=12.5, fontweight="bold", color=TEAL)
fig.text(0.8725, 0.90, "ASUINALUE", ha="center", va="bottom",
         fontsize=12.5, fontweight="bold", color=TEAL)

# Thin vertical separators between panels.
for xline in (0.325, 0.700):
    fig.add_artist(plt.Line2D([xline, xline], [0.08, 0.86],
                   color=GRID, lw=1.0, transform=fig.transFigure))

fig.savefig(PNG, dpi=220, facecolor=CREAM, bbox_inches="tight", pad_inches=0.06)
plt.close(fig)
print("chart ->", PNG)

# ---------------------------------------------------------------- pptx
PX_CREAM = RGBColor(0xF4, 0xEF, 0xE6)
PX_INK = RGBColor(0x2B, 0x2B, 0x2B)
PX_TEAL = RGBColor(0x13, 0x61, 0x5E)
PX_MUTED = RGBColor(0x6E, 0x68, 0x5C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
slide = prs.slides.add_slide(prs.slide_layouts[6])

# background
bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = PX_CREAM
bg.line.fill.background(); bg.shadow.inherit = False
slide.shapes._spTree.remove(bg._element)
slide.shapes._spTree.insert(2, bg._element)


def textbox(l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=0, line_spacing=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if line_spacing:
            p.line_spacing = line_spacing
        for txt, sz, col, bold in line:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col
            r.font.name = "Liberation Sans"
    return tb


# Accent bar top-left
acc = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42),
                             Inches(0.10), Inches(0.92))
acc.fill.solid(); acc.fill.fore_color.rgb = PX_TEAL
acc.line.fill.background(); acc.shadow.inherit = False

# Title (Finnish key message) + sub
textbox(Inches(0.80), Inches(0.40), Inches(12.0), Inches(1.0),
        [[("Vastaajat jakautuvat tasaisesti sukupuolen, iän ja alueen mukaan",
           23, PX_INK, True)],
         [("Otos edustaa väestöä kattavasti: sukupuolijakauma on lähes "
           "tasan ja ikäryhmät jakautuvat tasaisesti",
           13.5, PX_MUTED, False)]],
        space_after=3)

# Section label above panels
textbox(Inches(0.80), Inches(1.55), Inches(12.0), Inches(0.32),
        [[("VASTAAJAPROFIILI  ·  osuus vastaajista (%)", 11, PX_TEAL, True)]])

# Chart image (three panels)
iw, ih = Image.open(PNG).size
ar = iw / ih
pic_w = Inches(12.35)
pic_h = Emu(int(pic_w / ar))
max_h = Inches(4.75)
if pic_h > max_h:
    pic_h = max_h
    pic_w = Emu(int(pic_h * ar))
pic_l = Emu(int((SW - pic_w) / 2))
pic_t = Inches(2.10)
slide.shapes.add_picture(str(PNG), pic_l, pic_t, width=pic_w, height=pic_h)

# Footer: question caption + base n
textbox(Inches(0.80), Inches(7.00), Inches(9.4), Inches(0.45),
        [[("Taustamuuttujat: ", 9.5, PX_MUTED, True),
          ("sukupuoli, ikäryhmä ja asuinalue", 9.5, PX_MUTED, False)]])
textbox(Inches(9.5), Inches(7.00), Inches(3.3), Inches(0.45),
        [[(base, 9.5, PX_MUTED, True)]],
        align=PP_ALIGN.RIGHT)

prs.save(OUT)
print("slide ->", OUT)
