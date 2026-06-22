#!/usr/bin/env python3
"""Build 6 aided-awareness slides, each a DIFFERENT chart style (nSight house style).

Styles:
  1. Slope chart        Toukokuu 2024 -> Marraskuu 2025 per provider
  2. Dumbbell           first measured wave vs current wave per provider
  3. Lollipop           current-wave ranking (dot + stem), sorted
  4. Bump / rank        each provider's RANK across the 4 waves
  5. Bar + avg line     current vs across-provider average (above/below)
  6. Vertical columns   current-wave awareness as columns

All numbers EXACT from original_data.json. Charts rendered with matplotlib at
high DPI, embedded into one pptx via python-pptx.
"""
import json
from pathlib import Path

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
from matplotlib.lines import Line2D

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ---------------------------------------------------------------- paths
ROOT = Path("/home/johan/Projects/nsight/proto")
DATA = ROOT / "chart_lab" / "original_data.json"
IMGDIR = ROOT / "chart_lab"
OUT = ROOT / "work" / "agent_var_awareness.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------- fonts
for f in ["/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
          "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"]:
    fm.fontManager.addfont(f)
plt.rcParams["font.family"] = "Liberation Sans"

# ---------------------------------------------------------------- data
d = json.loads(DATA.read_text())
cats_all = d["categories"]
series = d["series"]

WAVES = ["Toukokuu 2024", "Marraskuu 2024", "Toukokuu 2025", "Marraskuu 2025"]
WAVE_SHORT = {"Toukokuu 2024": "5/24", "Marraskuu 2024": "11/24",
              "Toukokuu 2025": "5/25", "Marraskuu 2025": "11/25"}
CURRENT = "Marraskuu 2025"
FIRST = "Toukokuu 2024"

# Exclude "En mitään näistä" (not a provider).
prov_idx = [i for i, c in enumerate(cats_all) if c != "En mitään näistä"]
PROV = [cats_all[i] for i in prov_idx]

def val(wave, name):
    i = cats_all.index(name)
    return series[wave][i]

# ---------------------------------------------------------------- colours
CREAM   = "#F4EFE6"
CREAM2  = "#F7F3EC"
INK     = "#2B2B2B"
MUTED   = "#6E6A63"
TEAL    = "#13615E"
TEAL_MD = "#5E9C9A"
TEAL_LT = "#9CC6C4"
GRID    = "#DAD3C7"
RISE    = "#13615E"   # teal = up / good
FALL    = "#B5602F"   # warm clay = down
GOLD    = "#C99A3B"

DPI = 220

def style_ax(ax):
    ax.set_facecolor(CREAM)
    for s in ax.spines.values():
        s.set_visible(False)
    ax.tick_params(length=0)

# ================================================================
# 1. SLOPE CHART  — Toukokuu 2024 -> Marraskuu 2025
# ================================================================
def chart_slope():
    # providers measured in both waves (Validia has null in 5/24 -> skip)
    rows = [(p, val(FIRST, p), val(CURRENT, p)) for p in PROV
            if val(FIRST, p) is not None]
    fig = plt.figure(figsize=(11.6, 6.0), dpi=DPI)
    fig.patch.set_facecolor(CREAM)
    ax = fig.add_axes([0.16, 0.085, 0.70, 0.85])
    style_ax(ax)
    xL, xR = 0.0, 1.0
    for p, a, b in rows:
        col = RISE if b >= a else FALL
        lw = 2.6 if abs(b - a) >= 2 else 1.6
        ax.plot([xL, xR], [a, b], color=col, lw=lw, zorder=3,
                solid_capstyle="round")
        ax.scatter([xL, xR], [a, b], s=46, color=col, zorder=4,
                   edgecolor=CREAM, linewidth=1.0)
    # de-collide left/right labels
    def place_labels(side, items):
        items = sorted(items, key=lambda t: t[1])
        ys = [t[1] for t in items]
        minsep = 2.7
        for i in range(1, len(ys)):
            if ys[i] - ys[i-1] < minsep:
                ys[i] = ys[i-1] + minsep
        # shift down if overflow
        for p, _, _ in []:
            pass
        return [(items[i][0], items[i][1], ys[i], items[i][2]) for i in range(len(items))]
    left_items  = [(p, a, RISE if b >= a else FALL) for p, a, b in rows]
    right_items = [(p, b, RISE if b >= a else FALL) for p, a, b in rows]
    for name, yreal, ydraw, col in place_labels("L", left_items):
        ax.annotate(f"{name}  {int(yreal)}", xy=(xL - 0.02, yreal),
                    xytext=(xL - 0.04, ydraw), ha="right", va="center",
                    fontsize=10.5, color=INK,
                    fontweight="bold" if name == "Attendo" else "normal")
    for name, yreal, ydraw, col in place_labels("R", right_items):
        ax.annotate(f"{int(yreal)}  {name}", xy=(xR + 0.02, yreal),
                    xytext=(xR + 0.04, ydraw), ha="left", va="center",
                    fontsize=10.5, color=col,
                    fontweight="bold")
    ax.set_xlim(-0.42, 1.42)
    ymin = min(min(a, b) for _, a, b in rows) - 6
    ymax = max(max(a, b) for _, a, b in rows) + 6
    ax.set_ylim(ymin, ymax)
    ax.set_xticks([])
    ax.set_yticks([])
    # wave headers
    ax.text(xL, ymax + 1.5, "Toukokuu 2024", ha="center", va="bottom",
            fontsize=12, fontweight="bold", color=MUTED)
    ax.text(xR, ymax + 1.5, "Marraskuu 2025", ha="center", va="bottom",
            fontsize=12, fontweight="bold", color=TEAL)
    # vertical anchor lines
    ax.axvline(xL, color=GRID, lw=1.0, zorder=1)
    ax.axvline(xR, color=GRID, lw=1.0, zorder=1)
    # legend
    leg = [Line2D([0],[0], color=RISE, lw=3, label="Nousi / ennallaan"),
           Line2D([0],[0], color=FALL, lw=3, label="Laski")]
    lg = ax.legend(handles=leg, loc="lower center", bbox_to_anchor=(0.5, -0.04),
                   ncol=2, frameon=False, fontsize=10, handlelength=1.6)
    for t in lg.get_texts():
        t.set_color(MUTED)
    p = IMGDIR / "var_slope.png"
    fig.savefig(p, dpi=DPI, facecolor=CREAM, bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    return p

# ================================================================
# 2. DUMBBELL — first measured wave vs current wave
# ================================================================
def chart_dumbbell():
    # first MEASURED wave per provider (handle Validia null in 5/24)
    rows = []
    for p in PROV:
        first_w = next(w for w in WAVES if val(w, p) is not None)
        a = val(first_w, p)
        b = val(CURRENT, p)
        rows.append((p, first_w, a, b))
    rows.sort(key=lambda t: t[3], reverse=True)  # by current desc
    fig = plt.figure(figsize=(11.6, 6.0), dpi=DPI)
    fig.patch.set_facecolor(CREAM)
    ax = fig.add_axes([0.155, 0.10, 0.80, 0.86])
    style_ax(ax)
    n = len(rows)
    ypos = list(range(n))[::-1]
    for y, (p, fw, a, b) in zip(ypos, rows):
        col = RISE if b >= a else FALL
        ax.plot([a, b], [y, y], color=col, lw=3.0, zorder=2,
                solid_capstyle="round")
        # start dot (muted/open), end dot (solid)
        ax.scatter(a, y, s=120, facecolor=CREAM, edgecolor=MUTED,
                   linewidth=1.8, zorder=3)
        ax.scatter(b, y, s=150, facecolor=col, edgecolor=CREAM,
                   linewidth=1.2, zorder=4)
        # value labels: start label outside the LEFT dot, end label outside
        # the RIGHT dot, so they never collide even when a≈b.
        lo, hi = min(a, b), max(a, b)
        # start value (muted) goes to the left of the left-most point
        ax.text(lo - 1.6, y, f"{int(a)}", ha="right", va="center",
                fontsize=9.5, color=MUTED)
        # current value (bold, coloured) goes to the right of the right-most point
        ax.text(hi + 1.6, y, f"{int(b)}", ha="left", va="center",
                fontsize=10, color=col, fontweight="bold")
        # delta tag above the connector
        delta = b - a
        sign = "+" if delta > 0 else ("±" if delta == 0 else "−")
        ax.text((a + b) / 2, y + 0.34, f"{sign}{abs(int(delta))}",
                ha="center", va="bottom", fontsize=8.5, color=col,
                fontweight="bold")
    ax.set_yticks(ypos)
    ax.set_yticklabels([p for p, *_ in rows], fontsize=11.5, color=INK)
    for lbl in ax.get_yticklabels():
        if lbl.get_text() == "Attendo":
            lbl.set_fontweight("bold")
    ax.set_xlim(0, 100)
    ax.set_xticks([0, 20, 40, 60, 80, 100])
    ax.set_xticklabels(["0", "20", "40", "60", "80", "100 %"],
                       fontsize=9.5, color=MUTED)
    for x in [20, 40, 60, 80, 100]:
        ax.axvline(x, color=GRID, lw=0.8, zorder=1)
    ax.set_ylim(-0.7, n - 0.3)
    leg = [Line2D([0],[0], marker="o", color="none", markerfacecolor=CREAM,
                  markeredgecolor=MUTED, markersize=10, label="Ensimmäinen mittaus"),
           Line2D([0],[0], marker="o", color="none", markerfacecolor=TEAL,
                  markeredgecolor=CREAM, markersize=11, label="Marraskuu 2025")]
    lg = ax.legend(handles=leg, loc="lower right", bbox_to_anchor=(1.0, 0.0),
                   frameon=True, fontsize=9.5, handletextpad=0.4)
    lg.get_frame().set_facecolor("#FFFFFF")
    lg.get_frame().set_edgecolor(GRID)
    for t in lg.get_texts():
        t.set_color(INK)
    p = IMGDIR / "var_dumbbell.png"
    fig.savefig(p, dpi=DPI, facecolor=CREAM, bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    return p

# ================================================================
# 3. LOLLIPOP — current-wave ranking, sorted
# ================================================================
def chart_lollipop():
    rows = [(p, val(CURRENT, p)) for p in PROV]
    rows.sort(key=lambda t: t[1])  # ascending so biggest on top
    fig = plt.figure(figsize=(11.6, 6.0), dpi=DPI)
    fig.patch.set_facecolor(CREAM)
    ax = fig.add_axes([0.155, 0.10, 0.80, 0.86])
    style_ax(ax)
    n = len(rows)
    ypos = list(range(n))
    for y, (p, v) in zip(ypos, rows):
        is_top = (p == "Attendo")
        col = TEAL if v >= 40 else TEAL_MD
        ax.plot([0, v], [y, y], color=col, lw=2.6, zorder=2,
                solid_capstyle="round")
        ax.scatter(v, y, s=320 if is_top else 240, color=col, zorder=3,
                   edgecolor=CREAM, linewidth=1.5)
        ax.text(v, y, f"{int(v)}", ha="center", va="center", zorder=4,
                fontsize=8.5 if not is_top else 9.5, color="white",
                fontweight="bold")
    ax.set_yticks(ypos)
    ax.set_yticklabels([p for p, _ in rows], fontsize=11.5, color=INK)
    for lbl in ax.get_yticklabels():
        if lbl.get_text() == "Attendo":
            lbl.set_fontweight("bold")
    ax.set_xlim(0, 100)
    ax.set_xticks([0, 20, 40, 60, 80, 100])
    ax.set_xticklabels(["0", "20", "40", "60", "80", "100 %"],
                       fontsize=9.5, color=MUTED)
    for x in [20, 40, 60, 80, 100]:
        ax.axvline(x, color=GRID, lw=0.8, zorder=1)
    ax.set_ylim(-0.7, n - 0.3)
    p = IMGDIR / "var_lollipop.png"
    fig.savefig(p, dpi=DPI, facecolor=CREAM, bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    return p

# ================================================================
# 4. BUMP / RANK CHART — rank across 4 waves
# ================================================================
def chart_bump():
    # rank within each wave (1 = best known). Validia null in 5/24 -> no rank that wave.
    ranks = {p: {} for p in PROV}
    for w in WAVES:
        scored = [(p, val(w, p)) for p in PROV if val(w, p) is not None]
        scored.sort(key=lambda t: t[1], reverse=True)
        for r, (p, v) in enumerate(scored, start=1):
            ranks[p][w] = r
    n = len(PROV)
    fig = plt.figure(figsize=(11.6, 6.0), dpi=DPI)
    fig.patch.set_facecolor(CREAM)
    ax = fig.add_axes([0.13, 0.10, 0.74, 0.80])
    style_ax(ax)
    xs = list(range(len(WAVES)))
    # color by current rank: top emphasized teal, rest graded
    cur_order = sorted(PROV, key=lambda p: ranks[p][CURRENT])
    palette = [TEAL, "#2E7C77", TEAL_MD, "#7FB1AE", TEAL_LT, "#B7D3CF",
               GOLD, "#B5602F"]
    pcolor = {p: palette[i % len(palette)] for i, p in enumerate(cur_order)}
    for p in PROV:
        pts = [(x, ranks[p].get(w)) for x, w in zip(xs, WAVES)]
        pts = [(x, r) for x, r in pts if r is not None]
        xx = [x for x, _ in pts]; yy = [r for _, r in pts]
        col = pcolor[p]
        lw = 3.2 if p == "Attendo" else 2.2
        ax.plot(xx, yy, color=col, lw=lw, zorder=3, solid_capstyle="round")
        ax.scatter(xx, yy, s=210, color=col, zorder=4,
                   edgecolor=CREAM, linewidth=1.5)
        for x, r in pts:
            ax.text(x, r, f"{r}", ha="center", va="center", color="white",
                    fontsize=8.5, fontweight="bold", zorder=5)
        # label provider at right end
        lastx, lastr = pts[-1]
        ax.text(lastx + 0.10, lastr, p, ha="left", va="center",
                fontsize=10.5, color=col,
                fontweight="bold" if p == "Attendo" else "normal")
    ax.set_xlim(-0.35, len(WAVES) - 1 + 1.05)
    ax.set_ylim(n + 0.5, 0.4)  # invert: rank 1 at top
    ax.set_xticks(xs)
    ax.set_xticklabels([w for w in WAVES], fontsize=11, color=MUTED)
    ax.set_yticks(range(1, n + 1))
    ax.set_yticklabels([f"{i}." for i in range(1, n + 1)], fontsize=10, color=MUTED)
    ax.text(-0.30, 0.0, "Sija", ha="center", va="bottom", fontsize=9.5,
            color=MUTED, fontweight="bold")
    for x in xs:
        ax.axvline(x, color=GRID, lw=0.8, zorder=1)
    p = IMGDIR / "var_bump.png"
    fig.savefig(p, dpi=DPI, facecolor=CREAM, bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    return p

# ================================================================
# 5. BAR WITH MARKET-AVERAGE REFERENCE LINE
# ================================================================
def chart_avgline():
    rows = [(p, val(CURRENT, p)) for p in PROV]
    rows.sort(key=lambda t: t[1], reverse=True)
    avg = sum(v for _, v in rows) / len(rows)
    fig = plt.figure(figsize=(11.6, 6.0), dpi=DPI)
    fig.patch.set_facecolor(CREAM)
    ax = fig.add_axes([0.155, 0.10, 0.80, 0.86])
    style_ax(ax)
    n = len(rows)
    ypos = list(range(n))[::-1]
    for y, (p, v) in zip(ypos, rows):
        above = v >= avg
        col = TEAL if above else "#C9BFAE"
        ax.barh(y, v, height=0.62, color=col, edgecolor="none", zorder=3)
        ax.text(v + 1.2, y, f"{int(v)}", ha="left", va="center",
                fontsize=10, color=INK if above else MUTED,
                fontweight="bold" if above else "normal")
    ax.set_yticks(ypos)
    ax.set_yticklabels([p for p, _ in rows], fontsize=11.5, color=INK)
    for lbl in ax.get_yticklabels():
        if lbl.get_text() == "Attendo":
            lbl.set_fontweight("bold")
    ax.set_xlim(0, 100)
    ax.set_xticks([0, 20, 40, 60, 80, 100])
    ax.set_xticklabels(["0", "20", "40", "60", "80", "100 %"],
                       fontsize=9.5, color=MUTED)
    ax.set_ylim(-0.7, n - 0.3)
    # average reference line
    ax.axvline(avg, color=FALL, lw=2.0, ls=(0, (5, 3)), zorder=5)
    ax.text(avg, n - 0.35, f"Keskiarvo {avg:.0f} %", color=FALL,
            fontsize=10.5, fontweight="bold", ha="center", va="bottom")
    p = IMGDIR / "var_avgline.png"
    fig.savefig(p, dpi=DPI, facecolor=CREAM, bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    return p, avg

# ================================================================
# 6. VERTICAL COLUMN CHART — current wave
# ================================================================
def chart_columns():
    rows = [(p, val(CURRENT, p)) for p in PROV]
    rows.sort(key=lambda t: t[1], reverse=True)
    fig = plt.figure(figsize=(11.6, 5.7), dpi=DPI)
    fig.patch.set_facecolor(CREAM)
    ax = fig.add_axes([0.07, 0.16, 0.90, 0.78])
    style_ax(ax)
    n = len(rows)
    xpos = list(range(n))
    for x, (p, v) in zip(xpos, rows):
        col = TEAL if p == "Attendo" else (TEAL_MD if v >= 40 else TEAL_LT)
        ax.bar(x, v, width=0.64, color=col, edgecolor="none", zorder=3)
        ax.text(x, v + 1.6, f"{int(v)}", ha="center", va="bottom",
                fontsize=11, color=INK,
                fontweight="bold" if p == "Attendo" else "normal")
    ax.set_xticks(xpos)
    ax.set_xticklabels([p for p, _ in rows], fontsize=11, color=INK)
    for lbl in ax.get_xticklabels():
        if lbl.get_text() == "Attendo":
            lbl.set_fontweight("bold")
    ax.set_ylim(0, 100)
    ax.set_yticks([0, 20, 40, 60, 80, 100])
    ax.set_yticklabels(["0", "20", "40", "60", "80", "100 %"],
                       fontsize=9.5, color=MUTED)
    for yv in [20, 40, 60, 80, 100]:
        ax.axhline(yv, color=GRID, lw=0.8, zorder=1)
    ax.set_xlim(-0.7, n - 0.3)
    p = IMGDIR / "var_columns.png"
    fig.savefig(p, dpi=DPI, facecolor=CREAM, bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    return p

# ================================================================
# PPTX assembly
# ================================================================
PX_CREAM = RGBColor(0xF4, 0xEF, 0xE6)
PX_INK   = RGBColor(0x2B, 0x2B, 0x2B)
PX_TEAL  = RGBColor(0x13, 0x61, 0x5E)
PX_MUTED = RGBColor(0x6E, 0x6A, 0x63)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

def add_slide(png, title, subtitle, section, footer_q, footer_n,
              chart_top=1.92, chart_h=5.0, chart_w=12.15):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = PX_CREAM
    bg.line.fill.background(); bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    def textbox(l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                space_after=0):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = 0; tf.margin_right = 0
        tf.margin_top = 0; tf.margin_bottom = 0
        first = True
        for line in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.space_after = Pt(space_after); p.space_before = Pt(0)
            for txt, sz, col, bold in line:
                r = p.add_run(); r.text = txt
                r.font.size = Pt(sz); r.font.bold = bold
                r.font.color.rgb = col; r.font.name = "Liberation Sans"
        return tb

    acc = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42),
                                 Inches(0.10), Inches(0.92))
    acc.fill.solid(); acc.fill.fore_color.rgb = PX_TEAL
    acc.line.fill.background(); acc.shadow.inherit = False

    textbox(Inches(0.80), Inches(0.40), Inches(11.9), Inches(1.0),
            [[(title, 22, PX_INK, True)],
             [(subtitle, 13.5, PX_MUTED, False)]], space_after=3)
    textbox(Inches(0.80), Inches(1.54), Inches(11.9), Inches(0.32),
            [[(section, 11, PX_TEAL, True)]])

    iw, ih = Image.open(png).size
    ar = iw / ih
    pic_w = Inches(chart_w)
    pic_h = Emu(int(pic_w / ar))
    max_h = Inches(chart_h)
    if pic_h > max_h:
        pic_h = max_h
        pic_w = Emu(int(pic_h * ar))
    pic_l = int((SW - pic_w) / 2)
    slide.shapes.add_picture(str(png), pic_l, Inches(chart_top),
                             width=pic_w, height=pic_h)

    textbox(Inches(0.80), Inches(7.04), Inches(9.6), Inches(0.42),
            [[("Kysymys: ", 9, PX_MUTED, True),
              (footer_q, 9, PX_MUTED, False)]])
    textbox(Inches(10.4), Inches(7.04), Inches(2.4), Inches(0.42),
            [[(footer_n, 9, PX_MUTED, True)]], align=PP_ALIGN.RIGHT)

QUESTION = ("”Mitä seuraavista hoivapalveluiden tarjoajista tunnet "
            "vähintään nimeltä?”")
BASE = "Kaikki vastaajat, n = 1001"

# ---- build charts
p1 = chart_slope()
p2 = chart_dumbbell()
p3 = chart_lollipop()
p4 = chart_bump()
p5, avg = chart_avgline()
p6 = chart_columns()

# ---- slide 1: slope
add_slide(p1,
    "Attendo ja Esperi vahvistivat tunnettuuttaan – haastajat polkivat paikallaan",
    "Toukokuusta 2024 marraskuuhun 2025: kärki nousi, kun pienemmät toimijat jäivät ennalleen tai laskivat",
    "AUTETTU TUNNETTUUS  ·  muutos kahden mittauksen välillä (%)",
    QUESTION, BASE, chart_top=1.94, chart_h=4.95)

# ---- slide 2: dumbbell
add_slide(p2,
    "Tunnettuuden kasvu painottuu jo ennestään tunnetuimpiin toimijoihin",
    "Ensimmäisestä mittauksesta nykyhetkeen: Attendo +2 ja Esperi +3, häntäpäässä liike pientä",
    "AUTETTU TUNNETTUUS  ·  ensimmäinen mittaus vs. marraskuu 2025 (%)",
    QUESTION, BASE, chart_top=1.96, chart_h=4.95)

# ---- slide 3: lollipop
add_slide(p3,
    "Hoivamarkkinan tunnettuus jakautuu jyrkästi kärkeen ja häntään",
    "Marraskuu 2025: Attendo (86 %) ja Esperi (75 %) omassa sarjassaan, loput alle 45 %:n",
    "AUTETTU TUNNETTUUS  ·  marraskuu 2025, järjestys (%)",
    QUESTION, BASE, chart_top=1.96, chart_h=4.95)

# ---- slide 4: bump
add_slide(p4,
    "Tunnettuuden paremmuusjärjestys on pysynyt lähes muuttumattomana",
    "Sijoitukset neljässä mittauksessa: kärkikolmikko vakaa, vain häntäpäässä paikanvaihtoja",
    "AUTETTU TUNNETTUUS  ·  sija mittauksittain (1 = tunnetuin)",
    QUESTION, BASE, chart_top=1.96, chart_h=4.95)

# ---- slide 5: avg line
add_slide(p5,
    "Vain Attendo ja Esperi ylttävät markkinan keskitason yläpuolelle",
    f"Marraskuu 2025: toimijoiden keskimääräinen tunnettuus {avg:.0f} % – kaksi kärkeä sen yläpuolella, kuusi alle",
    "AUTETTU TUNNETTUUS  ·  marraskuu 2025 vs. toimijoiden keskiarvo (%)",
    QUESTION, BASE, chart_top=1.96, chart_h=4.95)

# ---- slide 6: columns
add_slide(p6,
    "Marraskuun 2025 tunnettuus: Attendo hallitsee hoivamarkkinaa",
    "Pystypylväät toimijoittain suuruusjärjestyksessä – ero kärjestä häntään on moninkertainen",
    "AUTETTU TUNNETTUUS  ·  marraskuu 2025 (%)",
    QUESTION, BASE, chart_top=2.05, chart_h=4.85)

prs.save(OUT)
print("slides ->", OUT, "count =", len(prs.slides._sldIdLst))
