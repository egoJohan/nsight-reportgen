"""Merge already-built slides from several source .pptx files into ONE ordered deck.

python-pptx has no slide-merge primitive, so we deep-copy each source slide part:
its full slide XML, every related part (images, charts, embeddings, chart data
workbooks, colors/style), and the rels that wire them together. Each cloned part
gets a fresh partname in the target package so nothing collides, and finally the
new slide id is appended to the target presentation's sldIdLst in order.

Run: uv run python chart_lab/agent_assemble_deck.py
"""
from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.opc.package import Part, _Relationship
from pptx.opc.constants import RELATIONSHIP_TARGET_MODE as RTM
from pptx.opc.packuri import PackURI
from pptx.oxml.ns import qn

LAB = Path(__file__).parent
WORK = LAB.parent / "work"
OUT = WORK / "attendo_agent_deck.pptx"

# Ordered list of (source file stem, 1-based slide number) to copy.
ORDER = [
    ("agent_structural", 1),
    ("agent_structural", 2),
    ("agent_structural", 3),
    ("agent_structural", 4),
    ("agent_structural", 6),
    ("agent_slide_awareness", 1),
    ("agent_slide_trend", 1),
    ("agent_slide_opinion", 1),
    ("agent_structural", 7),
    ("agent_slide_words", 1),
    ("agent_slide_radar", 1),
    ("agent_brand_images", 1),
    ("agent_brand_images", 2),
    ("agent_brand_images", 3),
    ("agent_brand_images", 4),
    ("agent_brand_images", 5),
    ("agent_brand_images", 6),
    ("agent_brand_images", 7),
    ("agent_brand_images", 8),
    ("agent_structural", 8),
]


def _next_partname(used: set, template: str) -> PackURI:
    """Return an unused PackURI like /ppt/media/image{n}.png, recording it in `used`.

    We track allocated partnames ourselves rather than scanning package.iter_parts()
    because a freshly-cloned part is not yet wired into the package rel graph (and so
    not yet visible to iter_parts) while its siblings are still being numbered.
    """
    n = 1
    while True:
        candidate = PackURI(template.format(n))
        if candidate not in used:
            used.add(candidate)
            return candidate
        n += 1


def _clone_part(src_part, target_package, cloned, used):
    """Deep-copy src_part (and transitively its related parts) into target_package.

    `cloned` maps source part -> cloned target part so shared parts are copied once.
    `used` is the set of partnames already allocated in the target package.
    Returns the cloned target part.
    """
    if src_part in cloned:
        return cloned[src_part]

    # Pick a fresh partname in the target package, keeping the same directory/ext.
    src_uri = str(src_part.partname)
    base = src_uri.rsplit("/", 1)[0]
    ext = src_uri.rsplit(".", 1)[-1]
    stem = src_uri.rsplit("/", 1)[-1].rsplit(".", 1)[0]
    # strip trailing digits from stem so template has a clean numbered base
    stem_base = stem.rstrip("0123456789") or "part"
    template = f"{base}/{stem_base}{{}}.{ext}"
    new_partname = _next_partname(used, template)

    # Copy the blob / element. XmlPart instances expose an lxml element we deep-copy;
    # plain Parts (media, embeddings) carry an opaque blob.
    element = getattr(src_part, "_element", None)
    if element is not None:
        new_part = src_part.__class__(
            new_partname,
            src_part.content_type,
            target_package,
            copy.deepcopy(element),
        )
    else:
        new_part = Part(
            new_partname,
            src_part.content_type,
            target_package,
            src_part.blob,
        )

    cloned[src_part] = new_part

    # Recurse into related parts and rebuild rels on the clone.
    for rel in src_part.rels.values():
        if rel.is_external:
            new_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            cloned_target = _clone_part(rel.target_part, target_package, cloned, used)
            # Preserve the original rId so the part XML (r:embed="rId3" etc.) stays valid.
            rels = new_part.rels
            rels._rels[rel.rId] = _Relationship(
                rels._base_uri,
                rel.rId,
                rel.reltype,
                target_mode=RTM.INTERNAL,
                target=cloned_target,
            )

    return new_part


def _copy_slide(src_slide, target_prs, cloned, used):
    """Copy one source slide part into target presentation, append to sldIdLst."""
    target_package = target_prs.part.package
    src_part = src_slide.part

    new_slide_part = _clone_part(src_part, target_package, cloned, used)

    # Wire the slide into the presentation: add a presentation->slide rel and a sldId.
    rId = target_prs.part.rels.get_or_add(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
        new_slide_part,
    )
    sldIdLst = target_prs.slides._sldIdLst
    sldId = sldIdLst.makeelement(qn("p:sldId"), {})
    # allocate a fresh sldId id (must be unique, >= 256)
    existing_ids = [int(s.get("id")) for s in sldIdLst.findall(qn("p:sldId"))]
    new_id = max(existing_ids + [255]) + 1
    sldId.set("id", str(new_id))
    sldId.set(qn("r:id"), rId)
    sldIdLst.append(sldId)


def build(order=ORDER, out_path: Path = OUT) -> Path:
    # Start from a blank deck so we don't inherit any unwanted source slides.
    # python-pptx default template is 4:3; force 16:9 to match all sources.
    target = Presentation()
    target.slide_width = 12191695
    target.slide_height = 6858000

    # Remove any default slides that ship with the blank template.
    sldIdLst = target.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        target.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    # Cache opened source presentations.
    sources: dict[str, Presentation] = {}
    cloned: dict = {}  # source part -> cloned target part (shared across all slides)
    # Partnames already present in the blank target (masters, layouts, theme, etc.).
    used = {p.partname for p in target.part.package.iter_parts()}

    for stem, num in order:
        if stem not in sources:
            sources[stem] = Presentation(WORK / f"{stem}.pptx")
        src_slide = sources[stem].slides[num - 1]
        _copy_slide(src_slide, target, cloned, used)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    target.save(out_path)
    return out_path


if __name__ == "__main__":
    p = build()
    final = Presentation(p)
    print(f"Wrote {p} with {len(final.slides._sldIdLst)} slides")
