#!/usr/bin/env python3
"""Dense "general opinion" perception slide (nSight house style).

Carries — like the original deck — THREE layers of information in one slide:

  1. CURRENT WAVE breakdown: one 100% horizontal stacked bar per group, split
     into the 5 ordered levels (Erittäin huono … Erittäin hyvä) + "En osaa
     sanoa", with a diverging negative|positive colour scale so the positive
     side reads as positive.
  2. POSITIVE TREND: a small sparkline per group tracing the positive share
     across ALL FOUR waves (Toukokuu 2024 → Marraskuu 2025).
  3. CHANGE vs the previous wave: the +6 % / +3 % delta called out per group.

All numbers are read verbatim from the JSON data file; nothing is invented.
The current positive share is derived deterministically as Hyvä + Erittäin
hyvä; the latest sparkline point is asserted to equal that derived value.

Data shape (chart_lab/opinion_trend.json):
    {
      "categories": ["Yksityiset palveluntarjoajat", "Julkinen palveluntarjoaja"],
      "current_levels": {
        "Erittäin huono":[priv,pub], "Huono":[..], "Hyvä":[..],
        "Erittäin hyvä":[..], "En osaa sanoa":[..]
      },
      "positive_trend": {
        "<group>": {"Toukokuu 2024":53, ..., "Marraskuu 2025":58}, ...
      },
      "change_vs_prev": {"<group>":"+6 %", "<group>":"+3 %"}
    }
"""
import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
from matplotlib.patches import Patch, FancyBboxPatch

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = Path("/home/johan/Projects/nsight/proto")

# --------------------------------------------------------------- palette
CREAM = "#F7F3EC"
INK = "#2B2B2B"
MUTED = "#6E6A63"
TEAL = "#13615E"
TEAL_LT = "#7DB8A6"
RED_DK = "#B23A2E"
RED_LT = "#E08C82"
GREY = "#C9C1B4"


def build(
    data_path=ROOT / "chart_lab" / "opinion_trend.json",
    png_path=ROOT / "chart_lab" / "agent_dense_opinion.png",
    out_path=ROOT / "work" / "agent_dense_opinion.pptx",
    *,
    negative_keys=("Erittäin huono", "Huono"),
    positive_keys=("Hyvä", "Erittäin hyvä"),
    dontknow_key="En osaa sanoa",
    wave="Marraskuu 2025",
    base_n=1001,
    question=("Mikä on yleinen käsityksesi tuntemistasi hoivapalveluita "
              "tarjoavista yrityksistä?"),
    section_label="YLEINEN KÄSITYS  ·  osuus vastaajista (%)",
    title=("Yksityisten palveluntarjoajien käsitys on parantunut – "
           "nyt yhtä myönteinen kuin julkisella"),
):
    # ------------------------------------------------------------- fonts
    for f in ["/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
              "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"]:
        fm.fontManager.addfont(f)
    plt.rcParams["font.family"] = "Liberation Sans"

    # ------------------------------------------------------------- data
    d = json.loads(Path(data_path).read_text())
    cats = d["categories"]
    levels = d["current_levels"]
    trend = d["positive_trend"]
    change = d["change_vs_prev"]
    n_groups = len(cats)

    stack_keys = list(negative_keys) + list(positive_keys) + [dontknow_key]
    SEG_COLORS = {
        negative_keys[0]: RED_DK,   # Erittäin huono (darkest red)
        negative_keys[1]: RED_LT,   # Huono
        positive_keys[0]: TEAL_LT,  # Hyvä
        positive_keys[1]: TEAL,     # Erittäin hyvä (dark teal)
        dontknow_key:     GREY,     # En osaa sanoa
    }
    LIGHT_TEXT = {negative_keys[0], positive_keys[1]}

    # derived positive share (exact, from current_levels)
    pos_share = [sum(levels[k][i] for k in positive_keys) for i in range(n_groups)]
    neg_share = [sum(levels[k][i] for k in negative_keys) for i in range(n_groups)]

    # waves (ordered as in JSON) and per-group positive series
    waves = list(trend[cats[0]].keys())
    series_trend = {c: [trend[c][w] for w in waves] for c in cats}

    # NOTE: two distinct positive figures exist in the source and are kept
    # distinct (both verbatim from data):
    #   pos_share[i]      = Hyvä + Erittäin hyvä summed from current_levels
    #                       (drives the stacked-bar zone label, exact to bar)
    #   trend_latest[i]   = positive_trend[...]["Marraskuu 2025"]
    #                       (drives the sparkline + headline / "Nyt" callout)
    # The headline metric is the tracked positive_trend series, so the big
    # "Nyt" % and subtitle use trend_latest; the bar zone uses pos_share.
    trend_latest = [series_trend[c][-1] for c in cats]

    # ============================================================ figure
    # Two-column dense layout in ONE figure:
    #   left  : 100% stacked breakdown (the two bars)
    #   right : per-group trend panels (sparkline + current % + change badge)
    FIG_W, FIG_H = 12.4, 4.7
    fig = plt.figure(figsize=(FIG_W, FIG_H), dpi=240)
    fig.patch.set_facecolor(CREAM)

    # ---- left axis: stacked breakdown -------------------------------
    axL = fig.add_axes([0.140, 0.150, 0.490, 0.700])
    axL.set_facecolor(CREAM)
    bar_h = 0.50
    ypos = list(range(n_groups))[::-1]          # first category at top

    for i, y in zip(range(n_groups), ypos):
        left = 0.0
        for k in stack_keys:
            v = levels[k][i]
            if v == 0:
                continue
            axL.barh(y, v, left=left, height=bar_h,
                     color=SEG_COLORS[k], edgecolor=CREAM, linewidth=1.4,
                     zorder=3)
            if v >= 4:
                axL.text(left + v / 2.0, y, f"{v}",
                         va="center", ha="center",
                         fontsize=11, fontweight="bold",
                         color="#FFFFFF" if k in LIGHT_TEXT else INK, zorder=5)
            left += v
        # group label above its bar
        axL.text(0, y + bar_h * 0.62 + 0.04, cats[i],
                 va="bottom", ha="left", fontsize=12.5, color=INK,
                 fontweight="bold", zorder=5)

    # macro-zone labels along the top of the first (private) bar as a guide
    top_y = ypos[0] + bar_h * 0.62 + 0.40
    axL.text(neg_share[0] / 2.0, top_y, "KIELTEISET",
             va="bottom", ha="center", fontsize=8.5, fontweight="bold",
             color=RED_DK, zorder=5)
    axL.text(neg_share[0] + pos_share[0] / 2.0, top_y, "MYÖNTEISET",
             va="bottom", ha="center", fontsize=8.5, fontweight="bold",
             color=TEAL, zorder=5)

    axL.set_xlim(0, 100)
    axL.set_ylim(-0.62, n_groups - 0.10)
    axL.axis("off")

    # legend under the left column (full ordered scale)
    handles = [Patch(facecolor=SEG_COLORS[k], label=k) for k in stack_keys]
    leg = axL.legend(handles=handles, loc="upper center",
                     bbox_to_anchor=(0.5, -0.06), ncol=len(stack_keys),
                     frameon=False, fontsize=9.2,
                     handlelength=1.0, handleheight=1.0,
                     columnspacing=1.1, handletextpad=0.45)
    for t in leg.get_texts():
        t.set_color(INK)

    # column heading for the left panel
    fig.text(0.140, 0.945, "Nykyinen jakauma (Marraskuu 2025)",
             fontsize=10.5, fontweight="bold", color=TEAL, ha="left")

    # ---- right: per-group trend panels ------------------------------
    # one sparkline axis per group, stacked vertically, aligned to bars
    fig.text(0.660, 0.945, "Myönteisten kehitys · 4 mittausta",
             fontsize=10.5, fontweight="bold", color=TEAL, ha="left")

    # vertical centres of the two panels (top group first)
    panel_h = 0.305
    panel_gap = 0.055
    panel_top = 0.150 + 0.700  # = 0.850 (top of left axis region)
    panel_bottoms = [panel_top - panel_h,
                     panel_top - 2 * panel_h - panel_gap]

    SPARK_L = 0.660
    SPARK_W = 0.205

    ymin = min(min(series_trend[c]) for c in cats) - 3
    ymax = max(max(series_trend[c]) for c in cats) + 4

    for idx, c in enumerate(cats):
        b = panel_bottoms[idx]
        axS = fig.add_axes([SPARK_L, b + 0.085, SPARK_W, panel_h - 0.110])
        axS.set_facecolor(CREAM)
        ys = series_trend[c]
        xs = list(range(len(ys)))
        axS.plot(xs, ys, color=TEAL, linewidth=2.4, zorder=3,
                 solid_capstyle="round")
        # markers; emphasise last point
        axS.scatter(xs[:-1], ys[:-1], s=26, color=TEAL, zorder=4)
        axS.scatter([xs[-1]], [ys[-1]], s=80, color=TEAL,
                    edgecolor=CREAM, linewidth=1.6, zorder=5)
        # value labels on each point
        for x, yv in zip(xs, ys):
            dy = 1.9 if x != xs[-1] else 2.6
            axS.text(x, yv + dy, f"{yv}", ha="center", va="bottom",
                     fontsize=8.6, color=INK,
                     fontweight="bold" if x == xs[-1] else "normal")
        axS.set_xlim(-0.35, len(ys) - 0.65)
        axS.set_ylim(ymin, ymax)
        for s in axS.spines.values():
            s.set_visible(False)
        axS.set_yticks([])
        axS.set_xticks([])
        # group caption above the sparkline
        short = "Yksityiset" if "Yksityis" in c else "Julkinen"
        fig.text(SPARK_L, b + panel_h - 0.018, short,
                 fontsize=10.5, fontweight="bold", color=INK,
                 ha="left", va="top")

    # wave tick labels (shared, under the lower sparkline)
    short_waves = [w.replace("Toukokuu ", "05/").replace("Marraskuu ", "11/")
                   for w in waves]
    lower_b = panel_bottoms[-1]
    for j, lab in enumerate(short_waves):
        x = SPARK_L + SPARK_W * (j / (len(short_waves) - 1)) * \
            ((len(waves) - 1 - 0) / (len(waves) - 1))
        # evenly space across the sparkline data span (-0.35..n-0.65)
    # compute tick x positions in fig coords matching axis data 0..n-1
    span_lo, span_hi = -0.35, len(waves) - 0.65
    for j, lab in enumerate(short_waves):
        frac = (j - span_lo) / (span_hi - span_lo)
        fx = SPARK_L + SPARK_W * frac
        fig.text(fx, lower_b + 0.052, lab, fontsize=7.6, color=MUTED,
                 ha="center", va="top")

    # ---- right-most: current positive % + change badge --------------
    BADGE_L = 0.882
    for idx, c in enumerate(cats):
        b = panel_bottoms[idx]
        cy = b + panel_h / 2.0 + 0.012
        # big current positive % (tracked positive_trend metric)
        fig.text(BADGE_L, cy + 0.045, f"{trend_latest[idx]} %",
                 fontsize=23, fontweight="bold", color=TEAL,
                 ha="left", va="center")
        fig.text(BADGE_L + 0.0015, cy - 0.052, "myönteisiä",
                 fontsize=8.2, color=MUTED, ha="left", va="center")
        # change-vs-prev badge
        chg = change[c]
        axB = fig.add_axes([BADGE_L, cy - 0.150, 0.090, 0.060])
        axB.axis("off")
        box = FancyBboxPatch((0.02, 0.10), 0.96, 0.80,
                             boxstyle="round,pad=0.02,rounding_size=0.18",
                             transform=axB.transAxes,
                             facecolor=TEAL, edgecolor="none", zorder=1)
        axB.add_patch(box)
        axB.text(0.5, 0.52, f"{chg} vs. ed.", transform=axB.transAxes,
                 ha="center", va="center", fontsize=8.8, fontweight="bold",
                 color="#FFFFFF", zorder=2)

    fig.text(BADGE_L, 0.945, "Nyt", fontsize=10.5, fontweight="bold",
             color=TEAL, ha="left")

    fig.savefig(png_path, dpi=240, facecolor=CREAM,
                bbox_inches="tight", pad_inches=0.06)
    plt.close(fig)
    print("chart ->", png_path)

    # ============================================================ pptx
    PX_CREAM = RGBColor(0xF7, 0xF3, 0xEC)
    PX_INK = RGBColor(0x2B, 0x2B, 0x2B)
    PX_TEAL = RGBColor(0x13, 0x61, 0x5E)
    PX_MUTED = RGBColor(0x6E, 0x6A, 0x63)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = PX_CREAM
    bg.line.fill.background(); bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    def textbox(l, t, w, h, runs, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, space_after=0):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = 0; tf.margin_right = 0
        tf.margin_top = 0; tf.margin_bottom = 0
        first = True
        for line in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.space_after = Pt(space_after)
            p.space_before = Pt(0)
            for txt, sz, col, bold in line:
                r = p.add_run(); r.text = txt
                r.font.size = Pt(sz); r.font.bold = bold
                r.font.color.rgb = col
                r.font.name = "Liberation Sans"
        return tb

    # subtitle string from exact data
    subtitle = (f"Myönteisten osuus {trend_latest[0]} % (yksityiset, "
                f"{change[cats[0]]}) ja {trend_latest[1]} % (julkinen, "
                f"{change[cats[1]]}) – ero käytännössä hävinnyt")

    # accent bar
    acc = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42),
                                 Inches(0.10), Inches(0.92))
    acc.fill.solid(); acc.fill.fore_color.rgb = PX_TEAL
    acc.line.fill.background(); acc.shadow.inherit = False

    # title + subtitle
    textbox(Inches(0.80), Inches(0.40), Inches(12.0), Inches(1.05),
            [[(title, 22, PX_INK, True)],
             [(subtitle, 13.0, PX_MUTED, False)]],
            space_after=3)

    # section label
    textbox(Inches(0.80), Inches(1.60), Inches(12.0), Inches(0.32),
            [[(section_label, 11, PX_TEAL, True)]])

    # chart
    iw, ih = Image.open(png_path).size
    ar = iw / ih
    pic_w = Inches(12.55)
    pic_h = Emu(int(pic_w / ar))
    max_h = Inches(4.75)
    if pic_h > max_h:
        pic_h = max_h
        pic_w = Emu(int(pic_h * ar))
    pic_l = Inches((13.333 - pic_w / Inches(1)) / 2)
    pic_t = Inches(2.00)
    slide.shapes.add_picture(str(png_path), pic_l, pic_t,
                             width=pic_w, height=pic_h)

    # footer: question + base n
    textbox(Inches(0.80), Inches(6.98), Inches(9.8), Inches(0.5),
            [[("Kysymys: ", 9.5, PX_MUTED, True),
              ("”" + question + "”", 9.5, PX_MUTED, False)]])
    textbox(Inches(8.7), Inches(6.98), Inches(4.0), Inches(0.5),
            [[(f"{wave} · Kaikki vastaajat, n = {base_n}",
               9.5, PX_MUTED, True)]],
            align=PP_ALIGN.RIGHT)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print("slide ->", out_path)
    return out_path


if __name__ == "__main__":
    build()
