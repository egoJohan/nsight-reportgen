#!/usr/bin/env python3
"""Build the "general opinion of providers" perception slide (nSight house style).

Reusable / parameterised generator. For the *current* wave it shows ONE diverging
100%-stacked horizontal bar per group (negatives to the LEFT of a zero line,
positives to the RIGHT, "En osaa sanoa" parked as a neutral grey tail). Alongside
each bar it adds a 4-wave positive-% sparkline and the change-vs-previous badge.

All numbers are read verbatim from the JSON data file; derived figures
(positive share, negative share) are computed in code.

Data shape (see opinion_trend.json):
    {
      "categories": ["Yksityiset ...", "Julkinen ..."],
      "current_levels": {
          "Erittäin huono":[priv,pub], "Huono":[..], "Hyvä":[..],
          "Erittäin hyvä":[..], "En osaa sanoa":[..] },   # one value per group, %
      "positive_trend": { "<group>": {"<wave>": %, ... } },  # 4 waves each
      "change_vs_prev": { "<group>": "+x %" }
    }
"""
import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
from matplotlib.patches import Patch

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = Path("/home/johan/Projects/nsight/proto")

# ----------------------------------------------------------------- palette
CREAM = "#F4EFE6"
INK = "#2B2B2B"
MUTED = "#8A857B"
TEAL = "#13615E"
# Diverging perception scale: dark red (worst) -> dark teal (best); grey = EOS.
SEG_COLORS = {
    "Erittäin huono": "#B23A2E",   # darkest red
    "Huono":          "#E08C82",   # light red
    "Hyvä":           "#7DB8A6",   # light teal
    "Erittäin hyvä":  "#13615E",   # dark teal
    "En osaa sanoa":  "#C9C1B4",   # neutral grey
}
LIGHT_TEXT = {"Erittäin huono", "Erittäin hyvä"}


def build(
    data_path=ROOT / "chart_lab" / "opinion_trend.json",
    png_path=ROOT / "chart_lab" / "agent_validate_opinion_chart.png",
    out_path=ROOT / "work" / "agent_validate_opinion.pptx",
    *,
    negative_keys=("Erittäin huono", "Huono"),
    positive_keys=("Hyvä", "Erittäin hyvä"),
    dontknow_key="En osaa sanoa",
    wave="Marraskuu 2025",
    base_n=1001,
    question=("Mikä on yleinen käsityksesi tuntemistasi hoivapalveluita "
              "tarjoavista yrityksistä?"),
    section_label="YLEINEN KÄSITYS  ·  osuus vastaajista (%)",
    title=("Yksityisten käsitys on parantunut – nyt yhtä myönteinen "
           "kuin julkisella"),
):
    # ------------------------------------------------------------- fonts
    for f in ["/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
              "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"]:
        fm.fontManager.addfont(f)
    plt.rcParams["font.family"] = "Liberation Sans"

    # ------------------------------------------------------------- data
    d = json.loads(Path(data_path).read_text())
    cats = d["categories"]
    levels = d["current_levels"]
    pos_trend = d["positive_trend"]
    change = d["change_vs_prev"]
    n_groups = len(cats)

    # derived totals (exact, from data)
    pos_share = [sum(levels[k][i] for k in positive_keys) for i in range(n_groups)]
    neg_share = [sum(levels[k][i] for k in negative_keys) for i in range(n_groups)]

    # subtitle from data: positive shares both groups
    subtitle = (f"Myönteisten osuus {pos_share[0]} % (yksityiset) ja "
                f"{pos_share[1]} % (julkinen) – ero käytännössä hävinnyt")

    # =============================================================== CHART
    # Two side-by-side panels: (left, wide) diverging stacked bars;
    # (right, narrow) per-group positive-% sparklines.
    FIG_W, FIG_H = 12.0, 4.35
    fig = plt.figure(figsize=(FIG_W, FIG_H), dpi=220)
    fig.patch.set_facecolor(CREAM)

    # gridspec: bars panel | sparkline panel
    gs = fig.add_gridspec(1, 2, width_ratios=[3.05, 1.0],
                          left=0.005, right=0.995, top=0.86, bottom=0.10,
                          wspace=0.04)
    axb = fig.add_subplot(gs[0, 0]); axb.set_facecolor(CREAM)
    axs = fig.add_subplot(gs[0, 1]); axs.set_facecolor(CREAM)

    bar_h = 0.40
    ypos = list(range(n_groups))[::-1]   # first category at top

    # -------- diverging bars: negatives left of 0, positives right, EOS tail
    # Layout along x (data units, %): EOS sits as a grey tail to the right
    # of the positives, separated by a small gap so the 0-line stays the
    # negative|positive divide.
    EOS_GAP = 3.0
    for i, y in zip(range(n_groups), ypos):
        # negatives stack leftward from 0 (worst furthest left)
        left = 0.0
        for k in negative_keys[::-1]:        # "Huono" nearest 0, "Erittäin huono" outer
            v = levels[k][i]
            if v:
                axb.barh(y, -v, left=left, height=bar_h, color=SEG_COLORS[k],
                         edgecolor=CREAM, linewidth=1.4, zorder=3)
                axb.text(left - v / 2.0, y, f"{v}", va="center", ha="center",
                         fontsize=11.5, fontweight="bold",
                         color="#FFFFFF" if k in LIGHT_TEXT else INK, zorder=5)
                left -= v
        # positives stack rightward from 0 ("Hyvä" nearest 0)
        right = 0.0
        for k in positive_keys:
            v = levels[k][i]
            if v:
                axb.barh(y, v, left=right, height=bar_h, color=SEG_COLORS[k],
                         edgecolor=CREAM, linewidth=1.4, zorder=3)
                axb.text(right + v / 2.0, y, f"{v}", va="center", ha="center",
                         fontsize=11.5, fontweight="bold",
                         color="#FFFFFF" if k in LIGHT_TEXT else INK, zorder=5)
                right += v
        # EOS grey tail, offset past the positives
        eos = levels[dontknow_key][i]
        if eos:
            eos_left = right + EOS_GAP
            axb.barh(y, eos, left=eos_left, height=bar_h,
                     color=SEG_COLORS[dontknow_key], edgecolor=CREAM,
                     linewidth=1.4, zorder=3)
            axb.text(eos_left + eos / 2.0, y, f"{eos}", va="center", ha="center",
                     fontsize=10.5, fontweight="bold", color=INK, zorder=5)

        # group label above the bar (left aligned to negative extent)
        axb.text(left, y + bar_h * 0.92, cats[i], va="bottom", ha="left",
                 fontsize=12.5, color=INK, fontweight="bold", zorder=5)
        # positive-share callout, just past EOS
        axb.text(eos_left + eos + 3.5, y, f"{pos_share[i]} %", va="center",
                 ha="left", fontsize=16, fontweight="bold", color=TEAL, zorder=6)

    # zero divide line
    axb.axvline(0, color=MUTED, lw=1.1, zorder=2)
    # zone labels along the top
    top_y = n_groups - 0.5 + 0.10
    axb.text(-neg_share[0] / 2.0, top_y, "KIELTEISET", va="bottom", ha="center",
             fontsize=9, fontweight="bold", color="#B23A2E", zorder=5)
    axb.text(pos_share[0] / 2.0, top_y, "MYÖNTEISET", va="bottom", ha="center",
             fontsize=9, fontweight="bold", color=TEAL, zorder=5)

    max_neg = max(neg_share)
    max_eos_end = max((max(0.0, (sum(levels[k][i] for k in positive_keys)))
                       + EOS_GAP + levels[dontknow_key][i]) for i in range(n_groups))
    axb.set_xlim(-max_neg - 8, max_eos_end + 16)
    axb.set_ylim(-0.55, n_groups - 0.20)
    axb.axis("off")

    # -------- EOS legend label near its tail (one small annotation)
    axb.text(pos_share[0] + EOS_GAP + levels[dontknow_key][0] / 2.0,
             n_groups - 0.5 + 0.10, "EOS", va="bottom", ha="center",
             fontsize=8.5, fontweight="bold", color=MUTED, zorder=5)

    # =============================================================== SPARKLINES
    axs.set_xlim(-0.4, 3.4)
    axs.set_ylim(-0.55, n_groups - 0.20)
    axs.axis("off")
    axs.text(1.5, n_groups - 0.5 + 0.10, "MYÖNTEISTEN OSUUS, 4 AALTOA",
             va="bottom", ha="center", fontsize=8.5, fontweight="bold",
             color=MUTED, zorder=5)

    for i, y in zip(range(n_groups), ypos):
        cat = cats[i]
        waves = list(pos_trend[cat].keys())
        vals = [pos_trend[cat][w] for w in waves]
        n_w = len(vals)
        # normalise spark y within a small band around the row centre
        vmin, vmax = min(vals), max(vals)
        span = max(vmax - vmin, 1)
        band = 0.30
        sx = [j / (n_w - 1) * 3.0 for j in range(n_w)]
        sy = [y - band + (v - vmin) / span * (2 * band) for v in vals]
        axs.plot(sx, sy, color=TEAL, lw=2.0, zorder=3, solid_capstyle="round")
        # markers; emphasise last (current) point
        axs.scatter(sx[:-1], sy[:-1], s=18, color=TEAL, zorder=4)
        axs.scatter([sx[-1]], [sy[-1]], s=46, color=TEAL,
                    edgecolor=CREAM, linewidth=1.2, zorder=5)
        # first + last value labels
        axs.text(sx[0] - 0.12, sy[0], f"{vals[0]}", va="center", ha="right",
                 fontsize=9, color=MUTED, zorder=5)
        axs.text(sx[-1] + 0.14, sy[-1], f"{vals[-1]} %", va="center", ha="left",
                 fontsize=10.5, color=TEAL, fontweight="bold", zorder=5)
        # change-vs-prev badge below the spark
        chg = change[cat]
        axs.text(1.5, y - band - 0.18, f"vs. edellinen  {chg}", va="top",
                 ha="center", fontsize=10, fontweight="bold",
                 color=TEAL, zorder=5)

    # =============================================================== legend
    stack_keys = list(negative_keys) + list(positive_keys) + [dontknow_key]
    handles = [Patch(facecolor=SEG_COLORS[k], label=k) for k in stack_keys]
    leg = fig.legend(handles=handles, loc="lower center",
                     bbox_to_anchor=(0.42, -0.01), ncol=len(stack_keys),
                     frameon=False, fontsize=10,
                     handlelength=1.0, handleheight=1.0,
                     columnspacing=1.3, handletextpad=0.5)
    for t in leg.get_texts():
        t.set_color(INK)

    fig.savefig(png_path, dpi=220, facecolor=CREAM,
                bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    print("chart ->", png_path)

    # =============================================================== PPTX
    PX_CREAM = RGBColor(0xF4, 0xEF, 0xE6)
    PX_INK = RGBColor(0x2B, 0x2B, 0x2B)
    PX_TEAL = RGBColor(0x13, 0x61, 0x5E)
    PX_MUTED = RGBColor(0x8A, 0x85, 0x7B)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = PX_CREAM
    bg.line.fill.background(); bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    def textbox(l, t, w, h, runs, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, space_after=0):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = 0; tf.margin_right = 0
        tf.margin_top = 0; tf.margin_bottom = 0
        first = True
        for line in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.space_after = Pt(space_after)
            p.space_before = Pt(0)
            for txt, sz, col, bold in line:
                r = p.add_run(); r.text = txt
                r.font.size = Pt(sz); r.font.bold = bold
                r.font.color.rgb = col
                r.font.name = "Liberation Sans"
        return tb

    # accent bar
    acc = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42),
                                 Inches(0.10), Inches(0.92))
    acc.fill.solid(); acc.fill.fore_color.rgb = PX_TEAL
    acc.line.fill.background(); acc.shadow.inherit = False

    # title + subtitle
    textbox(Inches(0.80), Inches(0.40), Inches(12.0), Inches(1.05),
            [[(title, 22, PX_INK, True)],
             [(subtitle, 13.5, PX_MUTED, False)]],
            space_after=3)

    # section label
    textbox(Inches(0.80), Inches(1.62), Inches(12.0), Inches(0.32),
            [[(section_label, 11, PX_TEAL, True)]])

    # chart
    iw, ih = Image.open(png_path).size
    ar = iw / ih
    pic_w = Inches(12.25)
    pic_h = Emu(int(pic_w / ar))
    max_h = Inches(4.65)
    if pic_h > max_h:
        pic_h = max_h
        pic_w = Emu(int(pic_h * ar))
    pic_l = Inches(0.55)
    pic_t = Inches(2.02)
    slide.shapes.add_picture(str(png_path), pic_l, pic_t,
                             width=pic_w, height=pic_h)

    # footer: question + base n
    textbox(Inches(0.80), Inches(6.98), Inches(9.8), Inches(0.5),
            [[("Kysymys: ", 9.5, PX_MUTED, True),
              ("”" + question + "”", 9.5, PX_MUTED, False)]])
    textbox(Inches(8.7), Inches(6.98), Inches(4.05), Inches(0.5),
            [[(f"{wave} · Kaikki vastaajat, n = {base_n}",
               9.5, PX_MUTED, True)]],
            align=PP_ALIGN.RIGHT)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print("slide ->", out_path)
    return out_path


if __name__ == "__main__":
    build()
