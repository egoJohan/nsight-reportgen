from __future__ import annotations

from pathlib import Path

from pptx import Presentation


class Template:
    """Open a .pptx and address shapes by (slide index, shape name)."""

    def __init__(self, path: Path) -> None:
        self._path = Path(path)
        self.prs = Presentation(str(self._path))

    def shape(self, *, slide_idx: int, name: str, occurrence: int = 0):
        slide = self.prs.slides[slide_idx]
        matches = [sh for sh in slide.shapes if sh.name == name]
        if occurrence < len(matches):
            return matches[occurrence]
        raise KeyError(
            f"shape {name!r} occurrence {occurrence} not found on slide {slide_idx} "
            f"({len(matches)} found)"
        )

    def save(self, out_path: Path) -> None:
        self.prs.save(str(out_path))
