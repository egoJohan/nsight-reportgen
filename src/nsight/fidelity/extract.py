from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation


@dataclass
class ChartData:
    name: str
    categories: list[str]
    series: dict[str, list[float]]


@dataclass
class TableData:
    name: str
    cells: dict[tuple[int, int], str]


@dataclass
class SlideData:
    idx: int
    charts: list[ChartData] = field(default_factory=list)
    tables: list[TableData] = field(default_factory=list)
    texts: list[str] = field(default_factory=list)


@dataclass
class DeckData:
    slides: list[SlideData]


def extract_deck(path: Path) -> DeckData:
    prs = Presentation(str(path))
    slides: list[SlideData] = []
    for i, slide in enumerate(prs.slides):
        sd = SlideData(idx=i)
        for sh in slide.shapes:
            if sh.has_chart:
                plot = sh.chart.plots[0]
                series = {s.name: list(s.values) for s in plot.series}
                sd.charts.append(ChartData(name=sh.name,
                                           categories=[str(c) for c in plot.categories],
                                           series=series))
            elif sh.has_table:
                cells = {}
                for r, row in enumerate(sh.table.rows):
                    for c, cell in enumerate(row.cells):
                        cells[(r, c)] = cell.text
                sd.tables.append(TableData(name=sh.name, cells=cells))
            elif sh.has_text_frame and sh.text_frame.text.strip():
                sd.texts.append(sh.text_frame.text)
        slides.append(sd)
    return DeckData(slides=slides)
