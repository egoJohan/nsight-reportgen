"""Objective number gates — fidelity layers 1 & 2 (design §10)."""
from __future__ import annotations
import re
import pdfplumber
from pptx import Presentation
from pptx.oxml.ns import qn
from reportbuilder.stats.series import SeriesResult

_NUM = re.compile(r"-?\d+(?:[.,]\d+)?")


def _xy_x_values(series_el) -> list[float]:
    """Extract x values from an XY scatter series element via c:xVal."""
    xVal = series_el.find(".//" + qn("c:xVal"))
    if xVal is None:
        return []
    return [float(pt.text) for pt in xVal.findall(".//" + qn("c:v")) if pt.text is not None]


def numbers_from_pptx(pptx_path: str) -> dict:
    """Read native chart series values: {series_name: [float, ...]}.

    For XY scatter charts, both x values and y values are included in the
    returned list (x values first, then y values) so that assert_series_match
    can verify the complete data payload.
    """
    prs = Presentation(pptx_path)
    out: dict[str, list[float]] = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_chart", False):
                continue
            for plot in shape.chart.plots:
                for series in plot.series:
                    y_vals = [float(v) for v in series.values]
                    # For XY scatter, also extract x values from the OOXML element.
                    x_vals = _xy_x_values(series._element)
                    out[series.name] = x_vals + y_vals if x_vals else y_vals
    return out

def numbers_from_pdf(pdf_path: str) -> list[float]:
    """Every numeric token in the rendered PDF text (data labels are real text)."""
    nums: list[float] = []
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            for tok in _NUM.findall(text):
                nums.append(float(tok.replace(",", ".")))
    return nums

def _expected_values(series: SeriesResult) -> list[float]:
    attr = {"pct": "pct", "count": "count", "mean": "mean"}[series.statistic]
    vals: list[float] = []
    for cat in series.categories:
        for seg in series.segments:
            v = getattr(series.cell(cat, seg), attr)
            if v is not None:
                vals.append(float(v))
    return vals

def assert_series_match(extracted: dict, series: SeriesResult, tol: float = 0.5) -> None:
    """Every SeriesResult value (per its statistic) appears in `extracted` within `tol`."""
    pool: list[float] = []
    for v in extracted.values():
        pool.extend(v if isinstance(v, (list, tuple)) else [v])
    missing = [
        exp for exp in _expected_values(series)
        if not any(abs(exp - got) <= tol for got in pool)
    ]
    assert not missing, f"values not found within tol={tol}: {missing} in {pool}"
