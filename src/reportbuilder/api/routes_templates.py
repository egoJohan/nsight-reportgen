"""Presentation templates: upload, bind, resolve.

nSight's analysts render into their CLIENTS' brand decks, so a template is
validated on upload and rejected with a reason rather than silently producing a
deck that ignores it (render/template_check.py).

Binding follows the card: a template can be set on an asiakas, a tutkimus or a
single report, and the lower level always wins. An already-delivered report
keeps the template it rendered with until someone asks for the update — see
Repository.resolve_template.
"""
import json
import logging
from dataclasses import replace

from fastapi import APIRouter, Depends, File, HTTPException, UploadFile
from fastapi.responses import Response
from pydantic import BaseModel

from reportbuilder.api.deps_auth import (
    require_case, require_case_in_customer, require_case_in_customer_write,
    require_case_write, require_customer, require_customer_write,
)
from reportbuilder.api.deps_store import get_auth, get_repository
from reportbuilder.auth.permissions import User
from reportbuilder.render.template_check import inspect_template
from reportbuilder.store.repository import Repository

log = logging.getLogger(__name__)
from reportbuilder.store.seam import AuthContext, NotFound

templates_router = APIRouter()

_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


class TemplateBinding(BaseModel):
    """None clears the binding, so the level above takes over again."""
    template_id: str | None = None


def _as_dict(t) -> dict:
    d = {"id": t.id, "name": t.name, "size": t.size,
         "layout_name": t.layout_name, "palette": list(t.palette),
         "heading_font": t.heading_font}
    # Recorded at upload; carried on every listing so the UI can keep flagging a
    # template whose fonts the render host cannot supply.
    fonts = list(getattr(t, "fonts", ()) or ())
    d["fonts"] = fonts
    # No recorded check (a template uploaded before this existed) is not the
    # same as a failed one, so absence reads as OK rather than as an alarm.
    d["fonts_ok"] = all(f.get("ok") for f in fonts) if fonts else True
    return d


def _live_font_status(stored: list[dict], families: list[str]) -> list[dict]:
    """Re-check *families* against the host as it is RIGHT NOW.

    The stored record is what was true at upload. It cannot know about a font
    installed since, or a stand-in chosen a minute ago, so a template row would
    keep saying "Missing font" after the admin had already dealt with it.

    Where a font is still unresolved the STORED reason wins: it explains the
    licence ("not an open-licence font, cannot be installed"), while a
    network-free re-check can only say it is absent. Cheap enough for a list —
    no network, and fontconfig's family list is cached in-process.

    Re-checks every family we RECORDED as well as the ones the caller knows
    about. Callers pass the theme's heading/body pair, but the upload check
    reads every family the FILE names — a slide master routinely sets fonts the
    theme never mentions. Driving the result from the caller's list alone
    dropped those: the upload warned that "Barlow Condensed" could not be
    installed and the template's own row then did not list it, leaving an admin
    a warning about a font the product denied needing.
    """
    from reportbuilder.render.fonts import check_template_fonts

    by_family = {f.get("family"): f for f in stored if isinstance(f, dict)}
    wanted = list(families)
    wanted += [f for f in by_family if f and f not in wanted]
    try:
        live = check_template_fonts(wanted, allow_network=False)
    except Exception:  # noqa: BLE001 — fall back to what we recorded
        return stored
    out = []
    for st in live:
        if st.ok:
            out.append(st.as_dict())
        else:
            out.append(by_family.get(st.family) or st.as_dict())
    return out


def _resolve_template_fonts(theme, pptx_path: str = "") -> list[dict]:  # theme OR Template record
    """Install the fonts this template names, or record why we cannot.

    Done at UPLOAD, because that is when someone is looking and can act: get a
    licence, install the font, or choose a different template. Discovering it at
    render time means finding out after the deck was sent.

    Never raises: a template with an unavailable font still uploads and still
    renders — it just renders in a substitute, which is exactly what the
    returned status says out loud.
    """
    from reportbuilder.render.fonts import check_template_fonts, families_in_template
    # Every font the FILE names, not just the theme's major/minor pair. On
    # Egoiq_x_Rahoo the theme says Arial twice while the master sets the title in
    # Bebas Neue and the body in Barlow Condensed Medium — both open-licence
    # Google families nSight could have installed, and never learned it needed,
    # so the customer's headline rendered in a substitute.
    families = [f for f in (getattr(theme, "heading_font", ""),
                            getattr(theme, "body_font", "")) if f]
    if pptx_path:
        families += families_in_template(pptx_path)
    try:
        return [st.as_dict() for st in check_template_fonts(families)]
    except Exception:  # noqa: BLE001 — a font check must not block an upload
        return []


@templates_router.post("/customers/{customer_id}/templates", status_code=201)
async def upload_template(customer_id: str, file: UploadFile = File(...),
                          auth: AuthContext = Depends(get_auth),
                          repo: Repository = Depends(get_repository),
                          user: User = Depends(require_customer_write)) -> dict:
    """Validate and store a template. 422 with the reason when it cannot work."""
    data = await file.read()
    if not data:
        raise HTTPException(422, "Empty file")

    # Validate from a temp copy: python-pptx needs a path, and a bad upload must
    # be rejected before anything is stored.
    import os
    import tempfile
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name
    try:
        report = inspect_template(tmp_path)
        # While the temp copy still exists: the font scan reads the FILE, and
        # this is the only moment it is on disk.
        font_status = _resolve_template_fonts(report.theme, tmp_path)
    finally:
        os.unlink(tmp_path)

    if not report.ok:
        raise HTTPException(422, "; ".join(report.problems))

    best = report.best
    summary = {
        "layout_index": best.index if best else -1,
        "layout_name": best.name if best else "",
        "palette": report.theme.palette,
        "heading_font": report.theme.heading_font,
        "body_font": report.theme.body_font,
        "slide_width_in": report.slide_width_in,
        "slide_height_in": report.slide_height_in,
        "fonts": font_status,
    }
    try:
        t = repo.upload_template(auth, customer_id, file.filename or "template.pptx",
                                 data, summary)
    except NotFound:
        raise HTTPException(404, f"Customer '{customer_id}' not found") from None
    # Warnings that did not block the upload (e.g. no theme colours) are worth
    # showing: the deck will render, it just will not carry the client's brand.
    # A font we cannot supply is not a reason to refuse the template. The .pptx
    # still NAMES the right font and looks correct on a machine that has it;
    # what degrades is the PDF and the previews, which are rasterised here.
    # Said plainly, because a warning that overstates its case gets ignored.
    font_problems = [f["reason"] for f in summary["fonts"] if not f["ok"]]
    return {**_as_dict(t), "warnings": report.problems + font_problems}


@templates_router.get("/customers/{customer_id}/templates")
def list_templates(customer_id: str, auth: AuthContext = Depends(get_auth),
                   repo: Repository = Depends(get_repository),
                   user: User = Depends(require_customer)) -> list[dict]:
    """List a customer's templates, checking any whose fonts we never checked.

    The backfill is here rather than on the render path for two reasons: a
    template uploaded before the check existed would otherwise never be flagged,
    and this runs when someone opens the template panel — a place where a short
    pause is acceptable and a network round trip during a render is not. It
    happens once per template; the result is stored.
    """
    from reportbuilder.api.routes_settings import load_substitutions

    load_substitutions(repo, auth)
    out = []
    for t in repo.list_templates(auth, customer_id):
        if not t.fonts and t.heading_font:
            fonts = _resolve_template_fonts(t)
            if fonts:
                try:
                    repo.record_template_fonts(auth, customer_id, t.id, fonts)
                except Exception:  # noqa: BLE001 — a cache miss, not a failure
                    pass
                t = replace(t, fonts=tuple(fonts))
        # Re-checked against the host as it is now, so installing a font or
        # choosing a stand-in clears the row instead of leaving it stale.
        families = [f for f in (t.heading_font, t.body_font) if f]
        t = replace(t, fonts=tuple(_live_font_status(list(t.fonts), families)))
        out.append(_as_dict(t))
    return out


@templates_router.get("/customers/{customer_id}/templates/{template_id}")
def template_detail(customer_id: str, template_id: str,
                    auth: AuthContext = Depends(get_auth),
                    repo: Repository = Depends(get_repository),
                    user: User = Depends(require_customer)) -> dict:
    """Everything known about one template, for the settings dialog.

    Font status is RE-RESOLVED rather than served from the stored record: a
    substitution chosen a minute ago has to show as resolved, and a font
    installed since the upload should stop being reported as missing.
    """
    from reportbuilder.api.routes_settings import load_substitutions
    from reportbuilder.render import fonts as F
    from reportbuilder.render import house_style as H

    for t in repo.list_templates(auth, customer_id):
        if t.id != template_id:
            continue
        load_substitutions(repo, auth)
        families = [f for f in (t.heading_font, t.body_font) if f]
        # allow_network=False: opening a dialog must not wait on Google Fonts.
        live = _live_font_status(list(t.fonts), families)
        return {**_as_dict(t), "body_font": t.body_font, "fonts": live,
                "fonts_ok": all(f["ok"] for f in live) if live else True,
                "available_fonts": H.available_chart_fonts()}
    raise HTTPException(404, f"Template '{template_id}' not found")


def _template_on_disk(repo, auth, customer_id: str, template_id: str) -> str:
    """The template's bytes as a file, content-addressed like the preview path."""
    import hashlib
    import os
    import uuid

    from reportbuilder import cache_dirs

    blob = repo.get_template_bytes(auth, customer_id, template_id)
    root = cache_dirs.template_root()
    f = root / f"{template_id}.{hashlib.sha256(blob).hexdigest()[:16]}.pptx"
    if not f.exists():
        tmp = root / f"{f.name}.{uuid.uuid4().hex[:8]}.part"
        tmp.write_bytes(blob)
        os.replace(tmp, f)
    return str(f)


def _in(emu) -> float:
    """EMU to inches, rounded to something a person would type."""
    return round(int(emu or 0) / 914400, 2)


def _content_rect(style, prs) -> tuple[int, int, int, int]:
    """Where a chart will ACTUALLY be drawn on this layout, in EMU.

    The layout's own content placeholder when we are taking it, and otherwise
    the renderer's own placement — under the title, inside the template's side
    margins. Reporting zeroes for the second case is what made most of the
    layouts in the picker show a content box collapsed into the top-left corner:
    "not visible, or invalid coordinates".

    Always inside the slide, and never smaller than an inch. A rectangle an
    author cannot see is one they cannot drag back.
    """
    sw, sh = int(prs.slide_width or 0), int(prs.slide_height or 0)
    slot = getattr(style, "chart_slot", None)
    if slot is not None and int(slot.width or 0) > 0 and int(slot.height or 0) > 0:
        rect = (int(slot.left), int(slot.top), int(slot.width), int(slot.height))
    else:
        profile = getattr(style, "profile", None)
        title = getattr(profile, "title", None) if profile else None
        if title is not None and getattr(title, "positioned", False):
            from reportbuilder.render.image.slide_chrome import harvested_chart_box

            rect = harvested_chart_box(profile, "", sw, sh)
        else:
            # No opinion anywhere: the house placement — a margin in from each
            # side, starting under where a title would sit.
            margin = int(0.05 * sw)
            rect = (margin, int(0.28 * sh), sw - 2 * margin, int(0.60 * sh))

    left, top, width, height = rect
    inch = 914400
    width = max(inch, min(width, sw))
    height = max(inch, min(height, sh))
    left = max(0, min(left, sw - width))
    top = max(0, min(top, sh - height))
    return left, top, width, height


@templates_router.get("/customers/{customer_id}/templates/{template_id}/layout")
def template_layout(customer_id: str, template_id: str, layout: int | None = None,
                    auth: AuthContext = Depends(get_auth),
                    repo: Repository = Depends(get_repository),
                    user: User = Depends(require_customer)) -> dict:
    """What we harvested from this template, what an author has overridden, and
    every layout they could choose instead.

    Harvesting is a guess made from a file nobody wrote for us. On the three
    customer templates we have it was wrong three different ways, so this is
    the surface where an author says what the answer is.
    """
    from pptx import Presentation

    from reportbuilder.render.style_spec import load_style_spec
    from reportbuilder.render.template_check import rank_layouts

    path = _template_on_disk(repo, auth, customer_id, template_id)
    stored = repo.template_layout(auth, customer_id, template_id)
    # Harvested FOR the layout in question: the caller's `layout` while they are
    # trying one in the dropdown, else the one they saved, else ours. Reporting
    # the automatic choice's numbers whatever was selected is what made the
    # dropdown look like it did nothing.
    chosen = layout if layout is not None else stored.get("layout_index")
    style = load_style_spec(path, force_layout=chosen if isinstance(chosen, int) else None)
    prs = Presentation(path)
    profile = getattr(style, "profile", None)
    title = getattr(profile, "title", None)
    slot_left, slot_top, slot_w, slot_h = _content_rect(style, prs)
    from reportbuilder.render import house_style as _H

    title_left = int(getattr(title, "left", 0) or slot_left)
    title_width = int(getattr(title, "width", 0) or slot_w)
    # Where the "n = 100" line really goes, from the function that puts it
    # there — a guess of "near the bottom" drew the box somewhere the text was
    # not. content_floor reads the template's own foot furniture, which is why
    # it needs a slide rather than just a page size.
    try:
        from reportbuilder.export.pptx_build import build_presentation
        from reportbuilder.render.image.slide_chrome import footer_top

        _report, _model, _df = _sample_report()
        _slide = build_presentation(_report, _model, _df, style=style).slides[0]
        footer_y = footer_top(_slide, int(prs.slide_height or 0), int(prs.slide_width or 0))
    except Exception:  # noqa: BLE001 — a template we cannot draw still opens
        footer_y = max(0, int(prs.slide_height or 0) - int(0.80 * 914400))
    ranked = {c.index: c for c in rank_layouts(prs)}
    # OUR pick, independent of what the caller is trying — the star in the
    # picker means "this is what we chose", and reporting the selected one as
    # chosen put a star on every option in turn.
    from collections import Counter

    from reportbuilder.render.template_profile import choose_layout

    usage: Counter = Counter()
    for slide in prs.slides:
        try:
            usage[slide.slide_layout.name] += 1
        except AttributeError:
            continue
    auto_index, _auto_geometry = choose_layout(list(ranked.values()), usage, len(prs.slides))
    slide_area = (int(prs.slide_width or 0) * int(prs.slide_height or 0)) or 1

    return {
        "slide": {"w": _in(prs.slide_width), "h": _in(prs.slide_height)},
        # What is in force for these numbers…
        "chosen_layout": getattr(profile, "layout_index", None),
        # …and what we would have chosen on our own.
        "auto_layout": auto_index,
        "content_is_chart_area": bool(
            getattr(profile, "layout_content_is_chart_area", False)),
        # Only the layouts that could hold a headline and a chart. Arla has 69
        # and 27 of them qualify; offering the other 42 — covers, dividers,
        # photo pages — is offering choices that cannot work.
        "layouts": [
            {"index": i, "name": lay.name,
             "content_pct": round(ranked[i].content_area_pct, 1),
             "suitable": True}
            for i, lay in enumerate(prs.slide_layouts) if i in ranked
        ],
        "harvested": {
            "title": {
                "x": _in(getattr(title, "left", 0)), "y": _in(getattr(title, "top", 0)),
                "w": _in(getattr(title, "width", 0)), "h": _in(getattr(title, "height", 0)),
                "font": getattr(title, "font", "") or "",
                "size": getattr(title, "size_pt", 0) or 0,
                "colour": getattr(title, "colour", "") or "",
            },
            "content": {
                "x": _in(slot_left), "y": _in(slot_top),
                "w": _in(slot_w), "h": _in(slot_h),
                "font": getattr(style, "body_font", "") or "",
                "size": (getattr(style, "fonts", {}) or {}).get("category_names", ("", 0))[1],
                "colour": "",
            },
            # Derived, not placed: the subtitle sits a fixed gap above the
            # chart sharing the title's left and width, and the footer a fixed
            # gap above the template's own foot. Reported so an author can SEE
            # where they land — and restyle them — without being offered a drag
            # that would do nothing.
            "subtitle": {
                "x": _in(title_left), "y": _in(max(0, slot_top - int(0.18 * 914400)
                                                   - int(0.55 * 914400))),
                "w": _in(title_width), "h": 0.55,
                "font": getattr(style, "body_font", "") or "",
                "size": getattr(style, "subtitle_size_pt", 0) or 13,
                "colour": getattr(style, "subtitle_colour", "") or "",
                "derived": True,
            },
            "footer": {
                "x": _in(title_left), "y": _in(footer_y), "w": _in(title_width), "h": 0.35,
                "font": (getattr(style, "fonts", {}) or {}).get("n_annotation", ("", 0))[0],
                "size": (getattr(style, "fonts", {}) or {}).get("n_annotation", ("", 0))[1],
                "colour": getattr(style, "footer_colour", "") or "",
                "derived": True,
            },
            "accent": getattr(style, "accent", "") or "",
            "background": getattr(style, "background", "") or "",
        },
        "available_fonts": _H.available_chart_fonts(),
        "overrides": repo.template_layout(auth, customer_id, template_id),
    }


@templates_router.get("/customers/{customer_id}/templates/{template_id}/ground.png")
def template_ground(customer_id: str, template_id: str, layout: int | None = None,
                    auth: AuthContext = Depends(get_auth),
                    repo: Repository = Depends(get_repository),
                    user: User = Depends(require_customer)) -> Response:
    """The customer's empty slide, as a picture, for the layout editor to draw on.

    The same image a preview is composed onto — their layout, or their harvested
    furniture redrawn — with no chart and none of our text. `layout` asks for a
    DIFFERENT one than we chose, which is what makes the dropdown show its
    effect before anything is saved.
    """
    from io import BytesIO

    from reportbuilder.render.image.fast_preview import ground_image
    from reportbuilder.render.style_spec import load_style_spec

    style = load_style_spec(_template_on_disk(repo, auth, customer_id, template_id))
    if layout is not None:
        style.chart_layout_index = layout
        if getattr(style, "profile", None) is not None:
            style.profile.layout_index = layout
    image = ground_image(style)
    if image is None:
        raise HTTPException(404, "this template has no ground to draw")
    buf = BytesIO()
    image.save(buf, format="PNG")
    return Response(content=buf.getvalue(), media_type="image/png",
                    headers={"Cache-Control": "no-store"})


def _sample_report():
    """A slide's worth of invented data, shaped like a real one.

    Five categories with a plausible spread, a Finnish question and a headline
    of the length these actually run to — enough that an author judging a layout
    is judging what their slides will look like rather than an empty rectangle.
    Invented on purpose: no customer's numbers are needed to show where a title
    sits, and a settings dialog is no place to load a study.
    """
    import pandas as pd

    from reportbuilder.model.question import (
        Question, QuestionModel, ValueLabel, Variable,
    )
    from reportbuilder.model.report import (
        ChartSpec, ElementToggles, NumberFormat, Report, SortSpec,
    )

    labels = ("Erittäin tyytyväinen", "Melko tyytyväinen", "Ei kumpaakaan",
              "Melko tyytymätön", "Erittäin tyytymätön")
    var = Variable(name="q1", label="Esimerkki", measurement="categorical",
                   value_labels=tuple(ValueLabel(float(i + 1), t)
                                      for i, t in enumerate(labels)),
                   missing_values=frozenset())
    model = QuestionModel(
        variables={"q1": var},
        questions=[Question(qid="q1", kind="single", variables=("q1",),
                            text="Kuinka tyytyväinen olet palveluun kokonaisuutena?")])
    counts = (34, 41, 12, 9, 4)
    df = pd.DataFrame({"q1": [float(i + 1) for i, n in enumerate(counts) for _ in range(n)]})
    spec = ChartSpec(
        question_ref="q1", chart_type="horizontal_bar", statistic="pct",
        classifying_var=None, number_format=NumberFormat(),
        sort=SortSpec(basis="data_order"), template_slot="s1",
        elements=ElementToggles(),
        slide_title="Enemmistö on tyytyväinen palveluun ja tyytymättömiä on vain harva",
    )
    return Report(name="sample", render_mode="image", template_ref="",
                  charts=(spec,)), model, df


@templates_router.get("/customers/{customer_id}/templates/{template_id}/sample.png")
def template_sample(customer_id: str, template_id: str, layout: int | None = None,
                    o: str | None = None,
                    auth: AuthContext = Depends(get_auth),
                    repo: Repository = Depends(get_repository),
                    user: User = Depends(require_customer)) -> Response:
    """A whole slide as this template would draw it — furniture, headline,
    question, chart and the n-line — for the layout editor to draw its areas on.

    The same path a preview takes, so what an author sees here is what a slide
    is. An empty ground told them where the boxes were and nothing about whether
    the result reads.
    """
    from io import BytesIO

    from reportbuilder.render.image.fast_preview import compose_from_slide, ground_image
    from reportbuilder.export.pptx_build import build_presentation
    from reportbuilder.render.style_spec import load_style_spec

    path = _template_on_disk(repo, auth, customer_id, template_id)
    stored = repo.template_layout(auth, customer_id, template_id)
    chosen = layout if layout is not None else stored.get("layout_index")
    # `o` is what the author is TRYING — the unsaved draft. Without it the
    # sample answers with what was last saved, so changing a font or a size
    # showed no effect until Save, which reads as the setting doing nothing.
    trying = stored
    if o:
        try:
            candidate = json.loads(o)
            if isinstance(candidate, dict):
                trying = candidate
                chosen = candidate.get("layout_index", chosen)
        except ValueError:
            pass
    style = load_style_spec(path, force_layout=chosen if isinstance(chosen, int) else None)
    from reportbuilder.render.style_spec import apply_template_overrides

    apply_template_overrides(style, trying)

    report, model, df = _sample_report()
    image = None
    try:
        image = compose_from_slide(
            style, build_presentation(report, model, df, style=style).slides[0])
    except Exception:  # noqa: BLE001 — a template we cannot draw still opens
        log.warning("sample slide failed for %s; falling back to the empty ground",
                    template_id, exc_info=True)
    if image is None:
        image = ground_image(style)
    if image is None:
        raise HTTPException(404, "this template has no slide to draw")
    buf = BytesIO()
    image.save(buf, format="PNG")
    return Response(content=buf.getvalue(), media_type="image/png",
                    headers={"Cache-Control": "no-store"})


class TemplateLayoutBody(BaseModel):
    layout_index: int | None = None
    title: dict = {}
    content: dict = {}
    accent: str = ""
    background: str = ""


@templates_router.put("/customers/{customer_id}/templates/{template_id}/layout")
def set_template_layout(customer_id: str, template_id: str,
                        body: TemplateLayoutBody,
                        auth: AuthContext = Depends(get_auth),
                        repo: Repository = Depends(get_repository),
                        user: User = Depends(require_customer_write)) -> dict:
    """Save an author's corrections. Blank fields are dropped, not stored as
    blanks, so "inherit" stays the absence of an opinion rather than a value
    that has to be recognised everywhere it is read."""
    def kept(area: dict) -> dict:
        out = {}
        for key in ("x", "y", "w", "h", "size"):
            value = area.get(key)
            if isinstance(value, (int, float)) and value > 0:
                out[key] = float(value)
        for key in ("font", "colour"):
            text = str(area.get(key) or "").strip()
            if text:
                out[key] = text
        return out

    layout = {k: v for k, v in {
        "layout_index": body.layout_index,
        "title": kept(body.title or {}),
        "content": kept(body.content or {}),
        "accent": str(body.accent or "").strip(),
        "background": str(body.background or "").strip(),
    }.items() if v not in (None, "", {})}
    repo.record_template_layout(auth, customer_id, template_id, layout)
    return {"saved": layout}


@templates_router.get("/customers/{customer_id}/templates/{template_id}/file")
def download_template(customer_id: str, template_id: str,
                      auth: AuthContext = Depends(get_auth),
                      repo: Repository = Depends(get_repository),
                      user: User = Depends(require_customer)) -> Response:
    try:
        return Response(repo.get_template_bytes(auth, customer_id, template_id),
                        media_type=_PPTX)
    except NotFound:
        raise HTTPException(404, f"Template '{template_id}' not found") from None


@templates_router.delete("/customers/{customer_id}/templates/{template_id}")
def delete_template(customer_id: str, template_id: str,
                    auth: AuthContext = Depends(get_auth),
                    repo: Repository = Depends(get_repository),
                    user: User = Depends(require_customer_write)) -> dict:
    return {"removed": repo.delete_template(auth, customer_id, template_id)}


@templates_router.put("/customers/{customer_id}/template")
def bind_customer_template(customer_id: str, body: TemplateBinding,
                           auth: AuthContext = Depends(get_auth),
                           repo: Repository = Depends(get_repository),
                           user: User = Depends(require_customer_write)) -> dict:
    try:
        repo.set_template(auth, body.template_id, customer_id=customer_id)
    except NotFound:
        raise HTTPException(404, f"Customer '{customer_id}' not found") from None
    return {"customer_id": customer_id, "template_id": body.template_id}


@templates_router.put("/customers/{customer_id}/cases/{case_id}/template")
def bind_case_template(customer_id: str, case_id: str, body: TemplateBinding,
                       auth: AuthContext = Depends(get_auth),
                       repo: Repository = Depends(get_repository),
                       user: User = Depends(require_case_in_customer_write)) -> dict:
    try:
        repo.set_template(auth, body.template_id, customer_id=customer_id,
                          case_id=case_id)
    except NotFound:
        raise HTTPException(404, f"Case '{case_id}' not found") from None
    return {"case_id": case_id, "template_id": body.template_id}


@templates_router.put(
    "/customers/{customer_id}/cases/{case_id}/reports/{report_id}/template")
def bind_report_template(customer_id: str, case_id: str, report_id: str,
                         body: TemplateBinding,
                         auth: AuthContext = Depends(get_auth),
                         repo: Repository = Depends(get_repository),
                         user: User = Depends(require_case_in_customer_write)) -> dict:
    """Set (or clear) a template on a single report.

    A report's choice lives in its own definition rather than beside it, because
    template_ref is already part of the report model and the renderer reads it
    from there. Clearing it drops the report back to inheriting from its
    tutkimus, which is what a user who unsets an override expects.

    Also clears any pin: choosing a template deliberately is not the same as
    keeping the one a previous render happened to use.

    Refused while somebody else has the report open. This writes into the
    report's own definition, so it is a change to a document another person is
    editing — and it restyles every slide they are looking at.
    """
    import json as _json

    held = repo._lock_state(auth, customer_id, case_id, report_id)
    if held and held.get("user_id") != getattr(user, "id", ""):
        raise HTTPException(
            409,
            f"{held.get('user_name') or 'Someone else'} is editing this report, "
            f"so its template cannot be changed.")
    try:
        raw = repo.load_report(auth, customer_id, case_id, report_id)
    except NotFound:
        raise HTTPException(404, f"Report '{report_id}' not found") from None
    try:
        report = _json.loads(raw)
    except ValueError:
        raise HTTPException(422, "Report definition is not valid JSON") from None

    report["template_ref"] = body.template_id or ""
    repo.save_report(auth, customer_id, case_id,
                     _json.dumps(report, ensure_ascii=False), report_id=report_id)
    repo.clear_pinned_template(auth, customer_id, case_id, report_id)

    template_id, level = repo.resolve_template(auth, customer_id, case_id, report_id)
    return {"template_id": template_id, "level": level}


@templates_router.get("/customers/{customer_id}/cases/{case_id}/template")
def case_template(customer_id: str, case_id: str,
                  auth: AuthContext = Depends(get_auth),
                  repo: Repository = Depends(get_repository),
                  user: User = Depends(require_case_in_customer)) -> dict:
    """What this tutkimus renders with, and where that came from."""
    template_id, level = repo.resolve_case_template(auth, customer_id, case_id)
    name = ""
    if template_id:
        name = next((t.name for t in repo.list_templates(auth, customer_id)
                     if t.id == template_id), "")
    return {"template_id": template_id, "level": level,
            "name": name or ("nSight default template" if not template_id else template_id)}


@templates_router.get(
    "/customers/{customer_id}/cases/{case_id}/reports/{report_id}/template")
def report_template(customer_id: str, case_id: str, report_id: str,
                    auth: AuthContext = Depends(get_auth),
                    repo: Repository = Depends(get_repository),
                    user: User = Depends(require_case_in_customer)) -> dict:
    """What this report renders with, and WHERE that came from.

    The level is what lets the UI say "inherited from Attendo" rather than
    showing a bare id the user cannot act on.
    """
    template_id, level = repo.resolve_template(auth, customer_id, case_id, report_id)
    name = ""
    if template_id:
        name = next((t.name for t in repo.list_templates(auth, customer_id)
                     if t.id == template_id), "")
    return {"template_id": template_id, "level": level,
            "name": name or ("nSight default template" if not template_id else template_id)}


@templates_router.post(
    "/customers/{customer_id}/cases/{case_id}/reports/{report_id}/template/refresh")
def refresh_report_template(customer_id: str, case_id: str, report_id: str,
                            auth: AuthContext = Depends(get_auth),
                            repo: Repository = Depends(get_repository),
                            user: User = Depends(require_case_in_customer_write)) -> dict:
    """The card's "päivitys pitää erikseen pyytää": move a delivered report onto
    whatever its tutkimus or asiakas now specifies."""
    repo.clear_pinned_template(auth, customer_id, case_id, report_id)
    template_id, level = repo.resolve_template(auth, customer_id, case_id, report_id)
    return {"template_id": template_id, "level": level}
