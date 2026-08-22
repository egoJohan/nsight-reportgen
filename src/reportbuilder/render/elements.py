"""Chart element profile: apply title/axes/legend/data-labels and N/filter annotations.

REQ-C-24a..i, C-25 (Task 5.3).
"""
from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.enum.chart import XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor

from reportbuilder.model.report import NumberFormat
from reportbuilder.render.base import RenderContext
from reportbuilder.render.house_style import furniture_colors
from reportbuilder.render.panels import panel_segments
import reportbuilder.stats.statistics  # noqa: F401 — ensure built-in registrations are loaded
import reportbuilder.stats.registry as _registry


def _native_ink(ctx: RenderContext) -> RGBColor:
    """Text colour for a NATIVE (python-pptx) chart's legend/data-label/axis
    text, derived from ctx's slide background rather than left to inherit the
    theme — the theme's own text colour can be dark, which on a dark template
    is exactly the invisible-on-dark defect this exists to fix. INK on a light
    background (today's default, unchanged); white on a dark one — same rule
    the image-mode builders apply via house_style.furniture_colors."""
    ink_hex, _muted, _grid = furniture_colors(getattr(ctx.style, "background", "") or "")
    return RGBColor.from_string(ink_hex.lstrip("#"))


# ---------------------------------------------------------------------------
# Number format
# ---------------------------------------------------------------------------

def number_format_code(fmt: NumberFormat, statistic: str) -> str:
    """Return an Excel-style number format string for use in pptx data labels.

    Delegates to the statistic registry so new statistics need no code here.
    """
    return _registry.statistic(statistic).fmt_code(fmt)


# ---------------------------------------------------------------------------
# Element application
# ---------------------------------------------------------------------------

def apply_elements(chart, ctx: RenderContext, title: str = "") -> None:
    """Apply chart element profile (title, data labels, legend, axis names) to *chart*.

    Each element is gated by the corresponding flag in ctx.spec.elements.
    Chart-type-specific failures (e.g. pie lacks value_axis) are silently skipped.
    """
    elements = ctx.spec.elements

    # --- Title ---
    if elements.title:
        chart.has_title = True
        tf = chart.chart_title.text_frame
        tf.text = title
        font_name, font_size = ctx.style.font_for("title")
        runs = tf.paragraphs[0].runs
        if runs:
            run = runs[0]
            run.font.name = font_name
            run.font.size = Pt(font_size)

    # --- Data labels ---
    if elements.data_labels:
        plot = chart.plots[0]
        try:
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.number_format = number_format_code(ctx.fmt, ctx.spec.statistic)
            dl.number_format_is_linked = False
            try:
                dl.position = XL_LABEL_POSITION.OUTSIDE_END
            except (ValueError, AttributeError):
                pass
            font_name, font_size = ctx.style.font_for("data_labels")
            dl.font.name = font_name
            dl.font.size = Pt(font_size)
            dl.font.color.rgb = _native_ink(ctx)
        except AttributeError:
            # Some plot types (e.g. XyPlot / CT_ScatterChart) don't support dLbls
            # in python-pptx; skip silently.
            pass

    # --- Legend ---
    if elements.legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        font_name, font_size = ctx.style.font_for("legend")
        chart.legend.font.name = font_name
        chart.legend.font.size = Pt(font_size)
        chart.legend.font.color.rgb = _native_ink(ctx)

    # --- Axis names / tick labels ---
    if elements.axis_names:
        try:
            vfont_name, vfont_size = ctx.style.font_for("axis_values")
            chart.value_axis.tick_labels.font.name = vfont_name
            chart.value_axis.tick_labels.font.size = Pt(vfont_size)
            chart.value_axis.tick_labels.font.color.rgb = _native_ink(ctx)

            cfont_name, cfont_size = ctx.style.font_for("category_names")
            chart.category_axis.tick_labels.font.name = cfont_name
            chart.category_axis.tick_labels.font.size = Pt(cfont_size)
            chart.category_axis.tick_labels.font.color.rgb = _native_ink(ctx)
        except (AttributeError, ValueError):
            # pie / doughnut / radar / scatter have no value_axis or category_axis;
            # python-pptx raises ValueError("chart has no value axis") for those types.
            pass


# ---------------------------------------------------------------------------
# Annotation helpers
# ---------------------------------------------------------------------------

def add_n_annotation(ctx: RenderContext) -> None:
    """Add a slide textbox near the slot bottom showing N=<Total>.

    Only added when ctx.spec.elements.n is True.
    """
    if not ctx.spec.elements.n:
        return

    slot = ctx.slot
    # Position: horizontally aligned to slot left, vertically just below slot bottom
    left = slot.left
    top = slot.top + slot.height - int(Inches(0.35))
    width = int(Inches(2.0))
    height = int(Inches(0.3))

    txBox = ctx.slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    # The SeriesResult contract guarantees base_n always contains "Total".
    # A KeyError here correctly signals a contract violation rather than hiding it.
    base_n = ctx.series.base_n["Total"]
    tf.text = f"N={base_n}"

    font_name, font_size = ctx.style.font_for("n_annotation")
    run = tf.paragraphs[0].runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)


# Only these chart types draw one panel per group and cap at three; every other
# chart type draws all its groups, so an omission clause there would be a lie.
_PANEL_CHART_TYPES = ("pie", "doughnut", "funnel")


def _omission_clause(ctx) -> str:
    """The footer's record of every classifier group the slide did NOT draw.

    The editor's warning stays in the editor; this line travels with the deck, so
    it is the authoritative account of what was omitted. The two reasons are kept
    apart because they mean different things to a reader: a group omitted for a
    thin base could not be reported at all, while a capped group fits the data but
    not the page. (spec 2026-08-22)

    Only pie/doughnut/funnel draw one panel per group and cap at three; every
    other chart type draws all its groups, so an omission clause there would be a
    false statement printed on a client slide. (ruling 2026-08-22)
    """
    if getattr(ctx.spec, "chart_type", "") not in _PANEL_CHART_TYPES:
        return ""
    sel = panel_segments(ctx.series)
    if not sel.split:
        return ""
    if sel.degraded:
        return " · Ryhmittelyä ei voitu piirtää"
    parts = []
    if sel.thin:
        parts.append("Ei raportoitu: " + ", ".join(sel.thin))
    if sel.capped:
        parts.append("Ei mahtunut sivulle: " + ", ".join(sel.capped))
    return (" · " + " · ".join(parts)) if parts else ""


def add_filter_annotation(ctx: RenderContext) -> None:
    """Add a slide textbox naming the classifying variable.

    Only added when ctx.spec.elements.filter_var is True AND
    ctx.spec.classifying_var is not None.
    """
    if not ctx.spec.elements.filter_var:
        return
    if not ctx.spec.classifying_var:
        return

    slot = ctx.slot
    # Position: below the N annotation (or at slot bottom if no N)
    left = slot.left + int(Inches(2.1))
    top = slot.top + slot.height - int(Inches(0.35))
    width = int(Inches(5.5))
    height = int(Inches(0.3))

    txBox = ctx.slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    # In the SEPARATE layout the slide is split by BOTH variables, side by side;
    # naming only the first would misdescribe the chart. (spec 2026-08-04)
    opts = getattr(ctx.spec, "options", None) or {}
    cv2 = getattr(ctx.spec, "classifying_var_2", None)
    if cv2 and opts.get("xtab_layout") == "separate":
        tf.text = f"{ctx.spec.classifying_var} · {cv2}"
    else:
        tf.text = f"{ctx.spec.classifying_var}{_omission_clause(ctx)}"

    font_name, font_size = ctx.style.font_for("filter_var")
    run = tf.paragraphs[0].runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
