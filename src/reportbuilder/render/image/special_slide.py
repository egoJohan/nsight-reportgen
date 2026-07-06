"""Render a "special" (non-chart) slide: Overview / Conclusion / Demographics.

Special slides carry no data series — their content is a heading (``slide_title``)
plus a list of bullet strings in ``spec.options["bullets"]``. They are drawn as
plain PowerPoint textboxes (house style: cream background, teal accent bar, bold
ink heading, muted bullet list), entirely independent of the chart pipeline —
no RenderContext, no SeriesResult, no plugin dispatch.
"""
from __future__ import annotations

import re

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from reportbuilder.model.report import ChartSpec
from reportbuilder.render.house_style import PX_CREAM, PX_INK, PX_TEAL
from reportbuilder.render.image.slide_chrome import _FONT, _slide_dims, TITLE_PT


# Inline markdown: **bold** / __bold__ and *italic* / _italic_ (non-nested).
_MD_RE = re.compile(r"(\*\*|__)(.+?)\1|(\*|_)(.+?)\3")


def _md_runs(text: str) -> list[tuple[str, bool, bool]]:
    """Split inline markdown into (text, bold, italic) runs so slide bullets can
    render emphasis. Returns the plain text as a single run when no markers."""
    runs: list[tuple[str, bool, bool]] = []
    pos = 0
    for m in _MD_RE.finditer(text):
        if m.start() > pos:
            runs.append((text[pos:m.start()], False, False))
        if m.group(1):  # **bold** / __bold__
            runs.append((m.group(2), True, False))
        else:  # *italic* / _italic_
            runs.append((m.group(4), False, True))
        pos = m.end()
    if pos < len(text):
        runs.append((text[pos:], False, False))
    return runs or [(text, False, False)]


def render_special_slide(slide, slot, style, spec: ChartSpec, heading: str = "") -> None:
    """Paint a heading + bullet list onto *slide* (house style).

    The heading is ``spec.slide_title`` when set, else the ``heading`` fallback
    (used by a themes slide, whose heading is the open-ended question text)."""
    sw, sh = _slide_dims(slide)

    # 1 — Cream background (full slide)
    bg = slide.shapes.add_shape(1, 0, 0, sw, sh)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PX_CREAM
    bg.line.fill.background()
    bg.shadow.inherit = False

    # 2 — Heading text (compute first so the accent bar can match its height).
    heading_text = (getattr(spec, "slide_title", None) or heading or "").strip()

    # 3 — Teal accent bar (top-left), sized to the heading's actual line height so
    #     it doesn't tower over a short one-line title (capped at the box height).
    _hlines = _heading_line_count(heading_text, sw) if heading_text else 1
    _line_h = Pt(_heading_size(heading_text or "x") * 1.25)
    bar_h = min(int(Inches(0.92)), _hlines * int(_line_h) + int(Inches(0.06)))
    acc = slide.shapes.add_shape(
        1, Inches(0.55), Inches(0.42), Inches(0.10), bar_h
    )
    acc.fill.solid()
    acc.fill.fore_color.rgb = PX_TEAL
    acc.line.fill.background()
    acc.shadow.inherit = False

    # 4 — Heading (slide_title, else the fallback — e.g. the question text)
    if heading_text:
        _heading_box(slide, sw, heading_text)

    # 4 — Bullet list. Each raw line is a markdown bullet: leading whitespace sets
    # the nesting level and a leading -,*,+,• marker is stripped. Tolerate a bare
    # string (don't iterate it into characters).
    raw = spec.options.get("bullets") or []
    if isinstance(raw, str):
        raw = [raw]
    # Drop degenerate "odd" bullets — empties, markdown code-fence lines
    # ("```question:yes_no"), or lines that are only markers / punctuation with no
    # real letters (defence in depth, in case options carry junk from a saved report).
    parsed: list[tuple[int, str]] = []
    for item in raw:
        for line in str(item).split("\n"):
            level, text = _bullet_level(line)
            if (text
                    and not text.startswith("```")
                    and not text.startswith("~~~")
                    and re.search(r"[^\s\-•*_:.,–—]", text)):
                parsed.append((level, text))
    if parsed:
        _bullet_box(slide, sw, sh, parsed)


def _heading_size(text: str) -> int:
    """The special-slide heading is a slide title just like a chart's key message,
    so it uses the one shared fixed title size (slide_chrome.TITLE_PT)."""
    return TITLE_PT


def _heading_line_count(text: str, sw: int) -> int:
    """Approximate how many lines the heading wraps to in its box (width sw - 1.0"),
    so the accent bar can match the title's actual height instead of a fixed box."""
    if not text:
        return 1
    size = _heading_size(text)
    box_pt = (sw / 914400 - 1.0) * 72              # heading box width in points
    chars_per_line = max(1, int(box_pt / (size * 0.55)))  # ~0.55·size pt per avg char
    lines = 0
    for seg in text.split("\n"):
        lines += max(1, -(-len(seg) // chars_per_line))   # ceil-divide
    return max(1, lines)


def _heading_box(slide, sw, text: str) -> None:
    tb = slide.shapes.add_textbox(
        Inches(0.80), Inches(0.42), sw - Inches(1.0), Inches(0.92)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(_heading_size(text))
    r.font.bold = True
    r.font.color.rgb = PX_INK
    r.font.name = _FONT


def _bullet_level(raw: str) -> tuple[int, str]:
    """(nesting level, text) for one raw markdown bullet line.

    Leading whitespace sets the level — a tab or two spaces per level, capped at 3 —
    and a leading list marker (-, *, +, •, ·, ◦) is stripped. So "  * Foo" → (1, "Foo")
    and "- Bar" → (0, "Bar"). A plain line (no marker) is a level-0 bullet."""
    expanded = raw.replace("\t", "  ")
    body = expanded.lstrip(" ")
    lead = len(expanded) - len(body)
    level = min(lead // 2, 3)
    text = re.sub(r"^[-*+•·◦]\s+", "", body).strip()
    return level, text


# Per-level bullet glyph + body font size (deeper levels are smaller / softer).
_LEVEL_GLYPH = {0: "•", 1: "–", 2: "·", 3: "·"}
_LEVEL_PT = {0: 16, 1: 14, 2: 13, 3: 13}


def _bullet_box(slide, sw, sh, bullets: list[tuple[int, str]]) -> None:
    tb = slide.shapes.add_textbox(
        Inches(0.85), Inches(1.55), sw - Inches(1.6), sh - Inches(2.1)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    # Hanging indent: the bullet glyph sits at the paragraph's left margin and
    # wrapped lines align EXACTLY under the first line's text. Nesting shifts the
    # whole paragraph right by _STEP per level; marL is the text start (offset by
    # level), indent = -_HANG pulls the glyph back, and a left TAB STOP at marL
    # snaps the first line's text to the same x as its wrapped continuation lines.
    _HANG = Inches(0.30)
    _STEP = Inches(0.34)
    for i, (level, text) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10 if level == 0 else 4)
        pPr = p._p.get_or_add_pPr()
        mar_l = int(_HANG) + level * int(_STEP)
        pPr.set("marL", str(mar_l))
        pPr.set("indent", str(-int(_HANG)))
        tab_lst = pPr.makeelement(qn("a:tabLst"), {})
        tab_lst.append(pPr.makeelement(qn("a:tab"), {"pos": str(mar_l), "algn": "l"}))
        pPr.append(tab_lst)
        pt = _LEVEL_PT.get(level, 13)
        # Teal bullet glyph + tab (snaps body text to the marL tab stop).
        dot = p.add_run()
        dot.text = f"{_LEVEL_GLYPH.get(level, '·')}\t"
        dot.font.size = Pt(pt)
        dot.font.bold = True
        dot.font.color.rgb = PX_TEAL
        dot.font.name = _FONT
        for seg, bold, italic in _md_runs(text):
            body = p.add_run()
            body.text = seg
            body.font.size = Pt(pt)
            body.font.bold = bold
            body.font.italic = italic
            body.font.color.rgb = PX_INK
            body.font.name = _FONT
