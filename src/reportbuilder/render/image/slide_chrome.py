"""House-style slide chrome for image-mode slides (REQ-C-24/25/27a, REQ-D-04).

`add_image_slide_chrome` adds the slide-level decorations — cream background,
teal accent bar, bold title, N annotation, and methodology footer — to the
slide *before* the chart picture is placed by the image builder.  Because
shapes are z-ordered by insertion order, the chart picture (added last) lands
on top of the chrome.

The function is intentionally generic: it reads title, statistic, and base-N
from the RenderContext (driven by ChartSpec + SeriesResult) and never
hard-codes any Attendo-specific content.

Slide-text polish (R2):
- Title = full question text (ctx.title), word-wrapped to 2 lines if long.
  (REQ-D-04)
- Methodology footer bottom-left: statistic label + "· n = N"  (REQ-C-24h)
  e.g. "Osuus vastaajista (%) · n = 1001"
"""
from __future__ import annotations

import logging

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from reportbuilder.render.base import RenderContext
from reportbuilder.render.elements import _omission_clause
from reportbuilder.render.template_profile import clone_furniture
from reportbuilder.render.house_style import (
    _relative_luminance,
    PX_CREAM, PX_INK, PX_TEAL, PX_MUTED, furniture_colors,
)


def _rgb(hex6: str):
    """'122D49' -> RGBColor. Empty/short input returns None."""
    from pptx.dml.color import RGBColor
    if not hex6 or len(hex6) != 6:
        return None
    try:
        return RGBColor(int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16))
    except ValueError:
        return None


def _furniture_px(style) -> tuple:
    """(ink, muted) RGBColor for THIS slide's own background — house_style.
    furniture_colors, converted from hex to pptx RGBColor. Same rule the image
    chart builders use (chart_furniture in render/image/_mpl.py), so the
    slide's own title/subtitle/footer text agrees with its charts about what
    "dark" means. A light/unstated background resolves to exactly PX_INK/
    PX_MUTED, unchanged."""
    from reportbuilder.render.resolved_style import furniture
    ink_hex, muted_hex, _grid = furniture(style)
    return (_rgb(ink_hex.lstrip("#")) or PX_INK,
            _rgb(muted_hex.lstrip("#")) or PX_MUTED)


def theme_colours(style):
    """(background, ink, accent) for a slide, preferring the template's own.

    A deck built on the client's template should not carry nSight's cream ground
    and teal accent bar — down to the bullet glyphs on a conclusion slide, which
    were house teal on an Attendo navy deck. Anything the template does not state
    falls back to the house value, so a template with a partial theme still
    renders.

    Ink the template doesn't state is derived from the background via
    `_furniture_px` (house_style.furniture_colors) rather than the fixed house
    PX_INK, so a dark background with no stated ink still gets legible
    near-white text instead of near-black on near-black. A light background —
    including the house default, an unstated one — resolves to exactly PX_INK,
    unchanged.
    """
    bg = _rgb(getattr(style, "background", "") or "") or PX_CREAM
    ink = _rgb(getattr(style, "ink", "") or "") or _furniture_px(style)[0]
    accent = PX_TEAL
    try:
        if getattr(style, "from_template", False):
            # style.accent is the template's brand accent, or the colour its own
            # slides are drawn in when its theme states no brand.
            accent = _rgb(getattr(style, "accent", "") or "") or _rgb(style.color_for(0)) or PX_TEAL
    except Exception:  # noqa: BLE001 — styling must not break a render
        pass
    return bg, ink, accent


def _theme_colours(ctx):
    return theme_colours(getattr(ctx, "style", None))


from reportbuilder.stats.engine import scale_endpoint_gloss

_FONT = "Liberation Sans"
# One fixed title size for EVERY slide (chart + special) so titles never vary in
# size between slides. (Shared by special_slide's heading.)
TITLE_PT = 18
_IN = Inches(1)

_STACKED_BAR_TYPES = frozenset({"stacked_horizontal_bar", "stacked_vertical_bar"})


# `scale_endpoint_gloss` lives in stats.engine so the questions API can offer the SAME
# default text to the frontend's Subtitle box (the subtitle owns the whole line).

# Statistic → Finnish methodology label (generic; no question-specific text)
_STAT_FOOTER: dict[str, str] = {
    "pct": "Osuus vastaajista (%)",
    "count": "Lukumäärä",
    "mean": "Keskiarvo",
    "median": "Mediaani",
    "sum": "Summa",
}


def _slide_dims(slide) -> tuple[int, int]:
    """Return (slide_width_emu, slide_height_emu) from the slide's parent presentation."""
    try:
        prs = slide.part.package.presentation_part.presentation
        return int(prs.slide_width), int(prs.slide_height)
    except Exception:
        # fallback: 13.333" × 7.5" (the deck default — 16:9 widescreen)
        return int(Inches(13.333)), int(Inches(7.5))


def wrapped_line_count(text: str, box_width_emu: int, size_pt: int) -> int:
    """Approximate how many lines *text* wraps to in a box *box_width_emu* wide at
    font *size_pt* (honours explicit '\\n'). Used to size the accent bar to the
    title/heading's actual height instead of a fixed box."""
    if not text:
        return 1
    box_pt = box_width_emu / 914400 * 72
    chars_per_line = max(1, int(box_pt / (size_pt * 0.55)))  # ~0.55·size pt per avg char
    lines = 0
    for seg in text.split("\n"):
        lines += max(1, -(-len(seg) // chars_per_line))       # ceil-divide
    return max(1, lines)


def _textbox(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font: str = ""):
    """Add a multi-run textbox.  `runs` is a list of (text, pt_size, rgb, bold) tuples."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    font_name = font or _FONT
    for txt, sz, col, bold in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = col
        r.font.name = font_name
    return tb


def body_font(style) -> str:
    """The typeface for text nSight draws on a templated slide.

    A subtitle in our font next to a title in the customer's is the kind of
    mismatch that reads as sloppy rather than as a missing font, so the
    template's own body face is used when it has one.

    Note this only affects how the .pptx NAMES the font. Whether it renders
    that way depends on the machine opening it — see the module docstring in
    render.fonts.
    """
    return getattr(style, "body_font", "") or ""


def _body_font(ctx) -> str:
    return body_font(getattr(ctx, "style", None))


def _from_template(ctx) -> bool:
    """True when this slide was built from the template's OWN chart layout.

    Then the layout already carries the customer's design — background, brand
    furniture, title styling — and anything we paint on top hides what they
    asked us to use. nSight contributes text and the chart, nothing else.
    """
    return getattr(ctx.style, "chart_layout_index", None) is not None


def harvested_profile(style):
    """The profile to draw this slide from, or None.

    Only when the design was harvested off a SLIDE. With a layout there is
    nothing to draw: the slide is built from that layout and inherits it.
    """
    profile = getattr(style, "profile", None)
    if profile is None or getattr(profile, "layout_index", None) is not None:
        return None
    if not (profile.furniture or profile.title.positioned):
        return None
    return profile


def slide_headline(spec, question: str) -> str:
    """The text that goes in the title — an authored headline, else the question.

    Shared with deck.py, which has to know how tall the title will be before it
    can decide where the chart starts.
    """
    return (getattr(spec, "slide_title", None) or "").strip() or (question or "").strip()


# How far the title box may grow before the type shrinks instead of the box.
# Twice the customer's own title height is roughly enough for our text to add
# one full wrapped line beyond theirs; past that, growing the box further is
# what produced the "elephant" title regression (a headline eating the slide),
# so the font steps down instead. A fixed multiple of the TEMPLATE's own
# height (rather than a fraction of the slide) keeps the rule proportional to
# how much room the customer actually drew for a title, whether that's a
# generous banner or a tight one-liner.
# The title may grow past the box the customer drew — their headline fits on one
# line because they wrote it to, ours is a survey question and wraps — but only
# this far. Past it the type scales down instead. At 2.0 a four-line question
# still "fit" and swallowed the slide, which is the complaint this bound exists
# to answer: there is a maximum height, and the title scales into it.
_TITLE_MAX_GROWTH = 1.35

# The font never steps down past this share of the template's own size. Below
# it the deck looks broken in a different way than an oversized title did, and
# the honest fix at that point is a shorter headline, not smaller type.
_TITLE_MIN_SCALE = 0.6

# Increment used to step the title size down. Coarse enough that the rendered
# result matches what wrapped_line_count predicted — a finer step just adds
# iterations without changing the outcome, since the line-count estimate is
# itself an approximation.
_TITLE_STEP_PT = 0.5


def _title_line_height(size_pt: float, st) -> int:
    """The height of one wrapped line of the title, at *size_pt*."""
    return int(Pt(size_pt * (st.line_spacing or 1.25)))


def fit_title_size(st, text: str, max_growth: float | None = None) -> float:
    """The size to actually render *text* at in title style *st*.

    Starts at the template's own size — the common case, a customer-length
    title, never needs to move off it. While the wrapped text is taller than
    `_TITLE_MAX_GROWTH` × the template's own box, the size steps down by
    `_TITLE_STEP_PT`, stopping at `_TITLE_MIN_SCALE` of the template's size
    even if the text still does not fit there.
    """
    base = st.size_pt or TITLE_PT
    if not text or not st.width or not st.height:
        return base
    floor = base * _TITLE_MIN_SCALE
    # `max_growth` lets a caller say how far past the template's own box the
    # headline may run. The harvested TEXTBOX may overrun a little (1.35) — it
    # owns the space under it. A layout PLACEHOLDER may not: the slide below it
    # is the template's, so a title that outgrows its box takes room that
    # belongs to the subtitle and the chart, and reads as oversized because it
    # is. See the caller in _fill_title_placeholder.
    max_height = int(st.height) * (
        _TITLE_MAX_GROWTH if max_growth is None else max_growth)
    size = base
    while True:
        # An all-caps title wraps sooner than the same string in mixed case,
        # and Holiday Club's is `cap="all"` — measure it as if it were a size
        # larger.
        # Measured, not estimated: the criterion is how much SPACE the text
        # takes, and uppercase is measured as uppercase rather than approximated
        # by inflating the size 12%.
        lines = measured_line_count(text, int(st.width), size, st)
        height = lines * _title_line_height(size, st)
        if height <= max_height or size <= floor:
            break
        size -= _TITLE_STEP_PT
    return max(size, floor)


def measured_line_count(text: str, box_width_emu: int, size_pt: float, st) -> int:
    """How many lines *text* really wraps to, measured with the host's own font.

    `wrapped_line_count` divides by an average character width (0.55em), which
    is a fine approximation for sizing an accent bar and a poor one for deciding
    whether a title fits: "WWW" and "iii" are the same length to it, and the
    caps case had to be fudged with a 1.12 multiplier. Fitting is about the
    SPACE the text takes, so measure the space — with the same PIL font
    resolution the preview compositor uses, so the two agree about wrapping.

    Falls back to the approximation when no font can be resolved (a headless
    host with no fontconfig match); a rough answer beats no title at all.
    """
    if not text:
        return 1
    body = text.upper() if getattr(st, "caps", False) else text
    try:
        from reportbuilder.render.image.fast_preview import _font

        # Any DPI works — only the ratio of text width to box width matters.
        px_per_pt = 96 / 72
        font = _font(getattr(st, "font", "") or "", size_pt * px_per_pt,
                     bold=bool(getattr(st, "bold", False)))
        if font is None:
            raise RuntimeError("no font")
        box_px = box_width_emu / 914400 * 96
        if box_px <= 0:
            raise RuntimeError("no box")
        lines = 0
        for para in body.split("\n"):
            words, cur = para.split(), ""
            if not words:
                lines += 1
                continue
            n = 1
            for w in words:
                trial = f"{cur} {w}".strip()
                if font.getlength(trial) <= box_px or not cur:
                    cur = trial
                else:
                    n += 1
                    cur = w
            lines += n
        return max(1, lines)
    except Exception:  # noqa: BLE001 — measurement is best-effort, never fatal
        return wrapped_line_count(text, box_width_emu, int(round(size_pt)))


def harvested_title_box(profile, text: str) -> tuple[int, int, int, int]:
    """The title's box, at the size it will actually render.

    The customer's own title fits on one line because they wrote it to. Ours is
    a question or an AI headline and often wraps; rather than growing the box
    to fit it without limit (that produced titles that ate the slide), the type
    shrinks — see `fit_title_size` — and the box grows only up to
    `_TITLE_MAX_GROWTH` × the template's own height, for the cases a shrunk
    title still does not fit in the template's exact box.
    """
    st = profile.title
    size = fit_title_size(st, text)
    measured = size * 1.12 if st.caps else size
    lines = wrapped_line_count(text, int(st.width), int(measured))
    wanted = lines * _title_line_height(size, st)
    max_height = int(st.height) * _TITLE_MAX_GROWTH if st.height else wanted
    height = min(max(int(st.height), wanted), int(max_height))
    return int(st.left), int(st.top), int(st.width), height


def content_floor(slide, sw: int, sh: int) -> int:
    """The lowest point our own content may reach on *slide*.

    A template's furniture usually includes something at the foot of the slide,
    and it arrives by inheritance so it is invisible to `slide.shapes`. Synsam's
    master puts its logo at 6.73in on a 7.5in slide, exactly where the "N = ..."
    footer goes, and the two printed on top of each other. A full-slide backdrop
    is not a floor, and neither is anything in the right half — the footer is
    left-aligned, and Attendo's brand icons live bottom-right.
    """
    floor = sh
    layout = getattr(slide, "slide_layout", None)
    for container in (layout, getattr(layout, "slide_master", None)):
        if container is None:
            continue
        try:
            shapes = list(container.shapes)
        except Exception:  # noqa: BLE001 — an unreadable layout sets no floor
            continue
        for shape in shapes:
            try:
                if shape.is_placeholder:
                    continue
                left, top = int(shape.left or 0), int(shape.top or 0)
                width, height = int(shape.width or 0), int(shape.height or 0)
            except (TypeError, ValueError, AttributeError):
                continue
            if width <= 0 or height <= 0:
                continue
            if width * height >= 0.75 * sw * sh:
                continue
            if top + height < sh * 0.8 or left > sw * 0.5:
                continue
            floor = min(floor, top)
    return floor


def footer_top(slide, sh: int, sw: int) -> int:
    """Where the methodology footer's box starts: above the template's own foot."""
    return content_floor(slide, sw, sh) - int(Inches(0.45))


def harvested_chart_box(profile, text: str, sw: int, sh: int,
                        floor: int | None = None) -> tuple[int, int, int, int]:
    """Where the chart goes on a harvested slide.

    Under the title, and inside the template's own side margins — Johan's rule
    is that the TITLE follows the template and nSight positions the rest, but a
    chart 0.62in from the edge below a title 0.70in from the edge reads as a
    mistake. The gap under the title is the subtitle's room, and the gap under
    the chart is the footer's.
    """
    left, top, width, height = harvested_title_box(profile, text)
    chart_top = top + height + int(Inches(0.70))
    bottom = (sh if floor is None else floor) - int(Inches(0.70))
    return left, chart_top, width, max(int(Inches(1.0)), bottom - chart_top)


def template_ground(slide, style) -> bool:
    """Does the TEMPLATE supply this slide's background? Redraws its furniture.

    True means the house cream ground and teal accent bar must not be painted:
    either the slide was built from the customer's layout and inherits their
    design, or their design was harvested off a slide and is redrawn here.
    Every kind of slide asks this — chart, bullet list, demographics grid — so
    a deck does not come out half in the customer's design and half in ours.
    """
    profile = harvested_profile(style)
    if profile is not None:
        clone_furniture(slide, profile.furniture)
        return True
    return getattr(style, "chart_layout_index", None) is not None


# The subtitle always sits this far above the chart. One constant, so the gap is
# the same on every slide whatever the question's length — the thing that made
# the old placement look arbitrary was that the gap moved with the text.
_SUBTITLE_GAP = Inches(0.18)
# How close to the top edge a headline may start. Small on purpose: the band
# above a title is dead space, and a two-line headline needs the room.
_MIN_TITLE_TOP = int(Inches(0.30))
# How much taller a rendered line really is than PIL measures it. LibreOffice
# renders the deck and lays Bebas Neue out this much taller; measured on
# Egoiq_x_Rahoo, where a title the box said ended at 1.34" was drawing at 1.55".
# One constant, used BOTH to reserve height and to fit the size, so the two
# cannot drift apart.
_RENDERER_LINE_BOX = 1.18
# Room for the question to grow UPWARD into. Four lines at the largest step;
# the box is bottom-anchored, so unused height is invisible.
_SUBTITLE_MAX_H = Inches(1.10)


def title_left_width(slide, style, title: str) -> tuple[int, int]:
    """(left, width) of the TITLE on this slide, so the subtitle can line up
    with it rather than with the chart.

    A subtitle indented differently from the headline above it reads as a
    mistake — and the two were indeed placed from different boxes: the title
    from the template's placeholder, the subtitle from the chart's slot.
    """
    try:
        ph = slide.shapes.title
        if ph is not None and ph.width:
            # The placeholder keeps its own insets while `_textbox` zeroes
            # every margin, so matching the BOX edges left the subtitle about
            # 0.1" to the left of the headline — visibly out of line. Match
            # where the TEXT starts instead.
            inset_l = int(ph.text_frame.margin_left or 0)
            inset_r = int(ph.text_frame.margin_right or 0)
            return int(ph.left or 0) + inset_l, max(1, int(ph.width) - inset_l - inset_r)
    except (AttributeError, KeyError):
        pass
    profile = harvested_profile(style)
    if profile is not None and profile.title.positioned:
        left, _top, width, _h = harvested_title_box(profile, title)
        return int(left), int(width)
    return 0, 0


def fit_subtitle_size(text: str, width_emu: int, height_emu: int, font: str,
                      *, max_pt: float = 15.0, min_pt: float = 11.0) -> float:
    """The largest size *text* fits its box at, between *min_pt* and *max_pt*.

    Replaces a ladder keyed on character count (15pt to 110 chars, then 13, 12,
    11), which could only guess: it knew nothing about the width it had or the
    height it was given, so on a template with a roomy header it shrank a
    question that would have fitted comfortably — Attendo's came out at 12pt in
    a box with space for more.

    Measured with the same font resolution and the same renderer line box the
    title uses, so the two agree about what fits.
    """
    from types import SimpleNamespace

    if not text or width_emu <= 0 or height_emu <= 0:
        return max_pt
    st = SimpleNamespace(size_pt=max_pt, width=width_emu, height=height_emu,
                         font=font, caps=False, line_spacing=0.0)
    size = max_pt
    while size > min_pt:
        lines = measured_line_count(text, width_emu, size, st)
        if lines * _title_line_height(size, st) * _RENDERER_LINE_BOX <= height_emu:
            break
        size -= 0.5
    return max(size, min_pt)


def _spec_title_pt(style, font: str) -> float:
    """The title size this template's spec states, or 0.0 if it cannot say.

    Read, not measured: `template_cache.resolve` settles the whole spec once per
    template, before any slide is built. The measuring path below is the
    fallback for a style that did not come through it — a house style, or a test
    building one directly.
    """
    spec = getattr(style, "resolved_spec", None)
    if spec is not None:
        return spec.title.size_pt
    try:
        from reportbuilder.render.resolved_style import build_spec
        return build_spec(style, title_font=font).title.size_pt
    except Exception:  # noqa: BLE001 — never fail a render over a font metric
        return 0.0


def _spec_subtitle_pt(style, font: str) -> float:
    """The subtitle size this template's spec states, or 0.0."""
    spec = getattr(style, "resolved_spec", None)
    if spec is not None:
        return spec.subtitle.size_pt
    try:
        from reportbuilder.render.resolved_style import build_spec
        return build_spec(style, subtitle_font=font).subtitle.size_pt
    except Exception:  # noqa: BLE001
        return 0.0


def _inherited_title_font(ph) -> str:
    """The face the template's own chain gives this title, if any."""
    try:
        from reportbuilder.render.image.fast_preview import (
            _inherited_placeholder_style,
        )
        return _inherited_placeholder_style(ph)[0]
    except Exception:  # noqa: BLE001
        return ""


def _rendered_title_height(ph, text: str) -> int:
    """How tall the fitted headline actually is in *ph*, in EMU.

    Measured with the size the run now states and the width the text really
    wraps inside — the box minus its own insets, which is what the renderer
    draws into and what the fitter (using the full box width) does not see.
    """
    try:
        from types import SimpleNamespace

        from reportbuilder.render.image.fast_preview import (
            _inherited_placeholder_style,
        )
        run = ph.text_frame.paragraphs[0].runs[0]
        font, inherited_pt, _c, _b, caps = _inherited_placeholder_style(ph)
        size_pt = run.font.size.pt if run.font.size is not None else (inherited_pt or TITLE_PT)
        inset = int(ph.text_frame.margin_left or 0) + int(ph.text_frame.margin_right or 0)
        width = max(1, int(ph.width or 0) - inset)
        st = SimpleNamespace(size_pt=size_pt, width=width, height=int(ph.height or 0),
                             font=font, caps=caps, line_spacing=0.0)
        lines = measured_line_count(text, width, size_pt, st)
        # 1.18: LibreOffice lays Bebas Neue's lines out taller than PIL measures
        # them, and it is LibreOffice that renders the deck. Measured on
        # Egoiq_x_Rahoo — a two-line headline the box said ended at 1.34" was
        # still drawing at 1.55" and came down across the subtitle. Reserving the
        # difference is what keeps the deck and the preview agreeing about where
        # the header ends.
        return int(lines * _title_line_height(size_pt, st) * _RENDERER_LINE_BOX)
    except Exception:  # noqa: BLE001 — a title must never fail a render
        logging.getLogger(__name__).warning("could not measure the title",
                                            exc_info=True)
        return 0


def draw_template_heading(slide, style, text: str) -> int:
    """Put *text* where the template says a title goes; return its bottom edge.

    0 means the template had no opinion and the caller should draw its own. On a
    layout slide the placeholder takes it and keeps every inherited property; on
    a harvested slide it is drawn in the box, font, size, weight and colour read
    off the customer's own slide.
    """
    if getattr(style, "chart_layout_index", None) is not None:
        # Also removes an empty placeholder, so PowerPoint shows no prompt.
        if not _fill_title_placeholder(slide, text, style) or not text:
            return 0
        profile = getattr(style, "profile", None)
        if profile is not None and profile.title.positioned:
            _left, top, _width, height = harvested_title_box(profile, text)
            return top + height
        ph = slide.shapes.title
        if ph is None:
            return 0
        # The BOX bottom is not the TEXT bottom. A customer's title box is drawn
        # for the headline they wrote; ours is often longer and wraps past it,
        # and the fitter only shrinks to `_TITLE_MIN_SCALE` — so on
        # Egoiq_x_Rahoo a two-line headline ran straight over the subtitle in
        # the deck and the preview both. Report where the text actually ends, so
        # whatever the caller places next starts below it.
        return int(ph.top or 0) + max(int(ph.height or 0), _rendered_title_height(ph, text))

    profile = harvested_profile(style)
    if not text or profile is None or not profile.title.positioned:
        return 0
    st = profile.title
    left, top, width, height = harvested_title_box(profile, text)
    # ONE source for the size: the template spec, sized from the font's own cap
    # height. This path used to fit against the template's stated size instead,
    # so a template whose title lands here rendered at a different physical size
    # than one whose title goes in a layout placeholder — the same deck, two
    # sizes, depending on a detail of how the template was built.
    size = _spec_title_pt(style, st.font) or fit_title_size(st, text)
    _textbox(slide, left, top, width, height,
             [(text, size, title_colour_for(st, style),
               True if st.bold is None else st.bold)],
             font=st.font or getattr(style, "heading_font", "") or "")
    return top + height


def header_furniture_floor(profile, chart_top: int, sw: int, sh: int) -> int:
    """The bottom of whatever the template draws between the title and the chart.

    Synsam rules a line under its titles and the agent deck sets a bar beside
    them; a subtitle placed by the title's height alone lands across it.
    Backdrops are ignored — they span the whole slide and would push everything
    to the floor.
    """
    floor = 0
    for item in getattr(profile, "furniture", []) or []:
        left, top, width, height = (int(v or 0) for v in item.box)
        if width <= 0 and height <= 0:
            continue
        if sw > 0 and sh > 0 and width * height >= 0.75 * sw * sh:
            continue
        bottom = top + height
        if top >= chart_top or bottom > chart_top:
            continue
        floor = max(floor, bottom)
    return floor


def title_colour_for(st, style):
    """The colour to draw the title in, taken from the TEMPLATE — but from the
    field of it that agrees with the background the same template states.

    A harvested title colour is read off one slide of the customer's deck, and a
    deck usually has both light and dark slides. Egoiq_x_Rahoo's harvest yields
    black on a ground the same template states as 37474F, so the headline came
    out dark-on-dark and all but invisible, while that template's own dk1 says
    white.

    So: the harvested colour when it is legible against the ground, otherwise
    the template's own ink for that ground. Both are the template's; this picks
    the one it did not contradict. Nothing here invents a colour.
    """
    harvested = _rgb(getattr(st, "colour", "") or "")
    bg = (getattr(style, "background", "") or "").strip()
    if harvested is None or not bg:
        return harvested or _furniture_px(style)[0]
    ground = bg if bg.startswith("#") else f"#{bg}"
    if abs(_relative_luminance(f"#{st.colour}") - _relative_luminance(ground)) < 0.25:
        return _furniture_px(style)[0]
    return harvested


def _fill_title_placeholder(slide, title: str, style=None) -> bool:
    """Put *title* in the layout's title placeholder. False if there isn't one.

    Only the TEXT is set: size, font, colour and position stay inherited, which
    is the whole point of using the customer's layout. LibreOffice resolves that
    chain correctly — the "slate" master states Bebas Neue 30pt for its title and
    that is what the deck shows. The preview compositor has to resolve the SAME
    chain (see fast_preview._inherited_placeholder_style); writing harvested
    values over the placeholder instead put Arial 14pt on a slide whose template
    asked for Bebas Neue 30pt.
    An unused placeholder is removed so PowerPoint does not show its "Click to
    add title" prompt.
    """
    try:
        ph = slide.shapes.title
    except (AttributeError, KeyError):
        return False
    if ph is None:
        return False
    if not title:
        ph._element.getparent().remove(ph._element)
        return True
    tf = ph.text_frame
    tf.word_wrap = True
    # TOP-anchored: a headline should start at the top of the space the template
    # gave it and grow down, not float in the middle of it. A middle-anchored box
    # drawn for one line leaves a band of empty slide above a two-line title and
    # pushes the second line toward the chart.
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.text = title

    # Pull the headline up into the top margin, and give the space it gains to
    # the box. Templates commonly park the title a third of the way down the
    # slide, which wastes the band above it AND leaves a long headline nowhere
    # to wrap but over the subtitle. Raising the box does both jobs at once and
    # costs nothing: the space above the title is empty by definition.
    try:
        left, top = int(ph.left or 0), int(ph.top or 0)
        width, height = int(ph.width or 0), int(ph.height or 0)
        if top > _MIN_TITLE_TOP and width and height:
            # All FOUR, explicitly. A placeholder inherits its position and size
            # from the layout, and writing one of them makes python-pptx emit an
            # <a:xfrm> carrying only that value — the others stop resolving and
            # the headline wrapped one word per line down the slide.
            ph.left, ph.width = left, width
            ph.height = height + (top - _MIN_TITLE_TOP)
            ph.top = _MIN_TITLE_TOP
    except (AttributeError, TypeError):
        pass

    # The size comes from the TEMPLATE SPEC — resolved once per template, from
    # the font's own cap height, so every deck's headline renders the same
    # physical size whatever face it is set in. It is stated explicitly because
    # an unstated placeholder size is resolved differently by LibreOffice,
    # PowerPoint and the preview compositor; font, colour and position are still
    # the template's.
    #
    # Nothing is fitted or shrunk here. A headline too long for its box
    # OVERFLOWS, deliberately: that is the signal to the author to write a
    # shorter one, and it is the same signal in the deck and in the preview.
    try:
        spec = getattr(style, "resolved_spec", None) if style else None
        if spec is None and style is not None:
            # A style that did not come through template_cache.resolve. Measure,
            # as this did for every slide before the spec was resolved once.
            from reportbuilder.render.resolved_style import build_spec
            spec = build_spec(style, title_font=_inherited_title_font(ph))
        if spec is not None and spec.title.size_pt:
            run = tf.paragraphs[0].runs[0]
            run.font.size = Pt(spec.title.size_pt)
    except Exception:  # noqa: BLE001 — a title must never fail a render
        logging.getLogger(__name__).warning("could not size the title",
                                            exc_info=True)
    return True


def add_image_slide_chrome(ctx: RenderContext) -> None:
    """Decorate an image-mode slide with house-style chrome.

    Adds (in z-order, bottom → top):
    1. Cream background rectangle (full slide)
    2. Teal vertical accent bar (top-left)
    3. Title textbox (ctx.title, bold INK, word-wrapped) — REQ-C-24a, REQ-D-04
    4. Methodology footer bottom-left (stat label + "· n = N") — REQ-C-24h
    5. N annotation textbox bottom-right (compact) — REQ-C-24h

    Call this *before* the image builder so the chart picture lands on top.
    """
    slide = ctx.slide
    sw, sh = _slide_dims(slide)

    _bg, _ink, _accent = _theme_colours(ctx)
    _muted = _furniture_px(ctx.style)[1]
    templated = _from_template(ctx)
    profile = harvested_profile(ctx.style)

    # 0 — The customer's own furniture, when their design was harvested off a
    #     slide rather than inherited from a layout: background, logo, the rule
    #     under the title. First, so everything else lands on top of it.
    if profile is not None:
        clone_furniture(slide, profile.furniture)

    # 1 — Cream background. Only when NO template laid this slide out: a
    #     full-slide rectangle would cover the customer's own background,
    #     brand furniture and everything else their layout provides.
    if not templated and profile is None:
        bg = slide.shapes.add_shape(1, 0, 0, sw, sh)
        bg.fill.solid()
        bg.fill.fore_color.rgb = _bg
        bg.line.fill.background()
        bg.shadow.inherit = False

    # The title block (accent bar + title + description) is gated on the
    # elements.title toggle. The live preview sets it False to render a
    # title-less PNG so the frontend can own the title region (progressive
    # "Generating title…" placeholder). The deck keeps it True (default).
    show_title = getattr(getattr(ctx.spec, "elements", None), "title", True)

    if show_title:
        # 2 — Title (top) + question subtitle (just above the chart)
        #     The top area is dedicated to the TITLE (AI key message when set,
        #     else the question). When a distinct headline is set, the actual
        #     QUESTION is a separate subtitle anchored to the BOTTOM of the header
        #     band — so the gap between the question and the chart stays constant
        #     no matter how many lines the question wraps to. It uses a lighter
        #     (non-bold) weight than the title. (REQ-C-24a, REQ-D-04)
        question = (ctx.title or "").strip()
        slide_title = (getattr(ctx.spec, "slide_title", None) or "").strip()
        slide_description = (getattr(ctx.spec, "slide_description", None) or "").strip()
        title = slide_title or question
        # Subtitle (the line just above the chart) is the editable slide_description; when
        # blank it defaults to the QUESTION — but only when the title is a DISTINCT headline
        # (otherwise the title already IS the question, so no redundant subtitle).
        has_distinct_title = bool(slide_title) and slide_title != question
        secondary = slide_description or (question if has_distinct_title else "")
        # On a STACKED bar the scale sits in the legend as bare numbers; move the endpoint
        # wording (1 = … · 7 = …) into the subtitle so the meaning isn't lost. (customer)
        # Only as the DEFAULT: an authored slide_description owns the whole line, so the
        # author can reword or drop a gloss that fits only one battery member (the levels
        # come from the first member with a parseable scale). The frontend prefills its
        # Subtitle box with this same default, so what you edit is what renders.
        if not slide_description and getattr(ctx.spec, "chart_type", "") in _STACKED_BAR_TYPES:
            gloss = scale_endpoint_gloss(ctx.series.categories)
            if gloss:
                secondary = f"{secondary}   {gloss}" if secondary else gloss
        # One fixed title size for every slide, from the template spec (the
        # house style is a template too, and is sized the same way).
        t_size = _spec_title_pt(ctx.style,
                               getattr(ctx.style, "heading_font", "")) or TITLE_PT

        # 3 — Teal accent bar. House furniture, so it is skipped on a templated
        #     slide: the customer's layout decides what sits beside a title.
        # Teal accent bar (thin vertical stripe, top-left), sized to the TITLE's
        #     actual height (its wrapped line count) so it doesn't tower over a short
        #     one-line headline. Capped at the title box height.
        if title:
            _n = wrapped_line_count(title, sw - Inches(1.0), t_size)
            bar_h = min(int(Inches(1.30)), _n * int(Pt(t_size * 1.25)) + int(Inches(0.06)))
        else:
            bar_h = int(Inches(0.30))
        if not templated and profile is None:
            acc = slide.shapes.add_shape(
                1, Inches(0.55), Inches(0.42), Inches(0.10), bar_h
            )
            acc.fill.solid()
            acc.fill.fore_color.rgb = _accent
            acc.line.fill.background()
            acc.shadow.inherit = False

        # 4 — Title. On a templated slide it goes into the layout's own title
        #     placeholder, so it inherits the customer's font, size, colour and
        #     position. Only a template without a title placeholder falls
        #     through to a box of our own.
        placed = bool(draw_template_heading(slide, ctx.style, title)) \
            if (templated or profile is not None) else False
        if title and not placed:
            # Tall, TOP-anchored box so the title can span up to ~4 lines (customers'
            # headlines are often 3) and honour manual line breaks ("\n") instead of
            # being clipped at 2. A short title still sits at the top (empty space
            # below is invisible); a long one grows DOWN toward the chart — if it meets
            # the question subtitle the author shortens the text.
            _textbox(
                slide,
                Inches(0.80), Inches(0.42),
                sw - Inches(1.0), Inches(1.30),
                [(title, t_size, _ink, True)],
                font=getattr(ctx.style, "heading_font", "") or "",
            )
        if secondary:
            # The question subtitle binds to the CHART: its box bottom sits just
            # above the chart (~1.84") and BOTTOM anchor makes multi-line questions
            # grow UPWARD toward the title. Its font steps down with length so a
            # long question always fits the box and is NEVER clipped at the top
            # (a bottom-anchored box clips overflow above it), while staying as
            # large as possible.
            s_size = 15.0   # replaced by the template spec below
            # With a template, the subtitle belongs to the TITLE — a line of
            # explanation directly under it, which is where every one of the
            # real client decks puts it. Left hanging off the chart instead it
            # floated in the middle of the slide with a band of empty cream
            # above it.
            # Bound to the CHART, never to the title: the box bottom sits a
            # fixed gap above the chart and BOTTOM anchor grows the text upward,
            # so the space between question and chart is identical on every
            # slide however long the question is. Placing it under the TITLE
            # instead made that gap depend on the headline's height, and a
            # two-line headline pushed the subtitle straight through it.
            anchor = MSO_ANCHOR.BOTTOM
            if templated or profile is not None:
                sub_bottom = int(ctx.slot.top) - int(_SUBTITLE_GAP)
                sub_h = min(int(_SUBTITLE_MAX_H), max(int(Inches(0.30)), sub_bottom))
                sub_top = max(0, sub_bottom - sub_h)
                # The BOTTOM is the fixed thing — a constant gap above the chart,
                # so that space never changes with the question's length. Growing
                # upward stops at whatever the header already occupies: the title
                # itself, and any rule the template draws under it (Synsam's).
                # Without the clamp a long question climbed through both.
                ceiling = 0
                title_profile = getattr(ctx.style, "profile", None)
                if title_profile is not None and title_profile.title.positioned:
                    _l, t_top, _w, t_h = harvested_title_box(title_profile, title)
                    ceiling = max(ceiling, int(t_top + t_h))
                    ceiling = max(ceiling, header_furniture_floor(
                        title_profile, int(ctx.slot.top), sw, sh))
                # The REAL headline, not the harvested box it was sized from:
                # the placeholder was raised and given the top margin's space, so
                # a two-line title ends well below where the harvest says. Using
                # the harvested bottom left the subtitle under the first line and
                # the second line came down across it.
                try:
                    tph = slide.shapes.title
                except (AttributeError, KeyError):
                    tph = None
                if tph is not None and title:
                    ceiling = max(ceiling, int(tph.top or 0)
                                  + _rendered_title_height(tph, title))
                if ceiling:
                    ceiling += int(Inches(0.06))     # a hair of clearance
                # Push the box DOWN to clear the header — `min` capped it instead,
                # which left it exactly where it had been overlapping.
                sub_top = max(sub_top, ceiling)
                sub_h = max(int(Inches(0.20)), sub_bottom - sub_top)
                # Left edge and width from the TITLE, so the two line up.
                t_left, t_width = title_left_width(slide, ctx.style, title)
                sub_left = t_left or int(ctx.slot.left)
                sub_w = t_width or int(ctx.slot.width)
            else:
                sub_h, sub_top = int(Inches(0.92)), int(Inches(0.92))
                sub_left, sub_w = int(Inches(0.80)), int(sw - Inches(1.0))
            s_font = _body_font(ctx)
            # From the spec, not fitted to the box: the question must be the
            # same size on every slide of every deck, which is the whole point
            # of sizing by cap height.
            # An author's own size wins: a question that runs three lines is
            # shrunk by hand, not by us guessing a different rule.
            s_font = getattr(ctx.style, "subtitle_font", "") or s_font
            s_size = (getattr(ctx.style, "subtitle_size_pt", 0.0)
                      or _spec_subtitle_pt(ctx.style, s_font) or 13.0)
            _textbox(
                slide,
                sub_left, sub_top, sub_w, sub_h,
                [(secondary, s_size,
                  f"#{getattr(ctx.style, 'subtitle_colour', '')}"
                  if getattr(ctx.style, "subtitle_colour", "") else _muted, False)],
                anchor=anchor,
                font=s_font,
            )

    # 4 — Methodology footer bottom-left (REQ-C-24h)
    #     Auto format: a simple "N = <base_n>". An author can override it per slide via
    #     spec.footer_note; "{n}" expands to the base and "{stat}" to the statistic label
    #     (e.g. "{stat} · n = {n}" restores the verbose form), so "N = {n}" keeps the
    #     count live.
    base_n = ctx.series.base_n.get("Total")
    stat_label = _STAT_FOOTER.get(ctx.series.statistic, ctx.series.statistic)
    override = (getattr(ctx.spec, "footer_note", None) or "").strip()
    if override:
        footer_text = override.replace("{n}", str(base_n if base_n is not None else "")) \
                              .replace("{stat}", stat_label)
    elif base_n is not None:
        footer_text = f"N = {base_n}"
    else:
        footer_text = stat_label
    # A pie/doughnut/funnel split into panels can drop a group (too thin a base,
    # or more groups than the page holds) — the editor's warning stays in the
    # editor, so this is the ONLY record of it that travels with the deck. Image
    # mode has no separate "classifying variable" box the way the native builder
    # does, so the disclosure rides on the same footer line. It names only the
    # omitted GROUPS, never the classifier's raw code (e.g. "var7") — RenderContext
    # carries no model to resolve that to a human label, and the panel titles
    # already show which variable the split used. (spec 2026-08-22, ruling 2026-08-22)
    if getattr(ctx.spec, "classifying_var", None):
        omission = _omission_clause(ctx).removeprefix(" · ")
        if omission:
            footer_text = f"{footer_text}   ·   {omission}"
    # Left margin follows the chart on a templated or harvested slide: those
    # margins are the customer's, and a footer 0.08in off from the chart above
    # it reads as a mistake rather than as a choice.
    foot_left = int(ctx.slot.left) if (templated or profile is not None) else int(Inches(0.62))
    foot_top = footer_top(slide, sh, sw)
    _textbox(
        slide,
        foot_left, foot_top,
        sw - Inches(4.0), Inches(0.40),
        [(footer_text, 9.5, _muted, False)],
        align=PP_ALIGN.LEFT,
        font=_body_font(ctx),
    )
    # Scale endpoint legend for a partially-labelled numeric scale (e.g. "1 = täysin
    # eri mieltä · 7 = täysin samaa mieltä") — a small caption just above the footer,
    # so the numeric axis (1..7) reads cleanly and the text isn't lost. (REQ-C-24c)
    caption = getattr(ctx.series, "caption", None)
    if caption:
        # Right-aligned on the footer row (below the plot) so it never overlaps the
        # chart's x-axis; shares the line with the left-aligned methodology footer.
        _textbox(
            slide,
            sw - Inches(6.4), foot_top,
            Inches(6.0), Inches(0.40),
            [(caption, 9.5, _muted, False)],
            align=PP_ALIGN.RIGHT,
            font=_body_font(ctx),
        )
    # n is shown once, in the methodology footer above (it already reads
    # "<stat label> · n = N"). The previous separate bottom-right "n = N"
    # annotation was redundant and has been removed.
