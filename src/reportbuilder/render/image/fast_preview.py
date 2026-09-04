"""A chart preview without running LibreOffice per chart.

The Design page asks for one PNG per chart, and each one used to be a whole
miniature deck: build a one-slide .pptx, hand it to LibreOffice, get a PDF back,
rasterize it. Measured on a real report that is 4.4s a chart, of which 3.3s is
LibreOffice — and the page asks for sixty of them.

But every one of those sixty conversions renders the SAME thing behind the
chart: the customer's ground, their band, their logo. Only the chart on top
changes. So LibreOffice runs once per template to produce that ground, the image
is cached, and each preview is the chart composited onto a copy of it.

LibreOffice stays the authority on what a customer's template looks like — this
is the same conversion the deck goes through, not a re-drawing of their design —
it just stops being paid sixty times over.

Falls back to the caller's slow path on anything unexpected: a preview that is
wrong is worse than a preview that is slow.
"""
from __future__ import annotations

import hashlib
import logging
import os
import subprocess
import tempfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

from reportbuilder import cache_dirs
from reportbuilder.export.pdf_convert import pptx_to_pdf
from reportbuilder.export.preview import rasterize_pages

log = logging.getLogger(__name__)

EMU_PER_INCH = 914400
# PowerPoint's own text-frame insets, applied when a shape does not state its own.
_DEFAULT_LR_INSET = 91440      # 0.1in
_DEFAULT_TB_INSET = 45720      # 0.05in
_CACHE = cache_dirs.ground_root()


def _key(style, dpi: int) -> str:
    """Identity of a ground: which template file, as it is right now, at what
    dpi, drawn the way this host currently draws text.

    The last part matters: the ground carries the template's own text, so a new
    font stand-in changes it while the template file is untouched.
    """
    from reportbuilder.render.fonts import rendering_fingerprint

    source = str(getattr(style, "spec_source", "") or "generic")
    stamp = ""
    try:
        stamp = str(os.path.getmtime(source))
    except OSError:
        pass
    # The LAYOUT and the author's corrections change what an empty slide looks
    # like while the template file is untouched — a different layout has a
    # different band and logo, and a corrected background is a different colour.
    # Without them a ground cached for one layout was served for every other,
    # and choosing a layout appeared to do nothing.
    slot = getattr(style, "chart_slot", None)
    shape = (f"{getattr(style, 'chart_layout_index', None)}"
             f"|{getattr(style, 'background', '')}|{getattr(style, 'ink', '')}"
             f"|{getattr(style, 'accent', '')}"
             f"|{'' if slot is None else f'{slot.left},{slot.top},{slot.width},{slot.height}'}")
    raw = (f"{source}|{stamp}|{dpi}|{getattr(style, 'slide_width', 0)}"
           f"|{shape}|{rendering_fingerprint()}")
    return hashlib.md5(raw.encode()).hexdigest()[:16]


def ground_image(style, dpi: int = 110):
    """The customer's empty slide, rendered once and cached. None if it cannot be.

    "Empty" means built exactly as a chart slide is — from their layout, or with
    their harvested furniture redrawn — and then left without a chart or any of
    our text.
    """
    _CACHE.mkdir(parents=True, exist_ok=True)
    cached = _CACHE / f"{_key(style, dpi)}.png"
    if cached.exists():
        try:
            return Image.open(cached).convert("RGB")
        except OSError:
            cached.unlink(missing_ok=True)

    from pptx import Presentation
    from pptx.util import Inches

    from reportbuilder.render.deck import _resolve_slot, _strip_slides
    from reportbuilder.render.image.slide_chrome import template_ground

    source = getattr(style, "spec_source", None)
    work = Path(tempfile.mkdtemp(prefix="nsight-ground-"))
    try:
        prs = None
        if source and source not in ("generic", "attendo-interim-proxy"):
            try:
                prs = Presentation(source)
                _strip_slides(prs)
            except Exception:  # noqa: BLE001 — unreadable template: no fast path
                prs = None
        if prs is None:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

        # The same slide the renderer would build, minus the chart and the text.
        slot = _resolve_slot(prs, style, "preview", "image", title="")
        slide = prs.slides[slot.slide_index]
        if not template_ground(slide, style):
            _house_ground(slide, int(prs.slide_width), int(prs.slide_height), style)
        # An empty title placeholder would print PowerPoint's own prompt text.
        for shape in list(slide.shapes):
            if shape.is_placeholder and shape.has_text_frame and not shape.text_frame.text:
                shape._element.getparent().remove(shape._element)

        path = str(work / "ground.pptx")
        prs.save(path)
        pdf = pptx_to_pdf(path, str(work))
        pngs = rasterize_pages(pdf, str(work / "png"), dpi=dpi, workers=1)
        if not pngs:
            return None
        image = Image.open(pngs[0]).convert("RGB")
        image.save(cached)
        return image
    except Exception:  # noqa: BLE001 — the slow path is always available
        log.warning("could not pre-render the template ground; previews stay slow",
                    exc_info=True)
        return None
    finally:
        import shutil

        shutil.rmtree(work, ignore_errors=True)


def _house_ground(slide, sw: int, sh: int, style) -> None:
    """nSight's own ground, for a report with no template."""
    from reportbuilder.render.image.slide_chrome import theme_colours

    bg, _ink, _accent = theme_colours(style)
    shape = slide.shapes.add_shape(1, 0, 0, sw, sh)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    shape.shadow.inherit = False


def _font(family: str, size_px: float, *, bold: bool = False,
          italic: bool = False):
    """A PIL font for *family*, resolved the way THIS HOST would render it.

    Through the configured substitutions first: a template naming Calibri on a
    host without it is drawn by LibreOffice in whatever stands in for Calibri,
    and a preview that picked something else instead would differ from the deck
    in the one place a preview is supposed to be trusted.
    """
    from reportbuilder.render import fonts as F

    wanted = (family or "").strip()
    try:
        wanted = F.substitutions().get(wanted, wanted)
    except Exception:  # noqa: BLE001
        pass
    path = _font_file(wanted, bold=bold, italic=italic)
    if path:
        try:
            # A FRACTIONAL size, deliberately: 11pt at 110 dpi is 16.81px, and
            # rounding that to 17 stretches every glyph by 1.2% — which is
            # invisible on one word and several pixels of drift by the end of a
            # long subtitle line.
            return ImageFont.truetype(path, size_px)
        except (OSError, TypeError):
            pass
    try:
        return ImageFont.truetype(path, int(round(size_px)))
    except (OSError, TypeError, ValueError):
        pass
    try:
        return ImageFont.load_default(size_px)
    except TypeError:
        return ImageFont.load_default()


_FONT_FILES: dict[tuple, str] = {}


def _font_file(family: str, *, bold: bool = False, italic: bool = False) -> str:
    """The font file this host would actually use for *family*.

    fontconfig, because that is who LibreOffice asks: on a host without Calibri,
    `fc-match Calibri` answers Liberation Sans and the deck is drawn in Liberation
    Sans, while matplotlib's own fallback would have said DejaVu — a preview in a
    different typeface from the deck it is previewing.
    """
    cache_key = (family, bold, italic)
    if cache_key in _FONT_FILES:
        return _FONT_FILES[cache_key]
    path = ""
    try:
        pattern = family or "sans-serif"
        if bold:
            pattern += ":bold"
        if italic:
            pattern += ":italic"
        out = subprocess.run(["fc-match", "-f", "%{file}", pattern],
                             capture_output=True, text=True, timeout=5)
        path = out.stdout.strip() if out.returncode == 0 else ""
    except (OSError, subprocess.SubprocessError):
        path = ""
    if not path:
        try:
            from matplotlib import font_manager

            path = font_manager.findfont(
                font_manager.FontProperties(family=family or None),
                fallback_to_default=True)
        except Exception:  # noqa: BLE001
            path = ""
    _FONT_FILES[cache_key] = path
    return path


def title_box_headers(style, text: str = "") -> dict[str, str]:
    """Where the template puts its title, as response headers for a caller that
    draws the title itself (routes_questions's fast preview path).

    The compositor never draws a title of its own on the fast path — the slide
    it is handed has none, because the caller only takes this path when the
    frontend owns that region. So the box has to travel out of band, and the
    template's own harvested profile (`TemplateStyleSpec.profile`) is the one
    place it is known: see style_spec.py and template_profile.py for how it
    got there. Nothing here re-derives or guesses a position.

    Empty when the template states no title box — no template, no profile, or
    a profile whose title was never positioned (`TextStyle.positioned`). The
    caller then sends no headers, and the frontend draws no title, same as
    today when a preview carries none.

    *text* is the headline this preview will carry. Given it, the size reported
    is the one the DECK would use — `fit_title_size` shrinks a long headline to
    fit rather than letting it overrun — and the box grows with the wrapped
    lines, exactly as `harvested_title_box` does for the real slide. Without it
    the frontend would draw a long title at the template's nominal size and
    overflow the box the deck keeps it inside, which is the whole defect the
    fitting exists to prevent.
    """
    from reportbuilder.render.image.slide_chrome import (
        TITLE_PT, fit_title_size, harvested_title_box)

    profile = getattr(style, "profile", None)
    title = getattr(profile, "title", None) if profile is not None else None
    if title is None or not title.positioned:
        return {}
    sw = int(getattr(style, "slide_width", 0) or 0)
    sh = int(getattr(style, "slide_height", 0) or 0)
    if sw <= 0 or sh <= 0:
        return {}
    colour = (title.colour or getattr(style, "ink", "") or "2B2B2B").lstrip("#").upper()
    # The box the DECK would give this headline, at the size the deck would use.
    # With no text we cannot fit anything, so fall back to the template's own.
    size_pt = fit_title_size(title, text) if text else (title.size_pt or TITLE_PT)
    if text:
        box_l, box_t, box_w, box_h = harvested_title_box(profile, text)
    else:
        box_l, box_t, box_w, box_h = title.left, title.top, title.width, title.height
    return {
        "X-Title-Box": "{:.4f},{:.4f},{:.4f},{:.4f}".format(
            box_l / sw, box_t / sh, box_w / sw, box_h / sh),
        "X-Title-Font": title.font or getattr(style, "heading_font", "") or "Arial",
        "X-Title-Size-Pt": f"{size_pt:g}",
        "X-Title-Color": colour,
        "X-Title-Align": title.align or "left",
        "X-Title-Caps": "1" if title.caps else "0",
        "X-Slide-Aspect": f"{sw / sh:.6f}",
    }


def compose_from_slide(style, slide, dpi: int = 110):
    """The finished preview: the cached ground with THIS slide's own shapes on it.

    Reads what the renderer actually put on the slide — the chart picture where
    it landed, the footer in the font and colour it was given — rather than
    re-deriving any of it. The .pptx is the layout description; this just draws
    it, for the one slide a preview needs, instead of starting LibreOffice.

    Everything else on the slide (ground, band, logo) is already in the cached
    image, because that image was built from the same template by the same code.

    Returns a PIL image, or None when there is no ground to build on.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    ground = ground_image(style, dpi)
    if ground is None:
        return None
    image = ground.convert("RGB").copy()
    slide_w = int(getattr(style, "slide_width", 0) or 0) or int(13.333 * EMU_PER_INCH)
    px = image.width / slide_w

    def to_px(emu: float) -> int:
        return int(round(emu * px))

    draw = ImageDraw.Draw(image)
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                _paste_picture(image, shape, to_px)
            elif shape.has_text_frame and (shape.text_frame.text or "").strip():
                _draw_text(draw, shape, to_px, dpi, style)
        except Exception:  # noqa: BLE001 — one shape must not lose the preview
            log.warning("preview: could not draw %s", getattr(shape, "name", "?"),
                        exc_info=True)
    return image


def _paste_picture(image, shape, to_px) -> None:
    """The chart, at the box the renderer gave it — letterboxing included."""
    import io

    with Image.open(io.BytesIO(shape.image.blob)) as chart:
        chart = chart.convert("RGBA")
        box = (max(1, to_px(int(shape.width or 0))), max(1, to_px(int(shape.height or 0))))
        resized = chart.resize(box, Image.LANCZOS)
    image.paste(resized, (to_px(int(shape.left or 0)), to_px(int(shape.top or 0))),
                resized)


def _line_step(shape, style, font) -> int:
    """Pixels from one baseline to the next, as PowerPoint computes it.

    `lnSpc` is a percentage of the FONT'S OWN line height — ascent plus descent,
    about 1.36em for Noto Sans — not of the point size. Multiplying the point
    size by Holiday Club's 90% made every line 26% too tight, so a two-line
    title came out crammed together while wrapping at exactly the right words.
    """
    ascent, descent = font.getmetrics()
    natural = ascent + descent
    if shape.is_placeholder:
        title = getattr(getattr(style, "profile", None), "title", None)
        if title is not None and title.line_spacing:
            return int(round(natural * title.line_spacing))
    return int(round(natural))


_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _caps(shape, run, style) -> bool:
    """Does this run render upper-case whatever was typed?

    `cap="all"` is a rendering instruction, not stored text: Holiday Club's
    title carries it, so LibreOffice prints the headline in capitals while the
    .pptx still holds mixed case. Drawing what the XML says instead of what
    PowerPoint shows made every Holiday Club preview visibly wrong.
    """
    rpr = run._r.find(f"{_A_NS}rPr")
    if rpr is not None and rpr.get("cap"):
        return rpr.get("cap") == "all"
    if shape.is_placeholder:
        title = getattr(getattr(style, "profile", None), "title", None)
        if title is not None:
            return bool(title.caps)
    return False


# Ground and legibility are NOT decided here — see render/resolved_style, which
# is the one place that answers "what does this template say?" for the deck and
# the preview alike.
from reportbuilder.render.resolved_style import (  # noqa: E402
    ground as chart_ground, legible_on as _legible_on,
)


def _inherited_placeholder_style(shape):
    """(font, size_pt, colour, bold, caps) a placeholder INHERITS, as PowerPoint
    and LibreOffice resolve it: the layout's matching placeholder first, then the
    master's, then the master's txStyles.

    A placeholder python-pptx writes states nothing, so everything about how it
    looks lives up that chain. Reading only the shape — or the harvested profile
    — is what drew Arial 14pt black on a deck whose master ("slate") asks for
    Bebas Neue 30pt: the preview showed a small sentence-case line where the deck
    showed a big condensed headline, on the same slide.

    Returns empty/zero fields where the chain says nothing, so the caller keeps
    its own fallbacks.
    """
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    try:
        ph_type = shape.placeholder_format.type
    except (AttributeError, ValueError):
        return "", 0.0, "", None, False

    def _read(el):
        """(font, sz, colour, bold, caps) from the first defRPr/rPr under `el`."""
        if el is None:
            return "", 0.0, "", None, False
        for tag in (f"{A}defRPr", f"{A}rPr"):
            for rpr in el.iter(tag):
                latin = rpr.find(f"{A}latin")
                fill = rpr.find(f"{A}solidFill")
                srgb = fill.find(f"{A}srgbClr") if fill is not None else None
                sz = rpr.get("sz")
                return (latin.get("typeface") if latin is not None else "",
                        int(sz) / 100.0 if sz else 0.0,
                        srgb.get("val") if srgb is not None else "",
                        {"1": True, "0": False}.get(rpr.get("b")),
                        rpr.get("cap") == "all")
        return "", 0.0, "", None, False

    def _matching(container):
        for cand in getattr(container, "placeholders", []):
            try:
                if cand.placeholder_format.type == ph_type:
                    return cand.element
            except (AttributeError, ValueError):
                continue
        return None

    sources = []
    try:
        layout = shape.part.slide.slide_layout
        sources.append(_matching(layout))
        master = layout.slide_master
        sources.append(_matching(master))
        # Last: the master's text styles, which state a size for every title.
        sources.append(master.element.find(f"{P}txStyles/{P}titleStyle")
                       if ph_type == 1 else
                       master.element.find(f"{P}txStyles/{P}bodyStyle"))
    except (AttributeError, KeyError):
        pass

    font, size_pt, colour, bold, caps = "", 0.0, "", None, False
    for src in sources:
        f, z, c, b, cp = _read(src)
        font = font or f
        size_pt = size_pt or z
        colour = colour or c
        bold = bold if bold is not None else b
        caps = caps or cp
        if font and size_pt and colour:
            break
    return font, size_pt, colour, bold, caps


def _run_style(shape, para, run, style):
    """(family, size_pt, colour, bold, italic) for one RUN.

    Per run, not per paragraph: a bullet line is a teal glyph followed by ink
    body text, and colouring the whole line from the first run painted the body
    teal. A title PLACEHOLDER states nothing and inherits from the layout and
    master — which is what the template profile already read off the template.
    """
    family = run.font.name or ""
    size_pt = float(run.font.size.pt) if run.font.size is not None else 0.0
    bold = bool(run.font.bold)
    italic = bool(run.font.italic)
    colour = ""
    try:
        if run.font.color is not None and run.font.color.type is not None:
            colour = str(run.font.color.rgb)
    except (AttributeError, ValueError):
        pass
    if not (family and size_pt and colour) and shape.is_placeholder:
        # The template's own chain first — it is what the DECK renders from.
        inh_font, inh_pt, inh_col, inh_bold, _caps = _inherited_placeholder_style(shape)
        family = family or inh_font
        size_pt = size_pt or inh_pt
        # The chain's colour only when it is legible on the ground this same
        # template states. Egoiq_x_Rahoo resolves black for its title over a
        # 37474F ground, so the headline came out invisible; its own dk1 says
        # white. Both are the template's — this takes the one it did not
        # contradict, and invents nothing.
        if inh_col and _legible_on(inh_col, chart_ground(style)):
            colour = colour or inh_col
        if run.font.bold is None and inh_bold is not None:
            bold = inh_bold
        title = getattr(getattr(style, "profile", None), "title", None)
        if title is not None:
            family = family or title.font
            size_pt = size_pt or title.size_pt
            # Same legibility rule as the chain above: a harvested colour is read
            # off ONE slide of the customer's deck, and a deck has light and dark
            # slides both.
            if title.colour and _legible_on(title.colour, chart_ground(style)):
                colour = colour or title.colour
            if run.font.bold is None:
                bold = bool(title.bold)
    if not colour:
        # Nothing legible stated anywhere: the slide's own furniture ink, which
        # is what the chart on the same slide uses, so text agrees with charts
        # about what this background needs.
        from reportbuilder.render.resolved_style import ink as _ink
        colour = _ink(style).lstrip("#")
    return family, size_pt or 11.0, colour, bold, italic


#: A tab in a run's text. Its own segment because it is not something to draw
#: but somewhere to move to — the tab stop this deck's bullets set at marL. An
#: object, not a "\t" string, so no piece of real text can ever be mistaken for
#: one.
_TAB = object()


def _indents(para, to_px) -> tuple[int, int]:
    """(first-line x, wrapped-line x) offsets from the shape's left edge.

    A bullet paragraph sets `marL` (where the text sits) and a negative `indent`
    (how far back the glyph hangs). Without them every bullet drew flush left and
    the nesting levels collapsed onto each other.
    """
    ppr = para._p.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr")
    if ppr is None:
        return 0, 0
    try:
        mar_l = int(ppr.get("marL") or 0)
        indent = int(ppr.get("indent") or 0)
    except ValueError:
        return 0, 0
    return to_px(mar_l + indent), to_px(mar_l)


def _draw_text(draw, shape, to_px, dpi: int, style) -> None:
    """A text shape on the slide — a title, a subtitle, a footer, a bullet list.

    Drawn run by run and line by line, so a paragraph made of several runs (a
    coloured bullet glyph then ink body text) keeps each run's own font, weight
    and colour, and a hanging indent puts the glyph where PowerPoint puts it.
    """
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    frame = shape.text_frame
    # A text frame insets its text from the shape's edges. nSight's own boxes set
    # all four to zero; a PLACEHOLDER inherits PowerPoint's defaults — 0.1in left
    # and right, 0.05in top and bottom — which at 110 dpi is the 10px left and
    # 5px down that every composited title was out by.
    l_ins = frame.margin_left if frame.margin_left is not None else _DEFAULT_LR_INSET
    t_ins = frame.margin_top if frame.margin_top is not None else _DEFAULT_TB_INSET
    r_ins = frame.margin_right if frame.margin_right is not None else _DEFAULT_LR_INSET
    b_ins = frame.margin_bottom if frame.margin_bottom is not None else _DEFAULT_TB_INSET
    left_emu = int(shape.left or 0) + int(l_ins)
    top_emu = int(shape.top or 0) + int(t_ins)
    width_px = max(1, to_px(int(shape.width or 0) - int(l_ins) - int(r_ins)))
    height_px = max(1, to_px(int(shape.height or 0) - int(t_ins) - int(b_ins)))
    y = to_px(top_emu)

    lines: list[tuple[list, object, int, int, bool]] = []
    line_steps: list[int] = []
    for para in frame.paragraphs:
        runs = [r for r in para.runs if r.text]
        if not runs:
            continue
        first_x, wrap_x = _indents(para, to_px)
        step, para_lines = 0, []
        # One drawable segment per run: (text, font, colour).
        segments = []
        for run in runs:
            family, size_pt, colour, bold, italic = _run_style(shape, para, run, style)
            font = _font(family, max(6.0, size_pt * dpi / 72), bold=bold, italic=italic)
            step = max(step, _line_step(shape, style, font))
            # The tab is kept, not spelled as spaces. A bullet is written as
            # "glyph TAB text" with a left tab stop at marL, which is what puts
            # the first line's text at the same x as its wrapped lines. Turning
            # it into two spaces started the text wherever the glyph happened to
            # end — about ten pixels short of its own continuation, on every
            # template, because the glyph's width has nothing to do with marL.
            text = run.text
            if _caps(shape, run, style):
                text = text.upper()
            for j, piece in enumerate(text.split("\t")):
                if j:
                    segments.append((_TAB, font, colour))
                if piece:
                    segments.append((piece, font, colour))
        # A first line that snaps to the tab stop has the wrapped lines' width,
        # not the glyph's head start — otherwise it takes one word too many.
        has_tab = any(t is _TAB for t, _f, _c in segments)
        avail = width_px - (wrap_x if has_tab else first_x)
        for i, chunk in enumerate(_wrap_segments(draw, segments, avail,
                                                 width_px - wrap_x)
                                  if frame.word_wrap is not False else [segments]):
            para_lines.append((chunk, para.alignment, first_x, wrap_x, i == 0))
        lines.extend(para_lines)
        line_steps.extend(step for _ in para_lines)

    if not lines:
        return
    if frame.vertical_anchor == MSO_ANCHOR.BOTTOM:
        y = to_px(top_emu) + height_px - sum(line_steps)
    for (segments, alignment, first_x, wrap_x, is_first), step in zip(lines, line_steps):
        x = to_px(left_emu) + (first_x if is_first else wrap_x)
        if alignment in (PP_ALIGN.RIGHT, PP_ALIGN.CENTER):
            width = sum(draw.textlength(t, font=f)
                    for t, f, _c in segments if t is not _TAB)
            if alignment == PP_ALIGN.RIGHT:
                x = to_px(left_emu) + width_px - int(width)
            else:
                x = to_px(left_emu) + (width_px - int(width)) // 2
        for text, font, colour in segments:
            if text is _TAB:
                # Left tab stop at marL — the same place the wrapped lines start.
                x = max(x, to_px(left_emu) + wrap_x)
                continue
            draw.text((x, y), text, font=font, fill=f"#{colour}")
            x += int(round(draw.textlength(text, font=font)))
        y += step


def _wrap_segments(draw, segments, first_width: int, wrap_width: int) -> list[list]:
    """Break a run-styled paragraph into lines, keeping each run's styling.

    Wrapping happens across the whole paragraph, not per run, because a line
    break can fall in the middle of a run — which is what makes this different
    from wrapping a plain string.
    """
    out: list[list] = []
    current: list = []
    used = 0.0
    limit = first_width
    for text, font, colour in segments:
        if text is _TAB:
            # Carried through, never broken on and never measured: it moves the
            # pen to the tab stop, and how far that is depends on where the line
            # starts, which is decided when the line is drawn.
            current.append((text, font, colour))
            continue
        for word in _tokens(text):
            w = draw.textlength(word, font=font)
            if current and used + w > limit and word.strip():
                out.append(current)
                current, used, limit = [], 0.0, wrap_width
            if (current and current[-1][0] is not _TAB
                    and current[-1][1] is font and current[-1][2] == colour):
                current[-1] = (current[-1][0] + word, font, colour)
            else:
                current.append((word, font, colour))
            used += w
    if current:
        out.append(current)
    return out or [[]]


def _tokens(text: str) -> list[str]:
    """Words with their trailing spaces kept, so re-joining preserves spacing."""
    out, buf = [], ""
    for ch in text:
        buf += ch
        if ch == " ":
            out.append(buf)
            buf = ""
    if buf:
        out.append(buf)
    return out
