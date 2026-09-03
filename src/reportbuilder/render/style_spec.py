"""Load a StyleSpec (fonts/colors/slots) from a template .pptx (REQ-C-25/27a)."""
from __future__ import annotations

import logging

from pptx import Presentation
from reportbuilder import config
from reportbuilder.render.base import Slot, StyleSpec

log = logging.getLogger(__name__)

_DEFAULT_FONTS: dict[str, tuple[str, int]] = {
    "title": ("Arial", 14),
    "axis_values": ("Arial", 10),
    "axis_names": ("Arial", 10),
    "category_names": ("Arial", 10),
    "data_labels": ("Arial", 10),
    "legend": ("Arial", 10),
    "n_annotation": ("Arial", 9),
    "filter_var": ("Arial", 9),
}
_DEFAULT_PALETTE = ["1F77B4", "FF7F0E", "2CA02C", "D62728", "9467BD", "8C564B", "E377C2", "7F7F7F"]

# PowerPoint's own default themes. A template whose accents are exactly one of
# these has never had its theme edited, so it states no brand — and using it
# would paint a client's deck in Office blue. `attendo_agent_deck.pptx` is that
# case: stock Office masters under slides drawn in cream and teal by hand.
_STOCK_THEMES: tuple[tuple[str, ...], ...] = (
    ("4F81BD", "C0504D", "9BBB59", "8064A2", "4BACC6", "F79646"),   # Office 2007-10
    ("5B9BD5", "ED7D31", "A5A5A5", "FFC000", "4472C4", "70AD47"),   # Office 2013+
)


def states_a_brand(palette) -> bool:
    """True when a theme's accents are a deliberate choice, not Office's default."""
    return bool(palette) and tuple(c.upper() for c in palette[:6]) not in _STOCK_THEMES


class TemplateStyleSpec(StyleSpec):
    # This style came from a real .pptx. Chart colours, background and ink
    # follow the template whichever way its slides are built, so this is the
    # flag for "a template applies" — chart_layout_index is not, since a deck
    # whose design lives on its slides has none.
    from_template: bool = True
    # Which layout new chart slides are built from, and where the chart goes on
    # it. None means the design was harvested off a slide instead (or no usable
    # layout was found), and the renderer builds the slide itself.
    chart_layout_index: int | None = None
    chart_slot: Slot | None = None
    # What to borrow from the template when its design is NOT in the layouts:
    # the title's style and box, and the furniture that repeats on every slide.
    profile: object | None = None
    # The template's type and colour, resolved ONCE — see template_cache.resolve,
    # which sets this. None on a style built without a template (or built
    # directly rather than through resolve), in which case the renderers measure
    # as they did before. Carried here because the renderers are handed a style,
    # not a ResolvedTemplate, at far too many call sites to thread a new one.
    resolved_spec: object | None = None
    # Hex, no leading '#'. Empty means "no template opinion", and the house
    # cream/ink apply.
    background: str = ""
    ink: str = ""
    # The theme's accents when the template actually states a brand — empty when
    # its theme is untouched Office. Charts use these in order when present.
    brand_palette: list[str] = []
    # The one colour this template leads with: its first brand accent, else the
    # colour its own slides are drawn in. Chart ramps are built from it.
    accent: str = ""
    # The template theme's own typefaces. Used for the text nSight draws itself
    # (subtitle, footer): those are textboxes, not placeholders, so they inherit
    # nothing and would otherwise sit on a customer's slide in our font.
    heading_font: str = ""
    body_font: str = ""
    #: Overrides for the two elements that have no box of their own. The
    #: subtitle sits a fixed gap above the chart and the footer a fixed gap
    #: above the template's own foot, so there is nothing to move — but their
    #: SIZE is the thing an author reaches for first: "kysymystekstin
    #: pienentäminen", shrinking a question that runs too long.
    title_size_pt: float = 0.0
    title_colour: str = ""
    subtitle_font: str = ""
    subtitle_size_pt: float = 0.0
    subtitle_colour: str = ""
    footer_colour: str = ""

    def __init__(self, slide_width, slide_height, slots, fonts, palette, spec_source="generic"):
        self.slide_width = slide_width
        self.slide_height = slide_height
        self._slots = slots
        self._fonts = fonts
        self._palette = palette
        self.spec_source = spec_source
        self.matches_client_spec = False

    def font_for(self, element_class: str) -> tuple[str, int]:
        return self._fonts.get(element_class, ("Arial", 10))

    def color_for(self, series_index: int) -> str:
        return self._palette[series_index % len(self._palette)]

    def slot(self, name: str) -> Slot:
        return self._slots[name]

    def slots(self) -> dict[str, Slot]:
        return dict(self._slots)


def load_style_spec(template_path: str,
                    force_layout: int | None = None) -> TemplateStyleSpec:
    prs = Presentation(template_path)
    fonts = dict(_DEFAULT_FONTS)
    slots: dict[str, Slot] = {}
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            name = shape.name or ""
            slots[name] = Slot(slide_index=i, left=int(shape.left or 0), top=int(shape.top or 0),
                               width=int(shape.width or 0), height=int(shape.height or 0), name=name)
            if name.startswith("style:"):
                cls = name.split(":", 1)[1]
                try:
                    run = shape.text_frame.paragraphs[0].runs[0]
                    fn = run.font.name or "Arial"
                    sz = int(run.font.size.pt) if run.font.size is not None else 10
                    fonts[cls] = (fn, sz)
                except (AttributeError, IndexError):
                    pass
    spec = TemplateStyleSpec(prs.slide_width, prs.slide_height, slots, fonts,
                             list(_DEFAULT_PALETTE), spec_source=str(template_path))

    # The layout a chart slide should be BUILT FROM, rather than a blank one we
    # then guess geometry on. The client's designer already decided where a
    # title and a content area belong on their slide; template_check finds that
    # layout by the size of its content placeholder, which survives the fact
    # that the same layout is called "1 layout area", "Innehåll" and "Title and
    # Content" in three real client decks.
    from reportbuilder.render.template_check import inspect_template
    from reportbuilder.render.template_profile import extract_profile

    report = inspect_template(str(template_path))
    spec.heading_font = report.theme.heading_font or ""
    spec.body_font = report.theme.body_font or ""
    if report.best is not None:
        spec.chart_layout_index = report.best.index
        layout = prs.slide_layouts[report.best.index]
        content = _largest_content_placeholder(layout)
        if content is not None:
            spec.chart_slot = Slot(
                slide_index=-1,  # -1: a slide is added per chart, not reused
                left=int(content.left or 0), top=int(content.top or 0),
                width=int(content.width or 0), height=int(content.height or 0),
                name="chart")
        # The template's own theme colours beat our defaults: accent1-6 is what
        # PowerPoint's charts use for series, so a deck in Attendo's template
        # gets Attendo's palette rather than nSight teal. Untouched Office
        # accents are not a brand and are not used — see states_a_brand.
        if states_a_brand(report.theme.palette):
            spec.brand_palette = list(report.theme.palette)
            spec._palette = list(report.theme.palette)
            spec.accent = report.theme.palette[0]
        spec.background = report.theme.background
        spec.ink = report.theme.ink

    # Where the design actually lives. A layout carries it only when the deck's
    # own slides are built on one — Attendo's are, Synsam's and the agent deck's
    # are not, and building those from their (stock Office) layouts is what made
    # a client deck come out looking like plain PowerPoint. When it does not,
    # the profile carries the title's style and box and the repeating furniture,
    # and the renderer draws the slide itself.
    try:
        spec.profile = extract_profile(str(template_path), force_layout)
    except Exception:  # noqa: BLE001 — a template we cannot harvest still renders
        # LOUDLY. Swallowed, this reverts to building on the template's layouts,
        # and for a stock-Office file that is the plain-PowerPoint deck this
        # whole module exists to prevent — with nothing in the log to say why.
        log.warning("could not harvest a style profile from %s; falling back",
                    template_path, exc_info=True)
        spec.profile = None
    if spec.profile is None and not spec.brand_palette:
        # Nothing harvested AND no brand in the theme: whatever layouts this
        # file has are Office's own, and a 44pt centred title placeholder on a
        # white slide is worse than nSight's house style. Take the house style.
        spec.chart_layout_index = None
        spec.chart_slot = None
    if spec.profile is not None:
        spec.chart_layout_index = spec.profile.layout_index
        spec.chart_slot = None
        # The title is drawn by US, on the ground the template gives us — so it
        # has to be the colour that template writes its titles in, whichever
        # branch below we take. This used to be set only when the design lived
        # on SLIDES, which is the one case where it rarely matters: Arla's
        # design lives in a LAYOUT whose title is white on a black band, so the
        # harvested white was discarded and the default black drawn onto black.
        # The headline was there on every slide and could not be read.
        if spec.profile.title is not None and spec.profile.title.colour:
            spec.ink = spec.profile.title.colour
        if spec.profile.layout_index is None:
            # Design on the slides means colours on the slides too. The theme of
            # such a file describes nothing that is on screen: attendo_agent_deck
            # renders cream and teal while its theme says white, black and Office
            # blue, and reading the theme is what made a client's deck come back
            # in colours that appear nowhere in their template.
            spec.background = spec.profile.background or spec.background
            spec.accent = spec.accent or spec.profile.accent
        if (spec.profile.layout_index is not None
                and spec.profile.layout_content_is_chart_area):
            content = _largest_content_placeholder(
                prs.slide_layouts[spec.profile.layout_index])
            if content is not None:
                spec.chart_slot = Slot(
                    slide_index=-1,  # -1: a slide is added per chart, not reused
                    left=int(content.left or 0), top=int(content.top or 0),
                    width=int(content.width or 0), height=int(content.height or 0),
                    name="chart")
    return spec


def _largest_content_placeholder(layout):
    """The placeholder a chart should occupy: the biggest content area."""
    from pptx.enum.shapes import PP_PLACEHOLDER

    wanted = {PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.BODY}
    best, area = None, 0
    for ph in layout.placeholders:
        if ph.placeholder_format.type not in wanted:
            continue
        size = int(ph.width or 0) * int(ph.height or 0)
        if size > area:
            best, area = ph, size
    return best


def attendo_interim_spec() -> TemplateStyleSpec:
    """Interim proxy style spec from the Attendo deck. Satisfies REQ-C-27a (renders
    against *a* spec); REQ-C-27b (match the *client's* spec) remains BLOCKED — marked
    via spec_source/matches_client_spec."""
    spec = load_style_spec(str(config.ATTENDO_TEMPLATE))
    spec.spec_source = "attendo-interim-proxy"
    return spec

# ---------------------------------------------------------------------------
# What an author says, on top of what we harvested
# ---------------------------------------------------------------------------

#: Geometry is given in INCHES, because that is what PowerPoint's own ruler
#: says and what somebody dragging a box on screen is really adjusting. EMU is
#: an implementation detail of the file format and belongs on this side of the
#: boundary only.
_EMU_PER_INCH = 914400


def _hex(value) -> str:
    """A colour an author typed, or "" for "no opinion"."""
    text = str(value or "").strip().lstrip("#").upper()
    return text if len(text) in (6, 8) and all(c in "0123456789ABCDEF" for c in text) else ""


def _num(value):
    try:
        out = float(value)
    except (TypeError, ValueError):
        return None
    return out if out > 0 or out == 0 else None


def apply_template_overrides(spec, overrides: dict | None) -> None:
    """Apply an author's corrections to a harvested style, in place.

    Harvesting is a guess made from a file nobody wrote for us, and each of the
    three customer templates we have is unusual in a different way. The rules
    are better now; the next template will find a new way. So an author can say
    what the answer is and be believed — per template, because everything we
    have seen is a property of the template rather than of one deck.

    BLANK MEANS INHERIT, everywhere. A template nobody has touched renders
    exactly as the harvester decided, which is what keeps this an override
    rather than a second place the truth has to be maintained.
    """
    if not overrides:
        return

    title = overrides.get("title") or {}
    content = overrides.get("content") or {}

    if _hex(title.get("colour")):
        spec.ink = _hex(title["colour"])
        spec.title_colour = _hex(title["colour"])
    if _num(title.get("size")):
        spec.title_size_pt = _num(title["size"])
    if _hex(overrides.get("accent")):
        spec.accent = _hex(overrides["accent"])
    if _hex(overrides.get("background")):
        spec.background = _hex(overrides["background"])

    profile = getattr(spec, "profile", None)
    if profile is not None and getattr(profile, "title", None) is not None:
        _apply_text(profile.title, title)
        _apply_box(profile.title, title)
    # AND the font role, which is where the drawn title actually gets its size:
    # `build_spec` reads `fonts["title"]`, not the profile, so setting only the
    # profile stored the number, showed it back in the editor, and left the
    # headline exactly as it was.
    if title.get("font") or _num(title.get("size")):
        had_family, had_size = spec.font_for("title")
        _set_fonts(spec, {"title": (str(title.get("font") or had_family),
                                    int(_num(title.get("size")) or had_size))})
    _apply_chart_text(spec, content)

    subtitle = overrides.get("subtitle") or {}
    if subtitle.get("font"):
        spec.subtitle_font = str(subtitle["font"])
    if _num(subtitle.get("size")):
        spec.subtitle_size_pt = _num(subtitle["size"])
    if _hex(subtitle.get("colour")):
        spec.subtitle_colour = _hex(subtitle["colour"])

    # The footer is the "n = 3144" line. It is a font role already, so its size
    # and face go where every other chart-text size goes.
    footer = overrides.get("footer") or {}
    if footer.get("font") or _num(footer.get("size")):
        had_family, had_size = spec.font_for("n_annotation")
        _set_fonts(spec, {"n_annotation": (
            str(footer.get("font") or had_family),
            int(_num(footer.get("size")) or had_size))})
    if _hex(footer.get("colour")):
        spec.footer_colour = _hex(footer["colour"])

    _apply_slot(spec, content)


#: The chart's own text. "Content font/size" means these: the row labels down
#: the side, the numbers in the bars, the legend and the axis. Naming them is
#: what makes the setting do the thing an author asked for — the first complaint
#: about a rendered slide was row labels overlapping each other, and the answer
#: to that is a smaller category font, not a different body font somewhere.
_CHART_TEXT_ROLES = ("category_names", "data_labels", "axis_values",
                     "axis_names", "legend")


def _set_fonts(spec, changes: dict[str, tuple[str, int]]) -> None:
    """Change what `font_for` answers.

    It reads `_fonts`, which the constructor set — assigning a `fonts` attribute
    instead made a new one that nothing reads, so a content or footer size was
    stored, echoed back by the API, and had no effect on a single pixel.

    The resolved spec is dropped with it: it was built from the sizes as they
    were, and anything reading it would go on answering with those.
    """
    fonts = dict(getattr(spec, "_fonts", {}) or {})
    fonts.update(changes)
    spec._fonts = fonts
    spec.resolved_spec = None


def _apply_chart_text(spec, given: dict) -> None:
    family = str(given.get("font") or "").strip()
    size = _num(given.get("size"))
    if not family and not size:
        return
    if family:
        spec.body_font = family
    _set_fonts(spec, {
        role: (family or spec.font_for(role)[0],
               int(size) if size else spec.font_for(role)[1])
        for role in _CHART_TEXT_ROLES
    })


def _apply_text(style, given: dict) -> None:
    if given.get("font"):
        style.font = str(given["font"])
    size = _num(given.get("size"))
    if size:
        style.size_pt = size
        # Not a starting point to shrink from: see TextStyle.size_locked.
        if hasattr(style, "size_locked"):
            style.size_locked = True


def _apply_box(style, given: dict) -> None:
    for key, attr in (("x", "left"), ("y", "top"), ("w", "width"), ("h", "height")):
        value = _num(given.get(key))
        if value is not None and hasattr(style, attr):
            setattr(style, attr, int(value * _EMU_PER_INCH))


def effective_content_rect(spec) -> tuple[int, int, int, int]:
    """Where the chart actually goes on this style, in EMU.

    The layout's own content placeholder when we are taking it, and otherwise
    the renderer's own placement — under the title, inside the template's side
    margins. There is no third answer, and every caller needs the same one: the
    editor draws this rectangle, and a drag has to amend THIS rectangle rather
    than an absence.
    """
    sw = int(getattr(spec, "slide_width", 0) or 0)
    sh = int(getattr(spec, "slide_height", 0) or 0)
    slot = getattr(spec, "chart_slot", None)
    if slot is not None and int(slot.width or 0) > 0 and int(slot.height or 0) > 0:
        rect = (int(slot.left), int(slot.top), int(slot.width), int(slot.height))
    else:
        profile = getattr(spec, "profile", None)
        title = getattr(profile, "title", None) if profile is not None else None
        if title is not None and getattr(title, "positioned", False):
            from reportbuilder.render.image.slide_chrome import harvested_chart_box

            rect = harvested_chart_box(profile, "", sw, sh)
        else:
            margin = int(0.05 * sw)
            rect = (margin, int(0.28 * sh), sw - 2 * margin, int(0.60 * sh))

    left, top, width, height = rect
    inch = _EMU_PER_INCH
    width = max(inch, min(width, sw or width))
    height = max(inch, min(height, sh or height))
    left = max(0, min(left, (sw or left + width) - width))
    top = max(0, min(top, (sh or top + height) - height))
    return left, top, width, height


def _apply_slot(spec, given: dict) -> None:
    """The chart's own box. Absent edges keep whatever they had."""
    edges = {k: _num(given.get(k)) for k in ("x", "y", "w", "h")}
    if all(v is None for v in edges.values()):
        return
    current = getattr(spec, "chart_slot", None)
    if current is None:
        # No slot does not mean no rectangle. Where a template offers us no
        # usable content area — Arla, whose every layout is two-column — the
        # renderer places the chart itself, and THAT is what a drag is amending.
        # Requiring all four numbers here meant dragging such a chart, which
        # sends only x and y, changed nothing at all.
        base_l, base_t, base_w, base_h = effective_content_rect(spec)
        current = Slot(slide_index=-1, left=base_l, top=base_t,
                       width=base_w, height=base_h, name="chart")
    spec.chart_slot = Slot(
        slide_index=current.slide_index,
        left=int(edges["x"] * _EMU_PER_INCH) if edges["x"] is not None else current.left,
        top=int(edges["y"] * _EMU_PER_INCH) if edges["y"] is not None else current.top,
        width=int(edges["w"] * _EMU_PER_INCH) if edges["w"] is not None else current.width,
        height=int(edges["h"] * _EMU_PER_INCH) if edges["h"] is not None else current.height,
        name=current.name)
