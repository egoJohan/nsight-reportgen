"""What to borrow from a customer's template, and nothing more.

Johan's spec: take the title's style and POSITION, the subtitle font, graphics
belonging to the title, and the furniture that repeats on every slide (logo,
footer, background). nSight then draws its own slide with those. The chart and
the rest of the text stay at sizes and positions nSight chooses.

The reason this exists rather than either of the two obvious approaches:

  * Building from the template's LAYOUTS works only when the brand is in the
    layouts. `Attendo Bränditutkimus Marraskuu 2025.pptx` is like that — 28
    layouts, Century Gothic, navy palette, and 42 of its 56 slides sit on "1
    layout area".
  * Copying the template's SLIDES hands us someone else's slide geometry, and
    Johan does not want the slides reused.

`attendo_agent_deck.pptx` is the case that breaks the first approach: stock
Office masters (Calibri, the 4F81BD/C0504D palette, English layout names) with
all 20 slides on the `Blank` layout carrying ~7 hand-drawn shapes each. Synsam's
deck is the same story in a brand template — 98 of 147 slides on "Tom". The
design is on the slides, not in the layouts.

So the profile is harvested from whichever place the design actually lives, and
the deck itself says which: the layout its own slides mostly sit on, when that
layout can hold a headline and a chart, otherwise its most representative slide.
The profile shape is the same either way, but how the renderer uses it is not —
see `layout_index`.
"""
from __future__ import annotations

import io
from collections import Counter
from copy import deepcopy
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from reportbuilder.render.template_check import rank_layouts, theme_root

_A = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
_P = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main",
      "a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

# A shape covering this much of the slide is background, not content.
_BACKDROP_AREA = 0.75
# Text this large is a title rather than body copy, when nothing else says so.
_TITLE_MIN_PT = 14.0
# A chart-capable layout has to be at least this much of the deck's mainstay
# before we believe the design lives in the layouts.
_MAINSTAY_SHARE = 0.5
# Furniture is what repeats: on this share of the body slides, and at least
# this many of them before repetition means anything at all.
_REPEAT_SHARE = 0.4
_MIN_REPEATS = 3
# A picture bigger than this much of the slide is a photo or a chart, not a mark.
_LOGO_AREA = 0.12


@dataclass
class TextStyle:
    """How a piece of text looks, and (for the title) where it sits."""

    font: str = ""
    size_pt: float = 0.0
    bold: bool | None = None
    #: PowerPoint's `cap="all"`: the text renders upper-case whatever we type,
    #: which makes a line about a tenth wider than the same string in mixed case.
    caps: bool = False
    colour: str = ""          # hex, no '#'
    left: int = 0
    top: int = 0
    width: int = 0
    height: int = 0

    @property
    def positioned(self) -> bool:
        return self.width > 0 and self.height > 0


@dataclass
class Furniture:
    """One borrowed shape, in a form that survives the move to another deck.

    A shape's XML is only half of it: a picture's `r:embed` points at an image
    part of the slide it came from, and that part goes when render_report strips
    the template's slides. So a picture travels as bytes and is redrawn, while
    everything else — bands, rules, backdrops — travels as XML.
    """

    element: object | None = None
    image: bytes = b""
    box: tuple = (0, 0, 0, 0)


@dataclass
class TemplateProfile:
    """The borrowed part of a template."""

    source: str = ""                  # "layout:<name>" | "slide:<n>" | ""
    #: Set when the design lives in a layout: build each slide FROM this layout
    #: and the client's furniture arrives by inheritance, no cloning. None means
    #: the design was harvested off a slide, so `furniture` has to be cloned onto
    #: a blank slide instead.
    layout_index: int | None = None
    title: TextStyle = field(default_factory=TextStyle)
    subtitle_font: str = ""
    #: Shape XML to clone onto every generated slide: background, logo, footer,
    #: and any graphic sitting with the title. Order preserved so z-order is.
    #: Empty for a layout source, where those shapes are inherited already.
    furniture: list = field(default_factory=list)
    slide_width: int = 0
    slide_height: int = 0

    @property
    def usable(self) -> bool:
        return bool(self.title.font or self.title.positioned or self.furniture)

    def describe(self) -> dict:
        """A summary safe to log or show in the UI."""
        return {"source": self.source, "layout_index": self.layout_index,
                "title_font": self.title.font,
                "title_pt": self.title.size_pt, "title_colour": self.title.colour,
                "title_box_in": [round(Emu(v).inches, 2) for v in
                                 (self.title.left, self.title.top,
                                  self.title.width, self.title.height)],
                "subtitle_font": self.subtitle_font,
                "furniture_shapes": len(self.furniture)}


# --- the theme, as a resolver -----------------------------------------------

class _Brand:
    """Resolves the references a template uses instead of naming things.

    A layout almost never says "Century Gothic, 28pt, navy". It says `+mj-lt`
    and `tx1`, and PowerPoint looks those up in the theme — which is why reading
    a title's run properties alone came back empty for all four real templates.
    """

    def __init__(self, prs):
        self._root = theme_root(prs)
        # The master remaps the slide-facing names onto theme slots, and the
        # remap is not always the identity: `tx1` is usually `dk1`, but a
        # reversed master swaps them.
        self._map = {}
        try:
            el = prs.slide_master._element.find(
                "{http://schemas.openxmlformats.org/presentationml/2006/main}clrMap")
            if el is not None:
                self._map = dict(el.attrib)
        except AttributeError:
            pass

    def font(self, typeface: str) -> str:
        """A typeface name, following `+mj-lt` / `+mn-lt` into the theme."""
        if not typeface:
            return ""
        if not typeface.startswith("+"):
            return typeface
        kind = "major" if typeface.startswith("+mj") else "minor"
        if self._root is None:
            return ""
        el = self._root.find(f".//a:fontScheme/a:{kind}Font/a:latin", _A)
        return el.get("typeface") or "" if el is not None else ""

    def colour(self, name: str) -> str:
        """A scheme colour name (`tx1`, `accent1`) as a hex string."""
        if not name or self._root is None:
            return ""
        slot = self._map.get(name, name)
        el = self._root.find(f".//a:clrScheme/a:{slot}", _A)
        if el is None:
            return ""
        srgb = el.find("a:srgbClr", _A)
        if srgb is not None and srgb.get("val"):
            return srgb.get("val").upper()
        sysc = el.find("a:sysClr", _A)
        if sysc is not None and sysc.get("lastClr"):
            return sysc.get("lastClr").upper()
        return ""

    def fill_colour(self, solid_fill) -> str:
        """The hex of an `<a:solidFill>`, whether it names an RGB or the theme."""
        if solid_fill is None:
            return ""
        srgb = solid_fill.find("a:srgbClr", _A)
        if srgb is not None and srgb.get("val"):
            return srgb.get("val").upper()
        scheme = solid_fill.find("a:schemeClr", _A)
        if scheme is not None:
            return self.colour(scheme.get("val") or "")
        return ""


# --- reading text style -----------------------------------------------------

def _first_run(shape):
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                return run, para
            return None, para
    except AttributeError:
        pass
    return None, None


def _style_of(shape) -> TextStyle:
    """The style stated on *shape*'s first run, plus the shape's own box.

    Most of these fields come back empty on a layout: what the run states is
    only the part the designer overrode. `_inherited` fills in the rest.
    """
    st = TextStyle(left=int(shape.left or 0), top=int(shape.top or 0),
                   width=int(shape.width or 0), height=int(shape.height or 0))
    run, para = _first_run(shape)
    holder = run.font if run is not None else (para.font if para is not None else None)
    if holder is None:
        return st
    st.font = holder.name or ""
    if holder.size is not None:
        st.size_pt = float(holder.size.pt)
    st.bold = holder.bold
    try:
        rpr = (run if run is not None else para)._rPr if run is not None else None
        if rpr is not None:
            st.caps = rpr.get("cap") == "all"
    except AttributeError:
        pass
    try:
        # An explicit RGB only. A theme-colour reference is resolved from the
        # XML below, where the scheme name is still readable.
        if holder.color is not None and holder.color.type is not None:
            st.colour = str(holder.color.rgb)
    except (AttributeError, ValueError):
        pass
    return st


def _from_rpr(rpr, brand: _Brand) -> TextStyle:
    """The style a `<a:defRPr>` or `<a:rPr>` element states."""
    st = TextStyle()
    if rpr is None:
        return st
    if rpr.get("sz"):
        try:
            st.size_pt = int(rpr.get("sz")) / 100.0
        except ValueError:
            pass
    if rpr.get("b") is not None:
        st.bold = rpr.get("b") in ("1", "true")
    st.caps = rpr.get("cap") == "all"
    latin = rpr.find("a:latin", _A)
    if latin is not None:
        st.font = brand.font(latin.get("typeface") or "")
    st.colour = brand.fill_colour(rpr.find("a:solidFill", _A))
    return st


def _merge(base: TextStyle, extra: TextStyle) -> TextStyle:
    """Fill what *base* does not state from *extra*. Position is never inherited."""
    for attr in ("font", "size_pt", "colour"):
        if not getattr(base, attr):
            setattr(base, attr, getattr(extra, attr))
    if base.bold is None:
        base.bold = extra.bold
    base.caps = base.caps or extra.caps
    return base


def _own_list_style(shape):
    """A shape's own `lstStyle` default, which is where a layout restyles a
    placeholder it inherited."""
    try:
        return shape._element.find(
            "./p:txBody/a:lstStyle/a:lvl1pPr/a:defRPr", _P)
    except AttributeError:
        return None


def _master_title_rpr(prs):
    """The master's title style — where "28pt, +mj-lt, tx1" is actually written."""
    try:
        return prs.slide_master._element.find(
            "./p:txStyles/p:titleStyle/a:lvl1pPr/a:defRPr", _P)
    except AttributeError:
        return None


def _text_shapes(container) -> list:
    out = []
    for sh in container.shapes:
        try:
            if sh.has_text_frame:
                out.append(sh)
        except AttributeError:
            continue
    return out


def _looks_like_title(shape) -> bool:
    if shape.is_placeholder:
        try:
            # 13 = TITLE, 0 = CENTER_TITLE in python-pptx's enum values.
            if shape.placeholder_format.idx == 0:
                return True
        except (AttributeError, ValueError):
            pass
    name = (shape.name or "").lower()
    return "title" in name or "otsik" in name or "rubrik" in name


def _pick_title(container) -> object | None:
    """The shape acting as the title.

    A real title placeholder settles it. Otherwise the biggest text near the
    top wins, with a nudge for a shape whose NAME says title in any of the
    three languages these templates come in — the name alone must not decide,
    because Synsam's slides carry four boxes all called "Rubrik 10" and the
    first of them is empty and sits at the foot of the slide.
    """
    shapes = _text_shapes(container)
    if not shapes:
        return None
    for sh in shapes:
        if sh.is_placeholder and _looks_like_title(sh):
            return sh
    with_text = [sh for sh in shapes if (sh.text_frame.text or "").strip()]
    candidates = with_text or shapes
    scored = []
    for sh in candidates:
        st = _style_of(sh)
        # Named-as-title first, then bigger type; ties break toward the top.
        scored.append((_looks_like_title(sh), st.size_pt or _TITLE_MIN_PT,
                       -int(sh.top or 0), sh))
    scored.sort(key=lambda t: t[:3], reverse=True)
    return scored[0][3]


def _title_style(prs, shape, brand: _Brand) -> TextStyle:
    """*shape*'s style, resolved through everything it inherits from.

    Run properties first (a hand-drawn title states them outright), then the
    shape's own list style, then — for a real title placeholder — the master's
    title style, and finally the theme's heading font. Without this a layout
    source yielded a title with no font, no size and no colour.
    """
    st = _style_of(shape)
    st = _merge(st, _from_rpr(_own_list_style(shape), brand))
    if _looks_like_title(shape):
        st = _merge(st, _from_rpr(_master_title_rpr(prs), brand))
    if not st.font:
        st.font = brand.font("+mj-lt")
    return st


def _pick_subtitle(container, title, brand: _Brand) -> str:
    """The body font: the next text shape below the title, else the theme's."""
    below = []
    for sh in _text_shapes(container):
        if sh is title:
            continue
        if int(sh.top or 0) >= int(getattr(title, "top", 0) or 0):
            below.append(sh)
    below.sort(key=lambda sh: int(sh.top or 0))
    for sh in below:
        st = _style_of(sh)
        if st.font:
            return st.font
    return brand.font("+mn-lt")


# --- choosing where to harvest from -----------------------------------------

def _decoration(container) -> list:
    """Shapes that are furniture rather than content placeholders."""
    out = []
    for sh in container.shapes:
        if sh.is_placeholder:
            continue
        out.append(sh)
    return out


def _design_layout(prs):
    """(index, layout) of the layout carrying this template's design, or None.

    Two questions, in this order. Which layouts COULD hold a headline and a
    chart — that is template_check's area ranking, which already identifies "1
    layout area", "Innehåll" and "Title and Content" across three languages.
    And then: is one of them what the deck itself is built on?

    The second question is the one that separates the two kinds of customer
    file. Attendo puts 42 of its 56 slides on "1 layout area", so its design is
    in the layouts and a slide built from that layout inherits it. Synsam puts
    98 of 147 on "Tom" and Johan's agent deck puts all 20 on `Blank` — a blank
    layout is not a candidate at all, which is exactly the signal that those
    decks were drawn by hand and must be harvested from a slide.

    Ranking on decoration count instead, as the first version did, picked
    Attendo's "Agenda slide" and Synsam's "Slutbild" (title box 0.14in tall,
    above the top edge of the slide). Decoration says a layout is pretty, not
    that it is the one the deck is made of.
    """
    candidates = {c.name: c for c in rank_layouts(prs)}
    if not candidates:
        return None

    usage = Counter()
    for slide in prs.slides:
        try:
            usage[slide.slide_layout.name] += 1
        except AttributeError:
            continue
    if not usage:
        # A template with no slides of its own — a real .potx, or one already
        # stripped. Nothing to learn from usage, so take the best candidate.
        best = min(candidates.values(), key=lambda c: (-c.score, c.index))
        return best.index, prs.slide_layouts[best.index]

    mainstay = usage.most_common(1)[0][1]
    for name, count in usage.most_common():
        if name not in candidates:
            continue
        # A candidate the deck barely touches is not where its design lives.
        if count < _MAINSTAY_SHARE * mainstay:
            break
        idx = candidates[name].index
        return idx, prs.slide_layouts[idx]
    return None


def _body_slides(prs) -> list:
    """The deck's ordinary content slides: the ones on its mainstay layout.

    A deck drawn by hand still says which layout it draws ON — "Tom" for
    Synsam, `Blank` for the agent deck — and its covers and photo pages sit
    elsewhere. Skipping that step harvested Synsam from slide 58, a full-bleed
    photo page whose title sits 5.8in down the slide.
    """
    if not len(prs.slides):
        return []
    usage = Counter()
    for slide in prs.slides:
        try:
            usage[slide.slide_layout.name] += 1
        except AttributeError:
            continue
    if not usage:
        return list(prs.slides)
    mainstay = usage.most_common(1)[0][0]
    return [s for s in prs.slides
            if getattr(s.slide_layout, "name", None) == mainstay]


def _representative_slide(prs):
    """The template's most typical content slide.

    Among the body slides, chosen by how many share its shape count: the
    ordinary content page outnumbers the one-off. Ties go to the later slide,
    which is more likely to be body content than a cover.
    """
    body = _body_slides(prs)
    if not body:
        return None
    counts = Counter(len(s.shapes) for s in body)
    # Ignore near-empty slides: they carry no design to harvest.
    usable = [(n, c) for n, c in counts.items() if n >= 2]
    if not usable:
        return None
    target = max(usable, key=lambda t: (t[1], t[0]))[0]
    match = [s for s in body if len(s.shapes) == target]
    return match[len(match) // 2] if match else None


def _area_share(shape, sw: int, sh_: int) -> float:
    if sw <= 0 or sh_ <= 0:
        return 0.0
    try:
        return int(shape.width or 0) * int(shape.height or 0) / (sw * sh_)
    except (TypeError, ValueError):
        return 0.0


def _is_backdrop(shape, sw: int, sh_: int) -> bool:
    return _area_share(shape, sw, sh_) >= _BACKDROP_AREA


def _off_slide(shape, sw: int, sh_: int) -> bool:
    """Wholly outside the slide, so it is not part of the design whatever it is.

    Synsam's body slides all carry a rectangle 7.7in down a 7.5in slide —
    invisible in the deck, and no reason to copy it into ours.
    """
    try:
        left, top = int(shape.left or 0), int(shape.top or 0)
        right = left + int(shape.width or 0)
        bottom = top + int(shape.height or 0)
    except (TypeError, ValueError):
        return False
    return right <= 0 or bottom <= 0 or left >= sw or top >= sh_


def _is_content_image(shape, sw: int, sh_: int) -> bool:
    """A picture too big to be a mark. The agent deck's own chart is a picture
    in the same box on all 20 slides, so repetition alone would clone last
    year's chart onto every new slide. A full-bleed backdrop is exempt: at that
    size a picture is the background, which is design."""
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    share = _area_share(shape, sw, sh_)
    return _LOGO_AREA < share < _BACKDROP_AREA


def _shape_key(shape) -> tuple:
    """What makes a shape "the same shape" from one slide to the next: its name
    and its box, to a hundredth of an inch."""
    def hundredths(v):
        return int(v or 0) // 9144
    return (shape.name or "",
            hundredths(shape.left), hundredths(shape.top),
            hundredths(shape.width), hundredths(shape.height))


def _shape_text(shape) -> str:
    try:
        if shape.has_text_frame:
            return (shape.text_frame.text or "").strip()
    except AttributeError:
        pass
    return ""


def _repeating(prs) -> set:
    """The shapes this deck puts in the same place slide after slide.

    That IS the definition of furniture: a logo, a footer rule, a background
    band. Taking one slide's decoration wholesale instead swept up its content
    — Synsam's chart slides carry a "5+6+7" label and an "n=1549 / 282 ..."
    line, both short enough to pass for a footer, neither of which belongs on
    next year's slide. A box that holds different words on different slides is
    a content slot, however reliably it appears, so it is excluded too.

    Returns an empty set when the deck has too few body slides to tell, and the
    caller then falls back to keeping the decoration it found.
    """
    body = _body_slides(prs)
    if len(body) < _MIN_REPEATS:
        return set()
    tally: dict = {}
    for slide in body:
        for key, text in {(_shape_key(sh), _shape_text(sh))
                          for sh in _decoration(slide)}:
            rec = tally.setdefault(key, [0, set()])
            rec[0] += 1
            rec[1].add(text)
    floor = max(_MIN_REPEATS, _REPEAT_SHARE * len(body))
    return {key for key, (count, texts) in tally.items()
            if count >= floor and len(texts - {""}) <= 1}


def _furniture(container, title, sw: int, sh_: int, keep: set | None = None) -> list:
    """Shape XML to clone onto every slide: background, logos, footer, rules.

    Excludes the title itself (nSight draws that), charts and tables (data, not
    design), and anything holding real body text — a chart placeholder's prompt
    or last year's bullet points must not travel onto this year's slide. A
    full-slide backdrop is kept even though it has no text, because it IS the
    design. When *keep* is given, a shape must also be one that repeats.
    """
    out = []
    for sh in _decoration(container):
        if sh is title:
            continue
        if getattr(sh, "has_chart", False) or getattr(sh, "has_table", False):
            continue
        if _off_slide(sh, sw, sh_) or _is_content_image(sh, sw, sh_):
            continue
        if keep and _shape_key(sh) not in keep:
            continue
        text = _shape_text(sh)
        if text and not _is_backdrop(sh, sw, sh_):
            # Short text is furniture (a footer, a page number, a brand line);
            # a paragraph is somebody's content.
            if len(text) > 60:
                continue
        box = (int(sh.left or 0), int(sh.top or 0),
               int(sh.width or 0), int(sh.height or 0))
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                out.append(Furniture(image=sh.image.blob, box=box))
            except (AttributeError, ValueError):
                continue
        elif sh._element.find(".//a:blip", _A) is not None:
            # A group or a picture-filled shape: its image lives in a part we
            # are not copying, and a dangling reference makes PowerPoint offer
            # to repair the file. Better to lose the shape than the deck.
            continue
        else:
            out.append(Furniture(element=deepcopy(sh._element), box=box))
    return out


def clone_furniture(slide, items) -> int:
    """Redraw harvested furniture onto *slide*, bottom of the z-order.

    Call before anything else is added: PowerPoint z-orders by document order,
    so the chart and our text have to come after the background they sit on.
    """
    placed = 0
    for item in items:
        try:
            if item.image:
                slide.shapes.add_picture(io.BytesIO(item.image), *item.box)
            elif item.element is not None:
                slide.shapes._spTree.append(deepcopy(item.element))
            else:
                continue
            placed += 1
        except Exception:  # noqa: BLE001 — one bad shape must not lose the slide
            continue
    return placed


# --- the entry point --------------------------------------------------------

def extract_profile(template_path: str) -> TemplateProfile:
    """Harvest the borrowable part of *template_path*.

    Never raises on a readable .pptx: a template we cannot read a title from
    still yields a profile with whatever furniture it had, and the renderer
    falls back to its own placement for the rest.
    """
    prs = Presentation(template_path)
    sw, sh_ = int(prs.slide_width or 0), int(prs.slide_height or 0)
    profile = TemplateProfile(slide_width=sw, slide_height=sh_)
    brand = _Brand(prs)

    found = _design_layout(prs)
    if found is not None:
        index, source = found
        profile.layout_index = index
        profile.source = f"layout:{source.name}"
    else:
        source = _representative_slide(prs)
        if source is None:
            return profile
        profile.source = f"slide:{list(prs.slides).index(source) + 1}"

    title = _pick_title(source)
    if title is not None:
        profile.title = _title_style(prs, title, brand)
        profile.subtitle_font = _pick_subtitle(source, title, brand)
    else:
        profile.subtitle_font = brand.font("+mn-lt")

    if profile.layout_index is None:
        # Only a slide's furniture needs cloning. A layout's own decoration, and
        # the master's, land on every slide built from it without our help —
        # cloning them there would draw the logo twice.
        profile.furniture = _furniture(source, title, sw, sh_, _repeating(prs))
    return profile
