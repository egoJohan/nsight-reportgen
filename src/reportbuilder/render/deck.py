"""Deck assembly entry point (Task 5.14 — design §C2).

render_report: open template, render each ChartSpec into its slot,
dispatch by render_mode.

render_to_file: convenience wrapper that saves to disk and returns the path
(REQ-C-29a).
"""
from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError
from pptx.util import Inches

from reportbuilder.model.report import Report, is_demographics_grid, renders_as_bullets
from reportbuilder.render.base import RenderContext, Slot
from reportbuilder.render.elements import apply_elements, add_n_annotation, add_filter_annotation
from reportbuilder.render.image.slide_chrome import (
    add_image_slide_chrome, content_floor, harvested_chart_box, harvested_profile,
    harvested_title_box, slide_headline,
)
from reportbuilder.render.image.special_slide import render_special_slide
from reportbuilder.render.image.demographics_grid import render_demographics_grid
from reportbuilder.render.image._mpl import render_empty_chart, series_is_empty
import reportbuilder.render.plugins as _plugins  # registers all plugins as side-effect


# ---------------------------------------------------------------------------
# Completeness and purity guards (REQ-C-18, REQ-C-23a)
# ---------------------------------------------------------------------------

class CompletenessError(Exception):
    """Generated deck doesn't match the report definition (REQ-C-18)."""


class NativePurityError(Exception):
    """A native-mode report has a picture shape in a chart slot (REQ-C-23a)."""


def _content_placeholders():
    """Placeholder types a chart occupies — imported lazily so this module keeps
    its light import cost."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    return {PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.BODY}


def _blank_layout(prs):
    """The emptiest layout in this presentation.

    `slide_layouts[6]` is blank only in python-pptx's OWN default template. In
    Attendo's it is "1 layoutarea + 1/3 image", whose picture placeholder lands
    on every synthesised slide — 6 charts produced 11 pictures and the
    completeness guard rejected the deck.

    So pick by inspection, and rank on the DECORATION a layout carries rather
    than on its placeholders. "End slide Thank You" declares zero placeholders
    yet holds four artwork shapes including a logo PICTURE — it won an earlier
    version of this ranking and would have stamped that logo on every chart
    slide. An unfilled placeholder is invisible in the export; a picture on the
    layout is not.

    Order: no real picture first, then least non-placeholder artwork, then
    fewest shapes overall.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    best, best_key = None, None
    for layout in prs.slide_layouts:
        try:
            artwork = [sh for sh in layout.shapes if not sh.is_placeholder]
            has_picture = any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE
                              for sh in artwork)
            key = (has_picture, len(artwork), len(layout.shapes))
        except Exception:  # noqa: BLE001 — a malformed layout is simply not chosen
            continue
        if best_key is None or key < best_key:
            best, best_key = layout, key
    # A presentation with no layouts at all cannot happen via python-pptx, but
    # index 6 is the historical behaviour and a safe last resort.
    return best if best is not None else prs.slide_layouts[6]


def _strip_slides(prs) -> int:
    """Remove every pre-existing slide from an opened template. Returns the count.

    Keeps masters, layouts and theme — everything that defines the look — and
    discards the example or finished slides that came with the file.
    """
    slide_ids = prs.slides._sldIdLst
    removed = 0
    for sld in list(slide_ids):
        rid = sld.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        slide_ids.remove(sld)
        if rid:
            try:
                prs.part.drop_rel(rid)
            except KeyError:
                pass
        removed += 1
    return removed


def _count_chart_shapes(prs: Presentation) -> tuple[int, int]:
    """Return (chart_count, picture_count) across all slides in *prs*."""
    charts = sum(
        1 for s in prs.slides for sh in s.shapes if getattr(sh, "has_chart", False)
    )
    pics = sum(
        1 for s in prs.slides for sh in s.shapes
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    return charts, pics


def assert_complete(prs: Presentation, report: Report,
                    expected_pics: int | None = None) -> None:
    """The deck contains exactly one rendered chart object per ChartSpec, nothing extra.

    Native mode: counts c:chart shapes.  Image mode: counts PICTURE shapes.
    Raises CompletenessError if the tally doesn't match. Bullet slides add no
    chart object; a demographics grid adds several pictures, so the caller passes
    ``expected_pics`` (computed where the series are known). (REQ-C-18)
    """
    charts, pics = _count_chart_shapes(prs)
    rendered = charts if report.render_mode == "native" else pics
    if expected_pics is not None and report.render_mode != "native":
        expected = expected_pics
    else:
        # Bullet/grid slides don't add exactly one picture; exclude them.
        expected = len([
            c for c in report.charts
            if not getattr(c, "excluded", False)
            and not renders_as_bullets(c) and not is_demographics_grid(c)
        ])
    if rendered != expected:
        raise CompletenessError(
            f"expected {expected} {report.render_mode} chart objects, found {rendered}"
        )


def assert_no_pictures_in_chart_slots(prs: Presentation, report: Report, style=None) -> None:
    """Native-mode reports must contain ZERO picture shapes (editability gate, REQ-C-23a).

    No-op for image mode (image reports legitimately use pictures).
    """
    if report.render_mode != "native":
        return
    _charts, pics = _count_chart_shapes(prs)
    if pics > 0:
        raise NativePurityError(
            f"native-mode report has {pics} picture shape(s) in chart slots"
        )


class RenderCancelled(Exception):
    """Raised to abort a deck render mid-way when the caller signals cancellation
    (e.g. the client aborted the request). Checked between slides so a long run
    (hundreds of slides) stops promptly instead of grinding to the end."""


def render_report(
    report: Report,
    series_by_ref: dict,
    style,
    titles: dict | None = None,
    cancel_check=None,
) -> Presentation:
    """Open the template, render each ChartSpec into its slot.

    Parameters
    ----------
    report:
        The Report definition (charts, render_mode, template_ref).
    series_by_ref:
        Maps ChartSpec.question_ref -> SeriesResult.
    style:
        A StyleSpec (base) or TemplateStyleSpec.  When a TemplateStyleSpec
        carries a template source file, that Presentation is used; otherwise
        a blank Presentation is created.
    titles:
        Optional mapping question_ref -> chart title text.  When omitted,
        chart titles default to "".
    """
    # --- Open or create the Presentation ---
    spec_source = getattr(style, "spec_source", None)
    # Only try to open if spec_source looks like a real path (not the
    # sentinel strings used by TemplateStyleSpec when built from load_style_spec).
    # A safe heuristic: try Presentation(spec_source) and fall back on any error.
    prs = None
    if spec_source and spec_source not in ("generic", "attendo-interim-proxy"):
        try:
            prs = Presentation(spec_source)
        except (FileNotFoundError, PackageNotFoundError):
            prs = None
        if prs is not None:
            # A client "template" is usually a FINISHED deck — Attendo's is 56
            # slides of last year's report. Using their template means
            # inheriting their look, not their content, so the slides go while
            # the masters, layouts, theme, fonts and page size stay.
            #
            # Dropping the relationship as well as the sldIdLst entry matters:
            # left related, the slide parts survive the save and every generated
            # deck would carry megabytes of somebody else's finished report.
            _strip_slides(prs)

    if prs is None:
        # Wizard/generic reports render on a blank deck at 16:9 (13.333"×7.5"), the
        # modern presentation standard. Template-based decks keep their own size.
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    _titles = titles or {}

    for spec in report.charts:
        if getattr(spec, "excluded", False):
            continue          # unticked in Select — kept in the report, off the deck
        # Cooperative cancellation: bail out promptly between slides when signalled.
        if cancel_check is not None and cancel_check():
            raise RenderCancelled()
        # --- Resolve slot and slide ---
        # The headline is passed in because a harvested slide's chart starts
        # under the title, and how far down that is depends on how many lines
        # the title wraps to.
        slot = _resolve_slot(prs, style, spec.template_slot, report.render_mode,
                             title=slide_headline(spec, _titles.get(spec.question_ref, "")))
        # slide_index may reference an existing slide or was just appended
        slide = prs.slides[slot.slide_index]

        # --- Demographics grid: several compact charts on one slide ---
        if is_demographics_grid(spec):
            render_demographics_grid(slide, slot, style, spec, series_by_ref, _titles)
            continue

        # --- Bullet slides (special slides + themes): render text, no series ---
        if renders_as_bullets(spec):
            render_special_slide(
                slide, slot, style, spec, heading=_titles.get(spec.question_ref, "")
            )
            continue

        # --- Build context ---
        series = series_by_ref[spec.question_ref]
        title = _titles.get(spec.question_ref, "")
        ctx = RenderContext(
            slide=slide,
            slot=slot,
            style=style,
            spec=spec,
            series=series,
            fmt=spec.number_format,
            title=title,
        )

        # --- Dispatch via ChartPlugin registry (REQ-C-13) ---
        p = _plugins.plugin(spec.chart_type)
        if report.render_mode == "native":
            gf = p.native_build(ctx)
            apply_elements(gf.chart, ctx, title)
            add_n_annotation(ctx)
            add_filter_annotation(ctx)
        else:
            # Add house-style slide chrome first so chart image lands on top
            # (REQ-C-24a/h, REQ-C-25, REQ-C-27a, REQ-D-04)
            add_image_slide_chrome(ctx)
            # A chart with nothing to plot (e.g. a scale variable with no value
            # labels) degrades to a placeholder instead of crashing the builder.
            if series_is_empty(series):
                render_empty_chart(ctx)
            else:
                p.image_build(ctx)

    # Expected pictures: 1 per normal chart, 0 per bullet slide, and one per grid
    # cell that actually has a series (computed above).
    expected_pics = 0
    for spec in report.charts:
        if getattr(spec, "excluded", False):
            continue
        if renders_as_bullets(spec):
            continue
        if is_demographics_grid(spec):
            expected_pics += sum(
                1
                for c in (spec.options.get("charts") or [])
                if series_by_ref.get(c.get("question_ref")) is not None
            )
        else:
            expected_pics += 1
    assert_complete(prs, report, expected_pics=expected_pics)
    assert_no_pictures_in_chart_slots(prs, report, style)
    return prs


def render_to_file(
    report: Report,
    series_by_ref: dict,
    style,
    out_path: str,
    titles: dict | None = None,
    cancel_check=None,
) -> str:
    """Render report to *out_path* and return the path (REQ-C-29a)."""
    prs = render_report(report, series_by_ref, style, titles, cancel_check=cancel_check)
    prs.save(out_path)
    return out_path


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _resolve_slot(prs: Presentation, style, slot_name: str,
                  render_mode: str = "native", title: str = "") -> Slot:
    """Return a Slot for *slot_name*, falling back to a new blank slide.

    Tries style.slot(slot_name) first.  If that raises KeyError (slot not in
    template) or AttributeError (base StyleSpec has no slot() method), a fresh
    blank slide is appended to *prs* and a synthesised Slot is returned.

    For image mode, the slot starts lower on the slide to leave room for the
    house-style title / accent chrome added by add_image_slide_chrome.
    """
    try:
        return style.slot(slot_name)
    except (KeyError, AttributeError):
        pass

    # Preferred: build the slide from the template's OWN chart layout, using the
    # content placeholder the client's designer positioned. Only when a template
    # supplies neither do we fall back to a blank slide with guessed geometry.
    chart_layout = getattr(style, "chart_layout_index", None)
    chart_slot = getattr(style, "chart_slot", None)
    if chart_layout is not None and chart_slot is not None:
        slide = prs.slides.add_slide(prs.slide_layouts[chart_layout])
        # The content placeholder has given us its rectangle; leaving it on the
        # slide would show PowerPoint's "Click to add text" prompt behind the
        # chart, so it is removed once its geometry is taken.
        for shape in list(slide.shapes):
            if shape.is_placeholder and shape.has_text_frame and not shape.text_frame.text:
                if shape.placeholder_format.type in _content_placeholders():
                    shape._element.getparent().remove(shape._element)
        # The layout's content area starts right under the title box the
        # customer drew for the headline THEY wrote. A question runs longer, and
        # the subtitle has to go somewhere, so the chart starts below whichever
        # is lower — the placeholder or the title as it actually falls. Its
        # bottom edge does not move, so the customer's margin is kept.
        top, height = int(chart_slot.top), int(chart_slot.height)
        profile = getattr(style, "profile", None)
        if profile is not None and profile.title.positioned:
            _l, t_top, _w, t_height = harvested_title_box(profile, title)
            wanted = t_top + t_height + int(Inches(0.70))
            if wanted > top:
                height -= wanted - top
                top = wanted
        return Slot(slide_index=len(prs.slides) - 1, left=chart_slot.left,
                    top=top, width=chart_slot.width,
                    height=max(int(Inches(1.0)), height), name=slot_name)

    # Fallback: add a new blank slide and synthesise a slot covering most of it
    layout = _blank_layout(prs)  # emptiest layout in THIS template
    slide = prs.slides.add_slide(layout)
    slide_index = len(prs.slides) - 1

    # Slots scale to the slide so they fill it at any aspect (4:3 or 16:9): fixed
    # side margins + top chrome, the rest of the width/height is the content area.
    sw, sh = int(prs.slide_width), int(prs.slide_height)
    # A template whose design lives on its slides gives us the title's box, so
    # the chart can start under the customer's title instead of at a guessed
    # 1.9in — and sit inside their side margins rather than ours.
    profile = harvested_profile(style)
    if render_mode == "image" and profile is not None and profile.title.positioned:
        left, top, width, height = harvested_chart_box(
            profile, title, sw, sh, floor=content_floor(slide, sw, sh))
        return Slot(slide_index=slide_index, left=left, top=top,
                    width=width, height=height, name=slot_name)
    if render_mode == "image":
        # Leave ~1.9" at top for house-style title chrome (REQ-C-24a, REQ-D-04)
        return Slot(
            slide_index=slide_index,
            left=int(Inches(0.62)),
            top=int(Inches(1.9)),
            width=sw - int(Inches(1.24)),
            height=sh - int(Inches(2.6)),
            name=slot_name,
        )
    return Slot(
        slide_index=slide_index,
        left=int(Inches(0.8)),
        top=int(Inches(1.0)),
        width=sw - int(Inches(1.6)),
        height=sh - int(Inches(2.5)),
        name=slot_name,
    )
