"""The house-style template.

Templates resolve report -> tutkimus -> asiakas -> THIS. Making the default a
real .pptx rather than a special case in the renderer means one code path: every
deck is rendered into some template, and "no template chosen" simply selects
this one. A branch for "no template" would be a second rendering path that only
the default exercises, and it would drift.

Deliberately plain: a white ground, one very common typeface, no decoration.
The typeface follows the customer templates in use (Calibri) rather than
asserting one of its own.
A deck with no template chosen is either a draft or headed for a customer whose
own template has not been set up yet, and in both cases nSight's own branding on
the slide is wrong — it is not the analyst's brand and not the client's. Plain
also carries anywhere: white and Arial survive being pasted into whatever deck
the slides end up in.

The chart SERIES keep the house palette. Those are not decoration — they have to
be distinguishable from each other, and a set of greys would make a chart harder
to read rather than more neutral.
"""
from __future__ import annotations

from lxml import etree
import os
import pathlib

from pptx import Presentation
from pptx.util import Inches, Pt

from reportbuilder.render import house_style as hs

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# The chart series palette, in the order a multi-series chart consumes it.
# Teal leads because a single-series chart is the common case and teal is the
# house colour; the alternating light variants keep adjacent bands legible.
_ACCENTS = [hs.TEAL, hs.TEAL_LT, hs.BLUE, hs.BLUE_LT, hs.RED, hs.RED_LT]

# Calibri: the typeface Attendo's own deck declares, and the one a slide with
# no template chosen should therefore match — a default that looks like the
# customer templates it sits beside is one less thing to notice. It is on
# every Windows and Mac that will open the deck.
#
# On the RENDER host it is a name like any other: without Calibri installed,
# fontconfig answers with whatever it considers closest (Noto Sans here),
# whose metrics differ — so line breaks and title fitting are computed against
# a slightly different face than PowerPoint will draw. Installing Calibri, or
# metric-compatible Carlito, through Settings > Fonts removes that difference.
_HEADING_FONT = "Calibri"
_BODY_FONT = "Calibri"

# Plain white, not the house cream: see the module docstring.
_BACKGROUND = "#FFFFFF"


_TITLE_PT = 18  # matches render.image.slide_chrome.TITLE_PT


def _rgb(hex6: str):
    """'1F9AA5' -> RGBColor, accepting a leading '#'."""
    from pptx.dml.color import RGBColor
    h = _hex(hex6)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _hex(value: str) -> str:
    return value.lstrip("#").upper()


def _set_scheme(theme_root) -> None:
    """Overwrite the theme's colour and font scheme with the house style."""
    scheme = theme_root.find(f".//{{{_A}}}themeElements/{{{_A}}}clrScheme")
    if scheme is not None:
        for i, colour in enumerate(_ACCENTS, start=1):
            el = scheme.find(f"{{{_A}}}accent{i}")
            if el is None:
                continue
            for child in list(el):
                el.remove(child)
            srgb = etree.SubElement(el, f"{{{_A}}}srgbClr")
            srgb.set("val", _hex(colour))
        # dk1/lt1 drive default text and background, so the deck reads as the
        # house style even before a chart is placed.
        for tag, colour in (("dk1", hs.INK), ("lt1", _BACKGROUND)):
            el = scheme.find(f"{{{_A}}}{tag}")
            if el is None:
                continue
            for child in list(el):
                el.remove(child)
            srgb = etree.SubElement(el, f"{{{_A}}}srgbClr")
            srgb.set("val", _hex(colour))

    fonts = theme_root.find(f".//{{{_A}}}themeElements/{{{_A}}}fontScheme")
    if fonts is not None:
        for kind, name in (("majorFont", _HEADING_FONT), ("minorFont", _BODY_FONT)):
            latin = fonts.find(f"{{{_A}}}{kind}/{{{_A}}}latin")
            if latin is not None:
                latin.set("typeface", name)


def _set_placeholder_style(ph, *, size_pt: int, bold: bool, colour: str) -> None:
    """Set a placeholder's inherited text style: left, sized, bold, ink.

    Written into the layout placeholder's ``a:lstStyle`` rather than onto its
    paragraphs. A slide's empty placeholder inherits from the LIST STYLE; run
    properties on the layout's own paragraphs are not what it reads, which is
    why the title kept arriving centred at the master's 44pt.
    """
    from pptx.oxml.ns import qn

    bodyPr = ph.text_frame._txBody
    lst = bodyPr.find(qn("a:lstStyle"))
    if lst is None:
        lst = bodyPr.makeelement(qn("a:lstStyle"), {})
        bodyPr.insert(1, lst)   # after a:bodyPr
    lvl = lst.find(qn("a:lvl1pPr"))
    if lvl is None:
        lvl = lst.makeelement(qn("a:lvl1pPr"), {})
        lst.append(lvl)
    lvl.set("algn", "l")
    rpr = lvl.find(qn("a:defRPr"))
    if rpr is None:
        rpr = lvl.makeelement(qn("a:defRPr"), {})
        lvl.append(rpr)
    rpr.set("sz", str(int(size_pt * 100)))
    rpr.set("b", "1" if bold else "0")
    fill = rpr.makeelement(qn("a:solidFill"), {})
    srgb = fill.makeelement(qn("a:srgbClr"), {"val": _hex(colour)})
    fill.append(srgb)
    rpr.append(fill)


def _bake_house_furniture(prs, layout) -> None:
    """Draw the default background INTO *layout*.

    The house look used to be painted onto every slide by
    `render.image.slide_chrome`. That is why a customer's template appeared
    unused: their design was there, underneath our cream rectangle. Now the
    LAYOUT owns the design in both cases — theirs carries it already, ours has
    to carry it too — and the renderer only ever adds text and the chart.

    python-pptx exposes no `add_shape` on a layout, so the shapes are built on
    a scratch slide and their XML is moved into the layout's tree. They are
    inserted at the FRONT of the tree so they sit behind the placeholders.
    """
    from copy import deepcopy

    scratch = prs.slides.add_slide(layout)
    sw, sh = int(prs.slide_width), int(prs.slide_height)

    bg = scratch.shapes.add_shape(1, 0, 0, sw, sh)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(_BACKGROUND)
    bg.line.fill.background()
    bg.shadow.inherit = False

    # No accent bar. A teal stripe beside every title is nSight's own branding,
    # and a deck with no template chosen belongs to whoever is making it.

    # Place the layout's own placeholders where the renderer expects them.
    # PowerPoint's stock positions are indented far to the right and would leave
    # the title adrift from the slide's left margin.
    sw_i, sh_i = sw, sh
    for ph in layout.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:            # title
            ph.left, ph.top = Inches(0.80), Inches(0.34)
            ph.width, ph.height = sw_i - Inches(1.60), Inches(1.10)
            _set_placeholder_style(ph, size_pt=_TITLE_PT, bold=True,
                                   colour=hs.INK)
        elif idx == 1:          # content — becomes the chart slot
            ph.left, ph.top = Inches(0.62), Inches(1.90)
            ph.width, ph.height = sw_i - Inches(1.24), sh_i - Inches(2.60)

    spTree = layout.shapes._spTree
    # After nvGrpSpPr + grpSpPr (the two mandatory heads), i.e. behind
    # everything the layout already defines.
    spTree.insert(2, deepcopy(bg._element))

    # Drop the scratch slide: it was only a place to build the XML.
    rid = prs.slides._sldIdLst[-1].rId
    prs.part.drop_rel(rid)
    del prs.slides._sldIdLst[-1]


def build_default_template(path: str) -> str:
    """Write the house-style template to *path*; return the path.

    Built from python-pptx's blank presentation, whose layout 6 IS blank and
    whose layout 1 is Title and Content — so the result satisfies
    template_check's requirement (a title plus a large content placeholder)
    the same way a client template does, rather than by exception.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    master = prs.slide_master.part
    for rel in master.rels.values():
        if "theme" in rel.reltype:
            root = etree.fromstring(rel.target_part.blob)
            _set_scheme(root)
            rel.target_part._blob = etree.tostring(root, xml_declaration=True,
                                                   encoding="UTF-8", standalone=True)
            break

    # Layout 1 is Title and Content — the one template_check ranks highest and
    # therefore the one chart slides are built from.
    _bake_house_furniture(prs, prs.slide_layouts[1])

    prs.save(path)
    return path


# The template nSight itself ships, and what a fresh hive starts on.
#
# A real .pptx, designed rather than constructed: `build_default_template`
# below assembles a serviceable deck from python-pptx's blank presentation, and
# that is the fallback, but the house design is a file a designer made. Shipping
# it as an asset is what makes "restore nSight's template" mean the same thing
# on every install.
#
# A tenant can replace it — Settings > Default template — and that upload is
# stored in the hive, not here; this is only ever what they get back.
_SHIPPED = pathlib.Path(__file__).with_name("assets") / "nsight_default.pptx"


def default_template_bytes() -> bytes:
    """The bytes a fresh hive is seeded with, and what "restore" restores."""
    if _SHIPPED.exists():
        return _SHIPPED.read_bytes()
    # No asset in this checkout (a source tree without it, or a trimmed
    # package): build one, so a hive always has SOMETHING to render on.
    import tempfile as _tf
    fd, tmp = _tf.mkstemp(prefix="nsight-default-", suffix=".pptx")
    os.close(fd)
    try:
        build_default_template(tmp)
        return pathlib.Path(tmp).read_bytes()
    finally:
        os.unlink(tmp)


def shipped_default_name() -> str:
    """What to call the built-in, so the UI can name what is in effect."""
    return _SHIPPED.name if _SHIPPED.exists() else "nSight's own template"
