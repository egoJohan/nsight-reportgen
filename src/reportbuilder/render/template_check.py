"""Can we render into this customer template?

nSight's analysts work in their CLIENTS' brand templates — Attendo's deck,
Holiday Club's deck — which nobody is going to re-author to suit us. So the
requirements have to be things a normal PowerPoint template already satisfies,
and the check has to say plainly what is wrong when one does not.

Layout NAMES are useless as an interface: the same layout is "1 layout area",
"Innehåll" and "Title and Content" in three real client templates. Placeholder
TYPES are not — they are a PowerPoint concept, identical across templates,
languages and vendors. All three of those layouts expose TITLE + OBJECT.

So a usable layout is defined by its placeholder signature, and the analyst
confirms the choice once per template. The heuristic picks the obvious
candidate; the confirmation absorbs the templates where "obvious" is wrong.
"""
from __future__ import annotations

from dataclasses import dataclass, field

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

_A = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

# Where a chart goes. OBJECT is PowerPoint's "content" placeholder (the one with
# the insert-chart/table icons); BODY is the plain text variant. Both are
# rectangles we can size an image into.
_CONTENT = {PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.BODY}
_TITLE = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
# Decoration that says "this layout is for a photo, not a chart".
_IMAGERY = {PP_PLACEHOLDER.PICTURE}


@dataclass
class LayoutCandidate:
    index: int
    name: str
    score: int
    has_title: bool
    content_count: int
    has_picture: bool
    content_area_pct: float = 0.0  # largest content placeholder, % of the slide


@dataclass
class TemplateTheme:
    """The client's brand, as the template states it.

    accent1-6 is what PowerPoint's own charts use for series colours, so it is
    the right palette for ours — a chart drawn in nSight teal on Attendo's navy
    template is the thing an analyst would have to fix by hand.

    Fonts are the theme's major (headings) and minor (body) latin typefaces.
    NOTE: naming a font is not having it. If the render host lacks it, matplotlib
    and LibreOffice both substitute silently — see export/pdf_convert.py, which
    already warns that a missing font shifts label metrics.
    """
    palette: list[str] = field(default_factory=list)
    heading_font: str = ""
    body_font: str = ""


@dataclass
class TemplateReport:
    ok: bool
    slide_width_in: float = 0.0
    slide_height_in: float = 0.0
    layout_count: int = 0
    candidates: list[LayoutCandidate] = field(default_factory=list)
    problems: list[str] = field(default_factory=list)
    theme: TemplateTheme = field(default_factory=TemplateTheme)

    @property
    def best(self) -> LayoutCandidate | None:
        return self.candidates[0] if self.candidates else None


def _theme(prs) -> TemplateTheme:
    """Read the colour and font scheme out of the template's theme part.

    python-pptx exposes no theme API, so this walks the master's relationships
    to the theme XML. Failing to read a theme is not a validation failure — the
    template is still usable, it just does not hand us its brand.
    """
    theme = TemplateTheme()
    try:
        part = None
        for rel in prs.slide_master.part.rels.values():
            if "theme" in rel.reltype:
                part = rel.target_part
                break
        if part is None:
            return theme
        root = etree.fromstring(part.blob)
    except Exception:  # noqa: BLE001
        return theme

    for i in range(1, 7):
        el = root.find(f".//a:clrScheme/a:accent{i}", _A)
        if el is None:
            continue
        srgb = el.find("a:srgbClr", _A)
        sysc = el.find("a:sysClr", _A)
        value = srgb.get("val") if srgb is not None else (
            sysc.get("lastClr") if sysc is not None else None)
        if value:
            theme.palette.append(value.upper())

    for kind, attr in (("major", "heading_font"), ("minor", "body_font")):
        el = root.find(f".//a:fontScheme/a:{kind}Font/a:latin", _A)
        if el is not None and el.get("typeface"):
            setattr(theme, attr, el.get("typeface"))
    return theme


def _score(has_title: bool, content_count: int, has_picture: bool,
           area_pct: float) -> int:
    """Rank a layout's fitness for one chart plus a headline.

    SIZE is the signal, not the count. Counting placeholders ranked Attendo's
    title-only "Endast rubrik" and Holiday Club's "Section Header" above their
    real chart layouts, because a divider's caption box counts the same as a
    full-slide content area. A chart needs a big rectangle, so the largest
    content placeholder's share of the slide dominates the score.

    Exactly one content area still beats several: a two-content layout is for
    comparing two things, and one chart in one half leaves the slide lopsided.
    A picture placeholder marks a photo layout, so it ranks below others rather
    than being excluded — a template of nothing but photo layouts should still
    produce a deck.
    """
    if not has_title or content_count == 0:
        return 0
    # 0-100 from area alone: a placeholder filling half the slide scores ~70.
    score = int(min(area_pct, 60) / 60 * 100)
    score -= (content_count - 1) * 15
    if has_picture:
        score -= 25
    return max(score, 1)


def inspect_template(path: str) -> TemplateReport:
    """Check an uploaded .pptx and rank its layouts. Never raises on bad input."""
    try:
        prs = Presentation(path)
    except Exception as exc:  # noqa: BLE001 — any unreadable file is one problem
        return TemplateReport(ok=False, problems=[f"Not a readable PowerPoint file: {exc}"])

    report = TemplateReport(
        ok=False,
        slide_width_in=round((prs.slide_width or 0) / 914400, 2),
        slide_height_in=round((prs.slide_height or 0) / 914400, 2),
        layout_count=len(prs.slide_layouts),
    )

    for i, layout in enumerate(prs.slide_layouts):
        try:
            kinds = [p.placeholder_format.type for p in layout.placeholders]
        except Exception:  # noqa: BLE001 — a malformed layout is skipped, not fatal
            continue
        has_title = any(k in _TITLE for k in kinds)
        content_phs = [p for p in layout.placeholders
                       if p.placeholder_format.type in _CONTENT]
        picture = any(k in _IMAGERY for k in kinds)

        slide_area = (prs.slide_width or 1) * (prs.slide_height or 1)
        largest = max((int(p.width or 0) * int(p.height or 0) for p in content_phs),
                      default=0)
        area_pct = largest / slide_area * 100 if slide_area else 0.0

        score = _score(has_title, len(content_phs), picture, area_pct)
        if score:
            report.candidates.append(LayoutCandidate(
                index=i, name=layout.name, score=score, has_title=has_title,
                content_count=len(content_phs), has_picture=picture,
                content_area_pct=round(area_pct, 1)))

    # Highest score first; ties keep template order so the result is stable.
    report.candidates.sort(key=lambda c: (-c.score, c.index))

    if report.layout_count == 0:
        report.problems.append("The template contains no slide layouts.")
    if not report.candidates:
        report.problems.append(
            "No layout has both a title and a content placeholder. nSight needs "
            "one layout it can put a headline and a chart on.")
    if report.slide_width_in and report.slide_width_in < 9:
        report.problems.append(
            f"Slides are only {report.slide_width_in}in wide; charts will be cramped.")

    report.theme = _theme(prs)
    if not report.theme.palette:
        # Not fatal: we fall back to house colours, but the deck will not look
        # like the client's, which is the whole point of using their template.
        report.problems.append(
            "The template declares no theme colours; charts will use nSight's "
            "default palette instead of the client's brand.")

    report.ok = not [p for p in report.problems if "will use nSight" not in p]
    return report
